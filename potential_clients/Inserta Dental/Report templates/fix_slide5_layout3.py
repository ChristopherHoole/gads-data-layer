"""Final Slide 5 redesign: compact table, give cards proper breathing room.

Vertical layout:
  y=0.30 .. 1.10  Title + subtitle (existing)
  y=1.10 .. 1.40  Daily CPA heading
  y=1.45 .. 3.95  Table (compact, 2.50" tall, 7 rows)
  y=1.10 .. 4.00  Narrative card (parallel to table)
  y=4.15 .. 4.45  What Changed heading
  y=4.55 .. 4.85  Card headings
  y=4.90 .. 6.90  Card bodies (2.00" tall = room for 5-6 lines)
  y=7.00 .. 7.30  Footer
"""
import sys
from pathlib import Path

if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt

FILE_PATH = Path(__file__).resolve().parents[3] / "potential_clients" / "Inserta Dental" / "End-of-week reports" / "DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v6.pptx"


def main():
    prs = Presentation(str(FILE_PATH))
    slide = list(prs.slides)[4]

    table_shape = None
    heading_what_changed = None
    card_heading_shapes = []
    card_body_shapes = []
    narrative_body = None
    narrative_heading = None

    for shape in slide.shapes:
        if shape.has_table and shape.table.cell(0, 0).text.strip() == "Day":
            table_shape = shape
        elif shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt == "What Changed & What's Next":
                heading_what_changed = shape
            elif txt in ("Bid Strategy Switch", "Flagged for Week 2", "Recommendation"):
                card_heading_shapes.append(shape)
            elif ("From: Max Conversion Value" in txt or
                  "Google still reports 'most asset" in txt or
                  "Hold the new bid strategy untouched" in txt):
                card_body_shapes.append(shape)
            elif txt == "Week 1 Story":
                narrative_heading = shape
            elif "Bid strategy switched Wed 15 April" in txt:
                narrative_body = shape

    # Table: compact y=1.45, h=2.50 (rows 0.357" each — tight but readable)
    if table_shape:
        table_shape.top = Inches(1.45)
        table_shape.height = Inches(2.50)
        # Set explicit row heights to force compact layout
        for row in table_shape.table.rows:
            row.height = Inches(0.35)
        print(" [OK] Table: y=1.45, h=2.50, rows 0.35\" each")

    # Narrative card (align with table)
    if narrative_heading:
        narrative_heading.top = Inches(1.20)
    if narrative_body:
        narrative_body.top = Inches(1.65)
        narrative_body.height = Inches(2.30)
    # Narrative BG rectangles
    for shape in slide.shapes:
        if shape.shape_type == 1:
            left = shape.left / 914400 if shape.left else 0
            w = shape.width / 914400 if shape.width else 0
            if abs(left - 7.00) < 0.05 and abs(w - 0.08) < 0.03:
                shape.top = Inches(1.10)
                shape.height = Inches(2.90)
            elif abs(left - 7.08) < 0.05 and w > 5.5:
                shape.top = Inches(1.10)
                shape.height = Inches(2.90)
    print(" [OK] Narrative card: y=1.10..4.00")

    if heading_what_changed:
        heading_what_changed.top = Inches(4.15)
        print(" [OK] 'What Changed' heading: y=4.15")

    for shape in card_heading_shapes:
        shape.top = Inches(4.55)
    print(f" [OK] {len(card_heading_shapes)} card headings: y=4.55")

    for shape in card_body_shapes:
        shape.top = Inches(4.90)
        shape.height = Inches(2.00)  # room for 5-6 lines
    print(f" [OK] {len(card_body_shapes)} card bodies: y=4.90, h=2.00")

    # Now the card BODY backgrounds. Find the shape rectangles that are behind each card.
    # Original v4/v5 cards had fill rects; these are the ones with rounded-rect or rectangle type
    # at the card positions (x=0.85, 4.95, 9.05)
    for shape in slide.shapes:
        if shape.shape_type == 1:
            left = shape.left / 914400 if shape.left else 0
            top = shape.top / 914400 if shape.top else 0
            w = shape.width / 914400 if shape.width else 0
            # Card background rects (width ~3.40-3.50, at x=0.85, 4.95, 9.05)
            if w > 3.0 and w < 4.0 and (abs(left - 0.85) < 0.1 or abs(left - 4.95) < 0.1 or abs(left - 9.05) < 0.1):
                # Move these card BGs to match new card positions
                shape.top = Inches(4.45)
                shape.height = Inches(2.55)  # covers heading y=4.55 to body end y=6.90
    # Left-edge colour borders (narrow, same x as card fill rects)
    for shape in slide.shapes:
        if shape.shape_type == 1:
            left = shape.left / 914400 if shape.left else 0
            w = shape.width / 914400 if shape.width else 0
            # Narrow coloured border strips (w<0.15) at x~0.75, ~4.85, ~8.95 (just to left of card fills)
            if w < 0.15 and (abs(left - 0.75) < 0.1 or abs(left - 4.85) < 0.1 or abs(left - 8.95) < 0.1):
                shape.top = Inches(4.45)
                shape.height = Inches(2.55)
    print(" [OK] Card backgrounds adjusted to y=4.45, h=2.55")

    prs.save(str(FILE_PATH))
    print(f"\nSaved: {FILE_PATH.name}")


if __name__ == "__main__":
    main()
