"""Final Slide 5 layout fix: absolute positioning so nothing overflows footer."""
import sys
from pathlib import Path

if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches

FILE_PATH = Path(__file__).resolve().parents[3] / "potential_clients" / "Inserta Dental" / "End-of-week reports" / "DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v6.pptx"


def main():
    prs = Presentation(str(FILE_PATH))
    slide = list(prs.slides)[4]

    # Vertical budget:
    # y=0.30 .. y=1.10  title + subtitle (0.80")
    # y=1.10 .. y=1.40  Daily CPA heading (0.30")
    # y=1.45 .. y=4.45  Table (3.00")        ← CHANGED (was 1.45..3.95 in v6 fix 1, then 1.45..4.95)
    # y=4.60 .. y=4.90  What Changed heading (0.30")
    # y=5.00 .. y=5.30  Card headings (0.30")
    # y=5.30 .. y=6.80  Card bodies (1.50")  ← BIGGER
    # y=7.00 .. y=7.30  Footer

    table_shape = None
    heading_what_changed = None
    card_heading_shapes = []
    card_body_shapes = []
    narrative_body = None
    narrative_heading = None

    for shape in slide.shapes:
        if shape.has_table and shape.table.cell(0, 0).text.strip() == "Day":
            table_shape = shape
        elif shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt == "What Changed & What's Next":
                heading_what_changed = shape
            elif txt in ("Bid Strategy Switch", "Flagged for Week 2", "Recommendation"):
                card_heading_shapes.append(shape)
            elif ("From: Max Conversion Value" in txt or
                  "Google still reports 'most asset" in txt or
                  "Hold the new bid strategy untouched" in txt):
                card_body_shapes.append(shape)
            elif txt == "Week 1 Story":
                narrative_heading = shape
            elif "Bid strategy switched Wed 15 April" in txt:
                narrative_body = shape

    # Set absolute positions/sizes
    if table_shape:
        table_shape.top = Inches(1.45)
        table_shape.height = Inches(3.00)
        print(" [OK] Table: y=1.45, h=3.00")

    if heading_what_changed:
        heading_what_changed.top = Inches(4.60)
        print(" [OK] 'What Changed' heading: y=4.60")

    for shape in card_heading_shapes:
        shape.top = Inches(5.00)
    print(f" [OK] {len(card_heading_shapes)} card headings: y=5.00")

    for shape in card_body_shapes:
        shape.top = Inches(5.30)
        shape.height = Inches(1.50)
    print(f" [OK] {len(card_body_shapes)} card bodies: y=5.30, h=1.50")

    # Narrative card — align with table top/bottom
    if narrative_heading:
        narrative_heading.top = Inches(1.20)
    if narrative_body:
        narrative_body.top = Inches(1.65)
        narrative_body.height = Inches(2.80)

    # Narrative card BG rectangles — need to extend to match table bottom (y=4.45)
    # Card covers y=1.10 .. y=4.45, so h=3.35
    for shape in slide.shapes:
        if shape.shape_type == 1:
            left = shape.left / 914400 if shape.left else 0
            w = shape.width / 914400 if shape.width else 0
            # Navy border strip (x~7.00, w~0.08)
            if abs(left - 7.00) < 0.05 and abs(w - 0.08) < 0.03:
                shape.top = Inches(1.10)
                shape.height = Inches(3.35)
            # White body rect (x~7.08, w>5.5)
            elif abs(left - 7.08) < 0.05 and w > 5.5:
                shape.top = Inches(1.10)
                shape.height = Inches(3.35)
    print(" [OK] Narrative card: y=1.10..4.45 (matches table)")

    prs.save(str(FILE_PATH))
    print(f"\nSaved: {FILE_PATH.name}")


if __name__ == "__main__":
    main()
