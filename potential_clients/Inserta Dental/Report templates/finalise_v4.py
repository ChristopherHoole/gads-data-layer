"""Final polish pass on v4.

- Remove leftover charts from the reference (slides 3, 5, 7, 8, 9)
- Fix slide 12 bullet #4 label (was 4., became 12.)
- Expand slide 8 table from 5 data rows to 8 data rows (full 8-week context)
"""
import sys, copy
from pathlib import Path

if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

FILE_PATH = Path(__file__).resolve().parents[3] / "potential_clients" / "Inserta Dental" / "End-of-week reports" / "DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v4.pptx"

NAVY = RGBColor(0x1A, 0x23, 0x7E)


def delete_charts_from_slide(slide):
    """Delete large embedded chart images (type=13 PICTURE, not charts).

    The reference report's charts are embedded PNG images, not native charts.
    Identify them by: PICTURE type + large size (w>3" and h>2") + NOT the footer logo.
    """
    to_remove = []
    for s in slide.shapes:
        if s.shape_type == 13:  # PICTURE
            w = s.width / 914400 if s.width else 0
            h = s.height / 914400 if s.height else 0
            if w > 3 and h > 2:  # large image = a chart, not the logo
                to_remove.append(s)
    for s in to_remove:
        s._element.getparent().remove(s._element)
    return len(to_remove)


def add_row_to_table(table, row_data, is_alt=False):
    """Clone the last row XML and populate it with row_data."""
    tbl_elem = table._tbl
    last_tr = tbl_elem.findall(qn('a:tr'))[-1]
    new_tr = copy.deepcopy(last_tr)
    last_tr.addnext(new_tr)
    # Get the new row via table.cell()
    new_row_idx = len(table.rows) - 1
    for c_idx, val in enumerate(row_data):
        cell = table.cell(new_row_idx, c_idx)
        tf = cell.text_frame
        # Preserve first-run formatting
        saved = None
        for p in tf.paragraphs:
            for r in p.runs:
                f = r.font
                saved = {
                    'name': f.name, 'size': f.size, 'bold': f.bold,
                    'align': p.alignment,
                }
                try:
                    saved['color'] = f.color.rgb if f.color and f.color.type else None
                except Exception:
                    saved['color'] = None
                break
            if saved: break
        tf.clear()
        p = tf.paragraphs[0]
        if saved and saved.get('align'):
            p.alignment = saved['align']
        r = p.add_run()
        r.text = val
        if saved:
            if saved['name']: r.font.name = saved['name']
            if saved['size']: r.font.size = saved['size']
            if saved['bold'] is not None: r.font.bold = saved['bold']
            if saved['color']: r.font.color.rgb = saved['color']


def main():
    prs = Presentation(str(FILE_PATH))
    slides = list(prs.slides)

    # Delete leftover charts
    for slide_num in [3, 5, 7, 8, 9]:
        n = delete_charts_from_slide(slides[slide_num - 1])
        print(f" [OK] Slide {slide_num}: removed {n} leftover chart(s)")

    # Slide 8: expand table from 5 data rows to 8 data rows
    slide_8 = slides[7]
    for shape in slide_8.shapes:
        if shape.has_table:
            # Check it's the 8-week table (4 columns: Week, Spend, Primary Conv, CPA)
            if len(shape.table.columns) == 4 and shape.table.cell(0, 0).text.strip() == 'Week':
                # Add 3 more rows (for 30 Mar, 6 Apr, 13 Apr)
                add_row_to_table(shape.table, ['30 Mar', '£9,997', '148.4', '£67'])
                add_row_to_table(shape.table, ['6 Apr', '£10,291', '134.1', '£77'])
                add_row_to_table(shape.table, ['13 Apr', '£8,462', '101', '£84'])
                print(" [OK] Slide 8: expanded table to 8 weeks (added 30 Mar, 6 Apr, 13 Apr)")
                break

    # Slide 12: fix bullet #4 — currently shows "12." instead of "4."
    slide_12 = slides[11]
    for shape in slide_12.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "12.  Conversion rate improvements" in text:
                # Replace run by run
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if "12." in r.text and "Conversion" in r.text:
                            r.text = r.text.replace("12.", "4.")
                        elif r.text.strip() == "12.":
                            r.text = "4."
                        elif r.text.strip() == "12":
                            # Check next run for Conversion
                            r.text = "4"
                print(" [OK] Slide 12: bullet #4 label restored")
                break

    # Slide 12: also check the big "12" in stat card (Compliance Actions)
    # This is correct — keep as "12"
    # And the bullet "12. Conversion rate improvements" should be "4."
    # The replace_text_in_shape_any(slide, "4", "12") was too aggressive.
    # For the stat card "12" (Compliance Actions), should be 12 — correct now.

    prs.save(str(FILE_PATH))
    print(f"\nSaved: {FILE_PATH.name}")


if __name__ == "__main__":
    main()
