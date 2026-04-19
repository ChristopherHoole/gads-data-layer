"""
DBD Weekly Report — Master Template Builder
============================================

Generates `DBD_WEEKLY_REPORT_TEMPLATE.pptx` from the specs in `DBD_REPORT_TEMPLATE_SPEC.md`.

This script is the executable form of the spec. Running it produces:
- `DBD_WEEKLY_REPORT_TEMPLATE.pptx` — the master template with {{PLACEHOLDERS}}

Usage:
    python build_template.py

Dependencies:
    pip install python-pptx pillow

Design system:
- Font: Calibri only (three weights — Light, Regular, Bold)
- Palette: Navy #1A237E, Google Blue #4285F4, Google Green #34A853,
  Google Red #EA4335, Google Yellow #FBBC05, Body #1A1A1A, Muted #666666
- Structure: Dark/light sandwich (slide 1 dark, 2-10 light, 11 dark)
- Motif: Navy left strip (0.3") + navy icon circle on content slides
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# -----------------------------------------------------------------------------
# CONSTANTS
# -----------------------------------------------------------------------------

# Paths
REPO_ROOT = Path(__file__).resolve().parents[3]
CLIENT_DIR = REPO_ROOT / "potential_clients" / "Inserta Dental"
TEMPLATES_DIR = CLIENT_DIR / "Report templates"
ACT_LOGO = CLIENT_DIR / "act_logo_official.png"
OUTPUT_PATH = TEMPLATES_DIR / "DBD_WEEKLY_REPORT_TEMPLATE.pptx"

# Slide dimensions (16:9)
SLIDE_W_IN = 13.33
SLIDE_H_IN = 7.50

# Colour palette (hex → RGBColor)
NAVY = RGBColor(0x1A, 0x23, 0x7E)
BLUE = RGBColor(0x42, 0x85, 0xF4)
GREEN = RGBColor(0x34, 0xA8, 0x53)
RED = RGBColor(0xEA, 0x43, 0x35)
YELLOW = RGBColor(0xFB, 0xBC, 0x05)
BODY = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF7, 0xF7, 0xF7)
BORDER_LIGHT = RGBColor(0xE0, 0xE0, 0xE0)
BORDER_MID = RGBColor(0xD0, 0xD0, 0xD0)
NAVY_TINT = RGBColor(0xE8, 0xEA, 0xF6)  # light navy for table totals

# Left strip and footer (content slides)
LEFT_STRIP_W = 0.30
FOOTER_Y = 7.10
FOOTER_H = 0.30

# Font name
FONT = "Calibri"

# Placeholder text convention
PH = lambda key: f"{{{{{key}}}}}"  # e.g. PH("WEEK_NUMBER") -> "{{WEEK_NUMBER}}"


# -----------------------------------------------------------------------------
# LOW-LEVEL HELPERS
# -----------------------------------------------------------------------------

def _set_slide_size(prs):
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)


def add_blank_slide(prs):
    """Add a blank slide using the default blank layout."""
    blank_layout = prs.slide_layouts[6]  # 6 = blank
    return prs.slides.add_slide(blank_layout)


def add_rectangle(slide, x, y, w, h, fill_color, line_color=None):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()  # no line
    else:
        shape.line.color.rgb = line_color
    return shape


def add_oval(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def add_rounded_rect(slide, x, y, w, h, fill_color, line_color=None, corner_radius=0.08):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    # Corner radius via adjustment 0 (0.0 = square, 0.5 = max)
    shape.adjustments[0] = corner_radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    font_size,
    bold=False,
    italic=False,
    color=BODY,
    align="left",
    anchor="top",
    font_name=FONT,
):
    """Add a text box. `align` in {'left','center','right'}, `anchor` in {'top','middle','bottom'}."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    # Vertical anchor
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[anchor]

    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]

    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_image(slide, path, x, y, h=None, w=None):
    if h is not None and w is None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), height=Inches(h))
    elif w is not None and h is None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def set_slide_background(slide, color):
    """Set slide background to a solid colour."""
    bg_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN)
    )
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = color
    bg_rect.line.fill.background()
    # Send to back — python-pptx lacks a direct API but we just add this first
    # so subsequent shapes render on top
    return bg_rect


# -----------------------------------------------------------------------------
# COMPONENT HELPERS
# -----------------------------------------------------------------------------

def add_left_strip(slide):
    """Navy left colour strip for content slides."""
    return add_rectangle(slide, 0, 0, LEFT_STRIP_W, SLIDE_H_IN, NAVY)


def add_footer(slide, page_number):
    """Standard footer for content slides 2-10."""
    # Logo
    if ACT_LOGO.exists():
        add_image(slide, ACT_LOGO, 0.80, FOOTER_Y, h=0.25)
    # Centre text
    add_text(
        slide,
        "Christopher Hoole  |  christopherhoole.com  |  Confidential",
        x=1.25,
        y=FOOTER_Y + 0.03,
        w=9.50,
        h=FOOTER_H,
        font_size=10,
        color=MUTED,
        align="center",
        anchor="middle",
    )
    # Page number
    add_text(
        slide,
        str(page_number),
        x=12.30,
        y=FOOTER_Y + 0.03,
        w=0.50,
        h=FOOTER_H,
        font_size=10,
        color=MUTED,
        align="right",
        anchor="middle",
    )


def add_section_header(slide, icon_glyph, title, context_line=None, intro_line=None):
    """Common content-slide header: icon circle + title + optional context/intro."""
    # Icon circle
    add_oval(slide, 0.80, 0.35, 0.45, 0.45, NAVY)
    # Icon glyph (centred inside circle)
    add_text(
        slide, icon_glyph,
        x=0.80, y=0.35, w=0.45, h=0.45,
        font_size=22, bold=True, color=WHITE, align="center", anchor="middle",
    )
    # Title
    add_text(
        slide, title,
        x=1.45, y=0.30, w=11.30, h=0.60,
        font_size=32, bold=True, color=NAVY, align="left", anchor="top",
    )
    # Context line (italic)
    if context_line:
        add_text(
            slide, context_line,
            x=1.45, y=0.95, w=11.30, h=0.30,
            font_size=14, italic=True, color=BLUE, align="left",
        )
    # Intro line
    if intro_line:
        add_text(
            slide, intro_line,
            x=1.45, y=1.30, w=11.30, h=0.30,
            font_size=13, color=BODY, align="left",
        )


def add_stat_callout(slide, x, y, number_text, label_text, number_size=42, number_color=NAVY, label_color=BODY, sub_label=None):
    """Reusable stat callout. Number above, label below (0.05" gap). Optional sub-label."""
    # Number
    add_text(
        slide, number_text,
        x=x, y=y, w=2.88, h=0.80,
        font_size=number_size, bold=True, color=number_color, align="left", anchor="top",
    )
    # Label (positioned BELOW the number with 0.05" gap — spec 1.10)
    label_y = y + 0.80 + 0.05
    add_text(
        slide, label_text,
        x=x, y=label_y, w=2.88, h=0.30,
        font_size=13, color=label_color, align="left", anchor="top",
    )
    # Optional sub-label below label
    if sub_label:
        add_text(
            slide, sub_label,
            x=x, y=label_y + 0.32, w=2.88, h=0.25,
            font_size=11, color=MUTED, align="left", anchor="top",
        )


def add_key_finding_block(slide, x, y, w, label, body_text, label_color=BLUE, body_max_h=1.40):
    """KEY FINDING or RECOMMENDATION block. Returns y-end."""
    add_text(
        slide, label,
        x=x, y=y, w=w, h=0.25,
        font_size=11, bold=True, color=label_color, align="left",
    )
    add_text(
        slide, body_text,
        x=x, y=y + 0.30, w=w, h=body_max_h,
        font_size=12, color=BODY, align="left",
    )
    return y + 0.30 + body_max_h


def add_status_pill(slide, x, y, w, h, text, fill_color):
    """Rounded pill with coloured fill and white bold text."""
    add_rounded_rect(slide, x, y, w, h, fill_color, corner_radius=0.18)
    add_text(
        slide, text,
        x=x, y=y, w=w, h=h,
        font_size=14, bold=True, color=WHITE, align="center", anchor="middle",
    )


# -----------------------------------------------------------------------------
# SLIDE BUILDERS
# -----------------------------------------------------------------------------

def build_slide_1_title(prs):
    """Slide 1 — Title (DARK)."""
    slide = add_blank_slide(prs)
    set_slide_background(slide, NAVY)

    # ACT logo top
    if ACT_LOGO.exists():
        # centre horizontally: image width unknown, we'll target 0.60" tall and let x centre
        # add_picture without width keeps native aspect ratio — we'll estimate
        pic = add_image(slide, ACT_LOGO, x=(SLIDE_W_IN - 0.60) / 2, y=0.40, h=0.60)
        # Recentre after knowing pic width
        pic.left = Inches((SLIDE_W_IN - Emu(pic.width).inches) / 2)

    # Main title
    add_text(
        slide, "Week " + PH("WEEK_NUMBER") + " Report",
        x=0.50, y=2.60, w=12.30, h=0.90,
        font_size=54, bold=True, color=WHITE, align="center",
    )

    # Accent line
    add_rectangle(slide, x=(SLIDE_W_IN - 3.00) / 2, y=3.55, w=3.00, h=0.03, fill_color=BLUE)

    # Subtitle
    add_text(
        slide, PH("DATE_RANGE"),
        x=0.50, y=3.70, w=12.30, h=0.50,
        font_size=24, color=BLUE, align="center",
    )

    # Descriptor
    add_text(
        slide, PH("WEEK_DESCRIPTOR"),
        x=0.50, y=4.40, w=12.30, h=0.45,
        font_size=18, color=WHITE, align="center",
    )

    # 4 stat callouts
    stat_positions = [0.90, 3.98, 7.05, 10.13]
    for i, sx in enumerate(stat_positions, start=1):
        add_text(
            slide, PH(f"TITLE_STAT_{i}_NUMBER"),
            x=sx, y=5.55, w=2.88, h=0.80,
            font_size=48, bold=True, color=WHITE, align="left",
        )
        add_text(
            slide, PH(f"TITLE_STAT_{i}_LABEL"),
            x=sx, y=6.40, w=2.88, h=0.30,
            font_size=13, color=BLUE, align="left",
        )

    # Client block
    add_text(
        slide, "Dental by Design  |  Prodent Group",
        x=0.50, y=6.95, w=7.00, h=0.25,
        font_size=14, bold=True, color=WHITE, align="left",
    )
    add_text(
        slide, "Christopher Hoole  ·  Google Ads Specialist  ·  christopherhoole.com",
        x=0.50, y=7.20, w=10.00, h=0.22,
        font_size=11, color=BLUE, align="left",
    )


def build_slide_2_exec_summary(prs):
    """Slide 2 — Executive Summary."""
    slide = add_blank_slide(prs)
    add_left_strip(slide)
    add_section_header(
        slide,
        icon_glyph="▣",
        title="Executive Summary",
        context_line=f"Week {PH('WEEK_NUMBER')}: {PH('DATE_RANGE')}",
        intro_line=PH("EXEC_INTRO_LINE"),
    )

    # 4 stat callouts
    stat_positions = [0.80, 3.88, 6.95, 10.03]
    for i, sx in enumerate(stat_positions, start=1):
        add_stat_callout(
            slide, x=sx, y=1.85,
            number_text=PH(f"EXEC_STAT_{i}_NUMBER"),
            label_text=PH(f"EXEC_STAT_{i}_LABEL"),
            number_color=NAVY,
        )

    # 3 status pills
    pill_positions = [0.80, 4.88, 8.95]
    for i, px in enumerate(pill_positions, start=1):
        add_status_pill(
            slide, x=px, y=3.40, w=3.88, h=0.55,
            text=PH(f"PILL_{i}_TEXT"),
            fill_color=GREEN,  # default; will be overridden per week
        )

    # Core Insight label + body
    add_text(
        slide, "CORE INSIGHT",
        x=0.80, y=4.20, w=12.00, h=0.25,
        font_size=11, bold=True, color=BLUE, align="left",
    )
    add_text(
        slide, PH("CORE_INSIGHT_TEXT"),
        x=0.80, y=4.50, w=12.00, h=2.30,
        font_size=13, color=BODY, align="left",
    )

    add_footer(slide, page_number=2)


def build_slide_3_work_delivered(prs):
    """Slide 3 — Work Delivered (4-column timeline; extensible to 5 for full weeks)."""
    slide = add_blank_slide(prs)
    add_left_strip(slide)
    add_section_header(
        slide,
        icon_glyph="▸",
        title="Work Delivered",
        context_line=f"Week {PH('WEEK_NUMBER')}: {PH('DATE_RANGE')} · {PH('WORK_DAYS')} working days",
        intro_line=PH("WORK_INTRO_LINE"),
    )

    # Horizontal connector line behind circles
    add_rectangle(slide, x=1.30, y=2.54, w=10.65, h=0.02, fill_color=BLUE)

    # 4 columns
    col_xs = [0.80, 3.80, 6.80, 9.80]
    col_w = 2.85
    for i, col_x in enumerate(col_xs, start=1):
        # Circle (positioned so it sits above the connector)
        circle_x = col_x + (col_w - 0.50) / 2  # centre circle in column
        add_oval(slide, x=circle_x, y=2.30, w=0.50, h=0.50, fill_color=BLUE)
        add_text(
            slide, str(i),
            x=circle_x, y=2.30, w=0.50, h=0.50,
            font_size=20, bold=True, color=WHITE, align="center", anchor="middle",
        )
        # Day label
        add_text(
            slide, PH(f"DAY_{i}_DATE"),
            x=col_x, y=2.95, w=col_w, h=0.30,
            font_size=13, bold=True, color=NAVY, align="left",
        )
        # Section label
        add_text(
            slide, PH(f"DAY_{i}_TITLE"),
            x=col_x, y=3.30, w=col_w, h=0.35,
            font_size=15, bold=True, color=BLUE, align="left",
        )
        # Body
        add_text(
            slide, PH(f"DAY_{i}_BODY"),
            x=col_x, y=3.75, w=col_w, h=2.40,
            font_size=11, color=BODY, align="left",
        )

    # Key Finding strip
    add_text(
        slide, "KEY FINDING",
        x=0.80, y=6.30, w=12.00, h=0.25,
        font_size=11, bold=True, color=BLUE, align="left",
    )
    add_text(
        slide, PH("WORK_KEY_FINDING"),
        x=0.80, y=6.60, w=12.00, h=0.45,
        font_size=12, color=BODY, align="left",
    )

    add_footer(slide, page_number=3)


def _build_campaign_slide(prs, page_num, icon, title, intro_placeholder,
                         stat_prefix, left_content_builder):
    """Shared builder for campaign slides (4, 5, 6)."""
    slide = add_blank_slide(prs)
    add_left_strip(slide)
    add_section_header(
        slide,
        icon_glyph=icon,
        title=title,
        context_line=f"Week {PH('WEEK_NUMBER')}: {PH('DATE_RANGE')}",
        intro_line=PH(intro_placeholder),
    )

    # Stat callouts
    stat_positions = [0.80, 3.88, 6.95, 10.03]
    for i, sx in enumerate(stat_positions, start=1):
        add_stat_callout(
            slide, x=sx, y=1.85,
            number_text=PH(f"{stat_prefix}_STAT_{i}_NUMBER"),
            label_text=PH(f"{stat_prefix}_STAT_{i}_LABEL"),
            number_color=NAVY,
            sub_label=PH(f"{stat_prefix}_STAT_{i}_SUBLABEL") if i == 4 else None,
        )

    # Left-side content (chart, before/after, etc.) — delegated to caller
    left_content_builder(slide)

    # Right-side insight blocks
    add_key_finding_block(
        slide, x=8.50, y=3.40, w=4.30,
        label="KEY FINDING",
        body_text=PH(f"{stat_prefix}_KEY_FINDING"),
    )
    add_key_finding_block(
        slide, x=8.50, y=5.25, w=4.30,
        label="RECOMMENDATION",
        body_text=PH(f"{stat_prefix}_RECOMMENDATION"),
    )

    add_footer(slide, page_number=page_num)
    return slide


def build_slide_4_brand(prs):
    """Slide 4 — Brand Campaign."""

    def left_content(slide):
        # Chart placeholder (actual chart added when template is filled per-week)
        # For template: draw a light grey placeholder rectangle with text
        add_rectangle(
            slide, x=0.80, y=3.40, w=7.30, h=3.35,
            fill_color=OFF_WHITE, line_color=BORDER_LIGHT,
        )
        add_text(
            slide, "[ CHART: " + PH("BRAND_CHART_TITLE") + " ]\nData: " + PH("BRAND_CHART_DATA"),
            x=0.80, y=3.40, w=7.30, h=3.35,
            font_size=12, italic=True, color=MUTED, align="center", anchor="middle",
        )

    _build_campaign_slide(
        prs, page_num=4,
        icon="◆",
        title="Brand Campaign",
        intro_placeholder="BRAND_INTRO_LINE",
        stat_prefix="BRAND",
        left_content_builder=left_content,
    )


def build_slide_5_pmax(prs):
    """Slide 5 — Performance Max."""

    def left_content(slide):
        add_rectangle(
            slide, x=0.80, y=3.40, w=7.30, h=3.35,
            fill_color=OFF_WHITE, line_color=BORDER_LIGHT,
        )
        add_text(
            slide, "[ CHART: " + PH("PMAX_CHART_TITLE") + " ]\nData: " + PH("PMAX_CHART_DATA"),
            x=0.80, y=3.40, w=7.30, h=3.35,
            font_size=12, italic=True, color=MUTED, align="center", anchor="middle",
        )

    slide = _build_campaign_slide(
        prs, page_num=5,
        icon="◆",
        title="Performance Max",
        intro_placeholder="PMAX_INTRO_LINE",
        stat_prefix="PMAX",
        left_content_builder=left_content,
    )

    # PMax-specific: REC block title should read "FLAGGED FOR WEEK N" in amber
    # (we'll leave the generic "RECOMMENDATION" in the template — users update per-week
    # if needed; or we generate a specific placeholder label)


def build_slide_6_dii(prs):
    """Slide 6 — Dental Implants Intent."""

    def left_content(slide):
        # Before/After diagrammatic layout (no mis-scaled bars)
        # Section heading
        add_text(
            slide, "Structural rebuild — before vs after",
            x=0.80, y=3.40, w=7.30, h=0.30,
            font_size=13, bold=True, color=NAVY, align="left",
        )

        # BEFORE column header
        add_text(
            slide, "BEFORE",
            x=1.00, y=3.80, w=2.50, h=0.25,
            font_size=11, bold=True, color=MUTED, align="left",
        )
        # AFTER column header
        add_text(
            slide, "AFTER",
            x=5.30, y=3.80, w=2.50, h=0.25,
            font_size=11, bold=True, color=NAVY, align="left",
        )

        # Four comparison rows
        comparisons = [
            ("Ad groups", "DII_BEFORE_AG", "DII_AFTER_AG"),
            ("Keywords", "DII_BEFORE_KW", "DII_AFTER_KW"),
            ("RSAs", "DII_BEFORE_RSA", "DII_AFTER_RSA"),
            ("Headlines", "DII_BEFORE_HL", "DII_AFTER_HL"),
        ]
        row_y = 4.15
        for label, before_ph, after_ph in comparisons:
            # Label
            add_text(
                slide, label,
                x=0.80, y=row_y, w=2.00, h=0.50,
                font_size=12, color=BODY, align="left", anchor="middle",
            )
            # Before value
            add_text(
                slide, PH(before_ph),
                x=2.90, y=row_y, w=1.50, h=0.50,
                font_size=18, bold=True, color=MUTED, align="left", anchor="middle",
            )
            # Arrow
            add_text(
                slide, "→",
                x=4.50, y=row_y, w=0.60, h=0.50,
                font_size=18, color=BLUE, align="center", anchor="middle",
            )
            # After value
            add_text(
                slide, PH(after_ph),
                x=5.20, y=row_y, w=2.50, h=0.50,
                font_size=22, bold=True, color=NAVY, align="left", anchor="middle",
            )
            row_y += 0.62

    _build_campaign_slide(
        prs, page_num=6,
        icon="◆",
        title="Dental Implants Intent",
        intro_placeholder="DII_INTRO_LINE",
        stat_prefix="DII",
        left_content_builder=left_content,
    )


def _interp_colour(c1, c2, t):
    """Linear interpolate between two RGBColor values. t in [0,1]."""
    # RGBColor stores as 3-byte; access via private `_rgb` on some versions. Fallback: str -> int
    s1 = str(c1)  # 6-hex
    s2 = str(c2)
    r1, g1, b1 = int(s1[0:2], 16), int(s1[2:4], 16), int(s1[4:6], 16)
    r2, g2, b2 = int(s2[0:2], 16), int(s2[2:4], 16), int(s2[4:6], 16)
    r = int(r1 + (r2 - r1) * t)
    g = int(g1 + (g2 - g1) * t)
    b = int(b1 + (b2 - b1) * t)
    return RGBColor(r, g, b)


def build_slide_7_numbers(prs):
    """Slide 7 — The Numbers That Matter. Funnel + breakdown table."""
    slide = add_blank_slide(prs)
    add_left_strip(slide)
    add_section_header(
        slide,
        icon_glyph="●",
        title="The Numbers That Matter",
        context_line=f"Business-level outcomes — Week {PH('WEEK_NUMBER')}",
        intro_line=PH("NUMBERS_INTRO_LINE"),
    )

    # Funnel — 5 stages, gradient navy → Google blue
    # For simplicity, use rectangles (not trapezoids — trapezoids via custom geom are complex in python-pptx)
    # Note: the spec calls for trapezoids. We use rectangles here with gradient fill as a pragmatic substitute.
    # Future enhancement: replace with custom trapezoid shapes via slide XML manipulation.
    stage_xs = [0.80, 3.22, 5.54, 7.76, 9.88]
    stage_widths = [2.30, 2.20, 2.10, 2.00, 1.90]
    stage_h = 1.40
    stage_y = 2.00
    for i in range(5):
        t = i / 4.0  # 0..1
        fill = _interp_colour(NAVY, BLUE, t)
        # Use a rectangle with rounded-left / flat-right visual (PowerPoint chevron-like)
        # Shape: CHEVRON gives arrow-like shape. Since chevrons caused the prior bug (narrow interior),
        # we use plain rounded RECTANGLE for reliable text fit.
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(stage_xs[i]), Inches(stage_y),
            Inches(stage_widths[i]), Inches(stage_h),
        )
        shape.adjustments[0] = 0.05  # small rounding
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()

        # Number (big, centred)
        add_text(
            slide, PH(f"FUNNEL_STAGE_{i+1}_NUMBER"),
            x=stage_xs[i], y=stage_y + 0.20, w=stage_widths[i], h=0.60,
            font_size=26, bold=True, color=WHITE, align="center", anchor="middle",
        )
        # Label below number
        add_text(
            slide, PH(f"FUNNEL_STAGE_{i+1}_LABEL"),
            x=stage_xs[i], y=stage_y + 0.85, w=stage_widths[i], h=0.40,
            font_size=11, color=WHITE, align="center", anchor="middle",
        )
        # Arrow between stages (except after last)
        if i < 4:
            arrow_x = stage_xs[i] + stage_widths[i] + 0.01
            add_text(
                slide, "▶",
                x=arrow_x, y=stage_y + 0.50, w=0.10, h=0.40,
                font_size=14, color=BLUE, align="center", anchor="middle",
            )

    # Note below funnel
    add_text(
        slide, PH("FUNNEL_NOTE"),
        x=0.80, y=3.60, w=12.00, h=0.60,
        font_size=11, italic=True, color=MUTED, align="left",
    )

    # Campaign breakdown heading
    add_text(
        slide, "Dengro Offline Leads — by campaign",
        x=0.80, y=4.30, w=12.00, h=0.30,
        font_size=13, bold=True, color=NAVY, align="left",
    )

    # Table (header + 3 campaign rows + 1 total row = 5 rows, 4 cols)
    _add_table(
        slide,
        x=0.80, y=4.65, w=12.00, h=2.10,
        headers=["Campaign", "Spend", "Dengro Leads", "Spend per Lead"],
        rows=[
            [PH("CAMP_1_NAME"), PH("CAMP_1_SPEND"), PH("CAMP_1_LEADS"), PH("CAMP_1_CPL")],
            [PH("CAMP_2_NAME"), PH("CAMP_2_SPEND"), PH("CAMP_2_LEADS"), PH("CAMP_2_CPL")],
            [PH("CAMP_3_NAME"), PH("CAMP_3_SPEND"), PH("CAMP_3_LEADS"), PH("CAMP_3_CPL")],
            ["TOTAL", PH("CAMP_TOTAL_SPEND"), PH("CAMP_TOTAL_LEADS"), PH("CAMP_TOTAL_CPL")],
        ],
        col_widths=[4.50, 2.50, 2.50, 2.50],
        col_aligns=["left", "right", "right", "right"],
        total_row_index=3,
    )

    add_footer(slide, page_number=7)


def build_slide_8_eight_week(prs):
    """Slide 8 — 8-Week Context."""
    slide = add_blank_slide(prs)
    add_left_strip(slide)
    add_section_header(
        slide,
        icon_glyph="◑",
        title="8-Week Context",
        context_line=PH("EIGHT_WEEK_DATE_RANGE"),
        intro_line=PH("CONTEXT_INTRO_LINE"),
    )

    # Two chart placeholder areas
    # Left
    add_rectangle(slide, x=0.80, y=1.90, w=5.90, h=3.80, fill_color=OFF_WHITE, line_color=BORDER_LIGHT)
    add_text(
        slide, "[ CHART: Weekly spend (£) ]\nData: " + PH("WEEKLY_SPEND_DATA"),
        x=0.80, y=1.90, w=5.90, h=3.80,
        font_size=11, italic=True, color=MUTED, align="center", anchor="middle",
    )
    # Right
    add_rectangle(slide, x=6.90, y=1.90, w=5.90, h=3.80, fill_color=OFF_WHITE, line_color=BORDER_LIGHT)
    add_text(
        slide, "[ CHART: Weekly primary conversions ]\nData: " + PH("WEEKLY_CONVERSIONS_DATA"),
        x=6.90, y=1.90, w=5.90, h=3.80,
        font_size=11, italic=True, color=MUTED, align="center", anchor="middle",
    )

    # Key Finding (capped height so it doesn't overflow bottom)
    add_text(
        slide, "KEY FINDING",
        x=0.80, y=6.00, w=12.00, h=0.25,
        font_size=11, bold=True, color=BLUE, align="left",
    )
    add_text(
        slide, PH("EIGHT_WEEK_KEY_FINDING"),
        x=0.80, y=6.28, w=12.00, h=0.72,
        font_size=12, color=BODY, align="left",
    )

    add_footer(slide, page_number=8)


def build_slide_9_clarifications(prs):
    """Slide 9 — Client Clarifications & Actions."""
    slide = add_blank_slide(prs)
    add_left_strip(slide)
    add_section_header(
        slide,
        icon_glyph="◉",
        title="Client Clarifications & Actions",
        context_line=PH("CLIENT_REPLY_CONTEXT"),
        intro_line=None,
    )

    # Table with placeholder rows
    _add_table(
        slide,
        x=0.80, y=1.80, w=12.00, h=5.00,
        headers=["Giulio confirmed", "Action taken / scheduled"],
        rows=[
            [PH(f"CLAR_{i}_TEXT"), PH(f"CLAR_{i}_ACTION")] for i in range(1, 10)
        ],
        col_widths=[6.00, 6.00],
        col_aligns=["left", "left"],
    )

    add_footer(slide, page_number=9)


def build_slide_10_plan(prs):
    """Slide 10 — Week N+1 Plan."""
    slide = add_blank_slide(prs)
    add_left_strip(slide)
    add_section_header(
        slide,
        icon_glyph="▶",
        title=f"Week {PH('NEXT_WEEK_NUMBER')} Plan",
        context_line=PH("NEXT_WEEK_DATE_RANGE"),
        intro_line=PH("PLAN_INTRO_LINE"),
    )

    _add_table(
        slide,
        x=0.80, y=1.80, w=12.00, h=5.00,
        headers=["Day", "Focus", "Expected output"],
        rows=[
            [PH(f"PLAN_DAY_{i}"), PH(f"PLAN_FOCUS_{i}"), PH(f"PLAN_OUTPUT_{i}")] for i in range(1, 6)
        ],
        col_widths=[1.80, 3.20, 7.00],
        col_aligns=["left", "left", "left"],
    )

    add_footer(slide, page_number=10)


def build_slide_11_closer(prs):
    """Slide 11 — Closer (DARK)."""
    slide = add_blank_slide(prs)
    set_slide_background(slide, NAVY)

    # "Questions?"
    add_text(
        slide, "Questions?",
        x=0.50, y=2.50, w=12.30, h=0.90,
        font_size=44, bold=True, color=WHITE, align="center",
    )
    # Accent line
    add_rectangle(slide, x=(SLIDE_W_IN - 2.00) / 2, y=3.50, w=2.00, h=0.03, fill_color=BLUE)
    # Subtitle
    add_text(
        slide,
        f"Happy to walk through any section of this report or align on Week {PH('NEXT_WEEK_NUMBER')} priorities over a call.",
        x=0.50, y=3.70, w=12.30, h=0.50,
        font_size=16, italic=True, color=BLUE, align="center",
    )
    # Contact name
    add_text(
        slide, "Christopher Hoole  ·  Google Ads Specialist",
        x=0.50, y=5.10, w=12.30, h=0.35,
        font_size=14, bold=True, color=WHITE, align="center",
    )
    # Contact details
    add_text(
        slide,
        "christopherhoole.com  ·  chrishoole101@gmail.com  ·  +44 7451 252857",
        x=0.50, y=5.55, w=12.30, h=0.30,
        font_size=12, color=BLUE, align="center",
    )
    # Logo centre-bottom
    if ACT_LOGO.exists():
        pic = add_image(slide, ACT_LOGO, x=(SLIDE_W_IN - 0.70) / 2, y=6.30, h=0.70)
        pic.left = Inches((SLIDE_W_IN - Emu(pic.width).inches) / 2)


# -----------------------------------------------------------------------------
# TABLE HELPER
# -----------------------------------------------------------------------------

def _add_table(slide, x, y, w, h, headers, rows, col_widths, col_aligns, total_row_index=None):
    """Add a styled table per spec 1.8.3."""
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    tbl = tbl_shape.table

    # Column widths
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)

    # Header row
    tbl.rows[0].height = Inches(0.40)
    for ci, header_text in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}[col_aligns[ci]]
        run = p.add_run()
        run.text = header_text
        run.font.name = FONT
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = WHITE

    # Body rows
    for ri, row_data in enumerate(rows):
        tbl.rows[ri + 1].height = Inches(0.35)
        is_total = (total_row_index is not None and ri == total_row_index)
        fill_color = NAVY_TINT if is_total else (OFF_WHITE if ri % 2 == 0 else WHITE)
        text_color = NAVY if is_total else BODY
        bold = is_total

        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}[col_aligns[ci]]
            run = p.add_run()
            run.text = cell_text
            run.font.name = FONT
            run.font.size = Pt(11)
            run.font.bold = bold
            run.font.color.rgb = text_color


# -----------------------------------------------------------------------------
# MAIN BUILD
# -----------------------------------------------------------------------------

def build_template():
    prs = Presentation()
    _set_slide_size(prs)

    build_slide_1_title(prs)
    build_slide_2_exec_summary(prs)
    build_slide_3_work_delivered(prs)
    build_slide_4_brand(prs)
    build_slide_5_pmax(prs)
    build_slide_6_dii(prs)
    build_slide_7_numbers(prs)
    build_slide_8_eight_week(prs)
    build_slide_9_clarifications(prs)
    build_slide_10_plan(prs)
    build_slide_11_closer(prs)

    TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"[OK] Template built: {OUTPUT_PATH}")
    print(f"     {len(prs.slides)} slides")


if __name__ == "__main__":
    build_template()
