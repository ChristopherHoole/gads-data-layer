"""
Fix Slide 5 (Performance Max) in v6:
- Widen daily CPA table and add Friday + Weekly Total rows
- Add a narrative card on the right showing the bid-strategy-switch story
- Replace em dashes in bottom cards with commas/hyphens (user preference)
- Apply user's v5 formatting preferences throughout slide 5 (no em dashes)
"""
import sys, copy
from pathlib import Path

if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

FILE_PATH = Path(__file__).resolve().parents[3] / "potential_clients" / "Inserta Dental" / "End-of-week reports" / "DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v6.pptx"

NAVY = RGBColor(0x1A, 0x23, 0x7E)
BLUE = RGBColor(0x42, 0x85, 0xF4)
RED = RGBColor(0xEA, 0x43, 0x35)
GREEN = RGBColor(0x34, 0xA8, 0x53)
AMBER = RGBColor(0xFB, 0xBC, 0x05)
BODY = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY_TINT = RGBColor(0xE8, 0xEA, 0xF6)


def add_row_to_table(table, row_data, bold=False, fill=None, text_color=None):
    """Clone last row and populate with row_data."""
    tbl_elem = table._tbl
    last_tr = tbl_elem.findall(qn('a:tr'))[-1]
    new_tr = copy.deepcopy(last_tr)
    last_tr.addnext(new_tr)
    new_row_idx = len(table.rows) - 1
    for c_idx, val in enumerate(row_data):
        cell = table.cell(new_row_idx, c_idx)
        # Save existing formatting
        saved = None
        tf = cell.text_frame
        for p in tf.paragraphs:
            for r in p.runs:
                f = r.font
                saved = {
                    'name': f.name, 'size': f.size, 'bold': f.bold,
                    'align': p.alignment,
                }
                try:
                    saved['color'] = f.color.rgb if f.color and f.color.type else None
                except Exception:
                    saved['color'] = None
                break
            if saved: break
        tf.clear()
        p = tf.paragraphs[0]
        if saved and saved.get('align'): p.alignment = saved['align']
        r = p.add_run()
        r.text = val
        if saved:
            if saved['name']: r.font.name = saved['name']
            if saved['size']: r.font.size = saved['size']
            r.font.bold = bold if bold else saved['bold']
            if text_color:
                r.font.color.rgb = text_color
            elif saved['color']:
                r.font.color.rgb = saved['color']
        # Apply cell fill if requested
        if fill:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill


def replace_em_dashes_in_shape(shape):
    """Replace em dashes in a shape's text with commas (user's v5 preference)."""
    if not shape.has_text_frame:
        return False
    changed = False
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if "—" in r.text:
                r.text = r.text.replace(" — ", ", ").replace("—", ",")
                changed = True
    return changed


def add_narrative_card_right_side(slide):
    """Add a right-side narrative card next to the daily table telling the bid-strategy-switch story."""
    # Card background — rounded rectangle with red-ish tint (like the 'Bid Strategy Switch' card)
    # Actually let's make it a subtle white card with a NAVY left border (info card)
    # Position: x=7.00, y=1.10, w=5.80, h=3.00
    card_x = 7.00
    card_y = 1.10
    card_w = 5.80
    card_h = 3.00

    # Left border strip (navy)
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(card_x), Inches(card_y), Inches(0.08), Inches(card_h),
    )
    border.fill.solid()
    border.fill.fore_color.rgb = NAVY
    border.line.fill.background()

    # Card body — white with subtle shadow
    body_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(card_x + 0.08), Inches(card_y), Inches(card_w - 0.08), Inches(card_h),
    )
    body_rect.fill.solid()
    body_rect.fill.fore_color.rgb = WHITE
    body_rect.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    # Heading
    heading = slide.shapes.add_textbox(
        Inches(card_x + 0.25), Inches(card_y + 0.10), Inches(card_w - 0.35), Inches(0.35)
    )
    heading.text_frame.word_wrap = True
    p = heading.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Week 1 Story"
    r.font.name = "Calibri"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = NAVY

    # Body copy
    body_text = slide.shapes.add_textbox(
        Inches(card_x + 0.25), Inches(card_y + 0.55), Inches(card_w - 0.35), Inches(card_h - 0.65)
    )
    tf = body_text.text_frame
    tf.word_wrap = True

    lines = [
        ("Bid strategy switched Wed 15 April, at the midpoint of the week. Each day tells a different part of the story:", False),
        ("", False),
        ("Mon 13 & Tue 14:", True),
        ("Operating under the old tROAS strategy with characteristic swings, £119 and £56 CPA.", False),
        ("", False),
        ("Wed 15:", True),
        ("First partial day under new Target CPA £60, delivered £25 CPA, the best single day in the account.", False),
        ("", False),
        ("Thu 16 & Fri 17:", True),
        ("Fully post-switch, settling toward the £60 target at £76 and £79 CPA.", False),
        ("", False),
        ("Week-total CPA: £73.13 vs prior week £85.86, a 15% improvement despite policy restrictions.", False),
    ]
    tf.clear()
    for i, (line, is_bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.name = "Calibri"
        r.font.size = Pt(11)
        r.font.bold = is_bold
        if is_bold:
            r.font.color.rgb = BLUE
        else:
            r.font.color.rgb = BODY


def main():
    prs = Presentation(str(FILE_PATH))
    slide = list(prs.slides)[4]  # slide 5

    # -------------------------------------------------------------------
    # 1. Find and widen the daily CPA table
    # -------------------------------------------------------------------
    table_shape = None
    for shape in slide.shapes:
        if shape.has_table and shape.table.cell(0, 0).text.strip() == "Day":
            table_shape = shape
            break

    if table_shape is None:
        print("ERROR: Could not find daily CPA table")
        return

    # Widen the table and move slightly: keep x=0.69, change width from 3.50 to 6.00
    table_shape.width = Inches(6.00)
    print(" [OK] Widened daily CPA table to 6.00\"")

    # Add Friday row
    add_row_to_table(
        table_shape.table,
        ['Fri 17 Apr', '£1,188', '15', '£79'],
    )
    print(" [OK] Added Fri 17 Apr row")

    # Add Weekly Total row (bold, with navy tint background)
    add_row_to_table(
        table_shape.table,
        ['Week total', '£5,960', '81.5', '£73.13'],
        bold=True,
        fill=NAVY_TINT,
        text_color=NAVY,
    )
    print(" [OK] Added Week total row")

    # -------------------------------------------------------------------
    # 2. Add narrative card on the right side
    # -------------------------------------------------------------------
    add_narrative_card_right_side(slide)
    print(" [OK] Added Week 1 Story narrative card")

    # -------------------------------------------------------------------
    # 3. Replace em dashes in all shapes on slide 5 (user's v5 preference)
    # -------------------------------------------------------------------
    n_shapes_fixed = 0
    for shape in slide.shapes:
        if replace_em_dashes_in_shape(shape):
            n_shapes_fixed += 1
    print(f" [OK] Replaced em dashes in {n_shapes_fixed} shape(s)")

    # Save
    prs.save(str(FILE_PATH))
    print(f"\nSaved v6: {FILE_PATH.name}")


if __name__ == "__main__":
    main()
