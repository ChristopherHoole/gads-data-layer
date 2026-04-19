"""Minimal v7 slide 5 fix: add Friday row only."""
import sys, copy
from pathlib import Path
if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.oxml.ns import qn

F = Path(__file__).resolve().parents[3] / "potential_clients" / "Inserta Dental" / "End-of-week reports" / "DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v7.pptx"

prs = Presentation(str(F))
table = list(prs.slides)[4].shapes[9].table

# Clone last row, insert after
tbl = table._tbl
last = tbl.findall(qn('a:tr'))[-1]
new = copy.deepcopy(last)
last.addnext(new)

# Populate: Fri 17 Apr / £1,188 / 15 / £79
data = ['Fri 17 Apr', '£1,188', '15', '£79']
ri = len(table.rows) - 1
for ci, val in enumerate(data):
    cell = table.cell(ri, ci)
    # preserve formatting from existing run
    saved = None
    for p in cell.text_frame.paragraphs:
        for r in p.runs:
            saved = {'name': r.font.name, 'size': r.font.size, 'bold': r.font.bold, 'align': p.alignment}
            try: saved['color'] = r.font.color.rgb if r.font.color.type else None
            except: saved['color'] = None
            break
        if saved: break
    cell.text_frame.clear()
    p = cell.text_frame.paragraphs[0]
    if saved and saved['align']: p.alignment = saved['align']
    r = p.add_run()
    r.text = val
    if saved:
        if saved['name']: r.font.name = saved['name']
        if saved['size']: r.font.size = saved['size']
        if saved['bold']: r.font.bold = saved['bold']
        if saved['color']: r.font.color.rgb = saved['color']

prs.save(str(F))
print("Done")
