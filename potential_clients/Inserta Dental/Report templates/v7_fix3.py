"""Give cards more vertical room so 'Flagged for Week 2' text doesn't overflow."""
import sys
from pathlib import Path
if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches

F = Path(__file__).resolve().parents[3] / "potential_clients" / "Inserta Dental" / "End-of-week reports" / "DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v7.pptx"
prs = Presentation(str(F))
s = list(prs.slides)[4]

# Tighten table (compact rows) to give cards more room below
for shape in s.shapes:
    if shape.has_table:
        shape.height = Inches(2.40)  # was 2.85
        for row in shape.table.rows:
            row.height = Inches(0.40)  # was 0.475
        continue

    if not shape.has_text_frame:
        # Card BGs at y=4.90 h=2.05 → move to y=4.15 h=2.80 (more room)
        t = shape.top/914400 if shape.top else 0
        h = shape.height/914400 if shape.height else 0
        if abs(t - 4.90) < 0.05 and abs(h - 2.05) < 0.05:
            shape.top = Inches(4.15)
            shape.height = Inches(2.80)
        continue

    txt = shape.text_frame.text.strip()
    if txt == "What Changed & What's Next":
        shape.top = Inches(3.95)
    elif txt in ("Bid Strategy Switch", "Flagged for Week 2", "Recommendation"):
        shape.top = Inches(4.25)
    elif txt.startswith("From: Max Conversion Value") or \
         txt.startswith("Google still reports") or \
         txt.startswith("Hold the new bid strategy"):
        shape.top = Inches(4.60)
        shape.height = Inches(2.30)

prs.save(str(F))
print("Done")
