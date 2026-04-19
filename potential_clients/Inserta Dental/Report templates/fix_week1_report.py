"""
Apply fixes to Week 1 v1 report → produces v3.

Targets the critical visual defects identified in the slide-by-slide review:
- Slide 1: Stat labels overlapping numbers (push labels down below numbers)
- Slide 4: "12.0" → "12" (remove trailing decimal)
- Slide 5: 4th stat: "48,349 Impressions" → "£3.22 Avg CPC" (consistency with Slide 4)
- Slide 6: Remove misleading before/after bars (different scales) + misleading "0 → 1,050" row
- Slide 7: Replace broken funnel chevrons with proper rounded rectangles
- Slide 7: Add TOTAL row to campaign table
- Slide 8: Fix Key Finding text cut-off (expand height, move up)
"""

import copy
import sys
from pathlib import Path

# Force UTF-8 output on Windows
if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from lxml import etree

# -----------------------------------------------------------------------------
# Setup
# -----------------------------------------------------------------------------

REPO_ROOT = Path(__file__).resolve().parents[3]
REPORTS_DIR = REPO_ROOT / "potential_clients" / "Inserta Dental" / "End-of-week reports"
SOURCE = REPORTS_DIR / "DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v1.pptx"
OUTPUT = REPORTS_DIR / "DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v3.pptx"

# Palette (matches spec)
NAVY = RGBColor(0x1A, 0x23, 0x7E)
BLUE = RGBColor(0x42, 0x85, 0xF4)
GREEN = RGBColor(0x34, 0xA8, 0x53)
RED = RGBColor(0xEA, 0x43, 0x35)
YELLOW = RGBColor(0xFB, 0xBC, 0x05)
BODY = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF7, 0xF7, 0xF7)
NAVY_TINT = RGBColor(0xE8, 0xEA, 0xF6)

FONT = "Calibri"


# -----------------------------------------------------------------------------
# Helpers
# -----------------------------------------------------------------------------

def set_textbox_text(shape, text, font_size=None, bold=None, color=None, font_name=None, align=None):
    """Rewrite a textbox's text while preserving other properties where possible."""
    tf = shape.text_frame
    # Clear existing paragraphs (keep one)
    tf.clear()
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = text
    if font_name:
        r.font.name = font_name
    if font_size is not None:
        r.font.size = Pt(font_size)
    if bold is not None:
        r.font.bold = bold
    if color is not None:
        r.font.color.rgb = color


def move_shape(shape, x=None, y=None, w=None, h=None):
    if x is not None:
        shape.left = Inches(x)
    if y is not None:
        shape.top = Inches(y)
    if w is not None:
        shape.width = Inches(w)
    if h is not None:
        shape.height = Inches(h)


def delete_shape(shape):
    """Remove a shape from its parent slide."""
    sp = shape._element
    sp.getparent().remove(sp)


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def add_rounded_rect(slide, x, y, w, h, fill_color, corner_radius=0.08):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.adjustments[0] = corner_radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_run(shape, text, font_size, bold=False, color=WHITE, align="center", anchor="middle"):
    """Set a shape's text (shape must already exist). Works for rounded rects or text boxes."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    tf.clear()
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        if isinstance(font_size, list):
            r.font.size = Pt(font_size[i] if i < len(font_size) else font_size[-1])
        else:
            r.font.size = Pt(font_size)
        r.font.bold = bold if not isinstance(bold, list) else bold[i] if i < len(bold) else bold[-1]
        r.font.color.rgb = color


def add_textbox(slide, text, x, y, w, h, font_size, bold=False, italic=False, color=BODY, align="left", anchor="top"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


# -----------------------------------------------------------------------------
# Fixes
# -----------------------------------------------------------------------------

def fix_slide_1(slide):
    """Fix label-overlap-with-number issue on title slide.

    Root cause: numbers are 72pt in 0.80"-tall boxes. 72pt visual height ≈ 1.0",
    so numbers overflow down into the labels at y=1.35 or even y=1.45.

    Fix: reduce number font to 54pt (fits cleanly in 0.80" box, more proportional
    on a title slide with 4 stats). Labels stay at y=1.45 with proper clearance.
    """
    shapes = list(slide.shapes)

    # Resize the number fonts from 72pt to 54pt on all 4 stat numbers
    # Number textboxes: indices 0, 2, 6, 8
    for num_idx in [0, 2, 6, 8]:
        tb = shapes[num_idx]
        for p in tb.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(54)

    # Move labels just below numbers with 0.05" gap
    # Top labels
    move_shape(shapes[1], y=1.45)
    move_shape(shapes[3], y=1.45)
    # Bottom labels
    move_shape(shapes[7], y=5.45)
    move_shape(shapes[9], y=5.45)


def fix_slide_4(slide):
    """Change '12.0' to '12' (remove trailing decimal)."""
    shapes = list(slide.shapes)
    # TextBox 6 (index 5) = "12.0"
    tb = shapes[5]
    # Preserve formatting by replacing text in existing run
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            if r.text == "12.0":
                r.text = "12"
                return
    # Fallback
    set_textbox_text(tb, "12", font_size=40, bold=True, color=GREEN, font_name=FONT, align="left")


def fix_slide_5(slide):
    """Change 4th stat from '48,349 Impressions' to '£3.22 Avg CPC'."""
    shapes = list(slide.shapes)
    # TextBox 11 (index 10) = "48,349"
    # TextBox 12 (index 11) = "Impressions"
    for p in shapes[10].text_frame.paragraphs:
        for r in p.runs:
            if r.text.strip() == "48,349":
                r.text = "£3.22"
    for p in shapes[11].text_frame.paragraphs:
        for r in p.runs:
            if r.text.strip() == "Impressions":
                r.text = "Avg CPC"


def fix_slide_6(slide):
    """Remove misleading before/after bars.

    The current layout has 4 proportional bars with different scales (one for each row),
    which creates a misleading visual comparison (27 → 675 "look" the same size as 4 → 25).

    Fix: delete all bar rectangles and bar-related labels. Keep just number arrows.
    Replace the bar area with a clean "X → Y" arrow-based comparison.
    """
    shapes = list(slide.shapes)

    # Shapes to delete (the before/after visualization mess + original heading)
    # Index 11 = "Structural rebuild — before / after" heading (original)
    # Indices 12 through 34 inclusive = the bars and their labels
    to_delete = []
    for i in range(11, 35):
        to_delete.append(shapes[i])

    for shape in to_delete:
        delete_shape(shape)

    # Now add a clean structural comparison (no bars — just numbers with arrows)
    # Position it where the old bars were (y=3.55 to 6.30)
    # Use 3 comparison rows only (drop "Headlines written: 0 → 1,050" which was misleading)

    # Section heading
    add_textbox(
        slide, "Structural rebuild — before → after",
        x=0.50, y=3.55, w=7.80, h=0.30,
        font_size=14, bold=True, color=NAVY, align="left",
    )

    # 3 comparison rows (dropped the misleading Headlines row)
    comparisons = [
        ("Ad groups", "4", "25"),
        ("Keywords", "27", "~675"),
        ("RSAs", "6", "75"),
    ]
    row_y = 4.10
    row_h = 0.75
    for i, (label, before_val, after_val) in enumerate(comparisons):
        y = row_y + i * row_h
        # Label
        add_textbox(
            slide, label,
            x=0.50, y=y, w=2.00, h=row_h,
            font_size=14, color=BODY, align="left", anchor="middle",
        )
        # Before value (grey, smaller)
        add_textbox(
            slide, before_val,
            x=2.70, y=y, w=1.20, h=row_h,
            font_size=24, bold=True, color=MUTED, align="center", anchor="middle",
        )
        # Arrow
        add_textbox(
            slide, "→",
            x=4.00, y=y, w=0.70, h=row_h,
            font_size=24, color=BLUE, align="center", anchor="middle",
        )
        # After value (navy, larger)
        add_textbox(
            slide, after_val,
            x=4.80, y=y, w=2.50, h=row_h,
            font_size=30, bold=True, color=NAVY, align="left", anchor="middle",
        )


def fix_slide_7(slide):
    """Replace broken funnel chevrons with proper rounded rectangles.

    The chevrons had narrow internal areas that caused text to wrap mid-word
    into unreadable fragments. Solution: delete chevrons, add rounded rectangles
    with proper dimensions. Also add a TOTAL row to the campaign table.
    """
    shapes = list(slide.shapes)

    # Delete chevrons (indices 3-7)
    chevrons_to_delete = [shapes[i] for i in range(3, 8)]
    for s in chevrons_to_delete:
        delete_shape(s)

    # Add proper funnel stages as rounded rectangles
    # Positions: 5 stages across width 12.30, starting x=0.50, y=1.90
    stages = [
        ("£8,462", "Spend"),
        ("2,339", "Clicks"),
        ("87", "Dengro Leads"),
        ("[awaiting audit]", "Bookings"),
        ("7", "Purchases"),
    ]
    # Colour gradient navy → blue
    stage_colors = [
        RGBColor(0x1A, 0x23, 0x7E),  # navy
        RGBColor(0x26, 0x3A, 0x9E),
        RGBColor(0x2F, 0x5C, 0xC0),
        RGBColor(0x36, 0x73, 0xDC),
        RGBColor(0x42, 0x85, 0xF4),  # blue
    ]
    # Use wide rectangles, no funnel narrowing — fix readability first
    stage_w = 2.30
    stage_h = 1.40
    gap = 0.15
    total_w = 5 * stage_w + 4 * gap  # = 12.10
    start_x = (13.33 - total_w) / 2  # centre

    for i, (number, label) in enumerate(stages):
        x = start_x + i * (stage_w + gap)
        rect = add_rounded_rect(slide, x=x, y=1.90, w=stage_w, h=stage_h, fill_color=stage_colors[i], corner_radius=0.08)
        # Add text inside the rectangle
        tf = rect.text_frame
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.12)
        tf.margin_bottom = Inches(0.12)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.clear()

        # Number paragraph (big)
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = number
        r1.font.name = FONT
        r1.font.size = Pt(22) if len(number) > 10 else Pt(26)
        r1.font.bold = True
        r1.font.color.rgb = WHITE

        # Label paragraph
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = label
        r2.font.name = FONT
        r2.font.size = Pt(12)
        r2.font.bold = False
        r2.font.color.rgb = WHITE

        # Arrow between stages (except after last)
        if i < 4:
            arrow_x = x + stage_w + gap * 0.1
            add_textbox(
                slide, "▶",
                x=arrow_x, y=2.30, w=gap * 0.8, h=0.60,
                font_size=14, color=BLUE, align="center", anchor="middle",
            )

    # Add TOTAL row to campaign table
    # Table is at index 10 in original (but we deleted 5 shapes above, so need fresh lookup)
    # Safer: find the table by iterating
    table_shape = None
    for shape in slide.shapes:
        if shape.has_table:
            table_shape = shape
            break
    if table_shape is None:
        return
    tbl = table_shape.table
    # Current table has 4 rows (header + 3 campaigns). Need to add a 5th row.
    # python-pptx doesn't directly support adding rows. Workaround: XML manipulation.
    # We'll add a new row by cloning the last row's XML
    tbl_elem = tbl._tbl
    last_tr = tbl_elem.findall(qn('a:tr'))[-1]
    new_tr = copy.deepcopy(last_tr)
    # Insert AFTER last row
    last_tr.addnext(new_tr)

    # Get the new row as a TableRow object
    new_row_idx = len(tbl.rows) - 1
    # Populate the new row with TOTAL values
    total_data = ["TOTAL", "£8,462", "87", "£97.27"]
    for ci, val in enumerate(total_data):
        cell = tbl.cell(new_row_idx, ci)
        # Set fill to light navy tint
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_TINT
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        # Match alignment of existing rows
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = val
        r.font.name = FONT
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = NAVY


def fix_slide_3(slide):
    """Fix context line overlap with day labels on Work Delivered slide.

    Context line was at y=1.05, h=0.35 (ends y=1.40).
    Day labels at y=1.15, h=0.40 (ends y=1.55).
    → Overlap from y=1.15 to y=1.40.

    Fix: Move context line up into the title's empty area.
    Title "Work Delivered" is 32pt bold in a 0.80"-tall box (y=0.30-1.10).
    32pt text renders y=0.30 to ~y=0.74, leaving y=0.74-1.10 empty.
    Move context line to y=0.90 with h=0.20 → ends at y=1.10, no overlap.
    """
    shapes = list(slide.shapes)
    # Shape index 1 = context line
    move_shape(shapes[1], y=0.90, h=0.20)


def fix_slide_8(slide):
    """Fix Key Finding text cut-off at bottom of slide.

    Root cause: charts extend to y=6.60 (y=2.10 + h=4.50), leaving only 0.50"
    before footer at y=7.10. Not enough room for 2-line Key Finding paragraph.

    Fix: shrink charts to end at y=5.80 (h=3.70), reposition Key Finding below.
    """
    shapes = list(slide.shapes)
    # Shape 3, 4 = two charts at y=2.10, h=4.50
    # Shape 5 = "Key Finding" label at y=6.55
    # Shape 6 = body at y=6.85, h=0.25

    # Shrink charts
    move_shape(shapes[3], h=3.70)  # left chart (was h=4.50)
    move_shape(shapes[4], h=3.70)  # right chart (was h=4.50)

    # Reposition Key Finding label + body
    move_shape(shapes[5], y=5.95)  # label
    move_shape(shapes[6], y=6.25, h=0.80)  # body — proper height for 2 lines


# -----------------------------------------------------------------------------
# Main
# -----------------------------------------------------------------------------

def main():
    prs = Presentation(str(SOURCE))
    slides = list(prs.slides)

    print(f"Loaded: {SOURCE.name}")
    print(f"Applying fixes...")

    fix_slide_1(slides[0]); print(" [OK] Slide 1 — label/number overlap fixed")
    fix_slide_3(slides[2]); print(" [OK] Slide 3 — context line / day label overlap fixed")
    fix_slide_4(slides[3]); print(" [OK] Slide 4 — '12.0' → '12'")
    fix_slide_5(slides[4]); print(" [OK] Slide 5 — 4th stat: Impressions → Avg CPC")
    fix_slide_6(slides[5]); print(" [OK] Slide 6 — before/after bars replaced with clean arrows")
    fix_slide_7(slides[6]); print(" [OK] Slide 7 — funnel rebuilt + TOTAL row added to table")
    fix_slide_8(slides[7]); print(" [OK] Slide 8 — Key Finding text cut-off fixed")

    prs.save(str(OUTPUT))
    print(f"\nSaved: {OUTPUT.name}")


if __name__ == "__main__":
    main()
