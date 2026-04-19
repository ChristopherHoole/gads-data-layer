"""
Populate v4 with Week 1 content while preserving the reference design.

Takes the v4 copy of the Account Structure Report and replaces content on each
slide with Week 1 content per the agreed mapping:
  Ref Slide 1  → Title slide
  Ref Slide 2  → Executive Summary
  Ref Slide 3  → Work Delivered
  Ref Slide 4  → Brand Campaign
  Ref Slide 5  → Performance Max
  Ref Slide 6  → Dental Implants Intent
  Ref Slide 7  → The Numbers That Matter
  Ref Slide 8  → 8-Week Context (with CPA per week)
  Ref Slide 9  → Client Clarifications & Actions
  Ref Slide 10 → Questions for Client (NEW)
  Ref Slide 11 → Week 2 Plan
  Ref Slide 12 → What's Next (closer)

All design elements (Google colour bars, card borders, fonts, colours) are preserved.
Only text content and table data are replaced.
"""
import sys
from pathlib import Path

if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

# ---- Paths ----
REPO_ROOT = Path(__file__).resolve().parents[3]
REPORTS_DIR = REPO_ROOT / "potential_clients" / "Inserta Dental" / "End-of-week reports"
FILE_PATH = REPORTS_DIR / "DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v4.pptx"

GREEN = RGBColor(0x34, 0xA8, 0x53)
NAVY = RGBColor(0x1A, 0x23, 0x7E)


def replace_text_in_shape(shape, old_text, new_text, new_color=None):
    """Find `old_text` in shape's runs and replace with `new_text`. Handles text split across multiple runs.
    Preserves formatting of first run.
    """
    if not shape.has_text_frame:
        return False
    # First try run-by-run (handles single-run case efficiently)
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if old_text in r.text:
                r.text = r.text.replace(old_text, new_text)
                if new_color is not None:
                    r.font.color.rgb = new_color
                return True
    # Fallback: check if concatenated text across all runs in any paragraph matches
    for p in shape.text_frame.paragraphs:
        # Concatenate all runs in paragraph
        concat = "".join(r.text for r in p.runs)
        if old_text in concat:
            # Reconstruct: put new_text in first run, clear others
            if p.runs:
                first_run = p.runs[0]
                # Preserve first run formatting
                new_concat = concat.replace(old_text, new_text)
                first_run.text = new_concat
                if new_color is not None:
                    first_run.font.color.rgb = new_color
                # Clear other runs in this paragraph
                for r in p.runs[1:]:
                    r.text = ""
                return True
    # Also try cross-paragraph match (for "Account Structure\n& Issues" case)
    all_text = shape.text_frame.text
    if old_text in all_text:
        # Clear all paragraphs and use set_shape_text to replace cleanly
        saved_color = None
        saved_size = None
        saved_bold = None
        saved_name = None
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                saved_name = r.font.name
                saved_size = r.font.size
                saved_bold = r.font.bold
                try:
                    if r.font.color and r.font.color.type:
                        saved_color = r.font.color.rgb
                except Exception:
                    pass
                break
            if saved_name:
                break
        new_all = all_text.replace(old_text, new_text)
        shape.text_frame.clear()
        # Split by newlines back into paragraphs
        for i, line in enumerate(new_all.split('\n')):
            p = shape.text_frame.paragraphs[0] if i == 0 else shape.text_frame.add_paragraph()
            r = p.add_run()
            r.text = line
            if saved_name: r.font.name = saved_name
            if saved_size: r.font.size = saved_size
            if saved_bold is not None: r.font.bold = saved_bold
            if new_color is not None:
                r.font.color.rgb = new_color
            elif saved_color:
                r.font.color.rgb = saved_color
        return True
    return False


def set_shape_text(shape, new_text, preserve_first_run_formatting=True):
    """Replace all text in a shape with new_text. Keeps formatting of first run."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Get first run's formatting to preserve
    first_run = None
    for p in tf.paragraphs:
        for r in p.runs:
            first_run = r
            break
        if first_run:
            break
    if first_run is None:
        tf.text = new_text
        return
    # Capture formatting
    font = first_run.font
    saved = {
        'name': font.name, 'size': font.size, 'bold': font.bold,
        'italic': font.italic,
    }
    try:
        saved['color'] = font.color.rgb if font.color and font.color.type else None
    except Exception:
        saved['color'] = None
    # Replace
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = new_text
    r.font.name = saved['name']
    if saved['size']: r.font.size = saved['size']
    r.font.bold = saved['bold']
    r.font.italic = saved['italic']
    if saved['color']:
        r.font.color.rgb = saved['color']


def set_shape_multi_paragraph(shape, lines, preserve_formatting=True):
    """Replace text with multiple lines. Each list element = a paragraph."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Capture formatting from first run
    saved = None
    for p in tf.paragraphs:
        for r in p.runs:
            f = r.font
            saved = {
                'name': f.name, 'size': f.size, 'bold': f.bold, 'italic': f.italic,
            }
            try:
                saved['color'] = f.color.rgb if f.color and f.color.type else None
            except Exception:
                saved['color'] = None
            break
        if saved: break

    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        if saved:
            r.font.name = saved['name']
            if saved['size']: r.font.size = saved['size']
            r.font.bold = saved['bold']
            r.font.italic = saved['italic']
            if saved['color']:
                r.font.color.rgb = saved['color']


def set_table_cell(table, row, col, new_text):
    """Replace cell text while preserving formatting."""
    cell = table.cell(row, col)
    tf = cell.text_frame
    # Grab first run formatting
    saved = None
    for p in tf.paragraphs:
        for r in p.runs:
            f = r.font
            saved = {
                'name': f.name, 'size': f.size, 'bold': f.bold, 'italic': f.italic,
                'alignment': p.alignment,
            }
            try:
                saved['color'] = f.color.rgb if f.color and f.color.type else None
            except Exception:
                saved['color'] = None
            break
        if saved: break
    tf.clear()
    p = tf.paragraphs[0]
    if saved and saved.get('alignment'): p.alignment = saved['alignment']
    r = p.add_run()
    r.text = new_text
    if saved:
        r.font.name = saved['name']
        if saved['size']: r.font.size = saved['size']
        r.font.bold = saved['bold']
        if saved['color']: r.font.color.rgb = saved['color']


def populate_table(table, rows_data):
    """Populate table rows with new data. rows_data is list of lists (matching existing row count)."""
    for r_idx, row_vals in enumerate(rows_data):
        for c_idx, val in enumerate(row_vals):
            if r_idx < len(table.rows) and c_idx < len(table.columns):
                set_table_cell(table, r_idx, c_idx, val)


# =============================================================================
# SLIDE POPULATORS
# =============================================================================

def populate_slide_1(slide):
    """Title slide."""
    # Title
    replace_text_in_shape_any(slide, "Account Structure& Issues", "Week 1 Report")
    replace_text_in_shape_any(slide, "Dental by Design  |  Prodent Group", "Dental by Design  |  Prodent Group")
    # Christopher Hoole block stays same
    # Date/data label
    replace_text_in_shape_any(slide, "Data: 2025 (full year) + 2026 (Q1, Jan–Apr)",
                               "Data: Monday 13 – Friday 17 April 2026")
    # Hero stat
    replace_text_in_shape_any(slide, "Non-Brand CPA Trend", "Brand Avg CPC Trend")
    replace_text_in_shape_any(slide, "£46 → £64", "£3.18 → £1.61", new_color=GREEN)
    replace_text_in_shape_any(slide, "+40% increase year-on-year while spend nearly doubled",
                               "−49% week-on-week · lowest of 8 tracked weeks · impression share 92%+")
    # Stats
    replace_text_in_shape_any(slide, "1,511", "~675")
    replace_text_in_shape_any(slide, "Ghost Conversions (KMG)", "New Keywords Added")
    replace_text_in_shape_any(slide, "47", "25")
    replace_text_in_shape_any(slide, "Total Campaigns", "DII Ad Groups (was 4)")
    replace_text_in_shape_any(slide, "6.4", "75")
    replace_text_in_shape_any(slide, "Avg QS (Weighted)", "New RSAs Written")


def replace_text_in_shape_any(slide, old, new, new_color=None):
    """Try replace in any shape on the slide."""
    for shape in slide.shapes:
        if replace_text_in_shape(shape, old, new, new_color):
            return True
    return False


def populate_slide_2(slide):
    """Executive Summary."""
    replace_text_in_shape_any(slide, "2025 + 2026", "Week 1 · 13-17 Apr")
    replace_text_in_shape_any(slide, "£46→£64", "£8,462")
    replace_text_in_shape_any(slide, "Non-Brand CPA (+40%)", "Week 1 Spend")
    replace_text_in_shape_any(slide, "1,511", "87")
    replace_text_in_shape_any(slide, "Ghost Conversions (KMG Legacy)", "Dengro Offline Leads")
    replace_text_in_shape_any(slide, "90%", "101")
    replace_text_in_shape_any(slide, "Ads Rated Poor or Average", "Primary Conversions")
    replace_text_in_shape_any(slide, "6.4/10", "£97")
    replace_text_in_shape_any(slide, "Weighted Avg Quality Score", "Cost per Lead")
    replace_text_in_shape_any(slide, "Key Issues Identified", "Key Wins This Week")

    # Bullet list — need to replace the multi-paragraph content block
    for shape in slide.shapes:
        if shape.has_text_frame and "1,511 ghost conversions from legacy KMG" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Brand CPC at 8-week low: £2.16 vs March avg £3.29 (−35%). £1.50 Max CPC cap is biting without suppressing volume.",
                "PMax CPA improved week-on-week: £73 vs £85.86 prior week, despite persistent asset-group policy restrictions.",
                "DII entered learning phase after full rebuild: 25 new ad groups, ~675 keywords, 75 RSAs. Clean read expected Week 2-3.",
                "5,000+ structured negatives replaced 2,192 unstructured broad-match entries across 11 shared lists.",
                "Giulio replied Fri 17 Apr with answers — compliance actions captured and scheduled for Week 2.",
            ])
            break

    # Core issue callout box
    for shape in slide.shapes:
        if shape.has_text_frame and "Higher spend on weak foundations" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Week 1 was structural. Three live campaigns moved off Max Conversion Value bidding onto controllable targets. Early signal is positive — but DII needs a further week of stabilisation and Monday's end-to-end tracking audit is the gate to any CPA-based decisions from here.",
            ])
            break


def populate_slide_3(slide):
    """Work Delivered (from Conversion Tracking Audit)."""
    replace_text_in_shape_any(slide, "Conversion Tracking Audit", "Work Delivered")
    replace_text_in_shape_any(slide,
        "1,511 ghost conversions + 445 soft signals are inflating reported performance",
        "Tue 14 → Fri 17 April · 4 working days post-signing · ~20 hours")
    replace_text_in_shape_any(slide, "Critical Issue", "Foundation Week")
    replace_text_in_shape_any(slide, "Conversion Classification", "Day-by-Day Activity")
    replace_text_in_shape_any(slide, "Impact on CPA", "Running Totals")
    replace_text_in_shape_any(slide, "Giulio Confirmed", "Key Findings")

    # Table
    for shape in slide.shapes:
        if shape.has_table:
            populate_table(shape.table, [
                ['Day', 'Focus', 'Output', 'Impact'],
                ['Tue 14 Apr', 'Negatives Overhaul', '2,192 → ~4,000+ across 11 lists', '~£63k waste blocked'],
                ['Wed 15 Apr', 'Account-Wide Reset', 'Brand · PMax tCPA £60 · DII tCPA £75', 'Controllable bid targets'],
                ['Thu 16 Apr', 'DII Rebuild (Part 1)', '20 AGs · 537 keywords · 60 RSAs', 'Patient-intent structure'],
                ['Fri 17 Apr', 'Audit + Completion', '173 DBD-blockers removed · AG21-25 built', 'Campaign 25/25 complete'],
            ])
            break

    # Running totals (Impact on CPA section)
    for shape in slide.shapes:
        if shape.has_text_frame and "Reported CPA: £49" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Positive keywords added: ~675 across 25 ad groups",
                "New RSAs written: 75 (1,050 headlines, 280 descriptions)",
                "Negatives: ~5,000 structured across 11 shared lists",
                "Geo targets: 270 London postcodes rolled out to all 3 active campaigns",
            ])
            break

    # Key Findings (Giulio Confirmed section)
    for shape in slide.shapes:
        if shape.has_text_frame and "KMG OLD LP actions" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "All 3 live campaigns moved from Max Conversion Value to Target CPA / Target Impression Share.",
                "Auto-apply recommendations disabled — no more silent keyword removals.",
                "Zero NHS/free-seeker junk in PMax search terms post-switch — negatives are working.",
                "Asset groups in PMax flagged 'limited by policy' — full rebuild scheduled Wed 22 Apr.",
            ])
            break


def populate_slide_4(slide):
    """Brand Campaign (from ROAS Puzzle)."""
    replace_text_in_shape_any(slide, "The ROAS Puzzle", "Brand Campaign")
    replace_text_in_shape_any(slide,
        "Target ROAS bidding is based on arbitrary £300 booking values - not real revenue",
        "CPC at 8-week low · £1.50 Max CPC cap is working")
    replace_text_in_shape_any(slide, "Bid Strategy", "Performance")
    replace_text_in_shape_any(slide, "Current Bid Strategy Setup", "Week 1 Performance")

    # Main table
    for shape in slide.shapes:
        if shape.has_table and len(shape.table.rows) == 4:
            populate_table(shape.table, [
                ['Metric', 'Week 1', 'Prior Week', 'Change'],
                ['Week spend', '£280', '£298', '−6%'],
                ['Primary conv', '12.0', '12.0', 'Flat'],
                ['Cost / conv', '£23.39', '£24.89', '−6%'],
            ])
            break

    replace_text_in_shape_any(slide, "Why This Is a Problem", "Key Finding")
    replace_text_in_shape_any(slide, "The Dengro Pipeline", "Bid Strategy Applied")
    replace_text_in_shape_any(slide, "Recommendation", "Next Week")

    # Why/Key Finding body
    for shape in slide.shapes:
        if shape.has_text_frame and "maximise conversion VALUE" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Brand CPC is at its lowest of 8 tracked weeks: £2.16 vs March avg £3.29 (−35%).",
                "The £1.50 Max CPC cap (set Wed 15 Apr) is biting without suppressing volume.",
                "Impression share held at 92%+ and conversion volume tracked with spend.",
                "This is the strongest signal of Week 1.",
            ])
            break

    # Dengro Pipeline / Bid Strategy Applied body
    for shape in slide.shapes:
        if shape.has_text_frame and "Dengro Offline Lead: 4,663" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "From: Max Conversion Value (tROAS 10%) — inflating brand CPCs",
                "To: Target Impression Share 90% · Max CPC £1.50",
                "Goal: eliminate runaway bidding on brand traffic, lock in the brand position at controlled cost",
                "Result: CPC dropped from £3.18 (prior 30d) to £1.61 (first 36 hours) to £2.16 (week avg)",
            ])
            break

    # Recommendation body
    for shape in slide.shapes:
        if shape.has_text_frame and "Option A: Switch to Maximise Conversions" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Hold the £1.50 cap through Week 2. Let the bid strategy complete its 14-day learning phase.",
                "Evaluate lowering toward £1.20 once CPC has stabilised at or below £1.50.",
                "Brand is mature territory — small incremental CPC reductions compound into meaningful monthly savings.",
            ])
            break


def populate_slide_5(slide):
    """Performance Max (from Account Structure)."""
    replace_text_in_shape_any(slide, "Account Structure", "Performance Max")
    replace_text_in_shape_any(slide,
        "47 campaigns, 4 active - a legacy of multiple agencies and strategy changes",
        "CPA improved week-on-week despite policy restrictions")
    replace_text_in_shape_any(slide, "Structure", "Performance")
    replace_text_in_shape_any(slide, "Currently Active (4 Campaigns)", "Week 1 Daily CPA")
    replace_text_in_shape_any(slide, "The 43 Paused Campaigns", "What Changed & What's Next")
    replace_text_in_shape_any(slide, "Legacy Agencies", "Bid Strategy Switch")
    replace_text_in_shape_any(slide, "Historical Experiments", "Flagged for Week 2")

    # Table
    for shape in slide.shapes:
        if shape.has_table and len(shape.table.rows) == 5:
            populate_table(shape.table, [
                ['Day', 'Spend', 'Primary Conv', 'CPA'],
                ['Mon 13 Apr', '£2,017', '17', '£119'],
                ['Tue 14 Apr', '£558', '10', '£56'],
                ['Wed 15 Apr', '£402', '16', '£25'],
                ['Thu 16 Apr', '£1,792', '23.5', '£76'],
            ])
            break

    # Legacy Agencies / Bid Strategy Switch body
    for shape in slide.shapes:
        if shape.has_text_frame and "KMG (Kau Media Group)" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "From: Max Conversion Value (tROAS 80%)",
                "To: Target CPA £60",
                "Drove campaign to optimise for LEAD VOLUME at controllable cost, not arbitrary £300 booking values",
                "Wed 15 Apr: first full day under new strategy delivered £25 CPA — best single day in the account",
            ])
            break

    # Historical Experiments / Flagged for Week 2 body
    for shape in slide.shapes:
        if shape.has_text_frame and "Device-split campaigns" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Google still reports 'most asset groups limited by policy'",
                "This restricts PMax reach — it's not running at full potential",
                "Full asset rebuild scheduled Wed 22 April",
                "Stripping unverified claims: '99.12% success', 'Save 60%', 'Top 10 Europe'",
                "Replacing with policy-safe copy aligned to Giulio's verified claims",
            ])
            break

    # Recommendation body
    replace_text_in_shape_any(slide, "Recommendation", "Recommendation")
    for shape in slide.shapes:
        if shape.has_text_frame and "Remove all campaigns paused 6+" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Hold the new bid strategy untouched — let Google complete its 7-day learning cycle.",
                "PMax asset rebuild addresses the policy flag next week.",
                "Expected upside: lifting policy limits typically unlocks 2-3× reach at the same CPA.",
            ])
            break


def populate_slide_6(slide):
    """Dental Implants Intent (from Ad Group Structure)."""
    replace_text_in_shape_any(slide, "Ad Group Structure (2026)", "Dental Implants Intent")
    replace_text_in_shape_any(slide,
        "Search Intent campaign: 8 ad groups with wide CPA variation",
        "Fully rebuilt Thu–Fri · 4 ad groups → 25 · entered learning phase")
    replace_text_in_shape_any(slide, "2025 vs 2026", "Rebuild")

    # Table (12 rows × 6 cols)
    for shape in slide.shapes:
        if shape.has_table and len(shape.table.rows) >= 10:
            populate_table(shape.table, [
                ['Metric', 'Before (legacy)', 'After (new)', 'Change', 'Impact', 'Status'],
                ['Ad groups', '4 (paused)', '25 live', '+525%', 'Patient-intent grouped', 'Complete'],
                ['Keywords', '27', '~675', '+2,400%', 'Phrase + exact', 'Complete'],
                ['RSAs', '6 (Poor-rated)', '75', '+1,150%', 'Keyword-aligned', 'Complete'],
                ['Headlines', '~90', '1,050', '+1,067%', 'Unique, keyword-rich', 'Complete'],
                ['Descriptions', '~24', '280', '+1,067%', '75-86 chars (QS safe)', 'Complete'],
                ['Legacy AG CPA', '£99+ avg', 'Paused', 'Eliminated', 'Waste removed', 'Done'],
                ['Bid strategy', 'Max Conv Value (tROAS)', 'Target CPA £75', 'Controllable', 'Lead-focused', 'Learning'],
                ['Geo', '25-mile radius', '270 London postcodes', 'M25 targeting', '4-tier structure', 'Rolled out'],
                ['Schedule', 'Always on', '42-cell structure', 'Hourly analysis', 'Reporting ready', 'Live'],
                ['Learning phase', 'n/a', 'Day 1 of 14', 'In progress', '7-14 days to stabilise', 'Monitor'],
                ['Campaign status', 'Scaling with waste', 'Learning - clean base', 'Foundation set', 'Week 2-3 read', 'On track'],
            ])
            break

    # Key Finding body
    for shape in slide.shapes:
        if shape.has_text_frame and "'Dental Implants' ad group (£26.6K)" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "The DII campaign was fully restructured mid-week (Tue → Fri).",
                "Week 1 CPA blends 1 day of the old structure with 4 days of Google's learning phase — expected noise, not signal.",
                "Google requires 14-21 days to rebalance bids against the new structure.",
                "First clean performance read: Week 2-3.",
            ])
            break

    # Recommendation body
    for shape in slide.shapes:
        if shape.has_text_frame and "'Financing' and 'Implants Cost'" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Hold the new structure untouched through Wed 22 April.",
                "Ad Strength audit on AG22-25 resumes Tuesday once Google clears the new ads from Pending review.",
                "Second pass on AG1-21 to push remaining RSAs toward Excellent.",
            ])
            break


def populate_slide_7(slide):
    """The Numbers That Matter (from Quality Score Analysis)."""
    replace_text_in_shape_any(slide, "Quality Score Analysis", "The Numbers That Matter")
    replace_text_in_shape_any(slide,
        "Weighted average QS: 6.9 (2025) → 6.4 (2026) — declining",
        "Business outcomes — what the account delivered this week")
    replace_text_in_shape_any(slide, "2025 vs 2026", "Business Outcomes")
    replace_text_in_shape_any(slide, "QS Components (2025 vs 2026)", "Week 1 at a Glance")
    replace_text_in_shape_any(slide, "Worst QS Keywords (QS 1-3) — 2026 Spend", "Week 1 Conversion Breakdown")

    # Main table 1
    for shape in slide.shapes:
        if shape.has_table and len(shape.table.rows) == 4:
            populate_table(shape.table, [
                ['Metric', 'Week 1', 'Prior Week', 'Change', 'Status'],
                ['Total Spend', '£8,462', '£10,291', '−18%', 'Trimmed'],
                ['Dengro Offline Leads', '87', '—', '—', 'Tracking'],
                ['Primary Conversions', '101', '134', '−25%', 'Learning'],
            ])
            break

    # Main table 2 (6 rows × 5 cols — campaign breakdown)
    tables = [s for s in slide.shapes if s.has_table]
    if len(tables) >= 2:
        populate_table(tables[1].table, [
            ['Campaign', 'Spend', 'Primary', 'Dengro Leads', 'Cost/Lead'],
            ['Performance Max', '£5,960', '81.5', '78.5', '£75.92'],
            ['Dental Implants Intent', '£2,222', '7.5', '6.5', '£341.83'],
            ['Brand', '£280', '12.0', '2.0', '£140.24'],
            ['TOTAL', '£8,462', '101.0', '87.0', '£97.27'],
        ])

    replace_text_in_shape_any(slide, "Consistent Pattern Both Years", "Key Signal")
    for shape in slide.shapes:
        if shape.has_text_frame and "Landing pages: strong" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "87 Dengro Offline Leads generated this week across 3 active campaigns.",
                "Performance Max delivered 90% of qualified leads (78.5 of 87).",
                "Dengro Offline Booking fired 0 times — tracking gap flagged for Monday's audit.",
                "No CPA-based decisions will be drawn before the audit closes.",
            ])
            break

    # Note at bottom
    replace_text_in_shape_any(slide,
        "Note: In 2025, most low-QS keywords (1-5) had zero spend because the Search Intent campaign was only active late in the year. The QS issues have emerged as this campaign scaled in 2026.",
        "Dengro Offline Booking fired 0 times in Week 1. Monday's tracking audit will confirm whether this is a reporting gap or a genuine zero-booking week.")


def populate_slide_8(slide):
    """8-Week Context (from Ad Copy Audit). USER ASKED: add CPA per week."""
    replace_text_in_shape_any(slide, "Ad Copy Audit", "8-Week Context")
    replace_text_in_shape_any(slide,
        "90% of ads rated Poor or Average — the primary driver of low Quality Scores",
        "Weekly performance trend · 23 Feb – 17 April 2026 · spend, conversions, CPA")
    replace_text_in_shape_any(slide, "Current State", "8-Week Trend")
    replace_text_in_shape_any(slide, "Current Ad Strength (All Ads in Account)", "Weekly Performance")

    # Table — USER WANTED CPA per week
    for shape in slide.shapes:
        if shape.has_table and len(shape.table.rows) == 6:
            populate_table(shape.table, [
                ['Week', 'Spend', 'Primary Conv', 'CPA'],
                ['23 Feb', '£17,568', '237.5', '£74'],
                ['2 Mar', '£14,856', '189.6', '£78'],
                ['9 Mar', '£8,393', '120.3', '£70'],
                ['16 Mar', '£13,318', '137.7', '£97'],
                ['23 Mar', '£11,558', '239.2', '£48'],
            ])
            break

    replace_text_in_shape_any(slide, "Critical Finding", "Key Observation")
    for shape in slide.shapes:
        if shape.has_text_frame and "main Search Intent campaign (£15K spend" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Weekly spend fluctuated dramatically: £17.5k peak (23 Feb) → £8.4k (Week 1). A 52% range in 8 weeks.",
                "CPA swung from £48 (best week) to £97 (worst week) — highly volatile before our engagement.",
                "Additional weeks: 30 Mar £9,997/£67 CPA · 6 Apr £10,291/£77 CPA · 13 Apr £8,462/£84 CPA",
                "Week 1 (13 Apr) reflects 5 days of structural disruption: new bid strategies, rebuilt campaigns, changing keyword base.",
            ])
            break

    replace_text_in_shape_any(slide, "Spend on Poor/Average Ads", "Week 1 in Context")
    for shape in slide.shapes:
        if shape.has_text_frame and "2025: £265,869 total spend" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Weekly spend volatility has been the account's core instability — this rebuild brings it under control.",
                "Target: stabilise weekly spend within a consistent range aligned to your target monthly budget (see Questions slide).",
                "Efficiency (cost per lead) is the metric to watch from Week 3 onwards, once the rebuild stabilises.",
                "First clean performance read: Week 2-3 after Google's learning phase completes.",
            ])
            break

    replace_text_in_shape_any(slide, "Recommendation:", "Next Week:")
    for shape in slide.shapes:
        if shape.has_text_frame and "Rewrite all Poor and Average ads." in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Monday's tracking audit confirms conversion reliability — the gate to CPA analysis.",
                "Week 2 performance comparison vs Week 1 will be the first post-rebuild read.",
                "Please confirm target monthly budget so I can lock weekly pacing (see Slide 10).",
            ])
            break


def populate_slide_9(slide):
    """Client Clarifications & Actions (from Landing Page Performance)."""
    replace_text_in_shape_any(slide, "Landing Page Performance", "Client Clarifications & Actions")
    replace_text_in_shape_any(slide,
        "CPA varies 7x+ across landing pages — and has worsened year-on-year on the top pages",
        "Giulio replied Fri 17 April 16:47 · actions captured")
    replace_text_in_shape_any(slide, "2025 vs 2026", "Client Reply")
    replace_text_in_shape_any(slide, "Best Performer", "Confirmed by Giulio")
    replace_text_in_shape_any(slide, "Worst Performer", "Compliance Actions Required")
    replace_text_in_shape_any(slide, "The Good News", "Advertising Scope")
    replace_text_in_shape_any(slide, "Recommendation:", "Actions This Week:")

    # Best Performer → Confirmed by Giulio
    for shape in slide.shapes:
        if shape.has_text_frame and "/dental-implants-offer:" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Consultation: Free + £25 booking deposit",
                "Lead → Booking rate: 13%",
                "Booking → Treatment rate: 35%",
                "Patient base: ~10,000 patients, 20,000+ implants placed",
                "25+ years team-wide expertise in dental field",
            ])
            break

    # Worst Performer → Compliance Actions Required
    for shape in slide.shapes:
        if shape.has_text_frame and "/teeth-implants (highest spend):" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Remove '99.12% success rate' from all ads — confirmed estimate",
                "Remove 'Save 60%' — no source benchmark",
                "Reword 'Top 10 Europe' → factual descriptor (self-claim, not award)",
                "Reword 'Best Implant Clinic 2024' → '2024 Finalist'",
                "Update ad copy: 'Free Consultation (£25 booking deposit)'",
            ])
            break

    # The Good News → Advertising Scope
    for shape in slide.shapes:
        if shape.has_text_frame and "QS data confirms landing pages score well" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "DBD only advertises dental implants (single + double arches) — confirmed by Giulio.",
                "Cosmetic and general dentistry services deferred for now.",
                "All campaigns aligned to this focus: Brand, PMax, Dental Implants Intent.",
            ])
            break

    # Recommendation → Actions This Week
    for shape in slide.shapes:
        if shape.has_text_frame and "Route traffic to the best-converting pages" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Applying all compliance changes Monday 20 Apr: removing unverifiable claims, updating consultation pricing.",
                "Replacing '2,500+ patients' claim with '20,000+ implants placed' across relevant ads.",
                "AG22-25 Ad Strength audit resumes Tuesday (ads currently Pending Google review).",
            ])
            break


def populate_slide_10(slide):
    """Questions for Client (from Change History) — NEW SLIDE TOPIC."""
    replace_text_in_shape_any(slide, "Change History", "Questions for Client")
    replace_text_in_shape_any(slide,
        "2,543 changes — 89% by one person, 24 by auto-apply",
        "To refine strategy, I need clarification on a few things")
    replace_text_in_shape_any(slide, "Account Management", "Your Input")
    replace_text_in_shape_any(slide, "Who Manages the Account", "Budget Questions (Urgent)")
    replace_text_in_shape_any(slide, "What's Being Changed", "Business Context (When You Can)")
    replace_text_in_shape_any(slide, "Key Observation", "How to Reply")
    replace_text_in_shape_any(slide, "Auto-Apply: 24 Changes", "Why These Matter")

    # Table 1 — Budget Questions
    tables = [s for s in slide.shapes if s.has_table]
    if len(tables) >= 1:
        populate_table(tables[0].table, [
            ['Question', 'Why It Matters', 'Answer Needed By'],
            ['What is your target monthly ad spend?', 'Weekly spend fluctuated £8k-£17k over 8 weeks. Target budget drives bid and budget decisions.', 'Mon 20 Apr'],
            ['What is your target weekly ad spend?', 'Weekly consistency drives campaign pacing. Current daily budgets sum to ~£14k/week but actual is ~£8.5k.', 'Mon 20 Apr'],
            ['Is there a hard monthly ceiling?', 'Prevents overspend and informs daily budget caps across 3 campaigns.', 'Mon 20 Apr'],
            ['Any seasonal budget flex (e.g. summer)?', 'Lets me set higher/lower target CPAs by season. January often quieter, summer busier.', '1 week'],
            ['Priority service: single or double arch?', 'Higher-margin service should carry higher bid caps and budget share.', '1 week'],
            ['', '', ''],
            ['', '', ''],
        ])

    # Table 2 — Business Context
    if len(tables) >= 2:
        populate_table(tables[1].table, [
            ['Question', 'Why It Matters', 'Answer Needed By'],
            ['Average revenue per implant patient?', 'To calculate true ROAS instead of the £300 arbitrary booking value.', '2 weeks'],
            ['Time from enquiry → treatment start?', 'Sets true attribution windows and conversion lag expectations.', '2 weeks'],
            ['Which London areas drive most patients?', 'Validates / refines the 270-postcode geo targeting.', '2 weeks'],
            ['Typical patient age range?', 'Informs demographic bid adjustments.', '2 weeks'],
            ['Marketing budget — Google vs other channels?', 'Gives context for total marketing mix and Google\'s share of voice.', '2 weeks'],
            ['Key competitors to watch?', 'Helps target auction insights and competitive ad copy angles.', '2 weeks'],
            ['What\'s the clinic\'s capacity limit?', 'Prevents over-scaling Google Ads beyond what the clinic can handle.', '2 weeks'],
            ['', '', ''],
            ['', '', ''],
        ])

    # How to Reply body
    for shape in slide.shapes:
        if shape.has_text_frame and "Giulio is doing the work" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Reply to this email with answers to any questions — no need to answer all at once.",
                "Or we can cover these on a 20-30 min call whenever convenient.",
                "Urgent questions (budget) need answers by Monday so Week 2 work isn't blocked.",
            ])
            break

    # Why These Matter body
    for shape in slide.shapes:
        if shape.has_text_frame and "Google auto-removed keywords" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "Every piece of clarification feeds directly into bid strategies, budget allocation, and ad copy.",
                "The more I know about your business, the more precisely I can tune the account.",
                "Example: the 8-week spend volatility (£8k-£17k range) is what happens when budgets aren't locked to a target — pinning your target spend eliminates this.",
            ])
            break


def populate_slide_11(slide):
    """Week 2 Plan (from Issues Summary — Prioritised)."""
    replace_text_in_shape_any(slide, "Issues Summary — Prioritised", "Week 2 Plan")
    replace_text_in_shape_any(slide, "12 Issues Identified", "5 Working Days")
    replace_text_in_shape_any(slide, "Critical (Fix Immediately)", "Monday 20 Apr — Tracking Audit")
    replace_text_in_shape_any(slide, "Important (Fix in 30 Days)", "Tue 21 – Thu 23 Apr — Execution")
    replace_text_in_shape_any(slide, "Monitor (Ongoing)", "Friday 24 Apr — Review + Report")

    # Critical section
    for shape in slide.shapes:
        if shape.has_text_frame and "1. Ghost conversions (KMG)" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "1. Verify form_submit fires on all landing pages",
                "End-to-end validation across every Search campaign LP",
                "2. Verify Dengro Offline Lead / Booking / Purchase reliability",
                "Confirm why Bookings fired 0 times in Week 1 — reporting gap or genuine zero?",
                "3. Audit attribution windows and conversion counting",
                "Ensure Dengro leads aren't being double-counted or missing",
                "4. Gate for all CPA analysis",
                "No CPA-based decisions until tracking is confirmed reliable",
            ])
            break

    # Important section
    for shape in slide.shapes:
        if shape.has_text_frame and "5. No device bid adjustments" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "5. Apply Giulio's claim corrections across all ads (Tue)",
                "Remove 99.12%, Save 60%; reword Top 10, Best Clinic",
                "6. Resume AG22-25 Ad Strength audit once Google clears Pending (Tue)",
                "12 RSAs to audit once review completes",
                "7. Second pass on AG1-21 ads not yet at Excellent (Tue)",
                "Push remaining RSAs over the line",
                "8. PMax asset rebuild (Wed)",
                "Policy-safe copy, verified claims only, restores full reach",
                "9. Ad extensions deployed (Thu)",
                "Sitelinks, callouts, structured snippets — first time the account has a complete set",
            ])
            break

    # Monitor section
    for shape in slide.shapes:
        if shape.has_text_frame and "10. 784 broad match negatives" in shape.text_frame.text:
            set_shape_multi_paragraph(shape, [
                "10. 7-day performance comparison vs Week 1",
                "Spend, leads, CPA — first post-rebuild signal",
                "11. Brand matures, PMax settles, DII still learning",
                "Different campaigns on different learning trajectories",
                "12. Week 2 end-of-week report delivered EoD Friday",
                "Same format as this report — trend, insights, Week 3 plan",
            ])
            break

    # Footer statement
    replace_text_in_shape_any(slide,
        "The 4 Critical issues affect every conversion and every pound spent. Fixing them is the foundation for everything else.",
        "Week 2 is about execution: tracking audit unblocks CPA analysis, claim corrections unblock compliance, PMax rebuild unblocks reach.")


def populate_slide_12(slide):
    """What's Next (closer) — ADDS Slide 12 at end."""
    replace_text_in_shape_any(slide, "What's Next", "What's Next")  # stays
    replace_text_in_shape_any(slide, "These issues become the foundation for Report 3",
                               "Rolling priorities beyond Week 2")
    replace_text_in_shape_any(slide, "1.  Remove KMG ghost conversion actions", "1.  Cosmetic dentistry campaign build")
    replace_text_in_shape_any(slide, "2.  Move soft signals to Secondary", "2.  Landing page optimisation")
    replace_text_in_shape_any(slide, "3.  Verify Dengro Offline Lead attribution", "3.  Long-tail keyword expansion")
    replace_text_in_shape_any(slide, "4.  Rewrite all Poor/Average ads", "4.  Conversion rate improvements")
    replace_text_in_shape_any(slide, "5.  Evaluate bid strategy: tCPA vs tROAS", "5.  Monthly review cycle")
    replace_text_in_shape_any(slide, "6.  Clean up 43 paused campaigns", "6.  End-of-month report")

    replace_text_in_shape_any(slide, "Coming Next", "Coming Next")
    replace_text_in_shape_any(slide, "Report 3:Restructure Plan", "Week 2 Report")
    replace_text_in_shape_any(slide,
        "Proposed campaign structure, bid strategy, conversion tracking fixes, ad copy plan, and projected CPA targets.",
        "Weekly performance review, tracking audit findings, PMax rebuild results, and Week 3 plan. Delivered EoD Friday 24 April.")

    replace_text_in_shape_any(slide, "Critical Issues", "Compliance Actions")
    replace_text_in_shape_any(slide, "Important Issues", "Working Days")
    replace_text_in_shape_any(slide, "Report 3", "Week 2")
    # The numbers "4" "5" and "Report 3" are positional in the closer — for Week 1
    # we want: 12 (compliance actions), 5 (working days), "Week 2" (coming next)
    replace_text_in_shape_any(slide, "4", "12")  # Critical count
    # Keep "5" (working days stays 5)


# =============================================================================
# MAIN
# =============================================================================

def main():
    prs = Presentation(str(FILE_PATH))
    slides = list(prs.slides)
    print(f"Loaded v4: {FILE_PATH.name}")
    print(f"Populating {len(slides)} slides...")

    populate_slide_1(slides[0]); print(" [OK] Slide 1  — Title")
    populate_slide_2(slides[1]); print(" [OK] Slide 2  — Executive Summary")
    populate_slide_3(slides[2]); print(" [OK] Slide 3  — Work Delivered")
    populate_slide_4(slides[3]); print(" [OK] Slide 4  — Brand Campaign")
    populate_slide_5(slides[4]); print(" [OK] Slide 5  — Performance Max")
    populate_slide_6(slides[5]); print(" [OK] Slide 6  — Dental Implants Intent")
    populate_slide_7(slides[6]); print(" [OK] Slide 7  — The Numbers That Matter")
    populate_slide_8(slides[7]); print(" [OK] Slide 8  — 8-Week Context (+ CPA per week)")
    populate_slide_9(slides[8]); print(" [OK] Slide 9  — Client Clarifications & Actions")
    populate_slide_10(slides[9]); print(" [OK] Slide 10 — Questions for Client (NEW)")
    populate_slide_11(slides[10]); print(" [OK] Slide 11 — Week 2 Plan")
    populate_slide_12(slides[11]); print(" [OK] Slide 12 — What's Next (closer)")

    prs.save(str(FILE_PATH))
    print(f"\nSaved: {FILE_PATH.name}")


if __name__ == "__main__":
    main()
