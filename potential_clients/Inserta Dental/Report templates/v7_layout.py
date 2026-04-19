"""v7 slide 5 layout fix — use full width, no wasted space."""
import sys
from pathlib import Path
if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches

F = Path(__file__).resolve().parents[3] / "potential_clients" / "Inserta Dental" / "End-of-week reports" / "DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v7.pptx"
prs = Presentation(str(F))
s = list(prs.slides)[4]

# Layout plan (12.13" content width between 0.60 and 12.73):
# y=0.30..1.10  Title + subtitle (keep)
# y=1.10..1.45  Daily CPA heading (keep at y=1.10)
# y=1.45..4.30  Table (full width 12.13", h=2.85, 6 rows at 0.475")
# y=4.55..4.85  "What Changed & What's Next" heading
# y=5.00..6.90  3 cards (each full height 1.90" — room for 5-6 text lines)

for i, shape in enumerate(s.shapes):
    l = shape.left/914400 if shape.left else 0
    t = shape.top/914400 if shape.top else 0
    w = shape.width/914400 if shape.width else 0
    h = shape.height/914400 if shape.height else 0

    # Table
    if shape.has_table:
        shape.left = Inches(0.60)
        shape.top = Inches(1.45)
        shape.width = Inches(12.13)
        shape.height = Inches(2.85)
        for row in shape.table.rows:
            row.height = Inches(0.475)
        continue

    if not shape.has_text_frame:
        # Card backgrounds and borders (no text frame)
        # BGs at y=4.65, h=1.70 → move to y=5.00, h=1.90
        if abs(t - 4.65) < 0.05 and abs(h - 1.70) < 0.05:
            shape.top = Inches(5.00)
            shape.height = Inches(1.90)
        continue

    txt = shape.text_frame.text.strip()

    # "Daily CPA" heading — widen to match table
    if txt == "Week 1 Daily CPA":
        shape.width = Inches(12.13)

    # "What Changed & What's Next" section heading
    elif txt == "What Changed & What's Next":
        shape.top = Inches(4.55)

    # Card headings
    elif txt in ("Bid Strategy Switch", "Flagged for Week 2", "Recommendation"):
        shape.top = Inches(5.10)

    # Card body text (identify by first words)
    elif txt.startswith("From: Max Conversion Value") or \
         txt.startswith("Google still reports") or \
         txt.startswith("Hold the new bid strategy"):
        shape.top = Inches(5.45)
        shape.height = Inches(1.45)

prs.save(str(F))
print("Done")
