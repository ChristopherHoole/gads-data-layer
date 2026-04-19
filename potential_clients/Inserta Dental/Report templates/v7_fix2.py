"""v7 final: fix table columns + card heights."""
import sys
from pathlib import Path
if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches

F = Path(__file__).resolve().parents[3] / "potential_clients" / "Inserta Dental" / "End-of-week reports" / "DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v7.pptx"
prs = Presentation(str(F))
s = list(prs.slides)[4]

for shape in s.shapes:
    if shape.has_table:
        # Table is wide now (12.13) — set column widths
        cols = shape.table.columns
        cols[0].width = Inches(2.50)  # Day
        cols[1].width = Inches(3.20)  # Spend
        cols[2].width = Inches(3.20)  # Primary Conv
        cols[3].width = Inches(3.23)  # CPA
        continue

    if not shape.has_text_frame:
        # Card BG (not text) at y=5.00 — give them more height
        t = shape.top/914400 if shape.top else 0
        h = shape.height/914400 if shape.height else 0
        if abs(t - 5.00) < 0.05 and abs(h - 1.90) < 0.05:
            shape.top = Inches(4.90)
            shape.height = Inches(2.05)
        continue

    txt = shape.text_frame.text.strip()
    # What Changed heading up
    if txt == "What Changed & What's Next":
        shape.top = Inches(4.55)
    # Card heading up
    elif txt in ("Bid Strategy Switch", "Flagged for Week 2", "Recommendation"):
        shape.top = Inches(5.00)
    # Card body — more room
    elif txt.startswith("From: Max Conversion Value") or \
         txt.startswith("Google still reports") or \
         txt.startswith("Hold the new bid strategy"):
        shape.top = Inches(5.35)
        shape.height = Inches(1.60)

prs.save(str(F))
print("Done")
