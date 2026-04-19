"""Clean redo of Slide 5 from v5 baseline.

Plan (v5 layout preserved, additions only):
1. Keep table at its v5 position (x=0.69, y=1.45, w=3.50, h=2.50) but add:
   - Friday row
   - Weekly Total row (bold, navy tint)
   Expand table h to 3.50 to fit all 7 rows.
2. Push 'What Changed' heading from y=4.30 to y=5.15
3. Push 3 card BGs + borders from y=4.65 to y=5.50 (h stays 1.70)
4. Push card heading textboxes from y=4.75 to y=5.55 (but they overflow? Need to verify position)
5. Push card body textboxes from y=5.05 to y=5.80 (shrink height to 1.10)
6. Add narrative card on RIGHT (x=4.30, y=1.45, w=8.40, h=3.50 to match table)
7. Replace em dashes everywhere on slide 5
"""
import sys, copy
from pathlib import Path

if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

FILE_PATH = Path(__file__).resolve().parents[3] / "potential_clients" / "Inserta Dental" / "End-of-week reports" / "DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v6.pptx"

NAVY = RGBColor(0x1A, 0x23, 0x7E)
BLUE = RGBColor(0x42, 0x85, 0xF4)
BODY = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY_TINT = RGBColor(0xE8, 0xEA, 0xF6)
BORDER_LIGHT = RGBColor(0xE0, 0xE0, 0xE0)


def add_row_clone(table, row_data, bold=False, fill=None, text_color=None):
    """Clone last row, insert after it, populate with row_data."""
    tbl_elem = table._tbl
    last_tr = tbl_elem.findall(qn('a:tr'))[-1]
    new_tr = copy.deepcopy(last_tr)
    last_tr.addnext(new_tr)
    new_row_idx = len(table.rows) - 1
    for c_idx, val in enumerate(row_data):
        cell = table.cell(new_row_idx, c_idx)
        saved = None
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                f = r.font
                saved = {'name': f.name, 'size': f.size, 'bold': f.bold, 'align': p.alignment}
                try:
                    saved['color'] = f.color.rgb if f.color and f.color.type else None
                except Exception:
                    saved['color'] = None
                break
            if saved: break
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        if saved and saved.get('align'): p.alignment = saved['align']
        r = p.add_run()
        r.text = val
        if saved:
            if saved['name']: r.font.name = saved['name']
            if saved['size']: r.font.size = saved['size']
            r.font.bold = bold if bold else saved['bold']
            if text_color: r.font.color.rgb = text_color
            elif saved['color']: r.font.color.rgb = saved['color']
        if fill:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill


def replace_em_dashes(shape):
    if not shape.has_text_frame:
        return False
    changed = False
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if "—" in r.text:
                r.text = r.text.replace(" — ", ", ").replace("—", ",")
                changed = True
    return changed


def main():
    prs = Presentation(str(FILE_PATH))
    slide = list(prs.slides)[4]

    # -------------------------------------------------------------------
    # STEP 1: Classify shapes (before any changes)
    # -------------------------------------------------------------------
    table_shape = None
    what_changed_heading = None
    card_bgs = []          # card body rectangles (pink/yellow/green tinted)
    card_borders = []      # 0.06"-wide coloured strips on left of each card
    card_heading_tbs = []  # "Bid Strategy Switch" / "Flagged for Week 2" / "Recommendation" text boxes
    card_body_tbs = []     # card body text boxes

    for shape in slide.shapes:
        if shape.has_table and shape.table.cell(0, 0).text.strip() == "Day":
            table_shape = shape
        elif shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt == "What Changed & What's Next":
                what_changed_heading = shape
            elif txt in ("Bid Strategy Switch", "Flagged for Week 2", "Recommendation"):
                card_heading_tbs.append(shape)
            elif ("From: Max Conversion Value" in txt or
                  "Google still reports 'most asset" in txt or
                  "Hold the new bid strategy untouched" in txt):
                card_body_tbs.append(shape)
        elif shape.shape_type == 1:
            left = shape.left / 914400 if shape.left else 0
            top = shape.top / 914400 if shape.top else 0
            w = shape.width / 914400 if shape.width else 0
            h = shape.height / 914400 if shape.height else 0
            # Card BGs (pink/yellow/green tinted, width ~3.80-3.90, h~1.70, y=4.65)
            if 3.5 < w < 4.1 and 1.5 < h < 1.9 and 4.0 < top < 5.0:
                card_bgs.append(shape)
            # Card border strips (w ~0.06, same y-range)
            elif w < 0.12 and h > 1.0 and 4.0 < top < 5.0:
                card_borders.append(shape)

    print(f"Found: table={table_shape is not None}, What Changed heading={what_changed_heading is not None}")
    print(f"Cards: {len(card_bgs)} BGs, {len(card_borders)} borders, {len(card_heading_tbs)} headings, {len(card_body_tbs)} bodies")

    # -------------------------------------------------------------------
    # STEP 2: Expand table to fit 7 rows
    # -------------------------------------------------------------------
    table_shape.height = Inches(3.50)
    for row in table_shape.table.rows:
        row.height = Inches(0.50)

    add_row_clone(table_shape.table, ['Fri 17 Apr', '£1,188', '15', '£79'])
    add_row_clone(
        table_shape.table,
        ['Week total', '£5,960', '81.5', '£73.13'],
        bold=True, fill=NAVY_TINT, text_color=NAVY,
    )
    # Set explicit heights on the new rows too
    for row in table_shape.table.rows:
        row.height = Inches(0.50)
    print(" [OK] Table expanded to 7 rows, height 3.50, Friday+Weekly Total added")

    # -------------------------------------------------------------------
    # STEP 3: Shift 'What Changed' section down to clear table
    # Table now ends at y=1.45+3.50=4.95
    # Give 0.15 gap, heading at y=5.10, card BGs at y=5.40
    # Card BGs h=1.60, ends y=7.00 (touches footer)
    # -------------------------------------------------------------------
    what_changed_heading.top = Inches(5.10)
    what_changed_heading.height = Inches(0.25)

    for bg in card_bgs:
        bg.top = Inches(5.40)
        bg.height = Inches(1.55)  # reduce slightly to not touch footer
    for border in card_borders:
        border.top = Inches(5.40)
        border.height = Inches(1.55)

    for th in card_heading_tbs:
        th.top = Inches(5.50)
        th.height = Inches(0.25)
    for tb in card_body_tbs:
        tb.top = Inches(5.80)
        tb.height = Inches(1.15)

    print(" [OK] 'What Changed' heading + 3 cards shifted down")

    # -------------------------------------------------------------------
    # STEP 4: Add right-side narrative card (fills empty space next to table)
    # Card at x=4.30, y=1.45, w=8.40, h=3.50 (matches table bounds)
    # -------------------------------------------------------------------
    card_x, card_y, card_w, card_h = 4.30, 1.45, 8.40, 3.50

    # Left-edge navy strip
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(card_x), Inches(card_y), Inches(0.08), Inches(card_h),
    )
    border.fill.solid()
    border.fill.fore_color.rgb = NAVY
    border.line.fill.background()

    # Card body
    body = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(card_x + 0.08), Inches(card_y), Inches(card_w - 0.08), Inches(card_h),
    )
    body.fill.solid()
    body.fill.fore_color.rgb = WHITE
    body.line.color.rgb = BORDER_LIGHT

    # Heading
    heading_tb = slide.shapes.add_textbox(
        Inches(card_x + 0.25), Inches(card_y + 0.15), Inches(card_w - 0.35), Inches(0.35)
    )
    p = heading_tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Week 1 Story"
    r.font.name = "Calibri"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = NAVY

    # Body — structured narrative
    body_tb = slide.shapes.add_textbox(
        Inches(card_x + 0.25), Inches(card_y + 0.60), Inches(card_w - 0.35), Inches(card_h - 0.75)
    )
    tf = body_tb.text_frame
    tf.word_wrap = True
    tf.clear()

    lines = [
        ("Bid strategy switched Wed 15 April, at the midpoint of the week. Each day tells a different part of the story:", False, BODY),
        ("", False, BODY),
        ("Mon 13 & Tue 14", True, BLUE),
        ("Operating under the old tROAS strategy with characteristic swings, £119 and £56 CPA.", False, BODY),
        ("", False, BODY),
        ("Wed 15", True, BLUE),
        ("First partial day under new Target CPA £60, delivered £25 CPA, the best single day in the account.", False, BODY),
        ("", False, BODY),
        ("Thu 16 & Fri 17", True, BLUE),
        ("Fully post-switch, settling toward the £60 target at £76 and £79 CPA.", False, BODY),
        ("", False, BODY),
        ("Week-total CPA: £73.13 vs prior week £85.86, a 15% improvement despite policy restrictions.", True, NAVY),
    ]
    for i, (line, is_bold, colour) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.name = "Calibri"
        r.font.size = Pt(11)
        r.font.bold = is_bold
        r.font.color.rgb = colour
    print(" [OK] Added 'Week 1 Story' narrative card (right side, aligned with table)")

    # -------------------------------------------------------------------
    # STEP 5: Replace em dashes across slide
    # -------------------------------------------------------------------
    n = 0
    for shape in slide.shapes:
        if replace_em_dashes(shape):
            n += 1
    print(f" [OK] Replaced em dashes in {n} shape(s)")

    prs.save(str(FILE_PATH))
    print(f"\nSaved: {FILE_PATH.name}")


if __name__ == "__main__":
    main()
