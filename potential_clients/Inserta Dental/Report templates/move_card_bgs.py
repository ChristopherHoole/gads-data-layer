"""Move card BGs + borders in slide 5 to match the repositioned card headings/bodies.

Current state in v6 slide 5 (AFTER redo_slide5.py):
- Table: y=1.45, h=3.50 ✓
- What Changed heading: y=5.10 ✓
- Card BGs (shapes 11, 15, 19): STILL at y=4.65, h=1.70 ← need to move to y=5.40, h=1.55
- Card borders (shapes 12, 16, 20): STILL at y=4.65, h=1.70 ← need to move to y=5.40, h=1.55
- Card headings (shapes 13, 17, 21): y=5.50 ✓
- Card bodies (shapes 14, 18, 22): y=5.80 ✓
"""
import sys
from pathlib import Path

if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches

FILE_PATH = Path(__file__).resolve().parents[3] / "potential_clients" / "Inserta Dental" / "End-of-week reports" / "DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v6.pptx"


def main():
    prs = Presentation(str(FILE_PATH))
    slide = list(prs.slides)[4]

    moved = 0
    for shape in slide.shapes:
        if shape.shape_type != 1:  # AUTO_SHAPE only
            continue
        top = shape.top / 914400 if shape.top else 0
        h = shape.height / 914400 if shape.height else 0
        w = shape.width / 914400 if shape.width else 0
        # Target: shapes at exactly y=4.65 with h=1.70 (the card BGs and borders)
        if abs(top - 4.65) < 0.02 and abs(h - 1.70) < 0.05:
            # Move them to y=5.40, h=1.55
            shape.top = Inches(5.40)
            shape.height = Inches(1.55)
            moved += 1

    print(f" [OK] Moved {moved} card BG/border shape(s) to y=5.40, h=1.55")
    prs.save(str(FILE_PATH))


if __name__ == "__main__":
    main()
