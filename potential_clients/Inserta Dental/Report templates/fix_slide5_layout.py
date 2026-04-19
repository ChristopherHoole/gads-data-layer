"""Adjust Slide 5 layout: expand table height, push 'What Changed' section down, align narrative card."""
import sys
from pathlib import Path

if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches

FILE_PATH = Path(__file__).resolve().parents[3] / "potential_clients" / "Inserta Dental" / "End-of-week reports" / "DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v6.pptx"


def main():
    prs = Presentation(str(FILE_PATH))
    slide = list(prs.slides)[4]

    # Find shapes by characteristic
    table_shape = None
    heading_what_changed = None
    cards_to_move = []
    narrative_heading = None
    narrative_body = None
    for shape in slide.shapes:
        if shape.has_table and shape.table.cell(0, 0).text.strip() == "Day":
            table_shape = shape
        elif shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt == "What Changed & What's Next":
                heading_what_changed = shape
            elif txt in ("Bid Strategy Switch", "Flagged for Week 2", "Recommendation"):
                cards_to_move.append(shape)
            elif txt == "Week 1 Story":
                narrative_heading = shape
            elif "Bid strategy switched Wed 15 April" in txt:
                narrative_body = shape

    # Also catch body text of 3 cards (not headings)
    card_y = 4.75
    card_bodies = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text
            if "From: Max Conversion Value" in t or \
               "Google still reports 'most asset" in t or \
               "Hold the new bid strategy untouched" in t:
                card_bodies.append(shape)

    # -------------------------------------------------------------------
    # Adjust table height
    # -------------------------------------------------------------------
    if table_shape:
        table_shape.height = Inches(3.50)  # was 2.50 — fits 7 rows at 0.50" each
        print(" [OK] Table height 2.50 -> 3.50 (to fit 7 rows)")

    # -------------------------------------------------------------------
    # Move 'What Changed' heading and 3 cards (heading + body) down
    # -------------------------------------------------------------------
    if heading_what_changed:
        heading_what_changed.top = Inches(5.15)  # was 4.30
        print(" [OK] 'What Changed' heading 4.30 -> 5.15")

    for shape in cards_to_move:
        shape.top = Inches(5.55)  # was 4.75
    print(f" [OK] Moved {len(cards_to_move)} card headings 4.75 -> 5.55")

    for shape in card_bodies:
        shape.top = Inches(5.85)  # was 5.05
    print(f" [OK] Moved {len(card_bodies)} card bodies 5.05 -> 5.85")

    # -------------------------------------------------------------------
    # Extend narrative card to match new table height
    # -------------------------------------------------------------------
    # Heading stays at y=1.20
    # Body expand height to cover the gap
    if narrative_body:
        narrative_body.height = Inches(3.30)  # extend down
        print(" [OK] Narrative body height extended to 3.30")

    # Also find & resize the card left-border strip and white body rect
    # These were added by fix_slide5_v6.py — find by distinctive colour/position
    # The NAVY strip at x=7.00 w=0.08 and WHITE rect at x=7.08 w=5.72
    for shape in slide.shapes:
        if shape.shape_type == 1:  # AUTO_SHAPE = rectangle
            left = shape.left / 914400 if shape.left else 0
            top = shape.top / 914400 if shape.top else 0
            w = shape.width / 914400 if shape.width else 0
            # The card border strip (navy, x~7.00, w~0.08)
            if abs(left - 7.00) < 0.05 and abs(w - 0.08) < 0.03:
                shape.height = Inches(3.80)
                print(" [OK] Narrative card navy border strip extended to 3.80")
            # The card body white rect (x~7.08, w~5.72)
            elif abs(left - 7.08) < 0.05 and w > 5.5:
                shape.height = Inches(3.80)
                print(" [OK] Narrative card body rect extended to 3.80")

    # Also extend narrative card heading + body to match
    # The narrative CARD covers y=1.10 to y=1.10+3.80=4.90.
    # Table covers y=1.45 to y=1.45+3.50=4.95 — roughly aligned.

    prs.save(str(FILE_PATH))
    print(f"\nSaved: {FILE_PATH.name}")


if __name__ == "__main__":
    main()
