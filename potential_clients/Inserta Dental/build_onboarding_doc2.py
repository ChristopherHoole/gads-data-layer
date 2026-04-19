"""
Dental by Design — Document 2: Technical Onboarding Checklist
Same design system as the audit reports.
DBD focus only — no partner clinic references.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
REPORTS = os.path.join(SCRIPT_DIR, 'reports')
LOGO = os.path.join(SCRIPT_DIR, 'act_logo_official.png')

BLUE=RGBColor(0x42,0x85,0xF4); RED=RGBColor(0xEA,0x43,0x35)
YELLOW=RGBColor(0xFB,0xBC,0x05); GREEN=RGBColor(0x34,0xA8,0x53)
NAVY=RGBColor(0x1A,0x23,0x7E); BLACK=RGBColor(0x1A,0x1A,0x1A)
GREY_BG=RGBColor(0xF5,0xF6,0xFA); WHITE=RGBColor(0xFF,0xFF,0xFF)
BORDER=RGBColor(0xE2,0xE8,0xF0)
BLUE_TINT=RGBColor(0xE8,0xF0,0xFE); RED_TINT=RGBColor(0xFC,0xE8,0xE6)
GREEN_TINT=RGBColor(0xE6,0xF4,0xEA); GREY_MID=RGBColor(0xCB,0xD5,0xE1)
FONT='Calibri'

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)

def rect(sl,l,t,w,h,fc,lc=None,lw=None):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h);s.fill.solid();s.fill.fore_color.rgb=fc
    if lc:s.line.color.rgb=lc;s.line.width=lw or Pt(1)
    else:s.line.fill.background()
def txt(sl,l,t,w,h,c,sz=11,co=BLACK,b=False,a=PP_ALIGN.LEFT):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    r=tf.paragraphs[0].add_run();r.text=c;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
    tf.paragraphs[0].alignment=a;return bx
def multipara(sl,l,t,w,h,pdata):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    for i,(tx,sz,co,b,sp) in enumerate(pdata):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run();r.text=tx;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
        if sp:p.space_after=Pt(sp)
def top_bar(sl,h=Inches(0.07)):
    for i,c in enumerate([BLUE,RED,YELLOW,GREEN]):rect(sl,Inches(i*3.333),Inches(0),Inches(3.333),h,c)
def bot_bar(sl,y=Inches(6.92)):
    for p,c in zip([Inches(0.60),Inches(3.63),Inches(6.67),Inches(9.70)],[BLUE,RED,YELLOW,GREEN]):rect(sl,p,y,Inches(3.03),Inches(0.03),c)
def bot_bar_title(sl):
    for i,c in enumerate([BLUE,RED,YELLOW,GREEN]):rect(sl,Inches(i*3.333),Inches(7.0),Inches(3.333),Inches(0.04),c)
def sidebar(sl): rect(sl,Inches(0),Inches(0.07),Inches(0.12),Inches(7.43),BLUE)
def logo(sl,l=Inches(0.60),t=Inches(0.50),s=Inches(0.65)):
    try:sl.shapes.add_picture(LOGO,l,t,s,s)
    except:pass
def footer(sl,n):
    bot_bar(sl)
    try:sl.shapes.add_picture(LOGO,Inches(0.60),Inches(7.0),Inches(0.22),Inches(0.22))
    except:pass
    txt(sl,Inches(0.90),Inches(7.0),Inches(6),Inches(0.25),"Christopher Hoole  |  christopherhoole.com  |  Confidential",11,NAVY,True)
    txt(sl,Inches(12.23),Inches(7.0),Inches(0.50),Inches(0.25),str(n),11,NAVY,a=PP_ALIGN.RIGHT)
def stitle(sl,t,sub=None):
    txt(sl,Inches(0.60),Inches(0.30),Inches(7),Inches(0.50),t,28,NAVY,True)
    if sub:txt(sl,Inches(0.60),Inches(0.85),Inches(9),Inches(0.30),sub,11,GREY_MID)
def badge(sl,t):
    rect(sl,Inches(9.13),Inches(0.30),Inches(3.60),Inches(0.45),WHITE,BLUE,Pt(1))
    txt(sl,Inches(9.23),Inches(0.32),Inches(3.40),Inches(0.40),t,11,BLUE,True,PP_ALIGN.CENTER)
def ibox(sl,l,t,w,h,ti,body,ac,bg=None):
    rect(sl,l,t,w,h,bg or GREY_BG);rect(sl,l,t,Inches(0.06),h,ac)
    txt(sl,l+Inches(0.25),t+Inches(0.10),w-Inches(0.40),Inches(0.30),ti,14,NAVY,True)
    txt(sl,l+Inches(0.25),t+Inches(0.40),w-Inches(0.40),h-Inches(0.50),body,11,BLACK)
def rbox(sl,l,t,w,h,c):
    rect(sl,l,t,w,h,BLUE_TINT,BLUE,Pt(1))
    txt(sl,l+Inches(0.25),t+Inches(0.05),w-Inches(0.50),h-Inches(0.10),c,11,NAVY,True)
def add_table(sl,l,t,w,h,hdrs,rows,cw=None):
    ts=sl.shapes.add_table(len(rows)+1,len(hdrs),l,t,w,h);tb=ts.table
    if cw:
        for i,c in enumerate(cw):tb.columns[i].width=Inches(c)
    for ci,hd in enumerate(hdrs):
        c=tb.cell(0,ci);c.text=hd;c.fill.solid();c.fill.fore_color.rgb=NAVY
        for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.font.color.rgb=WHITE;p.font.bold=True;p.font.name=FONT
    for ri,rd in enumerate(rows):
        for ci,v in enumerate(rd):
            c=tb.cell(ri+1,ci);c.text=str(v);c.fill.solid()
            c.fill.fore_color.rgb=GREY_BG if ri%2==0 else WHITE
            for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.font.color.rgb=BLACK;p.font.name=FONT

print("Building Document 2...")

# ── S1: TITLE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=GREY_BG
top_bar(sl);sidebar(sl);logo(sl)
txt(sl,Inches(0.60),Inches(1.40),Inches(5.74),Inches(1.80),"Technical Onboarding\nChecklist",44,NAVY,True)
rect(sl,Inches(0.60),Inches(3.30),Inches(2.50),Inches(0.05),BLUE)
txt(sl,Inches(0.60),Inches(3.55),Inches(5.50),Inches(0.50),"Dental by Design  |  Prodent Group",22,BLUE)
multipara(sl,Inches(0.60),Inches(4.30),Inches(5.50),Inches(1.00),[
    ("Christopher Hoole",11,BLACK,True,2),("Google Ads Specialist  |  April 2026",11,BLACK,False,2),
    ("christopherhoole.com",11,BLUE,False,0)])
rect(sl,Inches(0.60),Inches(5.50),Inches(4.50),Inches(0.50),WHITE,BLUE,Pt(1))
txt(sl,Inches(0.70),Inches(5.52),Inches(4.30),Inches(0.45),"Onboarding  |  Document 2 of 2",11,BLUE,True)

rect(sl,Inches(6.80),Inches(0.50),Inches(5.90),Inches(2.50),WHITE)
rect(sl,Inches(6.80),Inches(0.50),Inches(0.08),Inches(2.50),BLUE)
txt(sl,Inches(7.10),Inches(0.60),Inches(5.40),Inches(0.35),"Purpose",14,NAVY)
txt(sl,Inches(7.10),Inches(1.00),Inches(5.40),Inches(1.80),
    "This checklist covers all the technical access, tracking, and platform setup needed to begin implementation.\n\n"
    "Some items are quick confirmations. Others will require a separate technical session with Giulio.\n\n"
    "Each item has a status column \u2014 we'll update these as access is granted and audits are completed.",
    12,BLACK)

minicard_data = [("7","Sections",BLUE),("Access","Requests",RED),("Audits","To Complete",GREEN)]
for i,(v,l,c) in enumerate(minicard_data):
    left=Inches(6.80+i*2.03)
    rect(sl,left,Inches(3.20),Inches(1.83),Inches(1.15),WHITE)
    rect(sl,left,Inches(3.20),Inches(0.06),Inches(1.15),c)
    txt(sl,left+Inches(0.15),Inches(3.32),Inches(1.58),Inches(0.45),v,22,c,True)
    txt(sl,left+Inches(0.15),Inches(3.80),Inches(1.58),Inches(0.30),l,11,BLACK)

bot_bar_title(sl)


# ── S2: GOOGLE ADS ACCESS ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Google Ads Account Access","Already linked to MCC \u2014 need to confirm access level")
badge(sl,"Section 1 of 7")

add_table(sl,Inches(0.60),Inches(1.20),Inches(12.10),Inches(2.20),
    ["Item","Current Status","Action Required","Owner","Status"],
    [["MCC Link","Linked to Christopher Hoole MCC (152-796-4125)","Confirm link is active","Giulio","\u2705 Done"],
     ["Access Level","Read-only (granted for audit)","Upgrade to Admin access","Giulio","\u23F3 Pending"],
     ["Conversion Actions","1,511 KMG ghosts + 445 soft signals identified","Remove ghosts, move soft signals to Secondary","Christopher","Week 1"],
     ["Auto-Apply","5 recommendations enabled (last fired Feb 2026)","Disable all except ad rotation","Christopher","Week 1"],
     ["Bid Strategy","tROAS with arbitrary \u00a3300 values","Phased migration to tCPA","Christopher","Weeks 2-4"]],
    cw=[2.0,2.5,3.0,1.0,1.0])

ibox(sl,Inches(0.60),Inches(3.60),Inches(5.80),Inches(2.80),
    "Why Admin Access?",
    "Admin (not just Editor) is needed to:\n\n"
    "\u2022 Manage conversion actions (remove KMG ghosts)\n"
    "\u2022 Link/unlink GA4 and other properties\n"
    "\u2022 Manage audience lists\n"
    "\u2022 Access account-level settings\n"
    "\u2022 Manage negative keyword lists at account level\n\n"
    "This is essential for the Week 1 fixes identified in the audit reports.",
    RED,RED_TINT)

ibox(sl,Inches(6.80),Inches(3.60),Inches(5.80),Inches(2.80),
    "Questions for Giulio",
    "1. Can you upgrade my access from Read-only to Admin?\n\n"
    "2. Is the Google Ads account managed through a Majestico MCC, or directly? What is the relationship between Majestico and Prodent Group?\n\n"
    "3. Are there any other Google Ads accounts in the Prodent Group I should be aware of?",
    BLUE,BLUE_TINT)
footer(sl,2)


# ── S3: GA4 + SEARCH CONSOLE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"GA4 & Search Console Access","Need Editor access to both properties")
badge(sl,"Section 2 of 7")

add_table(sl,Inches(0.60),Inches(1.20),Inches(12.10),Inches(1.80),
    ["Platform","Property","Access Required","Grant To","Status"],
    [["GA4","dentalbydesign.co.uk property","Editor","chris@christopherhoole.com","\u23F3 Pending"],
     ["Search Console","dentalbydesign.co.uk (Domain)","Full Access","chris@christopherhoole.com","\u23F3 Pending"]],
    cw=[1.5,2.5,1.5,2.5,1.0])

ibox(sl,Inches(0.60),Inches(3.20),Inches(5.80),Inches(3.20),
    "GA4 Audits I Will Complete",
    "\u2022 Full GA4 \u2194 Google Ads linking audit \u2014 fix anything not connected\n\n"
    "\u2022 Full custom events and conversions audit \u2014 verify every event is firing correctly\n\n"
    "\u2022 Audience configuration \u2014 set up remarketing audiences for campaign targeting\n\n"
    "\u2022 Verify GA4 \u2194 Search Console link (enables organic data in GA4)\n\n"
    "\u2022 Check enhanced measurement settings (scroll, outbound clicks, file downloads)",
    GREEN,GREEN_TINT)

ibox(sl,Inches(6.80),Inches(3.20),Inches(5.80),Inches(3.20),
    "Questions for Giulio",
    "1. What is the GA4 property ID for dentalbydesign.co.uk?\n\n"
    "2. Is GA4 currently linked to the Google Ads account?\n\n"
    "3. Are there any custom events set up beyond defaults? (e.g. quiz completion steps, scroll depth, form interaction)\n\n"
    "4. Is Search Console verified as a Domain property or URL-prefix?\n\n"
    "5. Is Search Console linked to GA4?",
    BLUE,BLUE_TINT)
footer(sl,3)


# ── S4: GTM ACCESS ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"GTM Container Access","Need Publish access to audit and manage tracking tags")
badge(sl,"Section 3 of 7")

add_table(sl,Inches(0.60),Inches(1.20),Inches(12.10),Inches(1.20),
    ["Platform","Container","Access Required","Grant To","Status"],
    [["Google Tag Manager","dentalbydesign.co.uk container","Publish","chris@christopherhoole.com","\u23F3 Pending"]],
    cw=[2.0,2.5,1.5,2.5,1.0])

ibox(sl,Inches(0.60),Inches(2.60),Inches(5.80),Inches(2.50),
    "GTM Audits I Will Complete",
    "\u2022 Full tag audit \u2014 verify every tag is firing on every page type (homepage, promo pages, booking pages, thank you pages)\n\n"
    "\u2022 Verify Dengro tracking implementation \u2014 is it through GTM or hardcoded?\n\n"
    "\u2022 Check for redundant/orphaned tags from previous agencies\n\n"
    "\u2022 Validate conversion tracking tags match Google Ads conversion actions\n\n"
    "\u2022 Check if any other tracking exists (Facebook Pixel, LinkedIn Insight Tag, etc.)",
    GREEN,GREEN_TINT)

ibox(sl,Inches(6.80),Inches(2.60),Inches(5.80),Inches(2.50),
    "Questions for Giulio",
    "1. What is the GTM container ID?\n\n"
    "2. Who currently manages GTM? (Giulio, a developer, or someone else?)\n\n"
    "3. Is GTM installed on ALL pages including promo landing pages and the booking system?\n\n"
    "4. Are there any server-side GTM containers?\n\n"
    "5. How is Dengro tracking implemented? (GTM tag, hardcoded script, API integration?)",
    BLUE,BLUE_TINT)

rbox(sl,Inches(0.60),Inches(5.30),Inches(12.10),Inches(0.50),
    "Publish access (not just Read) is required because I'll be adding, editing, and fixing tracking tags as part of the Week 1 conversion tracking cleanup.")
footer(sl,4)


# ── S5: WEBSITE CMS ACCESS ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Website CMS & Landing Pages","Need to understand the website setup and agree a landing page workflow")
badge(sl,"Section 4 of 7")

ibox(sl,Inches(0.60),Inches(1.20),Inches(5.80),Inches(2.50),
    "What I Need to Know",
    "1. What CMS/platform is dentalbydesign.co.uk built on? (WordPress, Webflow, custom?)\n\n"
    "2. Can I get read-only CMS access to verify tracking is installed correctly?\n\n"
    "3. Who builds and maintains the /promo/google/ landing pages?\n\n"
    "4. How quickly can landing page changes be made?\n\n"
    "5. Is there a staging/test environment?",
    BLUE,BLUE_TINT)

ibox(sl,Inches(6.80),Inches(1.20),Inches(5.80),Inches(2.50),
    "Landing Page Workflow \u2014 Two Options",
    "Option A: I build HTML/CSS landing pages \u2192 hand to Giulio to upload to your server.\n\n"
    "Option B: Set up Vercel + GitHub with a CNAME on your domain (e.g. lp.dentalbydesign.co.uk). "
    "I build pages, get them signed off, and push live myself. Faster iteration cycle.\n\n"
    "Either way, new treatment-specific landing pages are planned for Month 2 of the restructure.",
    GREEN,GREEN_TINT)

ibox(sl,Inches(0.60),Inches(3.90),Inches(12.10),Inches(2.50),
    "Landing Pages Identified in Audit (2026 Performance)",
    "/teeth-implants: \u00a331K spend, \u00a3111 CPA, 8.5% CVR \u2014 highest spend, needs improvement\n"
    "/cheapest-dental-implants: \u00a322K spend, \u00a383 CPA, 7.1% CVR \u2014 second highest\n"
    "/all-on-4-offer: \u00a316K spend, \u00a3114 CPA, 6.2% CVR \u2014 high CPA for a high-value treatment\n"
    "/dental-implants-offer: \u00a33K spend, \u00a334 CPA, 24.7% CVR \u2014 best performer (Financing traffic)\n"
    "/invisalign: \u00a34K spend, \u00a345 CPA, 10.1% CVR \u2014 strong performer\n"
    "/dental-implants-hammersmith: \u00a32K spend, \u00a3142 CPA, 3.6% CVR \u2014 worst performer, needs rebuild",
    YELLOW)
footer(sl,5)


# ── S6: CONVERSION TRACKING AUDIT ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Conversion Tracking Audit","Full audit and cleanup \u2014 building on the findings from the pitch reports")
badge(sl,"Section 5 of 7")

add_table(sl,Inches(0.60),Inches(1.20),Inches(12.10),Inches(3.50),
    ["Conversion Action","Current Status","Action","Priority","Owner"],
    [["KMG OLD LP | Book A Free Consult","Primary (1,250 ghost cv)","Remove","Week 1","Christopher"],
     ["KMG OLD LP | Website Phone Call","Primary (57 ghost cv)","Remove","Week 1","Christopher"],
     ["KMG | Call Extension","Primary (204 ghost cv)","Remove","Week 1","Christopher"],
     ["Online Booking Click","Primary (264 soft signals)","Move to Secondary","Week 1","Christopher"],
     ["Phone Click / WhatsApp / Chat","Primary (181 soft signals)","Move to Secondary","Week 1","Christopher"],
     ["Dengro Offline Lead","Primary (4,663 cv)","Keep \u2014 verify ad-only attribution","Week 1","Giulio to confirm"],
     ["Dengro Offline Booking","Primary (\u00a3300 arbitrary)","Keep \u2014 discuss value approach","Week 1","Joint"],
     ["Calls from the ads","Primary (965 cv)","Keep \u2014 legitimate","No change","\u2014"],
     ["Submit Quiz / Request Callback","Primary","Keep \u2014 legitimate","No change","\u2014"]],
    cw=[2.5,2.0,2.0,1.0,1.5])

ibox(sl,Inches(0.60),Inches(4.90),Inches(5.80),Inches(1.50),
    "Bid Strategy Migration",
    "Phased migration from tROAS to tCPA:\n\n"
    "\u2022 PMax: tCPA \u00a350 (currently tROAS 30%)\n"
    "\u2022 Search campaigns: tCPA \u00a365, reduce to \u00a350 over 8 weeks\n"
    "\u2022 Brand: Max Clicks or low tCPA\n\n"
    "tCPA is the right approach for long acquisition cycles where real revenue can't reliably be pushed back into Google Ads months later.",
    GREEN,GREEN_TINT)

ibox(sl,Inches(6.80),Inches(4.90),Inches(5.80),Inches(1.50),
    "Questions for Giulio",
    "1. Can you confirm Dengro Offline Lead only counts ad-sourced leads? (Still pending from email)\n\n"
    "2. What other tracking exists outside Google Ads? (Facebook Pixel? LinkedIn?)\n\n"
    "3. Is there a phone tracking number on the website different from Google forwarding numbers?",
    BLUE,BLUE_TINT)
footer(sl,6)


# ── S7: CALL TRACKING ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Call Tracking","Understanding how phone calls work within the business")
badge(sl,"Section 6 of 7")

ibox(sl,Inches(0.60),Inches(1.20),Inches(5.80),Inches(5.20),
    "Questions for the Call",
    "1. Do you use a call tracking provider? (CallRail, Infinity, Mediahawk, ResponseTap \u2014 or just Google's forwarding numbers?)\n\n"
    "2. How are inbound calls routed? (Reception? IVR menu? Call centre?)\n\n"
    "3. Which phone numbers appear on the website? (One number sitewide, or different per page?)\n\n"
    "4. Are calls recorded?\n\n"
    "5. Is call quality tracked or scored with AI? (Some providers score calls as \u201Cqualified lead\u201D, \u201Cspam\u201D, \u201Cexisting patient\u201D)\n\n"
    "6. What happens when a call is missed? (Voicemail? Callback system? Lost?)\n\n"
    "7. What are the opening hours for phone enquiries? (If phones aren't answered outside hours but ads run 24/7, we're paying for unanswered calls)\n\n"
    "8. Does Dengro track which calls result in bookings?\n\n"
    "9. Are Google call extensions using forwarding numbers or your actual number?\n\n"
    "10. Giulio mentioned calls tracked only in CRM that don't report back to Google \u2014 can we get this data for the last 12 months?",
    BLUE,BLUE_TINT)

ibox(sl,Inches(6.80),Inches(1.20),Inches(5.80),Inches(5.20),
    "What We Know So Far",
    "From the Google Ads data:\n\n"
    "\u2022 \u201CCalls from the ads\u201D: 965 conversions (all-time)\n"
    "\u2022 Various call extensions: Brand (115), Implants PMax (24), Implants (21), Invisalign (35)\n"
    "\u2022 Google forwarding numbers are in use\n\n"
    "From Giulio's email:\n\n"
    "\u2022 \u201CSome calls are only tracked inside the CRM where the value does not get reported back to Google\u201D\n"
    "\u2022 He offered to share booking data if we give a timeframe\n"
    "\u2022 This means Google Ads is undercounting actual call conversions\n\n"
    "Gap: We don't know if there's a dedicated call tracking provider, whether calls are recorded, or whether call quality is being measured.",
    GREEN,GREEN_TINT)
footer(sl,7)


# ── S8: ACCESS SUMMARY ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=GREY_BG
top_bar(sl);sidebar(sl);logo(sl)
txt(sl,Inches(0.60),Inches(1.40),Inches(5.50),Inches(0.80),"Access Summary",44,NAVY,True)
rect(sl,Inches(0.60),Inches(2.20),Inches(2.50),Inches(0.05),BLUE)
txt(sl,Inches(0.60),Inches(2.40),Inches(5.50),Inches(0.40),"All access requests in one place",14,NAVY)

add_table(sl,Inches(0.60),Inches(2.90),Inches(5.50),Inches(3.30),
    ["Platform","Access Level","Grant To","Status"],
    [["Google Ads","Admin","chris@christopherhoole.com","\u23F3"],
     ["GA4","Editor","chris@christopherhoole.com","\u23F3"],
     ["GTM","Publish","chris@christopherhoole.com","\u23F3"],
     ["Search Console","Full","chris@christopherhoole.com","\u23F3"],
     ["Website CMS","Read-only","chris@christopherhoole.com","\u23F3"],
     ["Dengro","Export access (postcode data)","Christopher to receive","\u23F3"]],
    cw=[1.2,1.0,2.2,0.6])

rect(sl,Inches(6.80),Inches(0.50),Inches(5.90),Inches(2.50),WHITE)
rect(sl,Inches(6.80),Inches(0.50),Inches(0.08),Inches(2.50),GREEN)
txt(sl,Inches(7.10),Inches(0.60),Inches(5.40),Inches(0.35),"Once Access Is Granted",14,NAVY)
txt(sl,Inches(7.10),Inches(1.00),Inches(5.40),Inches(1.80),
    "Week 1 implementation can begin immediately:\n\n"
    "\u2022 Remove KMG ghost conversions\n"
    "\u2022 Move soft signals to Secondary\n"
    "\u2022 Negative keyword 9-list overhaul\n"
    "\u2022 Block NHS/free seekers (account-level)\n"
    "\u2022 Add PMax Display exclusions\n"
    "\u2022 Geo restructure (radius \u2192 postcodes)\n"
    "\u2022 Device + schedule bid adjustments\n"
    "\u2022 Disable auto-apply",
    12,BLACK)

rect(sl,Inches(6.80),Inches(3.20),Inches(5.90),Inches(1.00),BLUE_TINT,BLUE,Pt(1))
txt(sl,Inches(7.05),Inches(3.25),Inches(5.40),Inches(0.90),
    "Contact: chris@christopherhoole.com\nchristopherhoole.com  |  +44 7999 500 184",
    12,NAVY,True)

bot_bar_title(sl)


# ══════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════
output=os.path.join(REPORTS,'Onboarding Doc 2 - Technical Onboarding Checklist.pptx')
prs.save(output)
print(f"Saved to: {output}")
print(f"Slides: {len(prs.slides)}")
print("Done!")
