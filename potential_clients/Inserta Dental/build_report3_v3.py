"""
Dental by Design - Report 3: Restructure & Growth Plan (v3)
v3: Geo restructure + ad schedules moved to Week 1. Landing pages added to Month 2.
Timeline image replaced with built elements.
"""
import os, sys, warnings
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

warnings.filterwarnings('ignore')

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
CHARTS = os.path.join(SCRIPT_DIR, 'charts')
REPORTS = os.path.join(SCRIPT_DIR, 'reports')
LOGO = os.path.join(SCRIPT_DIR, 'act_logo_official.png')
os.makedirs(CHARTS, exist_ok=True)

# ── Design tokens ──
BLUE=RGBColor(0x42,0x85,0xF4); RED=RGBColor(0xEA,0x43,0x35)
YELLOW=RGBColor(0xFB,0xBC,0x05); GREEN=RGBColor(0x34,0xA8,0x53)
NAVY=RGBColor(0x1A,0x23,0x7E); BLACK=RGBColor(0x1A,0x1A,0x1A)
GREY_BG=RGBColor(0xF5,0xF6,0xFA); WHITE=RGBColor(0xFF,0xFF,0xFF)
BORDER=RGBColor(0xE2,0xE8,0xF0)
BLUE_TINT=RGBColor(0xE8,0xF0,0xFE); RED_TINT=RGBColor(0xFC,0xE8,0xE6)
GREEN_TINT=RGBColor(0xE6,0xF4,0xEA); GREY_MID=RGBColor(0xCB,0xD5,0xE1)
FONT='Calibri'
CB='#4285F4'; CR='#EA4335'; CY='#FBBC05'; CG='#34A853'; CN='#1A237E'

plt.rcParams.update({'font.family':'Calibri','font.size':11,
    'axes.spines.top':False,'axes.spines.right':False,
    'axes.edgecolor':'#E2E8F0','axes.linewidth':0.8,
    'figure.facecolor':'white','axes.facecolor':'white',
    'grid.color':'#F0F0F0','grid.linewidth':0.5})

# ── Verified numbers (hardcoded) ──
S25=265869; V25=6070; CPA25=44; NBCPA25=46
S26=175467; V26=2889; CPA26=61; NBCPA26=64
NB26_SPEND=171334; NB26_CONV=2660
MONTHLY_RUN=53990
TARGET_CPA=50

# PMax 30-day verified data
PM = {
    'spend': 30026, 'conv': 461, 'cpa': 65,
    'channels': {
        'Google Search': {'spend': 24537, 'pct': 81.7, 'conv': 311, 'cpa': 79},
        'Search Partners': {'spend': 3358, 'pct': 11.2, 'conv': 141, 'cpa': 24},
        'Display Network': {'spend': 1893, 'pct': 6.3, 'conv': 2, 'cpa': 947},
        'YouTube': {'spend': 169, 'pct': 0.6, 'conv': 6, 'cpa': 28},
        'Maps': {'spend': 69, 'pct': 0.2, 'conv': 1, 'cpa': 69},
    },
    'bookings_30d': 72, 'real_booking_cpa': 418, 'neg_keywords': 93,
}

# Financing ad group benchmark (proven Search performance)
FINANCING_CPA = 34
FINANCING_CVR = 24.7

print("Building Report 3 v2...")

# ══════════════════════════════════════════════════════════════════
# CHARTS
# ══════════════════════════════════════════════════════════════════
print("Generating charts...")

# Chart 1: PMax vs Search comparison
fig,ax=plt.subplots(figsize=(8,4))
metrics=['Reported CPA','Real Booking\nCPA','Search Term\nWaste Rate','Negative\nKeywords']
pmax_vals=[65,418,69.3,93]
search_vals=[94,0,40.5,2192]  # Search Intent doesn't have booking data
# Use 2 separate bar groups
x=np.arange(4)
ax.bar(x-0.2,[65,418,0,93],0.35,label='PMax (Current)',color=CR,alpha=0.8)
ax.bar(x+0.2,[94,0,0,2192],0.35,label='Search Intent (Current)',color=CB,alpha=0.8)
# This doesn't work well as a bar chart — let me do a table instead
plt.close()

# Chart 2: Projected timeline
fig,ax=plt.subplots(figsize=(11,3))
phases=[
    ('Week 1\nFoundations','Tracking fix\nNeg keywords\nNHS/free block',CR),
    ('Week 2-3\nBuild','New Search campaigns\nAd rewrites\nBid adjustments',CY),
    ('Week 4\nLaunch','Launch Search\nShift 30% PMax budget\nMonitor',CB),
    ('Month 2\nOptimise','Scale Search winners\nReduce PMax further\nGeo restructure',CB),
    ('Month 3\nScale','Search dominant\nPMax supplementary\nMulti-clinic plan',CG),
]
for i,(label,desc,color) in enumerate(phases):
    ax.barh(0,1,left=i,color=color,height=0.6,edgecolor='white',linewidth=2)
    ax.text(i+0.5,0.05,label,ha='center',va='center',fontsize=9,color='white',fontweight='bold')
    ax.text(i+0.5,-0.6,desc,ha='center',va='center',fontsize=8,color=CN)
ax.set_xlim(-0.1,5.1); ax.set_ylim(-1.1,0.5)
ax.axis('off')
ax.set_title('90-Day PMax \u2192 Search Migration Roadmap',fontsize=12,color=CN,fontweight='bold',loc='left')
plt.tight_layout()
plt.savefig(f"{CHARTS}/r3v2_timeline.png",dpi=200,bbox_inches='tight'); plt.close()

# Chart 3: Proposed campaign structure
fig,ax=plt.subplots(figsize=(8,4))
camps=['Single Implants\n(Search)','Full Arch / All-on-4\n(Search)','Financing\n(Search)','Invisalign\n(Search)','Brand\n(Search)','PMax\n(Supplementary)']
target_cpas=[60,70,35,45,18,50]
colors_c=[CG,CG,CG,CG,CB,CY]
bars=ax.barh(range(len(camps)),target_cpas,color=colors_c,height=0.5)
ax.set_yticks(range(len(camps))); ax.set_yticklabels(camps,fontsize=10)
ax.invert_yaxis()
for bar,cpa in zip(bars,target_cpas):
    ax.text(bar.get_width()+1,bar.get_y()+bar.get_height()/2,f'\u00a3{cpa} target CPA',va='center',fontsize=10,color=CN,fontweight='bold')
ax.set_xlabel('Target CPA (\u00a3)',fontsize=10)
ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f'\u00a3{x:.0f}'))
ax.set_xlim(0,90)
ax.set_title('Proposed Campaign Structure with Target CPAs',fontsize=12,color=CN,fontweight='bold',loc='left')
plt.tight_layout()
plt.savefig(f"{CHARTS}/r3v2_structure.png",dpi=200,bbox_inches='tight'); plt.close()

# Chart 4: Device bid adjustments
fig,ax=plt.subplots(figsize=(6,3))
devices=['Mobile','Desktop','Tablet']
current_cpa=[56,92,73]
colors_d=[CG,CR,CY]
bars=ax.barh(range(3),current_cpa,color=colors_d,height=0.5)
ax.set_yticks(range(3)); ax.set_yticklabels(devices,fontsize=10); ax.invert_yaxis()
for bar,cpa,a in zip(bars,current_cpa,['+0%','-25% to -35%','-20% to -30%']):
    ax.text(bar.get_width()+1,bar.get_y()+bar.get_height()/2,f'\u00a3{cpa} \u2192 {a}',va='center',fontsize=10,color=CN)
ax.set_xlabel('Current CPA (\u00a3)',fontsize=10)
ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f'\u00a3{x:.0f}'))
ax.set_xlim(0,140)
ax.set_title('Proposed Device Bid Adjustments',fontsize=11,color=CN,fontweight='bold',loc='left')
plt.tight_layout()
plt.savefig(f"{CHARTS}/r3v2_device.png",dpi=200,bbox_inches='tight'); plt.close()

print("Charts done.")

# ══════════════════════════════════════════════════════════════════
# PPTX HELPERS (same as Reports 1+2)
# ══════════════════════════════════════════════════════════════════
prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)

def rect(sl,l,t,w,h,fc,lc=None,lw=None):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h);s.fill.solid();s.fill.fore_color.rgb=fc
    if lc:s.line.color.rgb=lc;s.line.width=lw or Pt(1)
    else:s.line.fill.background()
def txt(sl,l,t,w,h,c,sz=11,co=BLACK,b=False,a=PP_ALIGN.LEFT):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    r=tf.paragraphs[0].add_run();r.text=c;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
    tf.paragraphs[0].alignment=a;return bx
def multirun(sl,l,t,w,h,runs,a=PP_ALIGN.LEFT):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=a
    for tx,sz,co,b in runs:
        r=p.add_run();r.text=tx;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
def multipara(sl,l,t,w,h,pdata):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    for i,(tx,sz,co,b,sp) in enumerate(pdata):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run();r.text=tx;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
        if sp:p.space_after=Pt(sp)
def bullets(sl,l,t,w,h,items,sz=11,co=BLACK):
    multipara(sl,l,t,w,h,[(i,sz,co,False,5) for i in items])
def top_bar(sl,h=Inches(0.07)):
    for i,c in enumerate([BLUE,RED,YELLOW,GREEN]):rect(sl,Inches(i*3.333),Inches(0),Inches(3.333),h,c)
def bot_bar(sl,y=Inches(6.92)):
    for p,c in zip([Inches(0.60),Inches(3.63),Inches(6.67),Inches(9.70)],[BLUE,RED,YELLOW,GREEN]):rect(sl,p,y,Inches(3.03),Inches(0.03),c)
def bot_bar_title(sl):
    for i,c in enumerate([BLUE,RED,YELLOW,GREEN]):rect(sl,Inches(i*3.333),Inches(7.0),Inches(3.333),Inches(0.04),c)
def sidebar(sl): rect(sl,Inches(0),Inches(0.07),Inches(0.12),Inches(7.43),BLUE)
def logo(sl,l=Inches(0.60),t=Inches(0.50),s=Inches(0.65)):
    try:sl.shapes.add_picture(LOGO,l,t,s,s)
    except:pass
def footer(sl,n):
    bot_bar(sl)
    try:sl.shapes.add_picture(LOGO,Inches(0.60),Inches(7.0),Inches(0.22),Inches(0.22))
    except:pass
    txt(sl,Inches(0.90),Inches(7.0),Inches(6),Inches(0.25),"Christopher Hoole  |  christopherhoole.com  |  Confidential",11,NAVY,True)
    txt(sl,Inches(12.23),Inches(7.0),Inches(0.50),Inches(0.25),str(n),11,NAVY,a=PP_ALIGN.RIGHT)
def stitle(sl,t,sub=None):
    txt(sl,Inches(0.60),Inches(0.30),Inches(7),Inches(0.50),t,28,NAVY,True)
    if sub:txt(sl,Inches(0.60),Inches(0.85),Inches(9),Inches(0.30),sub,11,RED)
def badge(sl,t):
    rect(sl,Inches(9.13),Inches(0.30),Inches(3.60),Inches(0.45),WHITE,BLUE,Pt(1))
    txt(sl,Inches(9.23),Inches(0.32),Inches(3.40),Inches(0.40),t,11,BLUE,True,PP_ALIGN.CENTER)
def mcard(sl,l,t,v,lab,ac,w=Inches(2.80),h=Inches(1.05)):
    rect(sl,l,t,w,h,WHITE,BORDER,Pt(0.75));rect(sl,l,t,Inches(0.06),h,ac)
    txt(sl,l+Inches(0.20),t+Inches(0.10),w-Inches(0.30),Inches(0.50),v,22,ac,True)
    txt(sl,l+Inches(0.20),t+Inches(0.60),w-Inches(0.30),Inches(0.35),lab,11,BLACK)
def minicard(sl,l,t,v,lab,ac):
    rect(sl,l,t,Inches(1.83),Inches(1.15),WHITE);rect(sl,l,t,Inches(0.06),Inches(1.15),ac)
    txt(sl,l+Inches(0.15),t+Inches(0.12),Inches(1.58),Inches(0.45),v,22,ac,True)
    txt(sl,l+Inches(0.15),t+Inches(0.60),Inches(1.58),Inches(0.30),lab,11,BLACK)
def ibox(sl,l,t,w,h,ti,body,ac,bg=None):
    rect(sl,l,t,w,h,bg or GREY_BG);rect(sl,l,t,Inches(0.06),h,ac)
    txt(sl,l+Inches(0.25),t+Inches(0.10),w-Inches(0.40),Inches(0.30),ti,14,NAVY,True)
    txt(sl,l+Inches(0.25),t+Inches(0.40),w-Inches(0.40),h-Inches(0.50),body,11,BLACK)
def rbox(sl,l,t,w,h,c):
    rect(sl,l,t,w,h,BLUE_TINT,BLUE,Pt(1))
    txt(sl,l+Inches(0.25),t+Inches(0.05),w-Inches(0.50),h-Inches(0.10),c,11,NAVY,True)
def add_table(sl,l,t,w,h,hdrs,rows,cw=None):
    ts=sl.shapes.add_table(len(rows)+1,len(hdrs),l,t,w,h);tb=ts.table
    if cw:
        for i,c in enumerate(cw):tb.columns[i].width=Inches(c)
    for ci,hd in enumerate(hdrs):
        c=tb.cell(0,ci);c.text=hd;c.fill.solid();c.fill.fore_color.rgb=NAVY
        for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.font.color.rgb=WHITE;p.font.bold=True;p.font.name=FONT
    for ri,rd in enumerate(rows):
        for ci,v in enumerate(rd):
            c=tb.cell(ri+1,ci);c.text=str(v);c.fill.solid();c.fill.fore_color.rgb=GREY_BG if ri%2==0 else WHITE
            for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.font.color.rgb=BLACK;p.font.name=FONT
    return ts
def img(sl,path,l,t,w,h):
    try:sl.shapes.add_picture(path,l,t,w,h)
    except Exception as e:print(f"  Warning: {e}")
def ncircle(sl,l,t,num,col):
    ns=sl.shapes.add_shape(MSO_SHAPE.OVAL,l,t,Inches(0.45),Inches(0.45))
    ns.fill.solid();ns.fill.fore_color.rgb=col;ns.line.fill.background()
    tf=ns.text_frame;tf.paragraphs[0].alignment=PP_ALIGN.CENTER;tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    r=tf.paragraphs[0].add_run();r.text=str(num);r.font.size=Pt(16);r.font.color.rgb=WHITE;r.font.bold=True;r.font.name=FONT

# ══════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════
print("Building slides...")

# ── S1: TITLE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=GREY_BG
top_bar(sl);sidebar(sl);logo(sl)
txt(sl,Inches(0.60),Inches(1.40),Inches(5.74),Inches(1.80),"Restructure &\nGrowth Plan",44,NAVY,True)
rect(sl,Inches(0.60),Inches(3.30),Inches(2.50),Inches(0.05),BLUE)
txt(sl,Inches(0.60),Inches(3.55),Inches(5.50),Inches(0.50),"Dental by Design  |  Prodent Group",22,BLUE)
multipara(sl,Inches(0.60),Inches(4.30),Inches(5.50),Inches(1.00),[
    ("Christopher Hoole",11,BLACK,True,2),("Google Ads Specialist  |  April 2026",11,BLACK,False,2),
    ("christopherhoole.com",11,BLUE,False,0)])
rect(sl,Inches(0.60),Inches(5.50),Inches(4.50),Inches(0.50),WHITE,BLUE,Pt(1))
txt(sl,Inches(0.70),Inches(5.52),Inches(4.30),Inches(0.45),"Based on data: 2025 + 2026 + PMax 30-day deep dive",11,BLUE,True)

rect(sl,Inches(6.80),Inches(0.50),Inches(5.90),Inches(2.50),WHITE)
rect(sl,Inches(6.80),Inches(0.50),Inches(0.08),Inches(2.50),GREEN)
txt(sl,Inches(7.10),Inches(0.60),Inches(5.40),Inches(0.35),"Target Non-Brand CPA",14,NAVY)
txt(sl,Inches(7.10),Inches(1.00),Inches(5.40),Inches(0.80),f"\u00a3{TARGET_CPA}",44,GREEN,True)
txt(sl,Inches(7.10),Inches(1.80),Inches(5.40),Inches(0.50),f"Down from \u00a3{NBCPA26} \u2014 {(NBCPA26-TARGET_CPA)/NBCPA26*100:.0f}% reduction within 90 days via PMax \u2192 Search migration",11,BLACK)

minicard(sl,Inches(6.80),Inches(3.20),"+29%","More Leads (Same Budget)",GREEN)
minicard(sl,Inches(8.83),Inches(3.20),"90 Days","PMax \u2192 Search Migration",BLUE)
minicard(sl,Inches(10.87),Inches(3.20),f"\u00a3{12000:,}/mo","Potential Monthly Saving",YELLOW)
bot_bar_title(sl)


# ── S2: EXECUTIVE SUMMARY ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06));stitle(sl,"Executive Summary");badge(sl,"Restructure Plan")

mcard(sl,Inches(0.60),Inches(1.20),f"\u00a3{TARGET_CPA}","Target Non-Brand CPA",GREEN)
mcard(sl,Inches(3.70),Inches(1.20),"+236/mo","Additional Leads (Same Budget)",BLUE)
mcard(sl,Inches(6.80),Inches(1.20),"PMax\u2192Search","Core Strategy Change",RED)
mcard(sl,Inches(9.90),Inches(1.20),f"\u00a3{MONTHLY_RUN:,.0f}","Current Monthly Spend",NAVY)

txt(sl,Inches(0.60),Inches(2.60),Inches(6),Inches(0.35),"The Plan",14,NAVY,True)
bullets(sl,Inches(0.70),Inches(3.00),Inches(11.50),Inches(2.60),[
    "Week 1: Fix conversion tracking, block NHS/free seekers, Display exclusions, negative keyword 9-list overhaul, geo restructure (radius \u2192 postcodes), ad schedule + device bid adjustments.",
    "Weeks 2-3: Build new treatment-specific Search campaigns (Single Implants, Full Arch, Financing, Invisalign) with proper ad copy, negatives, and tCPA bidding.",
    "Week 4: Launch new Search campaigns. Shift 30% of PMax budget. Monitor and compare real booking rates.",
    "Month 2: Scale Search winners. Reduce PMax further. Build dedicated landing pages per treatment type. Keyword optimization from first month of clean data.",
    "Month 3: Search dominant. PMax supplementary only. Per-treatment benchmarks established. Multi-clinic rollout planning.",
])

rbox(sl,Inches(0.60),Inches(5.85),Inches(12.13),Inches(0.70),
    f"The PMax deep-dive revealed the \u201Clow CPA\u201D is an illusion \u2014 real booking CPA is \u00a3{PM['real_booking_cpa']}. "
    f"The Financing ad group already proves Search can deliver \u00a3{FINANCING_CPA} CPA at {FINANCING_CVR}% CVR. The data supports this migration.")
footer(sl,2)


# ── S3: WHY PMax → SEARCH ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Why PMax \u2192 Search","The data from Reports 1 and 2 makes the case for migrating to structured Search campaigns")
badge(sl,"Strategic Decision")

# PMax problems (left)
ibox(sl,Inches(0.60),Inches(1.30),Inches(5.80),Inches(5.10),
    "PMax: The Problems (Last 30 Days)",
    f"\u2022 81.7% of spend is Google Search anyway \u2014 PMax isn't adding multi-channel value\n\n"
    f"\u2022 Display burns \u00a3{PM['channels']['Display Network']['spend']:,}/month on mobile games (Block Blast, Colour by Number) for 2 conversions\n\n"
    f"\u2022 Top converting terms are NHS/free seekers who never become patients\n\n"
    f"\u2022 Real booking CPA: \u00a3{PM['real_booking_cpa']} (only {PM['bookings_30d']} bookings from \u00a3{PM['spend']:,})\n\n"
    f"\u2022 Only {PM['neg_keywords']} negative keywords \u2014 vs 2,192 across the account\n\n"
    "\u2022 CPA doubled from \u00a326 (2025) to \u00a365 (last 30 days) as it scaled\n\n"
    "\u2022 Single asset group \u2014 no treatment segmentation",
    RED,RED_TINT)

# Search opportunity (right)
ibox(sl,Inches(6.80),Inches(1.30),Inches(5.80),Inches(5.10),
    "Search: The Opportunity",
    f"\u2022 The Financing ad group already delivers \u00a3{FINANCING_CPA} CPA at {FINANCING_CVR}% CVR \u2014 proof that well-structured Search works\n\n"
    "\u2022 Full control over keywords, match types, and negative keywords\n\n"
    "\u2022 Full visibility into every search term \u2014 no hidden waste\n\n"
    "\u2022 Treatment-specific ad copy matched to specific search intent\n\n"
    "\u2022 Device, schedule, and geo bid adjustments at campaign level\n\n"
    "\u2022 Route traffic to the right landing page per keyword intent\n\n"
    "\u2022 tCPA bidding gives direct control over cost per lead\n\n"
    "\u2022 Structured campaigns scale naturally to partner clinics",
    GREEN,GREEN_TINT)
footer(sl,3)


# ── S4: PROPOSED CAMPAIGN STRUCTURE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Proposed Campaign Structure","Treatment-specific Search campaigns with tCPA bidding")
badge(sl,"New Structure")

img(sl,f"{CHARTS}/r3v2_structure.png",Inches(0.60),Inches(1.30),Inches(7.50),Inches(3.50))

# Campaign detail table
add_table(sl,Inches(8.40),Inches(1.30),Inches(4.30),Inches(3.50),
    ["Campaign","Target CPA","Landing Page","Basis"],
    [["Single Implants","\u00a360","/dental-implants-offer","Best LP at \u00a334 CPA"],
     ["Full Arch / All-on-4","\u00a370","/all-on-4-offer","High value, higher CPA expected"],
     ["Financing","\u00a335","/dental-implants-offer","Proven: \u00a334 CPA, 24.7% CVR"],
     ["Invisalign","\u00a345","/invisalign","Currently \u00a345 CPA"],
     ["Brand","\u00a318","(homepage)","Currently \u00a318 \u2014 no change"],
     ["PMax (reduced)","\u00a350","Various","Supplementary only"]],
    cw=[1.2,0.7,1.2,1.2])

rbox(sl,Inches(0.60),Inches(5.10),Inches(12.10),Inches(0.60),
    "Each campaign gets proper negative keywords, treatment-specific ads, matched landing pages, and tCPA bidding. This is what PMax can't do.")

txt(sl,Inches(0.60),Inches(5.90),Inches(12.10),Inches(0.50),
    "Target CPAs are conservative starting points. The Financing ad group at \u00a334 CPA proves these targets are achievable with the right structure.",
    10,GREY_MID)
footer(sl,4)


# ── S5: CONVERSION TRACKING FIX ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Step 1: Fix Conversion Tracking","Remove ghosts, move soft signals, replace arbitrary ROAS values")
badge(sl,"Week 1")

add_table(sl,Inches(0.60),Inches(1.30),Inches(12.10),Inches(2.80),
    ["Conversion Action","Current","Proposed","Why"],
    [["KMG OLD LP | Book A Free Consult","Primary","Remove","Ghost \u2014 1,250 phantom conversions (Giulio confirmed)"],
     ["KMG OLD LP | Website Phone Call","Primary","Remove","Ghost \u2014 57 phantom conversions"],
     ["KMG | Call Extension","Primary","Remove","Legacy agency \u2014 204 phantom conversions"],
     ["Online Booking Click","Primary","Secondary","Click, not a booking \u2014 264 soft signals"],
     ["Phone Click / WhatsApp Click / Chat","Primary","Secondary","Clicks, not actions \u2014 181 soft signals"],
     ["Dengro Offline Booking","Primary (L300)","Keep (real value)","Replace L300 with actual treatment value"],
     ["Dengro Offline Lead","Primary","Keep (verify)","Verify only ad-sourced leads are counted"],
     ["Calls from the ads","Primary","Keep","Real calls \u2014 965 legitimate conversions"]],
    cw=[2.5,0.8,1.0,5.0])

ibox(sl,Inches(0.60),Inches(4.30),Inches(5.80),Inches(2.10),
    "Bid Strategy Recommendation",
    "Switch from Target ROAS to Target CPA on all Search campaigns.\n\n"
    "PMax: Start tCPA at \u00a350\n"
    "Search Intent: Start tCPA at \u00a365, reduce to \u00a350 over 8 weeks\n"
    "New campaigns: tCPA per treatment type (see slide 4)\n\n"
    "Target CPA gives direct control over cost per lead. The arbitrary \u00a3300 ROAS values have been misleading the algorithm.",
    GREEN,GREEN_TINT)

ibox(sl,Inches(6.80),Inches(4.30),Inches(5.80),Inches(2.10),
    "ROAS Alternative (If Preferred)",
    "If Tommaso prefers ROAS, replace arbitrary \u00a3300 with actual treatment revenue:\n\n"
    "\u2022 Single implant booking: \u00a31,695\n"
    "\u2022 All-on-4 booking: \u00a39,995\n"
    "\u2022 Vivo Bridge booking: \u00a315,990\n\n"
    "This lets the algorithm bid more for high-value treatments. Requires reliable Dengro attribution.",
    BLUE,BLUE_TINT)
footer(sl,5)


# ── S6: NEGATIVE KEYWORD OVERHAUL ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Step 2: Negative Keyword Overhaul","9-list system + immediate NHS/free/Display blocks for PMax")
badge(sl,"Week 1")

ibox(sl,Inches(0.60),Inches(1.30),Inches(5.80),Inches(2.00),
    "Immediate PMax Blocks (Day 1)",
    "Add account-level exact match negatives for:\n\n"
    "\u2022 [how to apply for dental implants on the nhs] (31 junk conversions/month)\n"
    "\u2022 [free dental implants near me] (30 junk conversions/month)\n"
    "\u2022 [how to get free dental implants on nhs] (14 junk conversions/month)\n"
    "\u2022 All 'nhs', 'free implant', 'clinical trial' variations\n\n"
    "These are PMax's top converting terms but they generate zero bookings.",
    RED,RED_TINT)

ibox(sl,Inches(6.80),Inches(1.30),Inches(5.80),Inches(2.00),
    "9-List Negative Keyword System",
    "Replace 784 broad match negatives with structured lists:\n\n"
    "\u2022 1 word \u2014 phrase + exact\n"
    "\u2022 2 words \u2014 phrase + exact\n"
    "\u2022 3 words \u2014 phrase + exact\n"
    "\u2022 4 words \u2014 exact only\n"
    "\u2022 5+ words \u2014 exact only\n\n"
    "No broad match. Weekly search term reviews.",
    GREEN,GREEN_TINT)

txt(sl,Inches(0.60),Inches(3.60),Inches(12),Inches(0.30),"Priority Negative Keyword Categories (Both Years Data)",14,NAVY,True)
add_table(sl,Inches(0.60),Inches(3.90),Inches(12.10),Inches(1.70),
    ["Category","Examples","2025 Waste","2026 Waste","Match Type","Priority"],
    [["NHS / Free seekers","nhs dental implants, free implants, clinical trials","\u00a31,378","\u00a3706 + junk conv","Phrase + Exact","Immediate"],
     ["Price Research (no location)","implant cost uk, how much is dental implants","\u00a320,540","\u00a314,211","Phrase","Week 1"],
     ["Wrong Treatment","implant dentures, dental bridge, dental crown","\u00a33,112","\u00a33,661","Phrase","Week 1"],
     ["Competitor / Tourism","brighton implant clinic, turkey teeth","\u00a31,724","\u00a31,390","Exact","Week 1"]],
    cw=[1.8,3.0,0.9,1.2,1.0,0.9])

rbox(sl,Inches(0.60),Inches(5.80),Inches(12.10),Inches(0.50),
    "Also: Add Display placement exclusions to PMax \u2014 remove mobile game apps, puzzle games, and irrelevant sites. This stops the \u00a31,893/month Display waste immediately.")
footer(sl,6)


# ── S7: AD COPY REWRITE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Step 3: Ad Copy Rewrite","90% of ads rated Poor/Average \u2014 the #1 lever for reducing CPA")
badge(sl,"Weeks 2-3")

txt(sl,Inches(0.60),Inches(1.20),Inches(12),Inches(0.30),"Proposed Ad Testing Framework (New Search Campaigns)",14,NAVY,True)
add_table(sl,Inches(0.60),Inches(1.55),Inches(12.10),Inches(2.20),
    ["Campaign","Angle 1 (Price)","Angle 2 (Speed)","Angle 3 (Quality)","Landing Page"],
    [["Single Implants","From \u00a31,695. 0% Finance Available.","Permanent Teeth in 48 Hours.","4.8\u2605 Rating. 399+ Reviews.","/dental-implants-offer"],
     ["Full Arch / All-on-4","Full Arch From \u00a39,995.","New Teeth in 48 Hours. Walk Out Smiling.","CQC Registered. Expert Surgeons.","/all-on-4-offer"],
     ["Financing","0% Finance. From \u00a347/month.","Book Free Consultation Today.","Trusted by 399+ Patients.","/dental-implants-offer"],
     ["Invisalign","Invisalign From \u00a32,500.","Free Scan. Results in 6 Months.","Diamond Apex Provider.","/invisalign"]],
    cw=[1.5,2.5,2.5,2.5,1.8])

ibox(sl,Inches(0.60),Inches(4.00),Inches(5.80),Inches(2.40),
    "Testing Plan",
    "2-3 RSA variants per ad group, each testing a different messaging angle.\n\n"
    "Run for 4-6 weeks, then keep winner and test new challenger.\n\n"
    "Every headline must match the keyword intent in that ad group. "
    "This is what the current ads fail to do \u2014 71% of keywords have Below Average ad relevance.\n\n"
    "Key USPs to leverage: \u201Cup to 60% less\u201D, Vivo Bridge 48 hours, 4.8\u2605 / 399+ reviews, 0% finance.",
    BLUE,BLUE_TINT)

ibox(sl,Inches(6.80),Inches(4.00),Inches(5.80),Inches(2.40),
    "Expected Impact",
    "Improving ad relevance from Below Average to Average on the top 20 keywords "
    "would lift QS by 2-3 points.\n\n"
    f"At QS 7 vs QS 3, CPC drops by an estimated 30-50%. "
    f"On \u00a3{MONTHLY_RUN:,}/month spend, that\u2019s \u00a316-27K/month in CPC savings.\n\n"
    f"Combined with the Financing ad group\u2019s proven \u00a3{FINANCING_CPA} CPA, "
    "well-written treatment-specific ads should bring the Search campaign CPA below \u00a360.",
    GREEN,GREEN_TINT)
footer(sl,7)


# ── S8: BID ADJUSTMENTS ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Step 4: Geo, Device & Schedule Adjustments","Quick wins that start building data immediately at this budget level")
badge(sl,"Week 1")

img(sl,f"{CHARTS}/r3v2_device.png",Inches(0.60),Inches(1.30),Inches(5.50),Inches(2.50))

txt(sl,Inches(6.50),Inches(1.20),Inches(6.20),Inches(0.30),"Day of Week (2026 Data)",14,NAVY,True)
add_table(sl,Inches(6.50),Inches(1.55),Inches(6.20),Inches(2.20),
    ["Day","2025 CPA","2026 CPA","Proposed"],
    [["Monday","\u00a341","\u00a364","-5%"],
     ["Thursday","\u00a344","\u00a358","+5%"],
     ["Friday","\u00a346","\u00a356","+10%"],
     ["Saturday","\u00a343","\u00a356","+10%"],
     ["Sunday","\u00a345","\u00a368","-15%"]],
    cw=[1.2,0.8,0.8,0.8])

ibox(sl,Inches(0.60),Inches(4.10),Inches(5.80),Inches(2.30),
    "Hour of Day",
    "Increase bids +10-15%:\n"
    "Hours 16-17 (4-5pm): \u00a346-48 CPA both years. High-volume, efficient.\n\n"
    "Reduce bids -15-20%:\n"
    "Hours 7-8 (7-8am): \u00a356-74 CPA. Expensive commute clicks.\n\n"
    "These adjustments apply to Search campaigns. PMax handles its own scheduling.",
    BLUE,BLUE_TINT)

ibox(sl,Inches(6.80),Inches(4.10),Inches(5.80),Inches(2.30),
    "Geo: Radius \u2192 Postcodes",
    "Replace radius targeting with London postcodes:\n\n"
    "\u2022 Bid +20% for W6, W12, W3 (near clinic, \u00a328-53 CPA)\n"
    "\u2022 No change for mid-range postcodes\n"
    "\u2022 Bid -20% for N13 \u00a3169, TW2 \u00a3137, WD23 \u00a3123\n\n"
    "All campaigns share the same postcode targets. Unified geographic view across the account.",
    GREEN,GREEN_TINT)
footer(sl,8)


# ── S9: MULTI-CLINIC SCALE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Multi-Clinic Growth","The optimised Dental by Design structure becomes the blueprint for partner clinics")
badge(sl,"Month 3+")

rect(sl,Inches(0.60),Inches(1.30),Inches(5.00),Inches(2.80),WHITE,BORDER,Pt(0.75))
rect(sl,Inches(0.60),Inches(1.30),Inches(0.08),Inches(2.80),GREEN)
txt(sl,Inches(0.90),Inches(1.45),Inches(4.50),Inches(0.40),"Dental by Design",22,GREEN,True)
txt(sl,Inches(0.90),Inches(1.90),Inches(4.50),Inches(0.25),"Flagship \u2014 Testing Ground",14,NAVY,True)
bullets(sl,Inches(0.90),Inches(2.25),Inches(4.50),Inches(1.50),[
    "Treatment-specific Search campaigns (proven model)",
    "Establish CPA benchmarks per treatment type",
    "Test ad copy angles \u2014 winning angles become the playbook",
    "Postcode-level geo data informs partner clinic targeting",
],11)

txt(sl,Inches(5.80),Inches(2.40),Inches(0.80),Inches(0.50),"\u25B6",28,BLUE,a=PP_ALIGN.CENTER)

rect(sl,Inches(6.80),Inches(1.30),Inches(5.80),Inches(2.80),WHITE,BORDER,Pt(0.75))
rect(sl,Inches(6.80),Inches(1.30),Inches(0.08),Inches(2.80),BLUE)
txt(sl,Inches(7.10),Inches(1.45),Inches(5.20),Inches(0.40),"Partner Clinics",22,BLUE,True)
txt(sl,Inches(7.10),Inches(1.90),Inches(5.20),Inches(0.25),"UK Network \u2014 Scaled Learnings",14,NAVY,True)
bullets(sl,Inches(7.10),Inches(2.25),Inches(5.20),Inches(1.50),[
    "Same campaign structure \u2014 localised per clinic",
    "Shared negative keyword lists (9-list system)",
    "Performance-based budgets \u2014 best clinics earn more spend",
    "Postcode targeting per clinic location",
],11)

txt(sl,Inches(0.60),Inches(4.40),Inches(12),Inches(0.30),"What Scales from Flagship to Network",14,NAVY,True)
mcard(sl,Inches(0.60),Inches(4.80),"9-List","Neg Keyword Framework",GREEN,Inches(2.40))
mcard(sl,Inches(3.20),Inches(4.80),"Ad Copy","Winning Angles",BLUE,Inches(2.40))
mcard(sl,Inches(5.80),Inches(4.80),"tCPA","Benchmarks per Treatment",YELLOW,Inches(2.40))
mcard(sl,Inches(8.40),Inches(4.80),"Reporting","Monthly Framework",GREEN,Inches(2.20))
mcard(sl,Inches(10.80),Inches(4.80),"Postcodes","Geo Structure",BLUE,Inches(1.90))
footer(sl,9)


# ── S10: PROJECTED RESULTS ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Projected Results","Conservative projections based on Reports 1, 2, and PMax deep-dive")
badge(sl,"90-Day Target")

mcard(sl,Inches(0.60),Inches(1.20),f"\u00a3{NBCPA26}\u2192\u00a3{TARGET_CPA}","Non-Brand CPA Target",GREEN,Inches(3.00))
mcard(sl,Inches(3.90),Inches(1.20),"+236/mo","Additional Leads (Same Budget)",BLUE,Inches(3.00))
mcard(sl,Inches(7.20),Inches(1.20),"\u00a312K/mo","Potential Monthly Saving",YELLOW,Inches(3.00))
mcard(sl,Inches(10.50),Inches(1.20),f"\u00a3{PM['real_booking_cpa']}\u2192?","Real Booking CPA (Improve)",RED,Inches(2.20))

ibox(sl,Inches(0.60),Inches(2.60),Inches(5.80),Inches(2.00),
    "Scenario 1: More Leads, Same Spend",
    f"Keep spend at \u00a3{MONTHLY_RUN:,.0f}/month.\n"
    f"Reduce non-brand CPA from \u00a3{NBCPA26} to \u00a3{TARGET_CPA}.\n"
    f"Result: ~236 additional patient leads per month (+29%).\n\n"
    "More importantly: these will be QUALIFIED leads (not NHS/free seekers) because the Search campaigns target paying patients specifically.",
    GREEN,GREEN_TINT)

ibox(sl,Inches(6.80),Inches(2.60),Inches(5.80),Inches(2.00),
    "Scenario 2: Same Leads, Less Spend",
    f"Maintain current lead volume.\n"
    f"Reduce CPA from \u00a3{NBCPA26} to \u00a3{TARGET_CPA}.\n"
    f"Result: \u00a312,000/month saving (\u00a3144K/year).\n\n"
    "The savings come from eliminating Display waste, blocking junk leads, and better keyword targeting.",
    BLUE,BLUE_TINT)

ibox(sl,Inches(0.60),Inches(4.85),Inches(12.10),Inches(1.60),
    "Why These Projections Are Achievable",
    f"\u2022 The Financing ad group already delivers \u00a3{FINANCING_CPA} CPA at {FINANCING_CVR}% CVR \u2014 with proper structure, other ad groups can approach this\n"
    f"\u2022 2025 non-brand CPA was \u00a3{NBCPA25} \u2014 the account has achieved below \u00a350 before\n"
    f"\u2022 \u00a3{PM['channels']['Display Network']['spend']:,}/month in Display waste is immediately recoverable\n"
    f"\u2022 Blocking NHS/free seekers stops the algorithm chasing junk conversions\n"
    "\u2022 Ad rewrites improving QS from 3 to 7 reduces CPCs by an estimated 30-50%",
    BLUE,BLUE_TINT)
footer(sl,10)


# ── S11: IMPLEMENTATION TIMELINE (built elements, not image) ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Implementation Timeline")
badge(sl,"90-Day Roadmap")

# Timeline bar built with elements (editable, not image)
timeline_phases=[
    ("Week 1",RED),("Weeks 2-3",YELLOW),("Week 4 + Month 2",BLUE),("Month 3",GREEN),
]
bar_w=Inches(2.90)
for i,(label,color) in enumerate(timeline_phases):
    cx=Inches(0.60+i*3.10)
    rect(sl,cx,Inches(1.20),bar_w,Inches(0.50),color)
    txt(sl,cx,Inches(1.22),bar_w,Inches(0.50),label,12,WHITE,True,PP_ALIGN.CENTER)
# Connecting line
rect(sl,Inches(0.60),Inches(1.68),Inches(12.40),Pt(2),BORDER)
# Arrow dots
for i in range(4):
    cx=Inches(0.60+i*3.10)+bar_w/2-Inches(0.10)
    circ=sl.shapes.add_shape(MSO_SHAPE.OVAL,cx,Inches(1.60),Inches(0.20),Inches(0.20))
    circ.fill.solid();circ.fill.fore_color.rgb=timeline_phases[i][1];circ.line.fill.background()

# Subtitle labels under timeline
subtitles=["Fix Foundations","Build + Rewrite","Launch + Optimise","Scale"]
for i,sub in enumerate(subtitles):
    cx=Inches(0.60+i*3.10)
    txt(sl,cx,Inches(1.85),bar_w,Inches(0.30),sub,11,NAVY,True,PP_ALIGN.CENTER)

phases=[
    ("Week 1","Fix Foundations",RED,[
        "Remove KMG ghost conversions",
        "Move soft signals to Secondary",
        "Block NHS/free seekers (account-level)",
        "Add PMax Display exclusions",
        "Negative keyword 9-list overhaul",
        "Geo: radius \u2192 postcodes",
        "Device + schedule bid adjustments",
        "Disable auto-apply",
    ]),
    ("Weeks 2-3","Build + Rewrite",YELLOW,[
        "Build 4 new Search campaigns",
        "Write treatment-specific RSAs",
        "Set tCPA targets per campaign",
        "Keyword strategy per ad group",
        "Landing page routing plan",
    ]),
    ("Week 4 + Month 2","Launch + Optimise",BLUE,[
        "Launch Search campaigns",
        "Shift 30% of PMax budget",
        "Monitor real booking rates",
        "Build dedicated landing pages",
        "Scale Search winners",
        "Keyword optimization (clean data)",
    ]),
    ("Month 3","Scale",GREEN,[
        "Search becomes dominant",
        "PMax \u2192 supplementary only",
        "Per-treatment benchmarks",
        "Multi-clinic planning",
        "Full reporting framework",
    ]),
]

for i,(label,title,color,items) in enumerate(phases):
    cx=Inches(0.60+i*3.10)
    rect(sl,cx,Inches(2.25),Inches(2.90),Inches(4.30),WHITE,BORDER,Pt(0.75))
    rect(sl,cx,Inches(2.25),Inches(2.90),Pt(4),color)
    rect(sl,cx,Inches(2.25),Inches(0.06),Inches(4.30),color)
    txt(sl,cx+Inches(0.15),Inches(2.35),Inches(2.60),Inches(0.25),label,11,color,True)
    txt(sl,cx+Inches(0.15),Inches(2.60),Inches(2.60),Inches(0.30),title,14,NAVY,True)
    bullets(sl,cx+Inches(0.15),Inches(3.00),Inches(2.60),Inches(3.30),items,9)
footer(sl,11)


# ── S12: NEXT STEPS ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=GREY_BG
top_bar(sl);sidebar(sl);logo(sl)
txt(sl,Inches(0.60),Inches(1.40),Inches(5.50),Inches(0.80),"Next Steps",44,NAVY,True)
rect(sl,Inches(0.60),Inches(2.20),Inches(2.50),Inches(0.05),BLUE)

steps=[
    ("1","Grant me editor access to the account","I can implement the Week 1 fixes in a single session",GREEN),
    ("2","Week 1: Tracking fix + negatives + geo + schedules + device bids","All foundational changes in one session \u2014 data starts building immediately",BLUE),
    ("3","New Search campaigns launched by end of Month 1","Treatment-specific structure with proper ads, negatives, tCPA",YELLOW),
    ("4","Monthly reporting from Day 1","Full transparency on what changed, why, and the measured impact",RED),
]

for i,(num,action,detail,color) in enumerate(steps):
    y=Inches(2.55+i*0.85)
    rect(sl,Inches(0.60),y,Inches(5.50),Inches(0.70),WHITE)
    rect(sl,Inches(0.60),y,Inches(0.06),Inches(0.70),color)
    ncircle(sl,Inches(0.80),y+Inches(0.10),num,color)
    txt(sl,Inches(1.40),y+Inches(0.05),Inches(4.50),Inches(0.30),action,12,NAVY,True)
    txt(sl,Inches(1.40),y+Inches(0.35),Inches(4.50),Inches(0.30),detail,10,BLACK)

rect(sl,Inches(6.80),Inches(0.50),Inches(5.90),Inches(2.50),WHITE)
rect(sl,Inches(6.80),Inches(0.50),Inches(0.08),Inches(2.50),BLUE)
txt(sl,Inches(7.10),Inches(0.60),Inches(5.40),Inches(0.35),"Ready to Start Immediately",14,NAVY)
txt(sl,Inches(7.10),Inches(1.00),Inches(5.40),Inches(0.80),"Christopher Hoole",40,BLUE,True)
multipara(sl,Inches(7.10),Inches(1.80),Inches(5.40),Inches(1.00),[
    ("chris@christopherhoole.com",14,BLUE,False,4),
    ("christopherhoole.com",14,BLUE,False,4),
    ("+44 7999 500 184",14,BLACK,False,0)])

minicard(sl,Inches(6.80),Inches(3.20),f"\u00a3{TARGET_CPA}","Target CPA",GREEN)
minicard(sl,Inches(8.83),Inches(3.20),"PMax\u2192Search","90-Day Migration",BLUE)
minicard(sl,Inches(10.87),Inches(3.20),"+29%","More Leads",RED)

rect(sl,Inches(6.80),Inches(4.60),Inches(5.90),Inches(1.20),BLUE_TINT,BLUE,Pt(1))
txt(sl,Inches(7.05),Inches(4.65),Inches(5.40),Inches(1.10),
    "\u201CThe data is clear. PMax is attracting the wrong patients and wasting budget on mobile games. "
    "Structured Search campaigns with proper targeting will deliver more real patient bookings "
    "for the same budget. I\u2019m ready to implement Week 1 the day I start.\u201D\n\n"
    "\u2014 Christopher Hoole",11,NAVY,True)
bot_bar_title(sl)


# ══════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════
output=os.path.join(REPORTS,'03_restructure_report_v3.pptx')
prs.save(output)
print(f"\nSaved to: {output}")
print(f"Slides: {len(prs.slides)}")
print("Done!")
