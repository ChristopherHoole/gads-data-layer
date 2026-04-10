"""
Christopher Hoole — Dental by Design Interview Presentation v2
Design system: exact match to Objection Experts winning decks
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ══════════════════════════════════════════════════════════════════
# DESIGN SYSTEM — exact colours from Objection Experts decks
# ══════════════════════════════════════════════════════════════════
BLUE    = RGBColor(0x42, 0x85, 0xF4)   # Google Blue — primary accent
RED     = RGBColor(0xEA, 0x43, 0x35)   # Google Red — warnings/negative
YELLOW  = RGBColor(0xFB, 0xBC, 0x05)   # Google Yellow — secondary
GREEN   = RGBColor(0x34, 0xA8, 0x53)   # Google Green — positive
NAVY    = RGBColor(0x1A, 0x23, 0x7E)   # Deep navy — titles, headers
BLACK   = RGBColor(0x1A, 0x1A, 0x1A)   # Near-black — body text
GREY_BG = RGBColor(0xF5, 0xF6, 0xFA)   # Light grey — backgrounds
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BORDER  = RGBColor(0xE2, 0xE8, 0xF0)   # Card border
BLUE_TINT = RGBColor(0xE8, 0xF0, 0xFE) # Light blue — recommendation boxes
RED_TINT  = RGBColor(0xFC, 0xE8, 0xE6) # Light red — warning boxes
GREEN_TINT = RGBColor(0xE6, 0xF4, 0xEA) # Light green — success boxes
GREY_MID  = RGBColor(0xCB, 0xD5, 0xE1) # Mid grey

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT = 'Calibri'

# Paths
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.normpath(os.path.join(SCRIPT_DIR, '..', '..'))
ACT_LOGO = os.path.join(PROJECT_ROOT, 'potential_clients', 'objection_experts', 'charts', 'act_logo.png')
ACT_SCREENSHOT = os.path.join(PROJECT_ROOT, 'act_dashboard', 'prototypes', 'screenshots', 'v6-light-full.png')

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ══════════════════════════════════════════════════════════════════
# REUSABLE LAYOUT COMPONENTS
# ══════════════════════════════════════════════════════════════════

def rect(slide, left, top, w, h, fill_color, line_color=None, line_width=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def text(slide, left, top, w, h, content, size=11, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box with a single paragraph."""
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = content
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT
    p.alignment = align
    return box


def multirun(slide, left, top, w, h, runs_data, align=PP_ALIGN.LEFT):
    """Add a text box with multiple styled runs in a single paragraph.
    runs_data = [(text, size, color, bold), ...]
    """
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for txt, sz, col, bld in runs_data:
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.color.rgb = col
        run.font.bold = bld
        run.font.name = FONT
    return box


def multipara(slide, left, top, w, h, paras_data):
    """Add text box with multiple paragraphs.
    paras_data = [(text, size, color, bold, space_after_pt), ...]
    """
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (txt, sz, col, bld, sp) in enumerate(paras_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.color.rgb = col
        run.font.bold = bld
        run.font.name = FONT
        if sp:
            p.space_after = Pt(sp)
    return box


def google_top_bar(slide, height=Inches(0.07)):
    """4-colour Google bar across top."""
    seg_w = Inches(3.333)
    for i, col in enumerate([BLUE, RED, YELLOW, GREEN]):
        rect(slide, Inches(i * 3.333), Inches(0), seg_w, height, col)


def google_bottom_bar(slide, y=Inches(6.92), height=Inches(0.03)):
    """4-colour Google bar near bottom (content slides)."""
    seg_w = Inches(3.03)
    positions = [Inches(0.60), Inches(3.63), Inches(6.67), Inches(9.70)]
    for pos, col in zip(positions, [BLUE, RED, YELLOW, GREEN]):
        rect(slide, pos, y, seg_w, height, col)


def google_bottom_bar_title(slide, y=Inches(7.0), height=Inches(0.04)):
    """4-colour Google bar at very bottom (title slides)."""
    seg_w = Inches(3.333)
    for i, col in enumerate([BLUE, RED, YELLOW, GREEN]):
        rect(slide, Inches(i * 3.333), y, seg_w, height, col)


def left_sidebar(slide):
    """Blue vertical sidebar on left edge (title slides)."""
    rect(slide, Inches(0), Inches(0.07), Inches(0.12), Inches(7.43), BLUE)


def add_logo(slide, left=Inches(0.60), top=Inches(0.50), size=Inches(0.65)):
    """Add ACT logo."""
    try:
        slide.shapes.add_picture(ACT_LOGO, left, top, size, size)
    except Exception:
        pass  # Skip if logo not found


def footer(slide, page_num):
    """Standard footer: 4-colour bar + logo + text + page number."""
    google_bottom_bar(slide)
    # Logo
    try:
        slide.shapes.add_picture(ACT_LOGO, Inches(0.60), Inches(7.0), Inches(0.22), Inches(0.22))
    except Exception:
        pass
    text(slide, Inches(0.90), Inches(7.0), Inches(6), Inches(0.25),
         "Christopher Hoole  |  christopherhoole.com  |  Confidential",
         size=11, color=NAVY, bold=True)
    text(slide, Inches(12.23), Inches(7.0), Inches(0.50), Inches(0.25),
         str(page_num), size=11, color=NAVY, align=PP_ALIGN.RIGHT)


def slide_title(slide, title_text, subtitle_text=None):
    """Standard content slide title at (0.60, 0.30)."""
    text(slide, Inches(0.60), Inches(0.30), Inches(7), Inches(0.50),
         title_text, size=28, color=NAVY, bold=True)
    if subtitle_text:
        text(slide, Inches(0.60), Inches(0.85), Inches(9), Inches(0.30),
             subtitle_text, size=11, color=RED)


def data_badge(slide, badge_text):
    """Top-right data/context badge (white pill with blue border)."""
    rect(slide, Inches(9.13), Inches(0.30), Inches(3.60), Inches(0.45),
         WHITE, line_color=BLUE, line_width=Pt(1))
    text(slide, Inches(9.23), Inches(0.32), Inches(3.40), Inches(0.40),
         badge_text, size=11, color=BLUE, bold=True, align=PP_ALIGN.CENTER)


def metric_card(slide, left, top, value, label, accent_color, w=Inches(2.80), h=Inches(1.05)):
    """White card with coloured left-edge bar, big value + small label."""
    rect(slide, left, top, w, h, WHITE, line_color=BORDER, line_width=Pt(0.75))
    rect(slide, left, top, Inches(0.06), h, accent_color)
    text(slide, left + Inches(0.20), top + Inches(0.10), w - Inches(0.30), Inches(0.50),
         value, size=22, color=accent_color, bold=True)
    text(slide, left + Inches(0.20), top + Inches(0.60), w - Inches(0.30), Inches(0.35),
         label, size=11, color=BLACK)


def mini_card(slide, left, top, value, label, accent_color, w=Inches(1.83), h=Inches(1.15)):
    """Smaller metric card for title slide (3-across on right side)."""
    rect(slide, left, top, w, h, WHITE)
    rect(slide, left, top, Inches(0.06), h, accent_color)
    text(slide, left + Inches(0.15), top + Inches(0.12), w - Inches(0.25), Inches(0.45),
         value, size=22, color=accent_color, bold=True)
    text(slide, left + Inches(0.15), top + Inches(0.60), w - Inches(0.25), Inches(0.30),
         label, size=11, color=BLACK)


def insight_box(slide, left, top, w, h, title_text, body_text, accent_color, bg_color=None):
    """Callout box with coloured left-edge bar."""
    bg = bg_color or GREY_BG
    rect(slide, left, top, w, h, bg)
    rect(slide, left, top, Inches(0.06), h, accent_color)
    text(slide, left + Inches(0.25), top + Inches(0.10), w - Inches(0.40), Inches(0.30),
         title_text, size=14, color=NAVY, bold=True)
    text(slide, left + Inches(0.25), top + Inches(0.40), w - Inches(0.40), h - Inches(0.50),
         body_text, size=11, color=BLACK)


def recommendation_box(slide, left, top, w, h, content_text):
    """Blue-tinted recommendation box with blue border."""
    rect(slide, left, top, w, h, BLUE_TINT, line_color=BLUE, line_width=Pt(1))
    text(slide, left + Inches(0.25), top + Inches(0.05), w - Inches(0.50), h - Inches(0.10),
         content_text, size=11, color=NAVY, bold=True)


def add_table(slide, left, top, w, h, headers, rows):
    """Add a styled table with navy header row and alternating grey/white rows."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, w, h)
    tbl = table_shape.table

    # Column widths — distribute evenly
    col_w = int(w / num_cols)
    for i in range(num_cols):
        tbl.columns[i].width = col_w

    # Header row
    for ci, header in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = FONT

    # Data rows
    for ri, row_data in enumerate(rows):
        for ci, val in enumerate(row_data):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = GREY_BG if ri % 2 == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = BLACK
                p.font.name = FONT

    return table_shape


def bullet_paragraphs(slide, left, top, w, h, items, size=11, color=BLACK):
    """Multiple paragraphs as bullet points."""
    paras = [(item, size, color, False, 4) for item in items]
    return multipara(slide, left, top, w, h, paras)


# ══════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = GREY_BG

google_top_bar(sl)
left_sidebar(sl)
add_logo(sl)

# Main title — left side
text(sl, Inches(0.60), Inches(1.40), Inches(5.74), Inches(1.80),
     "DRIVING PATIENT\nACQUISITION", size=44, color=NAVY, bold=True)

# Blue accent line
rect(sl, Inches(0.60), Inches(3.30), Inches(2.50), Inches(0.05), BLUE)

# Client name
text(sl, Inches(0.60), Inches(3.55), Inches(5.50), Inches(0.50),
     "Dental by Design  |  Prodent Group", size=22, color=BLUE)

# Author block
multipara(sl, Inches(0.60), Inches(4.30), Inches(5.50), Inches(1.00), [
    ("Christopher Hoole", 11, BLACK, True, 2),
    ("Google Ads Specialist  |  April 2026", 11, BLACK, False, 2),
    ("christopherhoole.com", 11, BLUE, False, 0),
])

# Data badge bottom-left
rect(sl, Inches(0.60), Inches(5.50), Inches(4.50), Inches(0.50), WHITE, line_color=BLUE, line_width=Pt(1))
text(sl, Inches(0.70), Inches(5.52), Inches(4.30), Inches(0.45),
     "Performance Marketing Specialist  |  Interview Presentation", size=11, color=BLUE, bold=True)

# ── Right side: Hero card ──
rect(sl, Inches(6.80), Inches(0.50), Inches(5.90), Inches(2.50), WHITE)
rect(sl, Inches(6.80), Inches(0.50), Inches(0.08), Inches(2.50), GREEN)

text(sl, Inches(7.10), Inches(0.70), Inches(5.40), Inches(0.35),
     "Specialist Focus", size=14, color=NAVY)
text(sl, Inches(7.10), Inches(1.10), Inches(5.40), Inches(0.80),
     "Dental & Aesthetics PPC", size=40, color=GREEN, bold=True)
text(sl, Inches(7.10), Inches(1.90), Inches(5.40), Inches(0.30),
     "Direct experience managing campaigns for dental clinics and cosmetic practices", size=11, color=BLACK)

# ── Mini metric cards ──
mini_card(sl, Inches(6.80), Inches(3.20), "16 Years", "Hands-On Google Ads", BLUE)
mini_card(sl, Inches(8.83), Inches(3.20), "100+", "Accounts Managed", YELLOW)
mini_card(sl, Inches(10.87), Inches(3.20), "£50M+", "Total Spend", GREEN)

google_bottom_bar_title(sl)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2: ABOUT ME
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

google_top_bar(sl, height=Inches(0.06))
slide_title(sl, "About Me")
data_badge(sl, "christopherhoole.com")

# 4 metric cards across
metric_card(sl, Inches(0.60), Inches(1.20), "16 Years", "In Paid Media (Since 2009)", BLUE)
metric_card(sl, Inches(3.70), Inches(1.20), "100+", "Accounts Managed", RED)
metric_card(sl, Inches(6.80), Inches(1.20), "£50M+", "Total Spend Managed", YELLOW)
metric_card(sl, Inches(9.90), Inches(1.20), "6 Levels", "ACT Optimization Engine", GREEN)

# Key points
text(sl, Inches(0.60), Inches(2.60), Inches(6), Inches(0.35),
     "Key Credentials", size=14, color=NAVY, bold=True)

bullet_paragraphs(sl, Inches(0.70), Inches(3.00), Inches(11.50), Inches(2.60), [
    "Hands-on platform specialist — I work inside Google Ads and Meta daily. Not a strategist who delegates to juniors.",
    "Direct dental & aesthetics PPC experience — managed campaigns for a cosmetic dental clinic and Medizen Aesthetics (medically-led cosmetic clinic, Birmingham).",
    "Built ACT (Ads Control Tower) — a proprietary Google Ads optimization engine with 30+ automated checks across Account, Campaign, Ad Group, Keyword, Ad, and Shopping levels.",
    "Notable clients: Lexus, Toyota, Virgin Media O2, Destinology, SixStarHolidays, Select Property Group.",
    "Managed budgets from £5,000 to £400,000 per month across competitive B2C markets.",
    "Key results: 63% CPA reduction at Destinology (£120 to £45). Grew B2B leads from 0 to 2,000/month at Virgin Media O2.",
])

# Recommendation box
recommendation_box(sl, Inches(0.60), Inches(5.85), Inches(12.13), Inches(0.70),
    "This is not a strategy-only profile. Every role I've held has been hands-on, in-platform, accountable for results.")

footer(sl, 2)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3: RELEVANT EXPERIENCE — DENTAL & AESTHETICS
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

google_top_bar(sl, height=Inches(0.06))
slide_title(sl, "Relevant Experience — Dental & Aesthetics",
            "Direct experience in the same market dynamics as Dental by Design")
data_badge(sl, "Competitive B2C Lead Gen")

# Left card: Cosmetic Dental
rect(sl, Inches(0.60), Inches(1.40), Inches(5.80), Inches(2.50), GREY_BG)
rect(sl, Inches(0.60), Inches(1.40), Inches(0.06), Inches(2.50), BLUE)
text(sl, Inches(0.85), Inches(1.50), Inches(5.20), Inches(0.30),
     "Cosmetic Dental Clinic (via PodDigital)", size=14, color=NAVY, bold=True)
bullet_paragraphs(sl, Inches(0.85), Inches(1.85), Inches(5.20), Inches(1.90), [
    "Google Ads for a cosmetic dental practice — implants, veneers, cosmetic treatments",
    "Generated patient enquiries in a competitive local market with high CPCs (£5-15+)",
    "Full search campaign management, bid strategy optimization, negative keyword framework",
    "Confidential — active agency client (can discuss approach and results in conversation)",
])

# Right card: Medizen
rect(sl, Inches(6.80), Inches(1.40), Inches(5.80), Inches(2.50), GREY_BG)
rect(sl, Inches(6.80), Inches(1.40), Inches(0.06), Inches(2.50), GREEN)
text(sl, Inches(7.10), Inches(1.50), Inches(5.20), Inches(0.30),
     "Medizen Aesthetics (Birmingham)", size=14, color=NAVY, bold=True)
bullet_paragraphs(sl, Inches(7.10), Inches(1.85), Inches(5.20), Inches(1.90), [
    "Medically-led cosmetic clinic — injectable treatments, laser procedures, aesthetic services",
    "High-value patient leads in a competitive B2C market",
    "Search-led campaigns targeting treatment-specific intent",
    "Focus on lead quality over volume — CPA reduction through audience and landing page optimization",
])

# Why this matters
text(sl, Inches(0.60), Inches(4.10), Inches(6), Inches(0.35),
     "Why This Matters for Dental by Design", size=14, color=NAVY, bold=True)

# 3 insight boxes
insight_box(sl, Inches(0.60), Inches(4.50), Inches(3.80), Inches(1.60),
    "Same Patient Journey",
    "Search intent → ad click → landing page → enquiry form → booking. I know what makes someone click and what makes them convert.",
    BLUE, BLUE_TINT)

insight_box(sl, Inches(4.70), Inches(4.50), Inches(3.80), Inches(1.60),
    "Same Cost Dynamics",
    "High CPCs (£5-15+), long decision cycle, price-sensitive patients, quality vs cost trade-offs. Every wasted click is expensive.",
    RED, RED_TINT)

insight_box(sl, Inches(8.80), Inches(4.50), Inches(3.80), Inches(1.60),
    "Same Success Metrics",
    "Patient enquiries that convert to bookings — not clicks, not impressions. Lead quality matters more than volume.",
    GREEN, GREEN_TINT)

footer(sl, 3)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4: CASE STUDY — ACCOUNT AUDIT & TURNAROUND
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

google_top_bar(sl, height=Inches(0.06))
slide_title(sl, "Case Study — Account Audit & Turnaround",
            "A UK professional services business — the methodology I would apply to your account")
data_badge(sl, "Anonymised Client")

# Hero metric
rect(sl, Inches(0.60), Inches(1.30), Inches(5.80), Inches(1.50), WHITE)
rect(sl, Inches(0.60), Inches(1.30), Inches(0.08), Inches(1.50), RED)
text(sl, Inches(0.90), Inches(1.35), Inches(4.00), Inches(0.80),
     "42.7%", size=44, color=RED, bold=True)
text(sl, Inches(0.90), Inches(2.10), Inches(5.30), Inches(0.40),
     "of budget wasted — £7,836 of £18,338 total spend identified as waste", size=14, color=BLACK)

# What I found
rect(sl, Inches(6.80), Inches(1.30), Inches(5.80), Inches(1.50), WHITE)
rect(sl, Inches(6.80), Inches(1.30), Inches(0.08), Inches(1.50), GREEN)
text(sl, Inches(7.10), Inches(1.35), Inches(4.00), Inches(0.80),
     "£20-25", size=44, color=GREEN, bold=True)
text(sl, Inches(7.10), Inches(2.10), Inches(5.30), Inches(0.40),
     "Target CPA achieved (down from £38) — a 35-48% reduction in cost per lead", size=14, color=BLACK)

# Findings table
add_table(sl, Inches(0.60), Inches(3.10), Inches(6.00), Inches(2.80),
    ["Issue Found", "Impact"],
    [
        ["No target CPA set", "Algorithm spending unchecked — CPA peaked at £65"],
        ["Conversion tracking inflated", "Micro-conversions counted as leads — misleading data"],
        ["Quality Scores of 2-3 / 10", "Paying 2-4x premium on every click"],
        ["Negative keyword gaps", "Irrelevant traffic consuming 42.7% of budget"],
        ["Auto-apply recommendations on", "Google making unsupervised bid strategy changes"],
    ])

# Recommendations
text(sl, Inches(6.80), Inches(3.10), Inches(5.80), Inches(0.30),
     "What I Recommended", size=14, color=NAVY, bold=True)

bullet_paragraphs(sl, Inches(7.00), Inches(3.45), Inches(5.40), Inches(2.00), [
    "Campaign restructure by service type — tighter ad groups for better Quality Scores",
    "Set proper target CPA — gave the algorithm a cost guardrail",
    "Negative keyword framework — 9 organised lists (phrase + exact match)",
    "Ad copy A/B testing — 3 RSA variants testing different messaging angles",
    "Disabled auto-apply — stopped Google making unsupervised changes",
])

# Bottom callout
rect(sl, Inches(6.80), Inches(5.55), Inches(5.80), Inches(1.05), BLUE_TINT, line_color=BLUE, line_width=Pt(1))
text(sl, Inches(7.05), Inches(5.60), Inches(5.30), Inches(0.90),
     "This is exactly the audit I would run on your Google Ads and Meta accounts in week 1. Same methodology, same rigour, same focus on measurable results.",
     size=11, color=NAVY, bold=True)

footer(sl, 4)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5: WEEK 1 AT DENTAL BY DESIGN
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

google_top_bar(sl, height=Inches(0.06))
slide_title(sl, "My Approach — Week 1 at Dental by Design",
            "I don't wait 3 months to show results. Week 1 delivers measurable improvements.")
data_badge(sl, "First 5 Days")

# Three columns
days_data = [
    ("Days 1-2", "Full Account Audit", BLUE, [
        "Deep-dive into Google Ads + Meta account structure",
        "Audit bid strategies, conversion tracking, search terms",
        "Analyse negative keywords, Quality Scores, ad copy",
        "Benchmark CPL, CPA, conversion rates per clinic",
    ]),
    ("Days 3-4", "Competitive Landscape", RED, [
        "Research competitor bidding — keywords, ad copy, landing pages",
        "Identify gaps in the dental implant keyword space",
        "Map high-intent search terms competitors are targeting",
        "Assess USP positioning (integrated model, pricing advantage)",
    ]),
    ("Day 5", "Quick Wins — Live", GREEN, [
        "Fix conversion tracking issues found in audit",
        "Add negative keywords to block wasted spend immediately",
        "Set proper target CPA / ROAS guardrails",
        "Pause underperforming campaigns or ad groups",
    ]),
]

for i, (day_label, day_title, color, items) in enumerate(days_data):
    col_left = Inches(0.60 + i * 4.10)
    col_w = Inches(3.80)

    # Day label bar
    rect(sl, col_left, Inches(1.40), Inches(1.40), Inches(0.35), color)
    text(sl, col_left, Inches(1.42), Inches(1.40), Inches(0.35),
         day_label, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Day title
    text(sl, col_left, Inches(1.85), col_w, Inches(0.35),
         day_title, size=16, color=NAVY, bold=True)

    # Content card
    rect(sl, col_left, Inches(2.25), col_w, Inches(3.20), GREY_BG)
    rect(sl, col_left, Inches(2.25), Inches(0.06), Inches(3.20), color)
    bullet_paragraphs(sl, col_left + Inches(0.20), Inches(2.40), col_w - Inches(0.35), Inches(2.90),
                      items, size=11, color=BLACK)

# Bottom deliverable bar
recommendation_box(sl, Inches(0.60), Inches(5.70), Inches(12.13), Inches(0.70),
    "Deliverable: Full audit report with findings, quick wins, and 90-day roadmap — by end of week 1")

footer(sl, 5)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6: MONTHS 1-3 ROADMAP
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

google_top_bar(sl, height=Inches(0.06))
slide_title(sl, "Months 1-3 Roadmap",
            "Targeting 20-30% CPL reduction within 90 days based on experience with similar accounts")
data_badge(sl, "90-Day Plan")

months_data = [
    ("Month 1", "Foundation & Restructure", BLUE, [
        "Restructure campaigns by treatment type (implants, cosmetic, general, emergency)",
        "A/B test ad copy — treatment-specific messaging angles",
        "Implement ad schedule + device bid adjustments based on data",
        "Launch negative keyword framework (9 shared lists per account)",
        "Fix conversion tracking — count real patient leads only",
    ]),
    ("Month 2", "Optimise & Expand", YELLOW, [
        "Landing page analysis + conversion rate optimization",
        "Refine audience targeting on Meta (lookalikes, remarketing)",
        "Expand keyword coverage from search term mining insights",
        "Quality Score improvement (ad relevance + landing page experience)",
        "Begin multi-clinic rollout — adapt messaging per practice",
    ]),
    ("Month 3", "Scale & Benchmark", GREEN, [
        "Scale winning campaigns — increase budget on best performers",
        "Cut underperformers — reallocate spend to proven winners",
        "Establish performance benchmarks per clinic (CPL, CVR, CPA)",
        "Full transparent reporting framework — monthly reports for each clinic",
        "Strategic roadmap for months 4-6 based on Q1 data",
    ]),
]

for i, (label, title_txt, color, items) in enumerate(months_data):
    col_left = Inches(0.60 + i * 4.10)
    col_w = Inches(3.80)

    # Full card
    rect(sl, col_left, Inches(1.40), col_w, Inches(5.00), GREY_BG)
    rect(sl, col_left, Inches(1.40), Inches(0.06), Inches(5.00), color)

    # Label + title
    text(sl, col_left + Inches(0.20), Inches(1.55), col_w - Inches(0.30), Inches(0.25),
         label, size=11, color=color, bold=True)
    text(sl, col_left + Inches(0.20), Inches(1.80), col_w - Inches(0.30), Inches(0.35),
         title_txt, size=16, color=NAVY, bold=True)

    bullet_paragraphs(sl, col_left + Inches(0.20), Inches(2.25), col_w - Inches(0.35), Inches(3.80),
                      items, size=11, color=BLACK)

footer(sl, 6)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7: UNDERSTANDING YOUR MARKET
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

google_top_bar(sl, height=Inches(0.06))
slide_title(sl, "Understanding Your Market",
            "The dental implant PPC landscape — and how your integrated model gives you an edge")
data_badge(sl, "Market Analysis")

# Left: Competitive dynamics table
add_table(sl, Inches(0.60), Inches(1.35), Inches(6.00), Inches(3.20),
    ["Market Factor", "Implication for Campaign Management"],
    [
        ["CPCs £5-15+ per click", "Every wasted click is expensive. Precision targeting and negative keywords are critical."],
        ["Long patient decision cycle", "Patients research extensively before committing to implants (£1,695-£15,990+). Remarketing is essential."],
        ["Lead quality > volume", "100 low-quality enquiries cost more than 20 high-intent patients. Structure must filter intent."],
        ["Intense local competition", "Harley Street clinics, NHS alternatives, dental tourism. Messaging must differentiate."],
        ["Price sensitivity", "Patients compare costs. Your pricing advantage must be prominent in ad copy and landing pages."],
    ])

# Right: Your competitive advantage
text(sl, Inches(6.80), Inches(1.20), Inches(5.80), Inches(0.35),
     "Your Competitive Advantage", size=14, color=NAVY, bold=True)

advantages = [
    ("Integrated Model", "Clinic + Lab + Manufacturer = control over quality AND pricing. No other London competitor can match this.", GREEN),
    ("\"Up to 60% Less\"", "A powerful ad copy hook. Price-sensitive patients will click — and your quality keeps them.", BLUE),
    ("Vivo Bridge — 48hrs", "\"Permanent teeth in 48 hours\" is an incredible USP. Speed + quality + price = unbeatable.", RED),
    ("4.8★ / 399+ Reviews", "Social proof in ad extensions and landing pages dramatically improves CTR and conversion rates.", YELLOW),
    ("Multi-Clinic Network", "Scale advantage — learnings from Dental by Design benefit every partner clinic nationwide.", GREEN),
]

for i, (title_txt, body_txt, color) in enumerate(advantages):
    y_pos = Inches(1.65 + i * 0.85)
    rect(sl, Inches(6.80), y_pos, Inches(5.80), Inches(0.75), WHITE, line_color=BORDER, line_width=Pt(0.75))
    rect(sl, Inches(6.80), y_pos, Inches(0.06), Inches(0.75), color)
    multirun(sl, Inches(7.10), y_pos + Inches(0.08), Inches(5.30), Inches(0.60), [
        (title_txt + "  —  ", 11, color, True),
        (body_txt, 11, BLACK, False),
    ])

# Bottom recommendation
recommendation_box(sl, Inches(0.60), Inches(5.85), Inches(12.13), Inches(0.70),
    "I would leverage your integrated model in every ad, every landing page, and every campaign. It's your strongest differentiator in a crowded market.")

footer(sl, 7)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8: MULTI-CLINIC CAMPAIGN MANAGEMENT
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

google_top_bar(sl, height=Inches(0.06))
slide_title(sl, "Multi-Clinic Campaign Management",
            "How I would structure campaigns for Dental by Design + your partner clinics")
data_badge(sl, "Network Strategy")

# Structure table
add_table(sl, Inches(0.60), Inches(1.35), Inches(6.00), Inches(3.00),
    ["Element", "Approach"],
    [
        ["Campaign structure", "Separate campaigns per clinic — tailored messaging, budgets, and targeting for each local market"],
        ["Ad group structure", "Treatment-level groups (Implants, Cosmetic, General, Emergency) — matching search intent precisely"],
        ["Negative keywords", "9 shared lists by word count + match type. One blocked term benefits every clinic instantly."],
        ["Reporting", "Centralised dashboard with clinic-level breakdown — one view of the entire network"],
        ["Budget allocation", "Performance-based — best-performing clinics earn more budget. Underperformers get optimised first."],
    ])

# Right side insights
insight_box(sl, Inches(6.80), Inches(1.35), Inches(5.80), Inches(1.70),
    "The Network Effect",
    "Each clinic gets individual attention, but benefits from the collective data. When one clinic discovers a winning ad copy angle, keyword, or negative keyword — it gets rolled out across the entire network.",
    BLUE, BLUE_TINT)

insight_box(sl, Inches(6.80), Inches(3.25), Inches(5.80), Inches(1.30),
    "Shared Learnings",
    "Dental by Design generates the most data — it will be the testing ground. Proven strategies then scale to partner clinics with localised messaging.",
    GREEN, GREEN_TINT)

# KPI targets
text(sl, Inches(0.60), Inches(4.60), Inches(6), Inches(0.35),
     "Indicative KPI Targets", size=14, color=NAVY, bold=True)

metric_card(sl, Inches(0.60), Inches(5.00), "20-30%", "CPL Reduction (90 days)", GREEN, w=Inches(3.80))
metric_card(sl, Inches(4.70), Inches(5.00), "Baseline →", "Clinic Benchmarks (Month 3)", BLUE, w=Inches(3.80))
metric_card(sl, Inches(8.80), Inches(5.00), "Weekly", "Reporting Cadence", YELLOW, w=Inches(3.80))

footer(sl, 8)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9: ACT — MY COMPETITIVE EDGE
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

google_top_bar(sl, height=Inches(0.06))
slide_title(sl, "ACT — My Competitive Edge",
            "A proprietary Google Ads optimization engine I built from the ground up")
data_badge(sl, "Proprietary Platform")

# Left: What ACT does
text(sl, Inches(0.60), Inches(1.30), Inches(5.80), Inches(0.35),
     "What ACT Does", size=14, color=NAVY, bold=True)

act_features = [
    ("30+ automated checks", "running overnight across all campaign elements", BLUE),
    ("6 optimization levels", "Account, Campaign, Ad Group, Keyword, Ad, Shopping", GREEN),
    ("7 bid strategy rule sets", "Manual CPC, tCPA, tROAS, Max Conv, Max Clicks, PMax, Shopping", RED),
    ("Negative keyword engine", "Auto-sorts into 9 standardised lists by word count + match type", YELLOW),
    ("Morning review dashboard", "Shows exactly what needs attention — nothing gets missed", BLUE),
    ("Human-in-the-loop", "Every action logged, explainable, and reversible", GREEN),
]

for i, (title_txt, body_txt, color) in enumerate(act_features):
    y_pos = Inches(1.75 + i * 0.60)
    rect(sl, Inches(0.60), y_pos, Inches(5.80), Inches(0.52), WHITE if i % 2 else GREY_BG, line_color=BORDER, line_width=Pt(0.5))
    rect(sl, Inches(0.60), y_pos, Inches(0.06), Inches(0.52), color)
    multirun(sl, Inches(0.85), y_pos + Inches(0.08), Inches(5.30), Inches(0.38), [
        (title_txt + "  —  ", 11, NAVY, True),
        (body_txt, 11, BLACK, False),
    ])

# Value prop box
rect(sl, Inches(0.60), Inches(5.45), Inches(5.80), Inches(1.10), BLUE_TINT, line_color=BLUE, line_width=Pt(1))
text(sl, Inches(0.85), Inches(5.50), Inches(5.30), Inches(1.00),
     "\"I don't just manage campaigns — I've built a system that continuously optimises them. Your campaigns are being refined 24/7, not just when I log in. No other candidate has this.\"",
     size=11, color=NAVY, bold=True)

# Right: ACT screenshot
text(sl, Inches(6.80), Inches(1.30), Inches(5.80), Inches(0.35),
     "ACT Morning Review Dashboard", size=14, color=NAVY, bold=True)

rect(sl, Inches(6.80), Inches(1.75), Inches(5.80), Inches(4.80), GREY_BG, line_color=BORDER, line_width=Pt(1))
try:
    sl.shapes.add_picture(ACT_SCREENSHOT, Inches(6.90), Inches(1.85), Inches(5.60), Inches(4.60))
except Exception:
    text(sl, Inches(7.50), Inches(3.80), Inches(4.50), Inches(1),
         "[ACT Dashboard Screenshot]", size=18, color=GREY_MID, align=PP_ALIGN.CENTER)

footer(sl, 9)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10: WHY ME
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = WHITE

google_top_bar(sl, height=Inches(0.06))
slide_title(sl, "Why Me")
data_badge(sl, "Summary")

why_items = [
    ("Hands-On Specialist",
     "You get direct access to the person making decisions and executing changes. I work inside the platforms daily — no juniors, no delegation, no layers.",
     BLUE),
    ("Direct Dental & Aesthetics Experience",
     "I've managed PPC for cosmetic dental clinics and aesthetic practices. I understand patient decision cycles, high CPCs, and the difference between a click and a booking.",
     GREEN),
    ("Proprietary Optimization Platform",
     "No other candidate has built their own Google Ads optimization engine. ACT automates repetitive optimization so I can focus on strategy and growth for your clinics.",
     RED),
    ("Proven Audit Methodology",
     "My audit process (waste analysis → structure review → restructure plan) has been proven with real clients. 42.7% waste identified, CPA reduced by 35-48%.",
     YELLOW),
    ("Results-Focused Commercial Thinking",
     "I measure success by patient enquiries and bookings, not impressions and clicks. Your integrated model is a competitive advantage I know how to leverage.",
     BLUE),
]

for i, (title_txt, body_txt, color) in enumerate(why_items):
    y_pos = Inches(1.20 + i * 1.10)
    rect(sl, Inches(0.60), y_pos, Inches(12.10), Inches(0.95), WHITE, line_color=BORDER, line_width=Pt(0.75))
    rect(sl, Inches(0.60), y_pos, Inches(0.06), Inches(0.95), color)

    # Number circle
    num_shape = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), y_pos + Inches(0.18), Inches(0.45), Inches(0.45))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = color
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = tf.paragraphs[0].add_run()
    run.text = str(i + 1)
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = FONT

    multirun(sl, Inches(1.50), y_pos + Inches(0.10), Inches(10.90), Inches(0.80), [
        (title_txt + "\n", 14, NAVY, True),
        (body_txt, 11, BLACK, False),
    ])

footer(sl, 10)


# ══════════════════════════════════════════════════════════════════
# SLIDE 11: NEXT STEPS (title-style slide)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = GREY_BG

google_top_bar(sl)
left_sidebar(sl)
add_logo(sl)

# Title
text(sl, Inches(0.60), Inches(1.40), Inches(5.50), Inches(0.80),
     "NEXT STEPS", size=44, color=NAVY, bold=True)

rect(sl, Inches(0.60), Inches(2.20), Inches(2.50), Inches(0.05), BLUE)

# Steps as cards
steps = [
    ("1", "Grant me access to your Google Ads and Meta accounts", "Read-only access is all I need to run the full audit", GREEN),
    ("2", "I'll deliver a full audit report within the first week", "Waste analysis, structure review, and restructure plan — same methodology as the case study", BLUE),
    ("3", "You'll see measurable improvements within 30 days", "Quick wins in week 1, structured optimization in months 1-3, benchmarks by month 3", RED),
]

for i, (num, main_txt, sub_txt, color) in enumerate(steps):
    y_pos = Inches(2.60 + i * 1.15)
    rect(sl, Inches(0.60), y_pos, Inches(5.50), Inches(1.00), WHITE)
    rect(sl, Inches(0.60), y_pos, Inches(0.06), Inches(1.00), color)

    # Number
    num_shape = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), y_pos + Inches(0.22), Inches(0.45), Inches(0.45))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = color
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = tf.paragraphs[0].add_run()
    run.text = num
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = FONT

    text(sl, Inches(1.50), y_pos + Inches(0.12), Inches(4.40), Inches(0.35),
         main_txt, size=14, color=NAVY, bold=True)
    text(sl, Inches(1.50), y_pos + Inches(0.50), Inches(4.40), Inches(0.40),
         sub_txt, size=11, color=BLACK)

# Right side — contact card
rect(sl, Inches(6.80), Inches(0.50), Inches(5.90), Inches(2.50), WHITE)
rect(sl, Inches(6.80), Inches(0.50), Inches(0.08), Inches(2.50), BLUE)

text(sl, Inches(7.10), Inches(0.60), Inches(5.40), Inches(0.35),
     "Ready to Start Immediately", size=14, color=NAVY)
text(sl, Inches(7.10), Inches(1.00), Inches(5.40), Inches(0.80),
     "Christopher Hoole", size=40, color=BLUE, bold=True)

multipara(sl, Inches(7.10), Inches(1.80), Inches(5.40), Inches(1.00), [
    ("chris@christopherhoole.com", 14, BLUE, False, 4),
    ("christopherhoole.com", 14, BLUE, False, 4),
    ("+44 7999 500 184", 14, BLACK, False, 0),
])

# Right side — key stats reminder
mini_card(sl, Inches(6.80), Inches(3.20), "16 Years", "Hands-On Experience", BLUE)
mini_card(sl, Inches(8.83), Inches(3.20), "Dental PPC", "Direct Experience", GREEN)
mini_card(sl, Inches(10.87), Inches(3.20), "ACT Engine", "Proprietary Platform", RED)

# Closing statement
rect(sl, Inches(6.80), Inches(4.60), Inches(5.90), Inches(1.20), BLUE_TINT, line_color=BLUE, line_width=Pt(1))
text(sl, Inches(7.05), Inches(4.65), Inches(5.40), Inches(1.10),
     "\"I don't just manage campaigns. I've built a system to continuously optimise them. I have direct dental PPC experience. And I'm ready to start immediately.\"\n\n— Christopher Hoole",
     size=11, color=NAVY, bold=True)

google_bottom_bar_title(sl)


# ══════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════
output_path = os.path.join(SCRIPT_DIR, 'Christopher_Hoole_Dental_by_Design_Presentation.pptx')
prs.save(output_path)
print(f"Saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
