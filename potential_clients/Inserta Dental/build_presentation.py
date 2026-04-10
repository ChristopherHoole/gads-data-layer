"""
Build Christopher Hoole — Dental by Design Interview Presentation
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

# ── Colour palette ──
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x00, 0xB4, 0xD8)
CHARCOAL = RGBColor(0x2C, 0x3E, 0x50)
LIGHT_GREY = RGBColor(0xF5, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xE7, 0x4C, 0x3C)
DARK_TEAL = RGBColor(0x00, 0x8B, 0xA8)
MID_GREY = RGBColor(0x95, 0xA5, 0xA6)
VERY_LIGHT = RGBColor(0xEB, 0xED, 0xF2)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper functions ──

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=CHARCOAL, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, font_size=13, color=CHARCOAL, line_space=1.2, bold_first=False):
    """lines = list of (text, bold, color_override) or just strings"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            txt, bld, col = line_data, False, color
        else:
            txt = line_data[0]
            bld = line_data[1] if len(line_data) > 1 else False
            col = line_data[2] if len(line_data) > 2 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = col
        p.font.bold = bld
        p.font.name = 'Calibri'
        p.space_after = Pt(font_size * (line_space - 1) + 2)
    return txBox

def add_footer(slide):
    add_text_box(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
                 "Christopher Hoole  |  christopherhoole.com  |  Confidential",
                 font_size=8, color=MID_GREY)

def add_slide_number(slide, num):
    add_text_box(slide, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.3),
                 str(num), font_size=8, color=MID_GREY, alignment=PP_ALIGN.RIGHT)

def add_accent_bar(slide, top=Inches(1.35)):
    add_shape_bg(slide, Inches(0.8), top, Inches(1.5), Pt(3), TEAL)

def add_slide_title(slide, title, subtitle=None):
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 title, font_size=28, color=NAVY, bold=True)
    add_accent_bar(slide, top=Inches(0.95))
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
                     subtitle, font_size=14, color=MID_GREY)

def add_metric_card(slide, left, top, value, label, color=TEAL, width=Inches(2.5)):
    card = add_shape_bg(slide, left, top, width, Inches(1.3), WHITE, border_color=VERY_LIGHT)
    # Colour accent bar at top of card
    add_shape_bg(slide, left, top, width, Pt(4), color)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.2), width - Inches(0.3), Inches(0.6),
                 value, font_size=26, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.8), width - Inches(0.3), Inches(0.4),
                 label, font_size=10, color=MID_GREY, alignment=PP_ALIGN.CENTER)

def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=CHARCOAL):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            txt, bld = item
        else:
            txt, bld = item, False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bld
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
        p.level = 0
    return txBox

def add_icon_bullet(slide, left, top, icon_char, title_text, body_text, width=Inches(5)):
    """Add a bullet with a teal circle icon + bold title + body"""
    # Teal circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Pt(2), Pt(22), Pt(22))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TEAL
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = icon_char
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.bold = True

    txBox = slide.shapes.add_textbox(left + Pt(30), top, width, Inches(0.6))
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    run1 = p.add_run()
    run1.text = title_text + "  "
    run1.font.size = Pt(13)
    run1.font.color.rgb = NAVY
    run1.font.bold = True
    run1.font.name = 'Calibri'
    run2 = p.add_run()
    run2.text = body_text
    run2.font.size = Pt(12)
    run2.font.color.rgb = CHARCOAL
    run2.font.name = 'Calibri'

# ═════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(sl, NAVY)

# Large teal accent shape at left
add_shape_bg(sl, Inches(0), Inches(0), Inches(0.15), SLIDE_H, TEAL)

# Title content
add_text_box(sl, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
             "Christopher Hoole", font_size=44, color=WHITE, bold=True)

add_shape_bg(sl, Inches(1.5), Inches(2.4), Inches(2), Pt(3), TEAL)

add_text_box(sl, Inches(1.5), Inches(2.7), Inches(10), Inches(0.6),
             "Performance Marketing Specialist", font_size=24, color=TEAL, bold=False)

add_text_box(sl, Inches(1.5), Inches(3.7), Inches(10), Inches(0.8),
             "Driving Patient Acquisition Through\nData-Led Paid Media",
             font_size=18, color=RGBColor(0xBB, 0xC5, 0xD5))

add_text_box(sl, Inches(1.5), Inches(5.3), Inches(5), Inches(0.3),
             "Prepared for Dental by Design  |  Prodent Group",
             font_size=13, color=MID_GREY)

add_text_box(sl, Inches(1.5), Inches(5.7), Inches(5), Inches(0.3),
             "christopherhoole.com  |  7 April 2026",
             font_size=12, color=MID_GREY)

add_slide_number(sl, 1)

# ═════════════════════════════════════════════════════════════════
# SLIDE 2: ABOUT ME
# ═════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
add_slide_title(sl, "About Me")

# Metric cards
add_metric_card(sl, Inches(0.8), Inches(1.5), "16", "Years in Paid Media", TEAL)
add_metric_card(sl, Inches(3.6), Inches(1.5), "100+", "Accounts Managed", NAVY)
add_metric_card(sl, Inches(6.4), Inches(1.5), "£50M+", "Total Spend Managed", TEAL)
add_metric_card(sl, Inches(9.2), Inches(1.5), "6", "ACT Optimization Levels", NAVY)

# Bullet points
bullets = [
    ("Hands-on platform specialist", True),
    "I work inside the platforms daily — not a strategist who delegates to juniors",
    "",
    ("Direct dental & aesthetics PPC experience", True),
    "Managed Google Ads for a cosmetic dental clinic and Medizen Aesthetics (medically-led cosmetic clinic)",
    "",
    ("Built ACT — my proprietary Google Ads optimization engine", True),
    "30+ automated checks across Account, Campaign, Ad Group, Keyword, Ad, and Shopping levels",
    "",
    ("Notable clients include:", True),
    "Lexus, Toyota, Virgin Media O2, Destinology, Select Property Group",
]
add_bullet_list(sl, Inches(0.8), Inches(3.1), Inches(11.5), Inches(4), bullets, font_size=13)

add_footer(sl)
add_slide_number(sl, 2)

# ═════════════════════════════════════════════════════════════════
# SLIDE 3: RELEVANT EXPERIENCE — DENTAL & AESTHETICS
# ═════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
add_slide_title(sl, "Relevant Experience — Dental & Aesthetics",
                "Direct experience in the markets that matter most for this role")

# Left card: Cosmetic Dental
card_w = Inches(5.5)
card_h = Inches(3.8)
add_shape_bg(sl, Inches(0.8), Inches(1.7), card_w, card_h, LIGHT_GREY, border_color=VERY_LIGHT)
add_shape_bg(sl, Inches(0.8), Inches(1.7), card_w, Pt(4), TEAL)
add_text_box(sl, Inches(1.1), Inches(1.9), Inches(5), Inches(0.4),
             "Cosmetic Dental Clinic", font_size=18, color=NAVY, bold=True)
add_text_box(sl, Inches(1.1), Inches(2.3), Inches(5), Inches(0.3),
             "via PodDigital Agency", font_size=11, color=MID_GREY)

dental_lines = [
    "▸  Google Ads for a cosmetic dental practice",
    "▸  Patient enquiries: implants, veneers, cosmetic treatments",
    "▸  Competitive local market with high CPCs (£5-15+)",
    "▸  Full search campaign management & optimization",
    "▸  Confidential — active agency client",
]
add_bullet_list(sl, Inches(1.1), Inches(2.8), Inches(4.8), Inches(2.5), dental_lines, font_size=12)

# Right card: Medizen Aesthetics
add_shape_bg(sl, Inches(7.0), Inches(1.7), card_w, card_h, LIGHT_GREY, border_color=VERY_LIGHT)
add_shape_bg(sl, Inches(7.0), Inches(1.7), card_w, Pt(4), NAVY)
add_text_box(sl, Inches(7.3), Inches(1.9), Inches(5), Inches(0.4),
             "Medizen Aesthetics", font_size=18, color=NAVY, bold=True)
add_text_box(sl, Inches(7.3), Inches(2.3), Inches(5), Inches(0.3),
             "Birmingham — medically-led cosmetic clinic", font_size=11, color=MID_GREY)

med_lines = [
    "▸  Injectable treatments, laser procedures, aesthetic services",
    "▸  High-value patient leads in competitive B2C market",
    "▸  Search-led campaigns targeting treatment-specific intent",
    "▸  Focus on lead quality over volume",
    "▸  CPA reduction through audience & landing page optimization",
]
add_bullet_list(sl, Inches(7.3), Inches(2.8), Inches(4.8), Inches(2.5), med_lines, font_size=12)

# Bottom quote
add_shape_bg(sl, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.0), NAVY)
add_text_box(sl, Inches(1.2), Inches(5.9), Inches(11), Inches(0.8),
             "\"I understand the dental and aesthetics patient journey — from search intent to enquiry\n"
             "to booking. I know what makes someone click and what makes them convert.\"",
             font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

add_footer(sl)
add_slide_number(sl, 3)

# ═════════════════════════════════════════════════════════════════
# SLIDE 4: CASE STUDY — ACCOUNT AUDIT & TURNAROUND
# ═════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
add_slide_title(sl, "Case Study — Account Audit & Turnaround",
                "A UK professional services business — methodology I would apply to your account")

# The challenge
add_text_box(sl, Inches(0.8), Inches(1.6), Inches(3), Inches(0.3),
             "THE CHALLENGE", font_size=11, color=TEAL, bold=True)
add_shape_bg(sl, Inches(0.8), Inches(1.95), Inches(5.5), Inches(1.6), LIGHT_GREY, border_color=VERY_LIGHT)
challenge_lines = [
    "▸  CPA increased 28% after agency change",
    "▸  Client frustrated with rising costs and falling lead quality",
    "▸  No visibility into what was going wrong",
    "▸  Budget: ~£1,500/month (similar scale to partner clinics)",
]
add_bullet_list(sl, Inches(1.0), Inches(2.05), Inches(5), Inches(1.4), challenge_lines, font_size=12)

# What I found
add_text_box(sl, Inches(0.8), Inches(3.75), Inches(3), Inches(0.3),
             "WHAT I FOUND", font_size=11, color=RED, bold=True)

# Metric cards for findings
add_metric_card(sl, Inches(0.8), Inches(4.1), "42.7%", "Budget Wasted", RED, width=Inches(2.7))
add_metric_card(sl, Inches(3.8), Inches(4.1), "£7,836", "Wasted Spend", RED, width=Inches(2.7))

findings_lines = [
    "▸  No target CPA set — algorithm spending unchecked",
    "▸  Conversion tracking inflated — micro-conversions counted as leads",
    "▸  Quality Scores of 2-3 (out of 10) — paying 2-4x premium per click",
    "▸  Negative keyword gaps letting irrelevant traffic through",
]
add_bullet_list(sl, Inches(0.8), Inches(5.55), Inches(6), Inches(1.3), findings_lines, font_size=11)

# Right side: Results / methodology
add_text_box(sl, Inches(7.0), Inches(1.6), Inches(3), Inches(0.3),
             "MY RECOMMENDATIONS", font_size=11, color=GREEN, bold=True)
add_shape_bg(sl, Inches(7.0), Inches(1.95), Inches(5.5), Inches(1.6), LIGHT_GREY, border_color=VERY_LIGHT)
rec_lines = [
    "▸  Campaign restructure by service type",
    "▸  Bid strategy fix: set proper target CPA",
    "▸  Negative keyword framework (9-list system)",
    "▸  Ad copy A/B testing programme",
]
add_bullet_list(sl, Inches(7.2), Inches(2.05), Inches(5), Inches(1.4), rec_lines, font_size=12)

add_text_box(sl, Inches(7.0), Inches(3.75), Inches(3), Inches(0.3),
             "PROJECTED RESULT", font_size=11, color=GREEN, bold=True)

add_metric_card(sl, Inches(7.0), Inches(4.1), "£20-25", "Target CPA (from £38)", GREEN, width=Inches(2.7))
add_metric_card(sl, Inches(10.0), Inches(4.1), "35-48%", "CPA Reduction", GREEN, width=Inches(2.5))

add_shape_bg(sl, Inches(7.0), Inches(5.55), Inches(5.5), Inches(1.1), NAVY)
add_text_box(sl, Inches(7.3), Inches(5.65), Inches(5), Inches(0.9),
             "This is exactly the audit I would run\non your account in week 1.",
             font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_footer(sl)
add_slide_number(sl, 4)

# ═════════════════════════════════════════════════════════════════
# SLIDE 5: WEEK 1 AT DENTAL BY DESIGN
# ═════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
add_slide_title(sl, "My Approach — Week 1 at Dental by Design",
                "I don't wait 3 months to show results. Week 1 delivers measurable improvements.")

# Day blocks
days = [
    ("Days 1-2", "Full Account Audit", [
        "Deep-dive into Google Ads & Meta account structure",
        "Audit bid strategies, conversion tracking, search terms",
        "Analyse negative keywords, Quality Scores, ad copy",
        "Benchmark current CPL, CPA, conversion rates per clinic",
    ]),
    ("Days 3-4", "Competitive Landscape Analysis", [
        "Research competitor bidding — keywords, ad copy, landing pages",
        "Identify gaps and opportunities in the dental implant space",
        "Map high-intent search terms competitors are targeting",
        "Assess your USP positioning (integrated model, pricing advantage)",
    ]),
    ("Day 5", "Quick Wins Implementation", [
        "Fix any conversion tracking issues found in audit",
        "Add negative keywords to block wasted spend immediately",
        "Adjust bid strategies with proper target CPA/ROAS guardrails",
        "Pause underperforming campaigns or ad groups burning budget",
    ]),
]

for idx, (day_label, day_title, items) in enumerate(days):
    col_left = Inches(0.8 + idx * 4.1)
    col_w = Inches(3.8)

    # Day label
    add_shape_bg(sl, col_left, Inches(1.6), Inches(1.2), Inches(0.35), TEAL)
    add_text_box(sl, col_left, Inches(1.62), Inches(1.2), Inches(0.35),
                 day_label, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Day title
    add_text_box(sl, col_left, Inches(2.05), col_w, Inches(0.4),
                 day_title, font_size=16, color=NAVY, bold=True)

    # Card
    add_shape_bg(sl, col_left, Inches(2.5), col_w, Inches(3.5), LIGHT_GREY, border_color=VERY_LIGHT)
    bullets_text = [f"▸  {item}" for item in items]
    add_bullet_list(sl, col_left + Inches(0.15), Inches(2.65), col_w - Inches(0.3), Inches(3.2),
                    bullets_text, font_size=11, color=CHARCOAL)

# Bottom bar
add_shape_bg(sl, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.5), NAVY)
add_text_box(sl, Inches(1.2), Inches(6.33), Inches(11), Inches(0.45),
             "Deliverable: Full audit report with findings, quick wins, and 90-day roadmap — by end of week 1",
             font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)

add_footer(sl)
add_slide_number(sl, 5)

# ═════════════════════════════════════════════════════════════════
# SLIDE 6: MONTH 1-3 ROADMAP
# ═════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
add_slide_title(sl, "My Approach — Months 1-3 Roadmap",
                "Targeting 20-30% CPL reduction within 90 days based on experience with similar accounts")

months = [
    ("Month 1", "Foundation & Restructure", TEAL, [
        "Restructure campaigns by treatment type (implants, cosmetic, general)",
        "A/B test ad copy — treatment-specific messaging angles",
        "Implement ad schedule & device bid adjustments based on data",
        "Launch negative keyword framework (9 shared lists per account)",
        "Set proper conversion tracking — count real leads only",
    ]),
    ("Month 2", "Optimise & Expand", NAVY, [
        "Landing page analysis & conversion rate optimization",
        "Refine audience targeting on Meta (lookalikes, remarketing)",
        "Expand keyword coverage from search term mining insights",
        "Quality Score improvement programme (ad relevance + landing pages)",
        "Begin multi-clinic rollout — adapt messaging per practice",
    ]),
    ("Month 3", "Scale & Benchmark", DARK_TEAL, [
        "Scale winning campaigns — increase budget on best performers",
        "Cut underperformers — reallocate spend to proven winners",
        "Establish performance benchmarks per clinic (CPL, CVR, CPA)",
        "Full reporting framework — transparent, data-driven, actionable",
        "Strategic roadmap for months 4-6",
    ]),
]

for idx, (label, title, color, items) in enumerate(months):
    col_left = Inches(0.8 + idx * 4.1)
    col_w = Inches(3.8)

    add_shape_bg(sl, col_left, Inches(1.5), col_w, Inches(5.2), LIGHT_GREY, border_color=VERY_LIGHT)
    add_shape_bg(sl, col_left, Inches(1.5), col_w, Pt(5), color)

    add_text_box(sl, col_left + Inches(0.15), Inches(1.7), col_w - Inches(0.3), Inches(0.3),
                 label, font_size=11, color=color, bold=True)
    add_text_box(sl, col_left + Inches(0.15), Inches(2.0), col_w - Inches(0.3), Inches(0.4),
                 title, font_size=16, color=NAVY, bold=True)

    bullets_text = [f"▸  {item}" for item in items]
    add_bullet_list(sl, col_left + Inches(0.15), Inches(2.5), col_w - Inches(0.3), Inches(4),
                    bullets_text, font_size=11, color=CHARCOAL)

add_footer(sl)
add_slide_number(sl, 6)

# ═════════════════════════════════════════════════════════════════
# SLIDE 7: UNDERSTANDING YOUR MARKET
# ═════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
add_slide_title(sl, "Understanding Your Market",
                "The dental implant PPC landscape — and how your model gives you an edge")

# Market dynamics section
add_text_box(sl, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.3),
             "COMPETITIVE DYNAMICS", font_size=11, color=TEAL, bold=True)

market_items = [
    ("High CPCs", "Dental implant keywords command £5-15+ per click in London. Every wasted click is expensive. Precision targeting and negative keyword management are critical."),
    ("Long Decision Cycle", "Patients research extensively before committing to implants (£1,695-£15,990+). Ad copy and landing pages must build trust at every touchpoint."),
    ("Lead Quality > Volume", "100 low-quality enquiries cost more than 20 high-intent patients. Campaign structure must filter intent, not just generate clicks."),
    ("Local Competition", "Harley Street clinics, NHS alternatives, Turkey dental tourism — your messaging must differentiate clearly."),
]

for i, (title, body) in enumerate(market_items):
    y_pos = Inches(2.0 + i * 1.15)
    add_shape_bg(sl, Inches(0.8), y_pos, Inches(5.8), Inches(1.0), LIGHT_GREY if i % 2 == 0 else WHITE, border_color=VERY_LIGHT)
    txBox = sl.shapes.add_textbox(Inches(1.0), y_pos + Inches(0.08), Inches(5.4), Inches(0.85))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = title + ":  "
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = NAVY
    r1.font.name = 'Calibri'
    r2 = p.add_run()
    r2.text = body
    r2.font.size = Pt(11)
    r2.font.color.rgb = CHARCOAL
    r2.font.name = 'Calibri'

# Right side: Your competitive advantage
add_text_box(sl, Inches(7.2), Inches(1.6), Inches(5.5), Inches(0.3),
             "YOUR COMPETITIVE ADVANTAGE", font_size=11, color=GREEN, bold=True)

add_shape_bg(sl, Inches(7.2), Inches(2.0), Inches(5.3), Inches(4.6), NAVY)

advantages = [
    ("Integrated Model", "Clinic + Lab + Manufacturer = control over quality AND pricing. No other competitor in London can say this."),
    ("\"Up to 60% Less\"", "This is a powerful ad copy hook. Price-sensitive patients searching for implants will click — and your quality keeps them."),
    ("Vivo Bridge — 48 Hours", "\"Permanent teeth in 48 hours\" is an incredible USP for ad copy. Speed + quality + price = unbeatable value proposition."),
    ("4.8★ / 399+ Reviews", "Social proof in ad extensions and landing pages dramatically improves CTR and conversion rates."),
    ("Multi-Clinic Network", "Scale advantage — learnings from Dental by Design can be applied across partner clinics nationwide."),
]

for i, (title, body) in enumerate(advantages):
    y_pos = Inches(2.2 + i * 0.85)
    txBox = sl.shapes.add_textbox(Inches(7.5), y_pos, Inches(4.8), Inches(0.75))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "✦  " + title + "  —  "
    r1.font.size = Pt(12)
    r1.font.bold = True
    r1.font.color.rgb = TEAL
    r1.font.name = 'Calibri'
    r2 = p.add_run()
    r2.text = body
    r2.font.size = Pt(11)
    r2.font.color.rgb = WHITE
    r2.font.name = 'Calibri'

add_footer(sl)
add_slide_number(sl, 7)

# ═════════════════════════════════════════════════════════════════
# SLIDE 8: MULTI-CLINIC MANAGEMENT
# ═════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
add_slide_title(sl, "Multi-Clinic Campaign Management",
                "How I would structure campaigns for Dental by Design + partner clinics")

# Left: structure approach
add_text_box(sl, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.3),
             "CAMPAIGN STRUCTURE", font_size=11, color=TEAL, bold=True)

structure_items = [
    "▸  Separate campaigns per clinic — each gets tailored messaging, budgets, and targeting for their local market",
    "▸  Treatment-level ad groups — Implants, Cosmetic, General, Emergency — matching search intent precisely",
    "▸  Shared negative keyword lists — 9 standardised lists (by word count + match type) applied across all clinics. Blocked irrelevant traffic benefits every clinic instantly",
    "▸  Centralised reporting with clinic-level breakdown — one dashboard showing performance across the entire network, drillable to individual clinics",
    "▸  Budget allocation based on performance — best-performing clinics earn more budget. Underperformers get optimised before scaling",
]

for i, item in enumerate(structure_items):
    y_pos = Inches(2.0 + i * 0.95)
    add_shape_bg(sl, Inches(0.8), y_pos, Inches(5.8), Inches(0.85), LIGHT_GREY if i % 2 == 0 else WHITE, border_color=VERY_LIGHT)
    add_text_box(sl, Inches(1.0), y_pos + Inches(0.05), Inches(5.4), Inches(0.75),
                 item, font_size=11, color=CHARCOAL)

# Right: the benefit
add_text_box(sl, Inches(7.2), Inches(1.6), Inches(5.5), Inches(0.3),
             "THE NETWORK EFFECT", font_size=11, color=NAVY, bold=True)

add_shape_bg(sl, Inches(7.2), Inches(2.0), Inches(5.3), Inches(2.5), LIGHT_GREY, border_color=VERY_LIGHT)
network_lines = [
    ("Each clinic gets individual attention,", True),
    ("but benefits from the collective data.", True),
    "",
    "When one clinic discovers a winning ad copy angle, keyword, or negative keyword — it gets rolled out across the network.",
    "",
    "This is exactly how I would use ACT's account-level budget allocation: performance-based distribution ensuring budget flows to the clinics generating the best patient enquiries.",
]
add_bullet_list(sl, Inches(7.4), Inches(2.15), Inches(4.9), Inches(2.2), network_lines, font_size=12, color=CHARCOAL)

# KPI targets
add_text_box(sl, Inches(7.2), Inches(4.7), Inches(5.5), Inches(0.3),
             "INDICATIVE MULTI-CLINIC KPIs", font_size=11, color=TEAL, bold=True)

add_shape_bg(sl, Inches(7.2), Inches(5.05), Inches(5.3), Inches(1.6), NAVY)
kpi_lines = [
    ("Flagship (Dental by Design):", True, WHITE),
    ("  Benchmark CPL, then target 20-30% reduction within 90 days", False, RGBColor(0xBB, 0xC5, 0xD5)),
    ("", False, WHITE),
    ("Partner Clinics:", True, WHITE),
    ("  Establish individual baselines → shared learnings → consistent improvement", False, RGBColor(0xBB, 0xC5, 0xD5)),
]
add_multiline_box(sl, Inches(7.5), Inches(5.15), Inches(4.8), Inches(1.4), kpi_lines, font_size=12, color=WHITE)

add_footer(sl)
add_slide_number(sl, 8)

# ═════════════════════════════════════════════════════════════════
# SLIDE 9: ACT — MY COMPETITIVE EDGE
# ═════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
add_slide_title(sl, "ACT — My Competitive Edge",
                "A proprietary Google Ads optimization engine I built from the ground up")

# Left content
add_text_box(sl, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.3),
             "WHAT ACT DOES", font_size=11, color=TEAL, bold=True)

act_items = [
    "▸  30+ automated optimization checks running overnight",
    "▸  Covers 6 levels: Account, Campaign, Ad Group, Keyword, Ad, Shopping",
    "▸  7 bid strategy rule sets (Manual CPC, tCPA, tROAS, Max Conv, Max Clicks, PMax, Shopping)",
    "▸  Automated negative keyword management — words auto-sorted into 9 standardised lists",
    "▸  Morning review dashboard — shows exactly what needs attention",
    "▸  Human-in-the-loop: every action is logged, explainable, and reversible",
]

for i, item in enumerate(act_items):
    y_pos = Inches(2.0 + i * 0.55)
    add_text_box(sl, Inches(0.8), y_pos, Inches(5.5), Inches(0.5), item, font_size=12, color=CHARCOAL)

# Value prop box
add_shape_bg(sl, Inches(0.8), Inches(5.4), Inches(5.5), Inches(1.2), NAVY)
add_text_box(sl, Inches(1.1), Inches(5.5), Inches(5), Inches(1.0),
             "\"I don't just manage campaigns — I've built\na system that continuously optimises them.\nYour campaigns are being refined 24/7,\nnot just when I log in.\"",
             font_size=13, color=WHITE, bold=False)

# Right: ACT screenshot
screenshot_path = os.path.join(os.path.dirname(__file__), '..', '..', 'act_dashboard', 'prototypes', 'screenshots', 'v6-light-full.png')
screenshot_path = os.path.normpath(screenshot_path)

add_text_box(sl, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.3),
             "ACT MORNING REVIEW DASHBOARD", font_size=11, color=NAVY, bold=True)

# Add screenshot with border
add_shape_bg(sl, Inches(6.9), Inches(1.95), Inches(5.7), Inches(4.75), VERY_LIGHT, border_color=MID_GREY)
try:
    sl.shapes.add_picture(screenshot_path, Inches(7.0), Inches(2.05), Inches(5.5), Inches(4.55))
except Exception:
    add_text_box(sl, Inches(7.0), Inches(3.5), Inches(5.5), Inches(1),
                 "[ACT Dashboard Screenshot]", font_size=18, color=MID_GREY, alignment=PP_ALIGN.CENTER)

add_footer(sl)
add_slide_number(sl, 9)

# ═════════════════════════════════════════════════════════════════
# SLIDE 10: WHY ME
# ═════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
add_slide_title(sl, "Why Me")

why_items = [
    ("1", "Hands-On Specialist",
     "You get direct access to the person making decisions and executing changes. I work inside the platforms daily — no juniors, no delegation, no layers.",
     TEAL),
    ("2", "Direct Dental & Aesthetics Experience",
     "I've managed PPC for cosmetic dental clinics and aesthetic practices. I understand patient decision cycles, high CPCs, and the difference between a click and a booking.",
     NAVY),
    ("3", "Proprietary Optimization Platform (ACT)",
     "No other candidate has built their own Google Ads optimization engine. ACT automates repetitive optimization so I can focus on strategy and growth.",
     TEAL),
    ("4", "Data-Driven Methodology",
     "Every recommendation is backed by data. My audit process (waste analysis → structure review → restructure plan) has been proven with real clients and real results.",
     NAVY),
    ("5", "Results-Focused Commercial Thinking",
     "I measure success by patient enquiries and bookings, not impressions and clicks. Your integrated model (clinic + lab + manufacturer) is a competitive advantage I know how to leverage in campaigns.",
     TEAL),
]

for i, (num, title, body, color) in enumerate(why_items):
    y_pos = Inches(1.4 + i * 1.1)

    # Number circle
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y_pos + Inches(0.08), Inches(0.45), Inches(0.45))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title + body
    txBox = sl.shapes.add_textbox(Inches(1.5), y_pos, Inches(11), Inches(1.0))
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    r1 = p.add_run()
    r1.text = title
    r1.font.size = Pt(16)
    r1.font.bold = True
    r1.font.color.rgb = NAVY
    r1.font.name = 'Calibri'
    p2 = tf2.add_paragraph()
    p2.text = body
    p2.font.size = Pt(12)
    p2.font.color.rgb = CHARCOAL
    p2.font.name = 'Calibri'
    p2.space_before = Pt(3)

add_footer(sl)
add_slide_number(sl, 10)

# ═════════════════════════════════════════════════════════════════
# SLIDE 11: NEXT STEPS
# ═════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, NAVY)

add_shape_bg(sl, Inches(0), Inches(0), Inches(0.15), SLIDE_H, TEAL)

add_text_box(sl, Inches(1.5), Inches(1.0), Inches(10), Inches(0.7),
             "Next Steps", font_size=36, color=WHITE, bold=True)

add_shape_bg(sl, Inches(1.5), Inches(1.8), Inches(2), Pt(3), TEAL)

steps = [
    ("1", "Grant me access to your Google Ads and Meta accounts"),
    ("2", "I'll deliver a full audit report within the first week"),
    ("3", "You'll see measurable improvements within 30 days"),
]

for i, (num, text) in enumerate(steps):
    y_pos = Inches(2.5 + i * 1.1)

    # Number
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y_pos, Inches(0.55), Inches(0.55))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TEAL
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(sl, Inches(2.3), y_pos + Inches(0.08), Inches(9), Inches(0.5),
                 text, font_size=20, color=WHITE)

# Ready to start
add_shape_bg(sl, Inches(1.5), Inches(5.7), Inches(4), Inches(0.55), TEAL)
add_text_box(sl, Inches(1.5), Inches(5.72), Inches(4), Inches(0.5),
             "Ready to start immediately.", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Contact info
add_text_box(sl, Inches(1.5), Inches(6.5), Inches(10), Inches(0.3),
             "chris@christopherhoole.com   |   christopherhoole.com   |   +44 7999 500 184",
             font_size=14, color=MID_GREY)

add_slide_number(sl, 11)

# ═════════════════════════════════════════════════════════════════
# SAVE
# ═════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), 'Christopher_Hoole_Dental_by_Design_Presentation.pptx')
prs.save(output_path)
print(f"Saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
