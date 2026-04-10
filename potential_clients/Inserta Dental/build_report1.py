"""
Dental by Design — Report 1: Waste Spend Analysis
Design system: exact match to Objection Experts winning reports
"""
import os, sys, warnings
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

warnings.filterwarnings('ignore')

# ── Paths ──
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
DATA = os.path.join(SCRIPT_DIR, 'data')
CHARTS = os.path.join(SCRIPT_DIR, 'charts')
REPORTS = os.path.join(SCRIPT_DIR, 'reports')
LOGO = os.path.join(SCRIPT_DIR, 'act_logo_official.png')
os.makedirs(CHARTS, exist_ok=True)
os.makedirs(REPORTS, exist_ok=True)

# ── Design system ──
BLUE = RGBColor(0x42, 0x85, 0xF4)
RED = RGBColor(0xEA, 0x43, 0x35)
YELLOW = RGBColor(0xFB, 0xBC, 0x05)
GREEN = RGBColor(0x34, 0xA8, 0x53)
NAVY = RGBColor(0x1A, 0x23, 0x7E)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GREY_BG = RGBColor(0xF5, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
BLUE_TINT = RGBColor(0xE8, 0xF0, 0xFE)
RED_TINT = RGBColor(0xFC, 0xE8, 0xE6)
GREEN_TINT = RGBColor(0xE6, 0xF4, 0xEA)
GREY_MID = RGBColor(0xCB, 0xD5, 0xE1)
FONT = 'Calibri'

# Chart colours (matplotlib)
C_BLUE = '#4285F4'
C_RED = '#EA4335'
C_YELLOW = '#FBBC05'
C_GREEN = '#34A853'
C_NAVY = '#1A237E'
C_GREY = '#F5F6FA'
C_BORDER = '#E2E8F0'

# ── Chart styling ──
plt.rcParams.update({
    'font.family': 'Calibri', 'font.size': 11,
    'axes.spines.top': False, 'axes.spines.right': False,
    'axes.edgecolor': '#E2E8F0', 'axes.linewidth': 0.8,
    'figure.facecolor': 'white', 'axes.facecolor': 'white',
    'grid.color': '#F0F0F0', 'grid.linewidth': 0.5,
})

def load(f, **kw):
    return pd.read_csv(os.path.join(DATA, f), skiprows=2, thousands=',', **kw)


# ══════════════════════════════════════════════════════════════════
# DATA LOADING & ANALYSIS
# ══════════════════════════════════════════════════════════════════

print("Loading data...")

# Campaign data
camp = load("01_campaign_monthly.csv")
camp = camp[camp['Campaign'].notna() & (camp['Campaign'] != ' --') & ~camp['Campaign status'].astype(str).str.startswith('Total')]
total_spend = camp['Cost'].sum()
total_conv = camp['Conversions'].sum()
total_cpa = total_spend / total_conv

# Search terms
st = load("04_search_terms.csv")
st = st[st['Search term'].notna() & (st['Search term'] != ' --')]
stg = st.groupby('Search term').agg({
    'Cost': 'sum', 'Clicks': 'sum', 'Impr.': 'sum', 'Conversions': 'sum',
    'Campaign': 'first', 'Campaign type': 'first'
}).reset_index()
zero_st = stg[stg['Conversions'] == 0]
zero_st_spend = zero_st['Cost'].sum()
big_waste = zero_st[zero_st['Cost'] >= 5].copy()

# Categorize waste
def categorize(term):
    t = str(term).lower()
    if any(c in t for c in ['brighton implant', 'damira', 'bupa', 'mydentist', 'dental prime', 'dental excellence', 'turkey', 'abroad', 'dental tourism']): return 'Competitor / Dental Tourism'
    if any(w in t for w in ['job', 'career', 'salary', 'hiring', 'vacancy', 'apprentice', 'course', 'training']): return 'Jobs & Training'
    if any(w in t for w in ['nhs', 'free dental', 'free implant', 'free consultation']): return 'NHS / Free Treatment'
    if any(w in t for w in ['how to', 'what is', 'can i', 'should i', 'is it', 'why do', 'how long', 'reviews', 'blog', 'reddit', 'forum']): return 'Informational / Research'
    if any(w in t for w in ['denture', 'crown', 'filling', 'extraction', 'root canal', 'braces', 'orthodont', 'retainer', 'wisdom', 'whitening']) and 'implant' not in t: return 'Wrong Treatment Type'
    if any(w in t for w in ['manchester', 'birmingham', 'leeds', 'liverpool', 'glasgow', 'edinburgh', 'cardiff', 'belfast', 'nottingham', 'sheffield', 'newcastle', 'brighton']) and 'implant clinic' not in t: return 'Wrong Location'
    if any(w in t for w in ['cost', 'price', 'cheap', 'cheapest', 'affordable', 'low cost']) and not any(w in t for w in ['london', 'near me', 'clinic', 'dentist', 'hammersmith']): return 'Price Research (No Location)'
    return 'Other Non-Converting'

big_waste['Category'] = big_waste['Search term'].apply(categorize)
cat_waste = big_waste.groupby('Category').agg({'Cost':'sum','Search term':'count','Clicks':'sum'}).sort_values('Cost', ascending=False)

# Keywords
kw = load("03_keywords_quarterly.csv")
kw = kw[kw['Keyword'].notna() & (kw['Keyword'] != ' --') & ~kw['Match type'].astype(str).str.startswith('Total')]
kwg = kw.groupby('Keyword').agg({'Cost':'sum','Conversions':'sum','Clicks':'sum'}).reset_index()
zero_kw = kwg[(kwg['Conversions']==0) & (kwg['Cost']>0)]

# Device
dv = load("08_device.csv")
dv = dv[dv['Device'].notna() & (dv['Device'] != ' --') & ~dv['Device'].astype(str).str.startswith('Total')]
dvg = dv.groupby('Device').agg({'Cost':'sum','Clicks':'sum','Conversions':'sum'}).sort_values('Cost', ascending=False)

# Day of week
dow = load("10_day_of_week.csv")
dow = dow[dow['Day of the week'].notna() & (dow['Day of the week'] != ' --') & ~dow['Day of the week'].astype(str).str.startswith('Total')]
dowg = dow.groupby('Day of the week').agg({'Cost':'sum','Clicks':'sum','Conversions':'sum'}).reset_index()
day_order = ['Monday','Tuesday','Wednesday','Thursday','Friday','Saturday','Sunday']
dowg['Day of the week'] = pd.Categorical(dowg['Day of the week'], categories=day_order, ordered=True)
dowg = dowg.sort_values('Day of the week')
dowg['CPA'] = dowg['Cost'] / dowg['Conversions']
dowg['CVR'] = dowg['Conversions'] / dowg['Clicks'] * 100

# Hour of day
hod = load("09_hour_of_day.csv")
hod = hod[hod['Hour of the day'].notna() & ~hod['Hour of the day'].astype(str).str.startswith('Total')]
hod['Hour of the day'] = pd.to_numeric(hod['Hour of the day'], errors='coerce')
hod = hod.dropna(subset=['Hour of the day'])
hodg = hod.groupby('Hour of the day').agg({'Cost':'sum','Clicks':'sum','Conversions':'sum'}).reset_index().sort_values('Hour of the day')
hodg['CPA'] = hodg.apply(lambda r: r['Cost']/r['Conversions'] if r['Conversions']>0 else 0, axis=1)

# Geographic
geo = load("07_geographic.csv")
geo = geo[geo['Location'].notna() & (geo['Location'] != ' --') & ~geo['Location'].astype(str).str.startswith('Total')]
geog = geo.groupby('Location').agg({'Cost':'sum','Clicks':'sum','Conversions':'sum'}).reset_index()
zero_geo = geog[geog['Conversions']==0]

# Ad groups
ag = load("02_adgroup_monthly.csv")
ag = ag[ag['Ad group'].notna() & (ag['Ad group'] != ' --') & ~ag['Ad group status'].astype(str).str.startswith('Total')]
agg_df = ag.groupby(['Campaign','Ad group']).agg({'Cost':'sum','Conversions':'sum','Clicks':'sum'}).reset_index()
zero_ag = agg_df[(agg_df['Conversions']==0) & (agg_df['Cost']>0)]

# Campaign waste (Demand Gen)
dg_spend = camp[camp['Campaign']=='Dental Implants - Demand Gen']['Cost'].sum()
dg_conv = camp[camp['Campaign']=='Dental Implants - Demand Gen']['Conversions'].sum()
dg_cpa = dg_spend / dg_conv if dg_conv > 0 else float('inf')

# KMG Max Conv Value waste
kmg_spend = camp[camp['Campaign']=='KMG | Dental Implants Max Conv Value']['Cost'].sum()
kmg_conv = camp[camp['Campaign']=='KMG | Dental Implants Max Conv Value']['Conversions'].sum()

# Total identifiable waste
waste_search_terms = zero_st_spend
waste_keywords = zero_kw['Cost'].sum()
waste_ad_groups = zero_ag['Cost'].sum()
waste_demand_gen = dg_spend  # all of it at £180 CPA
waste_geo = zero_geo['Cost'].sum()
total_waste = waste_search_terms + waste_keywords + waste_ad_groups + waste_demand_gen + waste_geo
waste_pct = total_waste / total_spend * 100

print(f"Total spend: £{total_spend:,.0f}")
print(f"Total waste: £{total_waste:,.0f} ({waste_pct:.1f}%)")
print(f"  Search terms: £{waste_search_terms:,.0f}")
print(f"  Keywords: £{waste_keywords:,.0f}")
print(f"  Ad groups: £{waste_ad_groups:,.0f}")
print(f"  Demand Gen: £{waste_demand_gen:,.0f}")
print(f"  Geographic: £{waste_geo:,.0f}")


# ══════════════════════════════════════════════════════════════════
# CHART GENERATION
# ══════════════════════════════════════════════════════════════════

print("Generating charts...")

# Chart 1: Search term waste categories
fig, ax = plt.subplots(figsize=(7, 4))
cats = cat_waste.index.tolist()
vals = cat_waste['Cost'].tolist()
colors = [C_RED if v > 5000 else C_YELLOW if v > 1000 else C_BLUE for v in vals]
bars = ax.barh(range(len(cats)), vals, color=colors, height=0.6)
ax.set_yticks(range(len(cats)))
ax.set_yticklabels(cats, fontsize=10)
ax.invert_yaxis()
ax.set_xlabel('Wasted Spend (£)', fontsize=10)
for bar, val in zip(bars, vals):
    ax.text(bar.get_width() + 200, bar.get_y() + bar.get_height()/2, f'£{val:,.0f}', va='center', fontsize=9, color=C_NAVY)
ax.set_xlim(0, max(vals) * 1.3)
ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'£{x:,.0f}'))
plt.tight_layout()
plt.savefig(os.path.join(CHARTS, 'waste_categories.png'), dpi=200, bbox_inches='tight')
plt.close()

# Chart 2: Day of week CPA
fig, ax = plt.subplots(figsize=(8, 4))
days = dowg['Day of the week'].astype(str).tolist()
cpas = dowg['CPA'].tolist()
avg = sum(cpas) / len(cpas)
bar_colors = [C_RED if c > avg * 1.05 else C_GREEN if c < avg * 0.95 else C_BLUE for c in cpas]
bars = ax.bar(range(len(days)), cpas, color=bar_colors, width=0.6)
ax.set_xticks(range(len(days)))
ax.set_xticklabels([d[:3] for d in days], fontsize=10)
ax.axhline(y=avg, color=C_NAVY, linestyle='--', linewidth=1, alpha=0.7)
ax.text(6.5, avg + 0.5, f'Avg: £{avg:.0f}', fontsize=9, color=C_NAVY, ha='right')
ax.set_ylabel('Cost Per Conversion (£)', fontsize=10)
for bar, cpa in zip(bars, cpas):
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.3, f'£{cpa:.0f}', ha='center', fontsize=9, color=C_NAVY, fontweight='bold')
ax.set_ylim(0, max(cpas) * 1.15)
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'£{x:.0f}'))
plt.tight_layout()
plt.savefig(os.path.join(CHARTS, 'day_of_week_cpa.png'), dpi=200, bbox_inches='tight')
plt.close()

# Chart 3: Hour of day CPA
fig, ax = plt.subplots(figsize=(10, 4))
hours = hodg['Hour of the day'].astype(int).tolist()
cpas_h = hodg['CPA'].tolist()
convs_h = hodg['Conversions'].tolist()
avg_h = hodg['Cost'].sum() / hodg['Conversions'].sum()
bar_colors = [C_RED if c > avg_h * 1.1 else C_GREEN if c < avg_h * 0.9 else C_BLUE for c in cpas_h]
ax.bar(hours, cpas_h, color=bar_colors, width=0.7)
ax.axhline(y=avg_h, color=C_NAVY, linestyle='--', linewidth=1, alpha=0.7)
ax.text(23, avg_h + 1, f'Avg: £{avg_h:.0f}', fontsize=9, color=C_NAVY, ha='right')
ax.set_xlabel('Hour of Day', fontsize=10)
ax.set_ylabel('Cost Per Conversion (£)', fontsize=10)
ax.set_xticks(range(0, 24))
ax.set_xticklabels([f'{h:02d}' for h in range(24)], fontsize=8)
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'£{x:.0f}'))
plt.tight_layout()
plt.savefig(os.path.join(CHARTS, 'hour_of_day_cpa.png'), dpi=200, bbox_inches='tight')
plt.close()

# Chart 4: Device CPA comparison
fig, ax = plt.subplots(figsize=(6, 3.5))
devices = []
device_cpas = []
device_colors = []
for d, r in dvg.iterrows():
    if r['Conversions'] > 0:
        devices.append(d)
        cpa = r['Cost']/r['Conversions']
        device_cpas.append(cpa)
        device_colors.append(C_GREEN if cpa < 50 else C_YELLOW if cpa < 70 else C_RED)
bars = ax.barh(range(len(devices)), device_cpas, color=device_colors, height=0.5)
ax.set_yticks(range(len(devices)))
ax.set_yticklabels(devices, fontsize=10)
ax.invert_yaxis()
ax.set_xlabel('Cost Per Conversion (£)', fontsize=10)
for bar, cpa in zip(bars, device_cpas):
    ax.text(bar.get_width() + 1, bar.get_y() + bar.get_height()/2, f'£{cpa:.0f}', va='center', fontsize=10, color=C_NAVY, fontweight='bold')
ax.set_xlim(0, max(device_cpas) * 1.25)
ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'£{x:.0f}'))
plt.tight_layout()
plt.savefig(os.path.join(CHARTS, 'device_cpa.png'), dpi=200, bbox_inches='tight')
plt.close()

# Chart 5: Monthly CPA trend
monthly = camp.groupby('Month').agg({'Cost':'sum','Conversions':'sum'}).reset_index()
# Parse month to sortable date
month_map = {'January':1,'February':2,'March':3,'April':4,'May':5,'June':6,
             'July':7,'August':8,'September':9,'October':10,'November':11,'December':12}
monthly['month_num'] = monthly['Month'].apply(lambda x: month_map.get(x.split()[0], 0))
monthly['year'] = monthly['Month'].apply(lambda x: int(x.split()[1]) if len(x.split()) > 1 else 2025)
monthly['sort_key'] = monthly['year'] * 100 + monthly['month_num']
monthly = monthly.sort_values('sort_key')
monthly['CPA'] = monthly['Cost'] / monthly['Conversions'].replace(0, float('nan'))
monthly['label'] = monthly['Month'].apply(lambda x: x[:3] + ' ' + x.split()[-1][2:] if len(x.split()) > 1 else x[:3])

fig, ax = plt.subplots(figsize=(10, 4))
ax.plot(range(len(monthly)), monthly['CPA'].values, color=C_BLUE, linewidth=2, marker='o', markersize=5)
ax.fill_between(range(len(monthly)), monthly['CPA'].values, alpha=0.1, color=C_BLUE)
ax.set_xticks(range(len(monthly)))
ax.set_xticklabels(monthly['label'].values, fontsize=8, rotation=45)
ax.set_ylabel('Cost Per Conversion (£)', fontsize=10)
# Annotate min and max
min_idx = monthly['CPA'].idxmin()
max_idx = monthly['CPA'].idxmax()
min_row = monthly.loc[min_idx]
max_row = monthly.loc[max_idx]
min_pos = monthly.index.get_loc(min_idx)
max_pos = monthly.index.get_loc(max_idx)
ax.annotate(f'£{min_row["CPA"]:.0f}', xy=(min_pos, min_row['CPA']), xytext=(min_pos, min_row['CPA']-8),
            fontsize=9, color=C_GREEN, fontweight='bold', ha='center')
ax.annotate(f'£{max_row["CPA"]:.0f}', xy=(max_pos, max_row['CPA']), xytext=(max_pos, max_row['CPA']+5),
            fontsize=9, color=C_RED, fontweight='bold', ha='center')
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'£{x:.0f}'))
ax.grid(axis='y', alpha=0.3)
plt.tight_layout()
plt.savefig(os.path.join(CHARTS, 'monthly_cpa_trend.png'), dpi=200, bbox_inches='tight')
plt.close()

# Chart 6: Waste summary waterfall
fig, ax = plt.subplots(figsize=(8, 4))
waste_items = [
    ('Search Terms\n(0 Conv)', waste_search_terms),
    ('Zero-Conv\nKeywords', waste_keywords),
    ('Zero-Conv\nAd Groups', waste_ad_groups),
    ('Demand Gen\nCampaign', waste_demand_gen),
    ('Geographic\nWaste', waste_geo),
]
names = [w[0] for w in waste_items]
values = [w[1] for w in waste_items]
colors_w = [C_RED, C_RED, C_YELLOW, C_RED, C_YELLOW]
bars = ax.bar(range(len(names)), values, color=colors_w, width=0.6)
ax.set_xticks(range(len(names)))
ax.set_xticklabels(names, fontsize=9)
for bar, val in zip(bars, values):
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 500, f'£{val:,.0f}', ha='center', fontsize=9, color=C_NAVY, fontweight='bold')
ax.set_ylabel('Wasted Spend (£)', fontsize=10)
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'£{x:,.0f}'))
ax.set_ylim(0, max(values) * 1.15)
plt.tight_layout()
plt.savefig(os.path.join(CHARTS, 'waste_summary.png'), dpi=200, bbox_inches='tight')
plt.close()

print("Charts generated.")


# ══════════════════════════════════════════════════════════════════
# PPTX HELPERS (same as v5 presentation)
# ══════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def rect(sl, l, t, w, h, fc, lc=None, lw=None):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if lc: s.line.color.rgb = lc; s.line.width = lw or Pt(1)
    else: s.line.fill.background()
    return s

def txt(sl, l, t, w, h, c, sz=11, co=BLACK, b=False, a=PP_ALIGN.LEFT):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = c; r.font.size = Pt(sz)
    r.font.color.rgb = co; r.font.bold = b; r.font.name = FONT
    tf.paragraphs[0].alignment = a
    return bx

def multirun(sl, l, t, w, h, runs, a=PP_ALIGN.LEFT):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = a
    for tx, sz, co, b in runs:
        r = p.add_run(); r.text = tx; r.font.size = Pt(sz); r.font.color.rgb = co; r.font.bold = b; r.font.name = FONT
    return bx

def multipara(sl, l, t, w, h, pdata):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, (tx, sz, co, b, sp) in enumerate(pdata):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = tx; r.font.size = Pt(sz); r.font.color.rgb = co; r.font.bold = b; r.font.name = FONT
        if sp: p.space_after = Pt(sp)
    return bx

def bullets(sl, l, t, w, h, items, sz=11, co=BLACK):
    return multipara(sl, l, t, w, h, [(i, sz, co, False, 4) for i in items])

def top_bar(sl, h=Inches(0.07)):
    for i, c in enumerate([BLUE, RED, YELLOW, GREEN]):
        rect(sl, Inches(i*3.333), Inches(0), Inches(3.333), h, c)

def bot_bar(sl, y=Inches(6.92)):
    for p, c in zip([Inches(0.60),Inches(3.63),Inches(6.67),Inches(9.70)], [BLUE,RED,YELLOW,GREEN]):
        rect(sl, p, y, Inches(3.03), Inches(0.03), c)

def bot_bar_title(sl):
    for i, c in enumerate([BLUE, RED, YELLOW, GREEN]):
        rect(sl, Inches(i*3.333), Inches(7.0), Inches(3.333), Inches(0.04), c)

def sidebar(sl):
    rect(sl, Inches(0), Inches(0.07), Inches(0.12), Inches(7.43), BLUE)

def logo(sl, l=Inches(0.60), t=Inches(0.50), s=Inches(0.65)):
    try: sl.shapes.add_picture(LOGO, l, t, s, s)
    except: pass

def footer(sl, n):
    bot_bar(sl)
    try: sl.shapes.add_picture(LOGO, Inches(0.60), Inches(7.0), Inches(0.22), Inches(0.22))
    except: pass
    txt(sl, Inches(0.90), Inches(7.0), Inches(6), Inches(0.25),
        "Christopher Hoole  |  christopherhoole.com  |  Confidential", 11, NAVY, True)
    txt(sl, Inches(12.23), Inches(7.0), Inches(0.50), Inches(0.25), str(n), 11, NAVY, a=PP_ALIGN.RIGHT)

def stitle(sl, t, sub=None):
    txt(sl, Inches(0.60), Inches(0.30), Inches(7), Inches(0.50), t, 28, NAVY, True)
    if sub: txt(sl, Inches(0.60), Inches(0.85), Inches(9), Inches(0.30), sub, 11, RED)

def badge(sl, t):
    rect(sl, Inches(9.13), Inches(0.30), Inches(3.60), Inches(0.45), WHITE, BLUE, Pt(1))
    txt(sl, Inches(9.23), Inches(0.32), Inches(3.40), Inches(0.40), t, 11, BLUE, True, PP_ALIGN.CENTER)

def mcard(sl, l, t, v, lab, ac, w=Inches(2.80), h=Inches(1.05)):
    rect(sl, l, t, w, h, WHITE, BORDER, Pt(0.75))
    rect(sl, l, t, Inches(0.06), h, ac)
    txt(sl, l+Inches(0.20), t+Inches(0.10), w-Inches(0.30), Inches(0.50), v, 22, ac, True)
    txt(sl, l+Inches(0.20), t+Inches(0.60), w-Inches(0.30), Inches(0.35), lab, 11, BLACK)

def minicard(sl, l, t, v, lab, ac):
    rect(sl, l, t, Inches(1.83), Inches(1.15), WHITE)
    rect(sl, l, t, Inches(0.06), Inches(1.15), ac)
    txt(sl, l+Inches(0.15), t+Inches(0.12), Inches(1.58), Inches(0.45), v, 22, ac, True)
    txt(sl, l+Inches(0.15), t+Inches(0.60), Inches(1.58), Inches(0.30), lab, 11, BLACK)

def ibox(sl, l, t, w, h, ti, body, ac, bg=None):
    rect(sl, l, t, w, h, bg or GREY_BG)
    rect(sl, l, t, Inches(0.06), h, ac)
    txt(sl, l+Inches(0.25), t+Inches(0.10), w-Inches(0.40), Inches(0.30), ti, 14, NAVY, True)
    txt(sl, l+Inches(0.25), t+Inches(0.40), w-Inches(0.40), h-Inches(0.50), body, 11, BLACK)

def rbox(sl, l, t, w, h, c):
    rect(sl, l, t, w, h, BLUE_TINT, BLUE, Pt(1))
    txt(sl, l+Inches(0.25), t+Inches(0.05), w-Inches(0.50), h-Inches(0.10), c, 11, NAVY, True)

def add_table(sl, l, t, w, h, hdrs, rows, col_widths=None):
    ts = sl.shapes.add_table(len(rows)+1, len(hdrs), l, t, w, h)
    tb = ts.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tb.columns[i].width = Inches(cw)
    else:
        cw = int(w/len(hdrs))
        for i in range(len(hdrs)): tb.columns[i].width = cw
    for ci, hd in enumerate(hdrs):
        c = tb.cell(0, ci); c.text = hd; c.fill.solid(); c.fill.fore_color.rgb = NAVY
        for p in c.text_frame.paragraphs: p.font.size=Pt(10); p.font.color.rgb=WHITE; p.font.bold=True; p.font.name=FONT
    for ri, rd in enumerate(rows):
        for ci, v in enumerate(rd):
            c = tb.cell(ri+1, ci); c.text = str(v); c.fill.solid()
            c.fill.fore_color.rgb = GREY_BG if ri%2==0 else WHITE
            for p in c.text_frame.paragraphs: p.font.size=Pt(10); p.font.color.rgb=BLACK; p.font.name=FONT
    return ts

def add_img(sl, path, l, t, w, h):
    try: sl.shapes.add_picture(path, l, t, w, h)
    except Exception as e: print(f"  Warning: could not add image {path}: {e}")


# ══════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════

print("Building slides...")

# ── SLIDE 1: TITLE ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = GREY_BG
top_bar(sl); sidebar(sl); logo(sl)

txt(sl, Inches(0.60), Inches(1.40), Inches(5.74), Inches(1.80),
    "Waste Spend\nAnalysis", 44, NAVY, True)
rect(sl, Inches(0.60), Inches(3.30), Inches(2.50), Inches(0.05), BLUE)
txt(sl, Inches(0.60), Inches(3.55), Inches(5.50), Inches(0.50),
    "Dental by Design  |  Prodent Group", 22, BLUE)
multipara(sl, Inches(0.60), Inches(4.30), Inches(5.50), Inches(1.00), [
    ("Christopher Hoole", 11, BLACK, True, 2),
    ("Google Ads Specialist  |  April 2026", 11, BLACK, False, 2),
    ("christopherhoole.com", 11, BLUE, False, 0),
])
rect(sl, Inches(0.60), Inches(5.50), Inches(4.50), Inches(0.50), WHITE, BLUE, Pt(1))
txt(sl, Inches(0.70), Inches(5.52), Inches(4.30), Inches(0.45),
    "Data period: Jan 2025 \u2013 Apr 2026 (16 months)", 11, BLUE, True)

# Hero stat
rect(sl, Inches(6.80), Inches(0.50), Inches(5.90), Inches(2.50), WHITE)
rect(sl, Inches(6.80), Inches(0.50), Inches(0.08), Inches(2.50), RED)
txt(sl, Inches(7.10), Inches(0.70), Inches(5.40), Inches(0.35), "Identified Waste", 14, NAVY)
txt(sl, Inches(7.10), Inches(1.10), Inches(5.40), Inches(0.80),
    f"\u00a3{total_waste:,.0f}", 44, RED, True)
txt(sl, Inches(7.10), Inches(1.90), Inches(5.40), Inches(0.30),
    f"{waste_pct:.1f}% of your \u00a3{total_spend:,.0f} total spend across all campaigns", 11, BLACK)

minicard(sl, Inches(6.80), Inches(3.20), f"{total_conv:,.0f}", "Total Conversions", GREEN)
minicard(sl, Inches(8.83), Inches(3.20), f"\u00a3{total_cpa:,.0f}", "Account CPA", YELLOW)
minicard(sl, Inches(10.87), Inches(3.20), "47", "Total Campaigns", BLUE)

bot_bar_title(sl)


# ── SLIDE 2: EXECUTIVE SUMMARY ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
top_bar(sl, Inches(0.06))
stitle(sl, "Executive Summary")
badge(sl, "Data: Jan 2025 \u2013 Apr 2026")

mcard(sl, Inches(0.60), Inches(1.20), f"\u00a3{total_spend:,.0f}", "Total Spend Analysed", NAVY)
mcard(sl, Inches(3.70), Inches(1.20), f"\u00a3{total_waste:,.0f}", "Waste Identified", RED)
mcard(sl, Inches(6.80), Inches(1.20), f"{waste_pct:.1f}%", "Of Budget Identifiable as Waste", YELLOW)
mcard(sl, Inches(9.90), Inches(1.20), f"{total_conv:,.0f}", "Total Conversions", GREEN)

txt(sl, Inches(0.60), Inches(2.60), Inches(6), Inches(0.35), "Key Findings", 14, NAVY, True)

findings = [
    f"\u00a3{waste_search_terms:,.0f} spent on {len(zero_st):,} search terms that generated zero conversions \u2014 the single largest source of waste in the account.",
    f"The Demand Gen campaign has spent \u00a3{dg_spend:,.0f} at \u00a3{dg_cpa:.0f} CPA \u2014 over 3x the account average. 17 conversions from \u00a3{dg_spend:,.0f} spend.",
    f"129 keywords have spent \u00a3{waste_keywords:,.0f} with zero conversions. The account has {kwg['Keyword'].nunique():,} unique keywords across 47 campaigns \u2014 many are redundant or competing.",
    f"784 account-level negative keywords use broad match \u2014 these may be blocking legitimate converting traffic. Broad match negatives are rarely appropriate.",
    f"Desktop CPA is \u00a367% higher than mobile (\u00a377 vs \u00a346) but there are no device bid adjustments in place anywhere in the account.",
]
bullets(sl, Inches(0.70), Inches(3.00), Inches(11.50), Inches(2.60), findings)

rbox(sl, Inches(0.60), Inches(5.85), Inches(12.13), Inches(0.70),
     f"Opportunity: Addressing the waste categories in this report could redirect \u00a3{total_waste/16:,.0f}+ per month toward patient-generating campaigns.")

footer(sl, 2)


# ── SLIDE 3: MONTHLY CPA TREND ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
top_bar(sl, Inches(0.06))
stitle(sl, "CPA Trend Over Time", "Cost per conversion has nearly tripled from its best month to its worst")
badge(sl, "Monthly Trend")

add_img(sl, os.path.join(CHARTS, 'monthly_cpa_trend.png'), Inches(0.60), Inches(1.30), Inches(8.00), Inches(3.40))

ibox(sl, Inches(8.90), Inches(1.30), Inches(3.80), Inches(1.50),
     "Best Month",
     f"September 2025: \u00a3{monthly.loc[monthly['CPA'].idxmin(), 'CPA']:.0f} CPA. PMax was ramping up and driving efficient conversions.",
     GREEN, GREEN_TINT)

ibox(sl, Inches(8.90), Inches(3.00), Inches(3.80), Inches(1.50),
     "Worst Month",
     f"April 2026: \u00a3{monthly.loc[monthly['CPA'].idxmax(), 'CPA']:.0f} CPA (partial month). CPA has been climbing since January as spend scaled to \u00a350K+/month.",
     RED, RED_TINT)

ibox(sl, Inches(0.60), Inches(5.00), Inches(12.10), Inches(1.40),
     "The Pattern",
     "CPA dropped dramatically when PMax launched (Sep-Oct 2025), but has been rising steadily since. "
     "Spend nearly quintupled from \u00a312K/month to \u00a354K/month. "
     "At this spend level, waste compounds faster \u2014 every percentage point of inefficiency costs more. "
     "The optimisations in this report become increasingly valuable as budget scales.",
     BLUE, BLUE_TINT)

footer(sl, 3)


# ── SLIDE 4: SEARCH TERM WASTE ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
top_bar(sl, Inches(0.06))
stitle(sl, "Search Term Waste", f"{len(zero_st):,} search terms spent \u00a3{zero_st_spend:,.0f} with zero conversions")
badge(sl, "Search Terms")

# Top 15 wasted terms table
top_waste = big_waste.nlargest(15, 'Cost')
add_table(sl, Inches(0.60), Inches(1.30), Inches(7.80), Inches(5.00),
    ["Search Term", "Spend", "Clicks", "Category"],
    [[r['Search term'][:45], f"\u00a3{r['Cost']:,.0f}", f"{r['Clicks']:.0f}", r['Category'][:25]]
     for _, r in top_waste.iterrows()],
    col_widths=[3.5, 0.8, 0.7, 2.8])

# Right side insight
ibox(sl, Inches(8.70), Inches(1.30), Inches(4.00), Inches(2.00),
     "Key Insight",
     f"These {len(zero_st):,} non-converting search terms represent {zero_st_spend/stg['Cost'].sum()*100:.1f}% of tracked search term spend. "
     f"The top 15 alone account for \u00a3{top_waste['Cost'].sum():,.0f} with zero patient enquiries.",
     RED, RED_TINT)

mcard(sl, Inches(8.70), Inches(3.60), f"\u00a3{zero_st_spend:,.0f}", "Total Zero-Conv Search Term Spend", RED, Inches(4.00))

ibox(sl, Inches(8.70), Inches(4.95), Inches(4.00), Inches(1.35),
     "Action Required",
     "Weekly search term review + structured negative keyword framework. Many of these terms could have been blocked months ago.",
     BLUE, BLUE_TINT)

footer(sl, 4)


# ── SLIDE 5: WASTE BY CATEGORY ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
top_bar(sl, Inches(0.06))
stitle(sl, "Waste by Category", "Non-converting search terms grouped by intent type")
badge(sl, "Search Term Categories")

add_img(sl, os.path.join(CHARTS, 'waste_categories.png'), Inches(0.60), Inches(1.30), Inches(7.00), Inches(4.00))

# Category insights on right
txt(sl, Inches(8.00), Inches(1.30), Inches(4.70), Inches(0.30), "Category Breakdown", 14, NAVY, True)

cat_items = []
for cat, r in cat_waste.iterrows():
    cat_items.append(f"{cat}: \u00a3{r['Cost']:,.0f} ({r['Search term']:.0f} terms)")
bullets(sl, Inches(8.00), Inches(1.70), Inches(4.70), Inches(3.60), cat_items, 11)

rbox(sl, Inches(0.60), Inches(5.55), Inches(12.13), Inches(0.80),
     "Price Research terms (\u00a316.6K) are people comparing costs but not in your area \u2014 "
     "negative keywords for generic price queries without location intent would block these. "
     "Competitor/tourism terms (\u00a31.4K) should be excluded immediately.")

footer(sl, 5)


# ── SLIDE 6: ZERO-CONVERSION KEYWORDS ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
top_bar(sl, Inches(0.06))
stitle(sl, "Zero-Conversion Keywords", f"129 keywords have spent \u00a3{waste_keywords:,.0f} with zero conversions")
badge(sl, "Keyword Waste")

# Top 10 wasted keywords
top_kw_waste = zero_kw.nlargest(10, 'Cost')
add_table(sl, Inches(0.60), Inches(1.30), Inches(7.00), Inches(3.50),
    ["Keyword", "Spend", "Clicks"],
    [[r['Keyword'][:50], f"\u00a3{r['Cost']:,.0f}", f"{r['Clicks']:.0f}"]
     for _, r in top_kw_waste.iterrows()],
    col_widths=[4.5, 1.2, 1.3])

ibox(sl, Inches(7.90), Inches(1.30), Inches(4.80), Inches(1.80),
     "The Problem",
     f"The account has {kwg['Keyword'].nunique():,} unique keywords across 47 campaigns. "
     f"Many are duplicated across campaigns (e.g. 'dental implants' appears in multiple campaigns competing against each other). "
     f"129 keywords have never converted.",
     RED, RED_TINT)

ibox(sl, Inches(7.90), Inches(3.30), Inches(4.80), Inches(1.50),
     "Recommendation",
     "Pause all keywords with zero conversions and meaningful spend. "
     "Consolidate duplicate keywords into the best-performing campaign. "
     "Stop internal keyword competition.",
     GREEN, GREEN_TINT)

# Zero-conv ad groups
txt(sl, Inches(0.60), Inches(5.10), Inches(7), Inches(0.30),
    f"Zero-Conversion Ad Groups: {len(zero_ag)} ad groups spent \u00a3{waste_ad_groups:,.0f} with zero conversions", 12, NAVY, True)

top_ag = zero_ag.nlargest(5, 'Cost')
add_table(sl, Inches(0.60), Inches(5.45), Inches(12.10), Inches(1.10),
    ["Campaign", "Ad Group", "Spend", "Clicks"],
    [[r['Campaign'][:30], r['Ad group'][:30], f"\u00a3{r['Cost']:,.0f}", f"{r['Clicks']:.0f}"]
     for _, r in top_ag.iterrows()],
    col_widths=[3.5, 3.5, 1.2, 1.2])

footer(sl, 6)


# ── SLIDE 7: CAMPAIGN-LEVEL WASTE ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
top_bar(sl, Inches(0.06))
stitle(sl, "Campaign-Level Waste", "Active campaigns with concerning performance")
badge(sl, "Campaign Analysis")

# Demand Gen highlight
rect(sl, Inches(0.60), Inches(1.30), Inches(5.80), Inches(2.50), WHITE, BORDER, Pt(0.75))
rect(sl, Inches(0.60), Inches(1.30), Inches(0.08), Inches(2.50), RED)
txt(sl, Inches(0.90), Inches(1.40), Inches(5.20), Inches(0.40),
    "Dental Implants \u2014 Demand Gen", 18, NAVY, True)
txt(sl, Inches(0.90), Inches(1.85), Inches(5.20), Inches(0.25),
    "Currently enabled \u2014 spending \u00a3100/day budget", 11, RED, True)
multipara(sl, Inches(0.90), Inches(2.20), Inches(5.20), Inches(1.50), [
    (f"\u00a3{dg_spend:,.0f} spent \u2014 only {dg_conv:.0f} conversions \u2014 \u00a3{dg_cpa:.0f} CPA", 14, RED, True, 6),
    ("This is 3.7x the account average CPA. The campaign runs on Maximise Conversions with no target.", 11, BLACK, False, 4),
    ("Recommendation: Pause immediately. Redirect budget to the PMax campaign which runs at \u00a336 CPA.", 11, NAVY, True, 0),
])

# KMG Max Conv Value
rect(sl, Inches(6.80), Inches(1.30), Inches(5.80), Inches(2.50), WHITE, BORDER, Pt(0.75))
rect(sl, Inches(6.80), Inches(1.30), Inches(0.08), Inches(2.50), YELLOW)
txt(sl, Inches(7.10), Inches(1.40), Inches(5.20), Inches(0.40),
    "KMG | Dental Implants Max Conv Value", 18, NAVY, True)
txt(sl, Inches(7.10), Inches(1.85), Inches(5.20), Inches(0.25),
    "Currently enabled \u2014 legacy campaign from previous agency", 11, YELLOW, True)
multipara(sl, Inches(7.10), Inches(2.20), Inches(5.20), Inches(1.50), [
    (f"\u00a3{kmg_spend:,.0f} spent \u2014 only {kmg_conv:.0f} conversions \u2014 \u00a3{kmg_spend/max(kmg_conv,0.1):,.0f} CPA", 14, RED, True, 6),
    ("A leftover campaign from Kau Media Group. Tiny volume but extremely high CPA.", 11, BLACK, False, 4),
    ("Recommendation: Pause. This campaign serves no purpose alongside the current structure.", 11, NAVY, True, 0),
])

# Structural mess
txt(sl, Inches(0.60), Inches(4.10), Inches(12), Inches(0.30),
    "Account Structure: 47 Campaigns", 14, NAVY, True)

ibox(sl, Inches(0.60), Inches(4.50), Inches(5.80), Inches(1.90),
     "The Problem",
     "47 campaigns exist in the account, but only 4 are currently active. "
     "The remaining 43 are paused relics from previous agencies and strategy changes. "
     "Campaign names include dates ('22), device splits (Mobile Only), agency names (KMG), "
     "and duplicate service targeting. This makes the account difficult to manage and audit.",
     RED, RED_TINT)

ibox(sl, Inches(6.80), Inches(4.50), Inches(5.80), Inches(1.90),
     "Recommendation",
     "Remove all paused campaigns that have been inactive for 6+ months. "
     "They serve no purpose and create confusion. A clean account structure "
     "is essential for effective management, especially when scaling across partner clinics.",
     GREEN, GREEN_TINT)

footer(sl, 7)


# ── SLIDE 8: DEVICE PERFORMANCE ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
top_bar(sl, Inches(0.06))
stitle(sl, "Device Performance", "No device bid adjustments in place despite significant CPA differences")
badge(sl, "Device Analysis")

add_img(sl, os.path.join(CHARTS, 'device_cpa.png'), Inches(0.60), Inches(1.30), Inches(5.50), Inches(3.00))

# Device table
device_rows = []
for d, r in dvg.iterrows():
    if r['Conversions'] > 0:
        cpa = r['Cost']/r['Conversions']
        cvr = r['Conversions']/r['Clicks']*100
        pct = r['Cost']/dvg['Cost'].sum()*100
        device_rows.append([d, f"\u00a3{r['Cost']:,.0f} ({pct:.0f}%)", f"{r['Conversions']:.0f}", f"\u00a3{cpa:.0f}", f"{cvr:.1f}%", "None"])

add_table(sl, Inches(6.50), Inches(1.30), Inches(6.20), Inches(1.80),
    ["Device", "Spend", "Conv", "CPA", "CVR", "Bid Adj."],
    device_rows,
    col_widths=[1.5, 1.3, 0.7, 0.7, 0.7, 0.7])

ibox(sl, Inches(6.50), Inches(3.30), Inches(6.20), Inches(1.50),
     "Key Finding",
     "Desktop CPA (\u00a377) is 67% higher than mobile (\u00a346), yet no bid adjustments exist. "
     "Tablets have the lowest conversion rate (2.3%) and should have reduced bids. "
     "Mobile drives 84% of spend and performs best on CPA.",
     RED, RED_TINT)

rbox(sl, Inches(0.60), Inches(5.10), Inches(12.10), Inches(0.60),
     "Recommendation: Implement device bid adjustments. Consider reducing desktop bids by 15-20% and tablet bids by 30-40% to align spend with conversion efficiency.")

footer(sl, 8)


# ── SLIDE 9: DAY OF WEEK ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
top_bar(sl, Inches(0.06))
stitle(sl, "Day of Week Performance", "Sunday CPA is 7% above average with the lowest conversion rate")
badge(sl, "Day of Week")

add_img(sl, os.path.join(CHARTS, 'day_of_week_cpa.png'), Inches(0.60), Inches(1.30), Inches(7.50), Inches(3.50))

ibox(sl, Inches(8.40), Inches(1.30), Inches(4.30), Inches(1.60),
     "Best Day",
     "Saturday: \u00a347 CPA. Monday: \u00a348 CPA, 7.0% CVR. "
     "Weekdays are generally more consistent than weekends.",
     GREEN, GREEN_TINT)

ibox(sl, Inches(8.40), Inches(3.10), Inches(4.30), Inches(1.60),
     "Worst Day",
     "Sunday: \u00a353 CPA, 4.1% CVR (lowest). "
     "Sunday spends \u00a3216K but converts at the lowest rate. "
     "Patients research on Sundays but book on weekdays.",
     RED, RED_TINT)

rbox(sl, Inches(0.60), Inches(5.10), Inches(12.10), Inches(0.60),
     "Recommendation: Consider reducing Sunday bids by 10-15%. The CPA difference is moderate but consistent. No day-of-week bid adjustments currently exist.")

footer(sl, 9)


# ── SLIDE 10: HOUR OF DAY ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
top_bar(sl, Inches(0.06))
stitle(sl, "Hour of Day Performance", "Early morning hours (7-8am) show significantly higher CPAs")
badge(sl, "Hour of Day")

add_img(sl, os.path.join(CHARTS, 'hour_of_day_cpa.png'), Inches(0.60), Inches(1.30), Inches(12.10), Inches(3.50))

ibox(sl, Inches(0.60), Inches(5.00), Inches(3.80), Inches(1.40),
     "Peak Efficiency",
     "Hours 3-4am and 16-17 (\u00a342-46 CPA). "
     "Late afternoon is both high-volume and efficient.",
     GREEN, GREEN_TINT)

ibox(sl, Inches(4.70), Inches(5.00), Inches(3.80), Inches(1.40),
     "Worst Efficiency",
     "Hour 7am: \u00a363 CPA. Hour 8am: \u00a357 CPA. "
     "Early morning commute clicks are expensive and convert poorly.",
     RED, RED_TINT)

ibox(sl, Inches(8.80), Inches(5.00), Inches(3.90), Inches(1.40),
     "Recommendation",
     "Implement ad schedule bid adjustments. "
     "Reduce bids during 7-8am by 15-20%. "
     "No scheduling adjustments currently exist.",
     BLUE, BLUE_TINT)

footer(sl, 10)


# ── SLIDE 11: GEOGRAPHIC WASTE ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
top_bar(sl, Inches(0.06))
stitle(sl, "Geographic Waste", f"{len(zero_geo)} locations spent \u00a3{waste_geo:,.0f} with zero conversions")
badge(sl, "Geographic Analysis")

# Top wasting postcodes
top_geo_waste = zero_geo.nlargest(12, 'Cost')
add_table(sl, Inches(0.60), Inches(1.30), Inches(6.50), Inches(4.00),
    ["Location", "Spend", "Clicks"],
    [[r['Location'][:45], f"\u00a3{r['Cost']:,.0f}", f"{r['Clicks']:.0f}"]
     for _, r in top_geo_waste.iterrows()],
    col_widths=[4.0, 1.2, 1.3])

ibox(sl, Inches(7.40), Inches(1.30), Inches(5.30), Inches(2.00),
     "Key Insight",
     f"Geographic waste is relatively low (\u00a3{waste_geo:,.0f}) compared to search term waste. "
     "The location targeting is focused on London postcodes and radius targets around the clinic. "
     "However, some outer London postcodes consistently generate clicks without conversions.",
     BLUE, BLUE_TINT)

ibox(sl, Inches(7.40), Inches(3.50), Inches(5.30), Inches(1.80),
     "Recommendation",
     "Review radius targeting settings. Consider tightening the radius for campaigns "
     "with high CPA in outer zones. The 35km and 20-mile radius targets have higher CPAs "
     "(\u00a393 and \u00a381) than the tighter 10-mile and 25-mile targets (\u00a323 and \u00a321).",
     GREEN, GREEN_TINT)

footer(sl, 11)


# ── SLIDE 12: NEGATIVE KEYWORD GAPS ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
top_bar(sl, Inches(0.06))
stitle(sl, "Negative Keyword Gaps", "784 account-level negatives use broad match \u2014 a risk to legitimate traffic")
badge(sl, "Negative Keywords")

mcard(sl, Inches(0.60), Inches(1.20), "2,192", "Total Negative Keywords", BLUE, Inches(3.00))
mcard(sl, Inches(3.90), Inches(1.20), "784", "Broad Match Negatives", RED, Inches(3.00))
mcard(sl, Inches(7.20), Inches(1.20), "770", "Phrase Match Negatives", GREEN, Inches(3.00))
mcard(sl, Inches(10.50), Inches(1.20), "425", "Exact Match Negatives", YELLOW, Inches(2.20))

ibox(sl, Inches(0.60), Inches(2.60), Inches(5.80), Inches(2.50),
     "The Broad Match Problem",
     "784 negative keywords at account level use broad match. "
     "Broad match negatives are dangerous because they can block legitimate searches. "
     "For example, a broad match negative for 'crown' would block 'dental implant crown' \u2014 "
     "a legitimate high-intent search. "
     "The OE account had a similar issue where broad match negatives were blocking converting traffic. "
     "This needs a full audit to identify which broad negatives are safe and which are blocking revenue.",
     RED, RED_TINT)

ibox(sl, Inches(6.80), Inches(2.60), Inches(5.80), Inches(2.50),
     "Recommended Framework",
     "Replace broad match negatives with phrase or exact match where possible. "
     "Implement a structured 9-list negative keyword system organised by word count and match type. "
     "This gives precise control \u2014 blocking what you want to block without accidentally "
     "blocking converting traffic. "
     "Weekly search term reviews ensure the lists stay current as new waste terms appear.",
     GREEN, GREEN_TINT)

ibox(sl, Inches(0.60), Inches(5.35), Inches(12.10), Inches(1.00),
     "Auto-Apply Risk",
     "5 auto-apply recommendations are enabled, including 'Remove conflicting negative keywords' which has been applied ~20 times. "
     "This means Google has been removing negative keywords it considers conflicting \u2014 potentially re-opening waste that was deliberately blocked. "
     "Last auto-applied: 19 Feb 2026. Recommend disabling this immediately.",
     YELLOW)

footer(sl, 12)


# ── SLIDE 13: TOTAL WASTE SUMMARY ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
top_bar(sl, Inches(0.06))
stitle(sl, "Total Waste Summary")
badge(sl, "All Categories")

# Hero stat
rect(sl, Inches(0.60), Inches(1.00), Inches(7.00), Inches(1.50), WHITE)
rect(sl, Inches(0.60), Inches(1.00), Inches(0.08), Inches(1.50), RED)
txt(sl, Inches(0.90), Inches(1.05), Inches(4.00), Inches(0.80),
    f"\u00a3{total_waste:,.0f}", 44, RED, True)
txt(sl, Inches(0.90), Inches(1.85), Inches(6.50), Inches(0.40),
    f"of your \u00a3{total_spend:,.0f} ({waste_pct:.1f}%) was spent on non-converting activity", 14, BLACK)

add_img(sl, os.path.join(CHARTS, 'waste_summary.png'), Inches(0.60), Inches(2.80), Inches(7.50), Inches(3.50))

# Breakdown cards on right
waste_items = [
    ("Search Terms (0 Conv)", f"\u00a3{waste_search_terms:,.0f}", f"{waste_search_terms/total_waste*100:.0f}%", RED),
    ("Zero-Conv Keywords", f"\u00a3{waste_keywords:,.0f}", f"{waste_keywords/total_waste*100:.0f}%", RED),
    ("Zero-Conv Ad Groups", f"\u00a3{waste_ad_groups:,.0f}", f"{waste_ad_groups/total_waste*100:.1f}%", YELLOW),
    ("Demand Gen Campaign", f"\u00a3{waste_demand_gen:,.0f}", f"{waste_demand_gen/total_waste*100:.0f}%", RED),
    ("Geographic Waste", f"\u00a3{waste_geo:,.0f}", f"{waste_geo/total_waste*100:.0f}%", YELLOW),
]

for i, (label, value, pct, color) in enumerate(waste_items):
    y = Inches(1.10 + i * 0.90)
    rect(sl, Inches(8.20), y, Inches(4.50), Inches(0.75), WHITE, BORDER, Pt(0.75))
    rect(sl, Inches(8.20), y, Inches(0.06), Inches(0.75), color)
    multirun(sl, Inches(8.50), y + Inches(0.08), Inches(2.80), Inches(0.30), [
        (label, 11, BLACK, False),
    ])
    multirun(sl, Inches(8.50), y + Inches(0.38), Inches(2.80), Inches(0.30), [
        (value, 14, color, True),
    ])
    txt(sl, Inches(11.50), y + Inches(0.15), Inches(1.00), Inches(0.45),
        pct, 14, NAVY, True, PP_ALIGN.RIGHT)

txt(sl, Inches(8.20), Inches(5.70), Inches(4.50), Inches(0.50),
    "Note: Some overlap exists between categories. This represents the upper bound of identifiable waste.", 11, NAVY)

footer(sl, 13)


# ── SLIDE 14: QUICK WINS ──
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = GREY_BG
top_bar(sl); sidebar(sl); logo(sl)

txt(sl, Inches(0.60), Inches(1.40), Inches(5.50), Inches(0.80),
    "Quick Wins", 44, NAVY, True)
rect(sl, Inches(0.60), Inches(2.20), Inches(2.50), Inches(0.05), BLUE)
txt(sl, Inches(0.60), Inches(2.40), Inches(5.50), Inches(0.40),
    "Immediate actions to reduce waste", 14, NAVY)

wins = [
    ("1", "Pause the Demand Gen campaign", f"Save ~\u00a3{dg_spend/16:,.0f}/month at \u00a3{dg_cpa:.0f} CPA", RED),
    ("2", "Pause 129 zero-conversion keywords", f"Stop \u00a3{waste_keywords/16:,.0f}/month in wasted keyword spend", RED),
    ("3", "Add negative keywords for top waste categories", "Block price research, competitor, and wrong treatment searches", BLUE),
    ("4", "Audit 784 broad match negative keywords", "Identify which are blocking legitimate converting traffic", YELLOW),
    ("5", "Disable auto-apply for negative keyword removal", "Stop Google from removing your negative keyword protection", GREEN),
    ("6", "Implement device bid adjustments", "Reduce desktop bids 15-20%, tablet bids 30-40%", BLUE),
]

for i, (num, action, detail, color) in enumerate(wins):
    y = Inches(2.90 + i * 0.65)
    rect(sl, Inches(0.60), y, Inches(5.50), Inches(0.55), WHITE)
    rect(sl, Inches(0.60), y, Inches(0.06), Inches(0.55), color)
    multirun(sl, Inches(0.85), y + Inches(0.08), Inches(5.00), Inches(0.40), [
        (f"{num}.  ", 12, color, True),
        (action, 12, NAVY, True),
        (f"  \u2014  {detail}", 11, BLACK, False),
    ])

# Right side: what's next
rect(sl, Inches(6.80), Inches(0.50), Inches(5.90), Inches(2.50), WHITE)
rect(sl, Inches(6.80), Inches(0.50), Inches(0.08), Inches(2.50), GREEN)
txt(sl, Inches(7.10), Inches(0.60), Inches(5.40), Inches(0.35), "What's Next", 14, NAVY)
txt(sl, Inches(7.10), Inches(1.00), Inches(5.40), Inches(0.80),
    "Report 2:\nAccount Structure", 40, GREEN, True)
txt(sl, Inches(7.10), Inches(1.90), Inches(5.40), Inches(0.60),
    "Deep-dive into bid strategies, conversion tracking, Quality Scores, ad copy, and the structural issues driving CPA higher.", 11, BLACK)

minicard(sl, Inches(6.80), Inches(3.20), f"\u00a3{total_waste:,.0f}", "Waste Identified", RED)
minicard(sl, Inches(8.83), Inches(3.20), "6", "Quick Win Actions", BLUE)
minicard(sl, Inches(10.87), Inches(3.20), "Report 2", "Coming Next", GREEN)

bot_bar_title(sl)


# ══════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════

output = os.path.join(REPORTS, '01_waste_spend_report.pptx')
prs.save(output)
print(f"\nSaved to: {output}")
print(f"Slides: {len(prs.slides)}")
print("Done!")
