"""Build the Sales Response Speed Analysis deck.
Matches the design of `DentalByDesign.co.uk - 400 Bookings Strategic Plan v1.pptx`.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- design tokens ----------
BG = RGBColor(0xF5, 0xF6, 0xFA)
NAVY = RGBColor(0x1A, 0x23, 0x7E)
BLUE = RGBColor(0x42, 0x85, 0xF4)
RED = RGBColor(0xEA, 0x43, 0x35)
YELLOW = RGBColor(0xFB, 0xBC, 0x05)
GREEN = RGBColor(0x34, 0xA8, 0x53)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xE5, 0xE7, 0xEB)

FONT = 'Calibri'

LOGO = r'C:\Users\User\Desktop\gads-data-layer\potential_clients\Inserta Dental\act_logo_official.png'

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   'DentalByDesign.co.uk - Maximising Bookings on Current Spend v1.pptx')


def set_bg(slide, rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_rect(slide, left, top, width, height, fill_rgb, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if not line:
        shp.line.fill.background()
    return shp


def add_text(slide, left, top, width, height, text, *,
             font_name=FONT, font_size=12, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_google_bars_top(slide):
    add_rect(slide, 0.00, 0.00, 3.33, 0.07, BLUE)
    add_rect(slide, 3.33, 0.00, 3.33, 0.07, RED)
    add_rect(slide, 6.67, 0.00, 3.33, 0.07, YELLOW)
    add_rect(slide, 10.00, 0.00, 3.33, 0.07, GREEN)


def add_google_bars_bottom(slide):
    add_rect(slide, 0.00, 7.00, 3.33, 0.04, BLUE)
    add_rect(slide, 3.33, 7.00, 3.33, 0.04, RED)
    add_rect(slide, 6.67, 7.00, 3.33, 0.04, YELLOW)
    add_rect(slide, 10.00, 7.00, 3.33, 0.04, GREEN)


def add_left_navy_strip(slide):
    add_rect(slide, 0.00, 0.07, 0.12, 7.43, BLUE)


def add_card(slide, left, top, width, height, accent_rgb, accent_w=0.08):
    """Card with white fill + coloured left border."""
    add_rect(slide, left, top, width, height, WHITE)
    add_rect(slide, left, top, accent_w, height, accent_rgb)


# =====================================================================
# SLIDE 1 — Cover
# =====================================================================
def slide_1_cover(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, BG)
    add_google_bars_top(slide)
    add_google_bars_bottom(slide)
    add_left_navy_strip(slide)
    # Logo top-left (cover position)
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(0.60), Inches(0.50),
                                 height=Inches(0.65))

    # Title block (left side)
    add_text(slide, 0.60, 1.40, 5.74, 1.80,
             'Maximising Bookings on Current Spend',
             font_size=40, bold=True, color=NAVY)
    add_rect(slide, 0.60, 3.30, 2.50, 0.05, BLUE)
    add_text(slide, 0.60, 3.55, 5.50, 0.50,
             'Dental by Design  |  Prodent Group',
             font_size=22, color=BLUE)
    add_text(slide, 0.60, 4.30, 5.50, 0.76,
             'Christopher Hoole\nGoogle Ads Specialist  |  May 2026\nchristopherhoole.com',
             font_size=12, bold=True, color=TEXT)
    add_rect(slide, 0.60, 5.50, 4.50, 0.50, WHITE)
    add_text(slide, 0.70, 5.52, 4.30, 0.45,
             'Data: 5 Jan - 26 Apr 2026  |  16 weeks analysed  |  Paid Search excl Bing',
             font_size=12, bold=True, color=BLUE)

    # Right side — hero stat (top, big card)
    add_card(slide, 6.80, 0.50, 5.90, 2.50, RED)
    add_text(slide, 7.10, 0.60, 5.40, 0.33,
             'The headline finding',
             font_size=14, color=NAVY)
    add_text(slide, 7.10, 1.00, 5.40, 0.83,
             '5.5x',
             font_size=44, bold=True, color=GREEN)
    add_text(slide, 7.10, 1.80, 5.40, 1.10,
             'Higher booking rate when leads are called within 15 minutes vs waiting 12+ hours.\nCallback speed is the single biggest lever on conversion.',
             font_size=12, color=TEXT)

    # Bottom row — 3 small stat cards
    add_card(slide, 6.80, 3.20, 1.83, 1.15, RED)
    add_text(slide, 6.95, 3.32, 1.58, 0.46, '28.2%', font_size=22, bold=True, color=RED)
    add_text(slide, 6.95, 3.80, 1.58, 0.46, 'Book rate when called within 15 min', font_size=11, color=TEXT)

    add_card(slide, 8.83, 3.20, 1.83, 1.15, YELLOW)
    add_text(slide, 8.98, 3.32, 1.58, 0.46, '5.1%', font_size=22, bold=True, color=YELLOW)
    add_text(slide, 8.98, 3.80, 1.58, 0.46, 'Book rate when waiting 12+ hours', font_size=11, color=TEXT)

    add_card(slide, 10.87, 3.20, 1.83, 1.15, BLUE)
    add_text(slide, 11.02, 3.32, 1.58, 0.46, '57%', font_size=22, bold=True, color=BLUE)
    add_text(slide, 11.02, 3.80, 1.58, 0.46, 'Of GAds spend buys leads in dead zones', font_size=11, color=TEXT)

    # Big projected impact card spans bottom-right
    add_card(slide, 6.80, 4.55, 5.90, 1.95, GREEN)
    add_text(slide, 7.10, 4.65, 5.40, 0.33,
             'Projected impact',
             font_size=14, color=NAVY)
    add_text(slide, 7.10, 5.00, 5.40, 0.83,
             '+25-40%',
             font_size=44, bold=True, color=GREEN)
    add_text(slide, 7.10, 5.80, 5.40, 0.65,
             'More bookings on the same monthly spend by shifting GAds to in-hours and adopting an SDR + Closer split.',
             font_size=12, color=TEXT)


def add_slide_chrome(slide, slide_num, total=7):
    """Standard non-cover slide: bg, top+bottom bars, left navy strip, footer."""
    set_bg(slide, BG)
    add_google_bars_top(slide)
    add_google_bars_bottom(slide)
    add_left_navy_strip(slide)
    # Logo bottom-left of footer
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(0.60), Inches(7.07),
                                 height=Inches(0.35))
    # Footer left — author/site/confidential
    add_text(slide, 1.10, 7.10, 9.00, 0.30,
             'Christopher Hoole  |  christopherhoole.com  |  Confidential',
             font_size=9, color=TEXT)
    # Footer right — slide number (separate block)
    add_text(slide, 12.40, 7.10, 0.50, 0.30, str(slide_num),
             font_size=9, color=TEXT, align=PP_ALIGN.RIGHT)


def add_slide_title(slide, title, subtitle=None):
    add_text(slide, 0.60, 0.30, 12.00, 0.65, title,
             font_size=24, bold=True, color=NAVY)
    add_rect(slide, 0.60, 0.95, 1.20, 0.04, BLUE)
    if subtitle:
        add_text(slide, 0.60, 1.05, 12.00, 0.40, subtitle,
                 font_size=12, color=TEXT)


def add_table(slide, left, top, width, headers, rows, *,
              col_widths=None, font_size=10, header_color=NAVY,
              header_text_color=WHITE, zebra=True, highlight_rows=None,
              row_height=0.30):
    """Manual table layout via shapes (gives us full styling control)."""
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [width / n_cols] * n_cols
    # Header
    x = left
    for i, h in enumerate(headers):
        cw = col_widths[i]
        add_rect(slide, x, top, cw, row_height, header_color)
        add_text(slide, x + 0.05, top, cw - 0.10, row_height, str(h),
                 font_size=font_size, bold=True, color=header_text_color,
                 anchor=MSO_ANCHOR.MIDDLE)
        x += cw
    # Body
    for ri, row in enumerate(rows):
        y = top + row_height * (ri + 1)
        bg = WHITE if (not zebra or ri % 2 == 0) else LIGHT_GREY
        if highlight_rows and ri in highlight_rows:
            bg = RGBColor(0xFF, 0xF4, 0xCC)  # soft yellow highlight
        x = left
        for i, cell in enumerate(row):
            cw = col_widths[i]
            add_rect(slide, x, y, cw, row_height, bg)
            add_text(slide, x + 0.05, y, cw - 0.10, row_height, str(cell),
                     font_size=font_size, color=TEXT,
                     anchor=MSO_ANCHOR.MIDDLE)
            x += cw


# =====================================================================
# SLIDE 2 — Headline finding (Table A)
# =====================================================================
def slide_2_headline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 2)
    add_slide_title(slide,
        '15-minute call response delivers 5.5x the bookings of a 12+h wait',
        'Paid Search excl Bing  |  Jan 5 – Apr 26 2026  |  16 weeks  |  3,297 leads')

    # Table A
    headers = ['Bucket', 'Leads', '% of total', 'Booked', 'Book %']
    col_widths = [1.40, 1.10, 1.20, 1.10, 1.20]  # total 6.00
    rows = [
        ['0-15m',   '393',   '11.9%', '111', '28.2%'],
        ['15-30m',  '176',   '5.3%',  '37',  '21.0%'],
        ['30-60m',  '162',   '4.9%',  '28',  '17.3%'],
        ['1-3h',    '227',   '6.9%',  '33',  '14.5%'],
        ['3-6h',    '200',   '6.1%',  '18',  '9.0%'],
        ['6-12h',   '268',   '8.1%',  '10',  '3.7%'],
        ['12+h',    '1,864', '56.5%', '95',  '5.1%'],
        ['Never',   '7',     '0.2%',  '0',   '0.0%'],
        ['TOTAL',   '3,297', '100%',  '332', '10.1%'],
    ]
    add_table(slide, 0.60, 1.55, 6.00, headers, rows,
              col_widths=col_widths, font_size=11,
              highlight_rows=[0, 6, 8])  # 0-15m, 12+h, TOTAL

    # Right column — 3 bullet takeaways as cards
    add_card(slide, 7.20, 1.55, 5.50, 1.60, GREEN)
    add_text(slide, 7.40, 1.65, 5.20, 0.40, 'The gold standard',
             font_size=14, bold=True, color=NAVY)
    add_text(slide, 7.40, 2.05, 5.20, 1.10,
             'Leads called within 15 minutes book at 28.2% — that is the conversion ceiling.',
             font_size=13, color=TEXT)

    add_card(slide, 7.20, 3.30, 5.50, 1.60, YELLOW)
    add_text(slide, 7.40, 3.40, 5.20, 0.40, 'The reality',
             font_size=14, bold=True, color=NAVY)
    add_text(slide, 7.40, 3.80, 5.20, 1.10,
             'Only 11.9% of leads currently get the gold-standard treatment.',
             font_size=13, color=TEXT)

    add_card(slide, 7.20, 5.05, 5.50, 1.60, RED)
    add_text(slide, 7.40, 5.15, 5.20, 0.40, 'The lost majority',
             font_size=14, bold=True, color=NAVY)
    add_text(slide, 7.40, 5.55, 5.20, 1.10,
             '56.5% of leads wait 12+ hours and book at just 5.1% — a 5.5x conversion gap.',
             font_size=13, color=TEXT)


# =====================================================================
# SLIDE 3 — Week-by-week GAds (Table B)
# =====================================================================
def slide_3_weekly(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 3)
    add_slide_title(slide,
        'Booking rate tracks team capacity — when 1 agent is overloaded, conversion collapses',
        'Paid Search excl Bing  |  week-by-week  |  16 weeks  |  3,297 leads')

    headers = ['week', 'leads', '% 0-30m', '% 30m-3h', '% 3h+', 'book %', 'Active agents', 'Notes']
    col_widths = [0.85, 0.55, 0.70, 0.80, 0.65, 0.70, 2.40, 5.50]  # total ~12.15
    rows = [
        ['Jan 5',  '200', '15.5%', '8.0%',  '76.5%', '17.5%', 'feria(17), patrycja(13), livi(3), tommaso(2)', '4 agents → highest book of period'],
        ['Jan 12', '281', '20.6%', '17.8%', '61.6%', '8.9%',  'feria(18), patrycja(7)',                       'Decent speed but only 2 agents on heavy volume'],
        ['Jan 19', '349', '9.5%',  '18.1%', '72.5%', '8.3%',  'patrycja(20), feria(4), livi(2)',              'High volume + slow speed → book drops'],
        ['Jan 26', '358', '4.7%',  '3.1%',  '92.2%', '4.5%',  'patrycja(6), feria(4), livi(1)',               '92% in 3h+ — full overload'],
        ['Feb 2',  '139', '15.8%', '12.9%', '71.2%', '15.1%', 'feria(19), patrycja(2)',                       'Low volume — feria carries it solo'],
        ['Feb 9',  '158', '26.6%', '10.1%', '63.3%', '10.1%', 'feria(11), patrycja(5)',                       'Fast speed, 2 agents — mid book'],
        ['Feb 16', '213', '19.7%', '15.0%', '65.3%', '12.7%', 'feria(20), patrycja(4), angelica(2), livi(1)', '4 agents active — book holds'],
        ['Feb 23', '229', '39.3%', '12.7%', '48.0%', '12.2%', 'feria(18), angelica(7), patrycja(3)',          'Best speed of any week (39% in 0-30m)'],
        ['Mar 2',  '200', '14.0%', '8.0%',  '78.0%', '13.5%', 'feria(13), patrycja(13), livi(1)',             'Slow speed but 2 senior agents kept book up'],
        ['Mar 9',  '149', '25.5%', '12.8%', '61.7%', '18.1%', 'feria(18), patrycja(7), livi(2)',              'Best book % — 3 agents + manageable volume'],
        ['Mar 16', '127', '25.2%', '15.0%', '59.8%', '11.0%', 'feria(14)',                                    'feria solo but low volume = manageable'],
        ['Mar 23', '240', '10.0%', '12.1%', '77.9%', '4.6%',  'feria(11)',                                    'feria alone vs 240 leads → bottleneck proven'],
        ['Mar 30', '124', '19.4%', '12.9%', '67.7%', '11.3%', 'feria(12), livi(1), patrycja(1)',              'Low volume + 3 agents → book holds'],
        ['Apr 6',  '134', '18.7%', '15.7%', '65.7%', '11.2%', 'feria(13), patrycja(1), livi(1)',              'Last "normal" week — low volume + 3 agents'],
        ['Apr 13', '176', '11.4%', '5.7%',  '83.0%', '6.8%',  'feria(11), birsu(1)',                          'birsu starts, 83% slip to 3h+'],
        ['Apr 20', '220', '19.5%', '10.9%', '69.5%', '6.8%',  'feria(7), birsu(5), livi(2), patrycja(1)',     'birsu takes 33% of bookings, book stays low'],
    ]
    add_table(slide, 0.60, 1.55, sum(col_widths), headers, rows,
              col_widths=col_widths, font_size=8, row_height=0.25,
              highlight_rows=[3, 11, 14, 15])  # 26 Jan, 23 Mar, 13 Apr, 20 Apr

    # Three takeaway cards along the bottom strip
    card_y = 6.10
    card_h = 0.85
    card_w = 3.95
    gap = 0.05

    # Card 1 — bottleneck (red)
    add_card(slide, 0.60, card_y, card_w, card_h, RED)
    add_text(slide, 0.75, card_y + 0.05, card_w - 0.20, 0.30,
             'Single-agent weeks crash', font_size=11, bold=True, color=NAVY)
    add_text(slide, 0.75, card_y + 0.35, card_w - 0.20, card_h - 0.40,
             'Mar 23: feria alone vs 240 leads → 4.6% book',
             font_size=10, color=TEXT)

    # Card 2 — works (green)
    add_card(slide, 0.60 + card_w + gap, card_y, card_w, card_h, GREEN)
    add_text(slide, 0.75 + card_w + gap, card_y + 0.05, card_w - 0.20, 0.30,
             '2+ agents holds book rate', font_size=11, bold=True, color=NAVY)
    add_text(slide, 0.75 + card_w + gap, card_y + 0.35, card_w - 0.20, card_h - 0.40,
             'Mar 2, Mar 9, Feb 16 all 13-18% book even with slower speed',
             font_size=10, color=TEXT)

    # Card 3 — caution (yellow)
    add_card(slide, 0.60 + 2*(card_w + gap), card_y, card_w, card_h, YELLOW)
    add_text(slide, 0.75 + 2*(card_w + gap), card_y + 0.05, card_w - 0.20, 0.30,
             'April reset', font_size=11, bold=True, color=NAVY)
    add_text(slide, 0.75 + 2*(card_w + gap), card_y + 0.35, card_w - 0.20, card_h - 0.40,
             'birsu ramping but team output not lifting → training / handover gap',
             font_size=10, color=TEXT)


# =====================================================================
# SLIDE 4 — Outbound queue trend (Table D)
# =====================================================================
def slide_4_outbound(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 4)
    add_slide_title(slide,
        'Outbound calling load doubled in 6 weeks — team did not grow',
        'Last 7 weeks  |  GAds (Paid Search excl Bing) + Paid Social — both need outbound calls')

    headers = ['week', 'GAds leads', 'GAds book %', 'Paid Social leads', 'Social book %', 'Outbound queue', 'Active agents']
    col_widths = [0.85, 1.10, 1.10, 1.45, 1.20, 1.40, 2.10]  # total ~9.20
    rows = [
        ['Mar 9',  '149', '18.1%', '15',  '0.0%', '164', '3'],
        ['Mar 16', '127', '11.0%', '28',  '7.1%', '155', '1'],
        ['Mar 23', '240', '4.6%',  '83',  '4.8%', '323', '1'],
        ['Mar 30', '124', '11.3%', '107', '1.9%', '231', '3'],
        ['Apr 6',  '134', '11.2%', '131', '0.8%', '265', '3'],
        ['Apr 13', '176', '6.8%',  '136', '0.7%', '312', '2 (birsu new)'],
        ['Apr 20', '220', '6.8%',  '148', '1.4%', '368', '4 (birsu ramp)'],
    ]
    add_table(slide, 0.60, 1.55, sum(col_widths), headers, rows,
              col_widths=col_widths, font_size=11, row_height=0.35,
              highlight_rows=[2, 5, 6])  # Mar 23, Apr 13, Apr 20

    # 3 takeaway cards along the right
    card_x = 9.95
    card_w = 2.80
    card_h = 1.80
    gap = 0.10

    add_card(slide, card_x, 1.55, card_w, card_h, RED)
    add_text(slide, card_x + 0.15, 1.65, card_w - 0.30, 0.30,
             'Paid Social scaled 6x', font_size=12, bold=True, color=NAVY)
    add_text(slide, card_x + 0.15, 1.95, card_w - 0.30, card_h - 0.40,
             'From ~25/wk in Jan-Mar to 148/wk by Apr 20',
             font_size=11, color=TEXT)

    add_card(slide, card_x, 1.55 + card_h + gap, card_w, card_h, YELLOW)
    add_text(slide, card_x + 0.15, 1.65 + card_h + gap, card_w - 0.30, 0.30,
             'Outbound queue +125%', font_size=12, bold=True, color=NAVY)
    add_text(slide, card_x + 0.15, 1.95 + card_h + gap, card_w - 0.30, card_h - 0.40,
             '~165/wk in early Mar → 368/wk by Apr 20 (2.2x in 6 weeks)',
             font_size=11, color=TEXT)

    add_card(slide, card_x, 1.55 + 2*(card_h + gap), card_w, card_h, BLUE)
    add_text(slide, card_x + 0.15, 1.65 + 2*(card_h + gap), card_w - 0.30, 0.30,
             'Same 1-4 agents', font_size=12, bold=True, color=NAVY)
    add_text(slide, card_x + 0.15, 1.95 + 2*(card_h + gap), card_w - 0.30, card_h - 0.40,
             'Owning all of it: outbound calls + inbound Metro + closing',
             font_size=11, color=TEXT)


# =====================================================================
# SLIDE 5 — Time-of-day factor (Table E)
# =====================================================================
def slide_5_timeofday(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 5)
    add_slide_title(slide,
        'When leads land matters as much as how fast they are called',
        'Paid Search excl Bing  |  16 weeks  |  3,297 leads')

    headers = ['Bucket', 'Leads', '% of total', 'Booked', 'Book %', '% 0-30m', 'Avg time']
    col_widths = [1.85, 0.90, 1.00, 0.90, 0.95, 0.95, 1.10]  # ~7.65
    rows = [
        ['Mon-Fri 00-08', '325',   '9.9%',  '21', '6.5%',  '0.0%',  '16h 24m'],
        ['Mon-Fri 08-12', '474',   '14.4%', '73', '15.4%', '27.2%', '9h 18m'],
        ['Mon-Fri 12-14', '287',   '8.7%',  '53', '18.5%', '40.8%', '9h 24m'],
        ['Mon-Fri 14-18', '655',   '19.9%', '98', '15.0%', '40.2%', '10h 42m'],
        ['Mon-Fri 18-20', '269',   '8.2%',  '19', '7.1%',  '7.4%',  '20h 06m'],
        ['Mon-Fri 20-24', '488',   '14.8%', '26', '5.3%',  '0.0%',  '20h 42m'],
        ['Weekend',       '799',   '24.2%', '42', '5.3%',  '5.0%',  '27h 42m'],
    ]
    add_table(slide, 0.60, 1.55, sum(col_widths), headers, rows,
              col_widths=col_widths, font_size=11, row_height=0.35,
              highlight_rows=[1, 2, 3, 6])  # in-hours + weekend

    # Right side — 3 factors block (the 3 levers)
    title_x = 8.55
    title_w = 4.20
    add_text(slide, title_x, 1.55, title_w, 0.40,
             '3 factors driving the recent collapse',
             font_size=14, bold=True, color=NAVY)

    factor_y = 2.05
    factor_h = 1.20
    factor_gap = 0.15

    add_card(slide, title_x, factor_y, title_w, factor_h, BLUE)
    add_text(slide, title_x + 0.15, factor_y + 0.05, title_w - 0.30, 0.35,
             '1. When leads land', font_size=12, bold=True, color=NAVY)
    add_text(slide, title_x + 0.15, factor_y + 0.40, title_w - 0.30, factor_h - 0.45,
             'We control this via GAds ad scheduling', font_size=11, color=TEXT)

    add_card(slide, title_x, factor_y + factor_h + factor_gap, title_w, factor_h, YELLOW)
    add_text(slide, title_x + 0.15, factor_y + factor_h + factor_gap + 0.05, title_w - 0.30, 0.35,
             '2. How fast leads are called', font_size=12, bold=True, color=NAVY)
    add_text(slide, title_x + 0.15, factor_y + factor_h + factor_gap + 0.40, title_w - 0.30, factor_h - 0.45,
             'Team capacity + structure (SDR role)', font_size=11, color=TEXT)

    add_card(slide, title_x, factor_y + 2*(factor_h + factor_gap), title_w, factor_h, RED)
    add_text(slide, title_x + 0.15, factor_y + 2*(factor_h + factor_gap) + 0.05, title_w - 0.30, 0.35,
             '3. Agent capacity vs lead volume', font_size=12, bold=True, color=NAVY)
    add_text(slide, title_x + 0.15, factor_y + 2*(factor_h + factor_gap) + 0.40, title_w - 0.30, factor_h - 0.45,
             'Team structure (SDR + Closer split)', font_size=11, color=TEXT)

    # Bottom strip — headline insight
    add_card(slide, 0.60, 6.05, 7.65, 0.85, RED)
    add_text(slide, 0.75, 6.10, 7.50, 0.30,
             'Working hours leads book at 15-18%. Out-of-hours leads book at 5-7%.',
             font_size=12, bold=True, color=NAVY)
    add_text(slide, 0.75, 6.45, 7.50, 0.40,
             '57% of all GAds spend currently buys leads landing in dead zones (eve, night, weekend).',
             font_size=11, color=TEXT)


# =====================================================================
# SLIDE 6 — Quick win: GAds schedule shift
# =====================================================================
def slide_6_quickwin(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 6)
    add_slide_title(slide,
        'Quick win — concentrate GAds spend in Mon-Fri 8am-6pm',
        'Phase 1 — fully within our control, deploy this week, no team changes needed')

    # Spend reallocation table (left)
    headers = ['Time bucket', 'Today', 'Proposed', 'Book %']
    col_widths = [2.40, 1.20, 1.40, 1.20]  # ~6.20
    rows = [
        ['Mon-Fri 00-08',  '9.9%',  'minimal',  '6.5%'],
        ['Mon-Fri 08-12',  '14.4%', 'maintain', '15.4%'],
        ['Mon-Fri 12-14',  '8.7%',  'maintain', '18.5%'],
        ['Mon-Fri 14-18',  '19.9%', 'maintain', '15.0%'],
        ['Mon-Fri 18-20',  '8.2%',  'reduce',   '7.1%'],
        ['Mon-Fri 20-24',  '14.8%', 'minimal',  '5.3%'],
        ['Weekend',        '24.2%', 'minimal',  '5.3%'],
    ]
    add_table(slide, 0.60, 1.55, sum(col_widths), headers, rows,
              col_widths=col_widths, font_size=11, row_height=0.34,
              highlight_rows=[1, 2, 3])  # the keepers

    # Headline summary card under the table
    add_card(slide, 0.60, 4.50, sum(col_widths), 1.10, GREEN)
    add_text(slide, 0.80, 4.60, 5.80, 0.30,
             'Result of shift', font_size=12, bold=True, color=NAVY)
    add_text(slide, 0.80, 4.90, 5.80, 0.70,
             '57% of spend in dead zones → ~10%. ~85% of spend lands in the 15-18% book-rate window.',
             font_size=11, color=TEXT)

    # Right column — what + why (no big stat block)
    add_text(slide, 7.20, 1.55, 5.40, 0.40,
             'What we will do (this week)',
             font_size=14, bold=True, color=NAVY)
    bullets_what = [
        '• Heavily reduce GAds bids/budget outside Mon-Fri 08-18',
        '• Concentrate spend in the proven 15.8% book window',
        '• Outside that window: minimal (brand-defence only) or pause',
    ]
    y = 2.00
    for b in bullets_what:
        add_text(slide, 7.20, y, 5.40, 0.40, b, font_size=11, color=TEXT)
        y += 0.36

    add_text(slide, 7.20, 3.30, 5.40, 0.40,
             'Why',
             font_size=14, bold=True, color=NAVY)
    bullets_why = [
        '• 57% of GAds spend buys leads in dead zones → 5.7% book',
        '• Mon-Fri 08-18 leads book at 15.8% — nearly 3x better',
        '• Fully within our control — no team or process changes needed',
    ]
    y = 3.75
    for b in bullets_why:
        add_text(slide, 7.20, y, 5.40, 0.40, b, font_size=11, color=TEXT)
        y += 0.36

    # Bottom strip — outcome + owner side by side
    add_card(slide, 0.60, 5.85, sum(col_widths), 1.05, BLUE)
    add_text(slide, 0.80, 5.95, 5.80, 0.30,
             'Expected outcome', font_size=12, bold=True, color=NAVY)
    add_text(slide, 0.80, 6.25, 5.80, 0.60,
             '+25-40% more bookings on the same monthly spend. 2-4 week feedback loop.',
             font_size=11, color=TEXT)

    add_card(slide, 7.20, 5.85, 5.40, 1.05, YELLOW)
    add_text(slide, 7.40, 5.95, 5.10, 0.30,
             'Owner & next step', font_size=12, bold=True, color=NAVY)
    add_text(slide, 7.40, 6.25, 5.10, 0.60,
             'Chris (Google Ads). Deploy this week, monitor weekly.',
             font_size=11, color=TEXT)


# =====================================================================
# SLIDE 7 — Structural recommendation: SDR + Closer + warm transfer
# =====================================================================
def slide_7_structural(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide, 7)
    add_slide_title(slide,
        'The structural fix — split sales into Lead Hunter (SDR) + Closer',
        'Phase 2 — implement once Phase 1 data confirms the in-hours capacity ceiling')

    # Two role cards side by side with an arrow between
    role_y = 1.55
    role_h = 2.40
    sdr_x = 0.60
    closer_x = 8.40
    role_w = 4.30

    # SDR card
    add_card(slide, sdr_x, role_y, role_w, role_h, BLUE)
    add_text(slide, sdr_x + 0.20, role_y + 0.10, role_w - 0.40, 0.35,
             'SDR — Lead Hunter', font_size=16, bold=True, color=NAVY)
    add_text(slide, sdr_x + 0.20, role_y + 0.50, role_w - 0.40, 0.30,
             '1 dedicated role', font_size=11, color=BLUE)
    sdr_lines = [
        '• Sole job: dial new leads within 15 minutes',
        '• Qualifies — fit, treatment, timeline, finance',
        '• Captures basics + warm-transfers live to Closer',
        '• Does NOT close the booking',
    ]
    y = role_y + 0.85
    for line in sdr_lines:
        add_text(slide, sdr_x + 0.20, y, role_w - 0.40, 0.32, line,
                 font_size=11, color=TEXT)
        y += 0.32

    # Arrow / handoff text in middle
    arrow_x = sdr_x + role_w + 0.10
    arrow_w = closer_x - arrow_x - 0.10
    add_card(slide, arrow_x, role_y, arrow_w, role_h, YELLOW)
    add_text(slide, arrow_x + 0.10, role_y + 0.10, arrow_w - 0.20, 0.35,
             'Warm transfer (live)', font_size=14, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_text(slide, arrow_x + 0.10, role_y + 0.55, arrow_w - 0.20, 0.30,
             '→', font_size=44, bold=True, color=YELLOW,
             align=PP_ALIGN.CENTER)
    add_text(slide, arrow_x + 0.10, role_y + 1.30, arrow_w - 0.20, 1.00,
             'SDR briefs Closer in real time while lead is on the phone — Closer starts mid-conversation.',
             font_size=10, color=TEXT, align=PP_ALIGN.CENTER)

    # Closer card
    add_card(slide, closer_x, role_y, role_w, role_h, GREEN)
    add_text(slide, closer_x + 0.20, role_y + 0.10, role_w - 0.40, 0.35,
             'Closer', font_size=16, bold=True, color=NAVY)
    add_text(slide, closer_x + 0.20, role_y + 0.50, role_w - 0.40, 0.30,
             'Senior agents (e.g. feria, patrycja)', font_size=11, color=GREEN)
    closer_lines = [
        '• Takes warm handoffs — already qualified',
        '• Owns the booking and the close',
        '• Also works booked-back appointments',
        '• Does NOT dial cold leads',
    ]
    y = role_y + 0.85
    for line in closer_lines:
        add_text(slide, closer_x + 0.20, y, role_w - 0.40, 0.32, line,
                 font_size=11, color=TEXT)
        y += 0.32

    # Why this works — bullets across the middle
    add_text(slide, 0.60, 4.20, 12.10, 0.40,
             'Why this works (and why warm transfer beats callback)',
             font_size=14, bold=True, color=NAVY)
    why_bullets = [
        '• Speed-to-lead is the highest-leverage variable but does not require closer skill — split the work to match the strengths',
        '• Lead is already on the phone — no risk of missing the second call. Closer starts primed, not from scratch',
        '• Industry data: warm transfers convert 2-3x higher than scheduled callbacks. Standard practice in insurance, mortgage, dental, aesthetics',
    ]
    y = 4.65
    for b in why_bullets:
        add_text(slide, 0.60, y, 12.10, 0.40, b, font_size=11, color=TEXT)
        y += 0.36

    # Bottom — phasing + decision needed
    add_card(slide, 0.60, 5.85, 6.05, 1.10, BLUE)
    add_text(slide, 0.80, 5.95, 5.65, 0.30,
             'Phasing', font_size=12, bold=True, color=NAVY)
    add_text(slide, 0.80, 6.25, 5.65, 0.60,
             'Phase 1 (Slide 6) deploys this week. The capacity ceiling it surfaces is the trigger for Phase 2.',
             font_size=11, color=TEXT)

    add_card(slide, 6.75, 5.85, 5.95, 1.10, RED)
    add_text(slide, 6.95, 5.95, 5.55, 0.30,
             'Decision needed (Tommaso / Giulio)', font_size=12, bold=True, color=NAVY)
    add_text(slide, 6.95, 6.25, 5.55, 0.60,
             'In-principle agreement to consider the split if Phase 1 data confirms the capacity ceiling.',
             font_size=11, color=TEXT)


# ---------- main ----------
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.50)
    slide_1_cover(prs)
    slide_2_headline(prs)
    slide_3_weekly(prs)
    slide_4_outbound(prs)
    slide_5_timeofday(prs)
    slide_6_quickwin(prs)
    slide_7_structural(prs)
    prs.save(OUT)
    print(f'saved: {OUT}')


if __name__ == '__main__':
    main()
