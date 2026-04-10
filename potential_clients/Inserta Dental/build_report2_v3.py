"""
Dental by Design - Report 2: Account Structure & Issues (v3)
v3: Fix slide 8 percentages (90% not 90%), fix slide 9 (YoY chart + homepage label),
fix slide 11 (3 priority cards not table wall).
"""
import os, sys, warnings, json
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

warnings.filterwarnings('ignore')

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
D25 = os.path.join(SCRIPT_DIR, 'data', '2025')
D26 = os.path.join(SCRIPT_DIR, 'data', '2026')
DOLD = os.path.join(SCRIPT_DIR, 'data', 'old')
CHARTS = os.path.join(SCRIPT_DIR, 'charts')
REPORTS = os.path.join(SCRIPT_DIR, 'reports')
LOGO = os.path.join(SCRIPT_DIR, 'act_logo_official.png')
os.makedirs(CHARTS, exist_ok=True)

# ── Design tokens ──
BLUE=RGBColor(0x42,0x85,0xF4); RED=RGBColor(0xEA,0x43,0x35)
YELLOW=RGBColor(0xFB,0xBC,0x05); GREEN=RGBColor(0x34,0xA8,0x53)
NAVY=RGBColor(0x1A,0x23,0x7E); BLACK=RGBColor(0x1A,0x1A,0x1A)
GREY_BG=RGBColor(0xF5,0xF6,0xFA); WHITE=RGBColor(0xFF,0xFF,0xFF)
BORDER=RGBColor(0xE2,0xE8,0xF0)
BLUE_TINT=RGBColor(0xE8,0xF0,0xFE); RED_TINT=RGBColor(0xFC,0xE8,0xE6)
GREEN_TINT=RGBColor(0xE6,0xF4,0xEA); GREY_MID=RGBColor(0xCB,0xD5,0xE1)
FONT='Calibri'
CB='#4285F4'; CR='#EA4335'; CY='#FBBC05'; CG='#34A853'; CN='#1A237E'

plt.rcParams.update({'font.family':'Calibri','font.size':11,
    'axes.spines.top':False,'axes.spines.right':False,
    'axes.edgecolor':'#E2E8F0','axes.linewidth':0.8,
    'figure.facecolor':'white','axes.facecolor':'white',
    'grid.color':'#F0F0F0','grid.linewidth':0.5})

def ld(folder,f): return pd.read_csv(os.path.join(folder,f),skiprows=2,thousands=',')
def ct(df,col): return df[df[col].notna()&~df[col].astype(str).str.startswith('Total')&(df[col]!=' --')]

# ══════════════════════════════════════════════════════════════════
# DATA
# ══════════════════════════════════════════════════════════════════
print("Loading data...")

c25=ct(ld(D25,"Campaign report 2025.csv"),'Campaign')
c26=ct(ld(D26,"Campaign report 2026.csv"),'Campaign')
kw25=ct(ld(D25,"Search keyword report 2025.csv"),'Keyword')
kw25=kw25[~kw25['Match type'].astype(str).str.startswith('Total')]
kw26=ct(ld(D26,"Search keyword report 2026.csv"),'Keyword')
kw26=kw26[~kw26['Match type'].astype(str).str.startswith('Total')]
ad26=ct(ld(D26,"Ad report 2026.csv"),'Ad status')
ad25=ct(ld(D25,"Ad report 2025.csv"),'Ad status')
ch=ld(DOLD,"Change history report.csv")

S25=c25['Cost'].sum(); V25=c25['Conversions'].sum(); CPA25=S25/V25
S26=c26['Cost'].sum(); V26=c26['Conversions'].sum(); CPA26=S26/V26
NB25=c25[c25['Campaign']!='Brand']['Cost'].sum(); NV25=c25[c25['Campaign']!='Brand']['Conversions'].sum(); NBCPA25=NB25/NV25
NB26=c26[c26['Campaign']!='Brand']['Cost'].sum(); NV26=c26[c26['Campaign']!='Brand']['Conversions'].sum(); NBCPA26=NB26/NV26

# QS data
kwu26=kw26.sort_values('Cost',ascending=False).drop_duplicates('Keyword',keep='first')
scored26=kwu26[kwu26['Quality Score']!=' --'].copy()
scored26['QS']=pd.to_numeric(scored26['Quality Score'],errors='coerce')
scored26=scored26.dropna(subset=['QS'])
avg_qs_w=((scored26['QS']*scored26['Cost']).sum()/scored26['Cost'].sum()) if scored26['Cost'].sum()>0 else 0

kwu25=kw25.sort_values('Cost',ascending=False).drop_duplicates('Keyword',keep='first')
scored25=kwu25[kwu25['Quality Score']!=' --'].copy()
scored25['QS']=pd.to_numeric(scored25['Quality Score'],errors='coerce')
scored25=scored25.dropna(subset=['QS'])
avg_qs_w25=((scored25['QS']*scored25['Cost']).sum()/scored25['Cost'].sum()) if scored25['Cost'].sum()>0 else 0

# Landing page data
lp26=ad26.groupby('Final URL').agg({'Cost':'sum','Conversions':'sum','Clicks':'sum'}).sort_values('Cost',ascending=False)
lp26=lp26[lp26['Cost']>0]
lp26['CPA']=lp26.apply(lambda r:r['Cost']/r['Conversions'] if r['Conversions']>0 else 999,axis=1)
lp26['CVR']=lp26.apply(lambda r:r['Conversions']/r['Clicks']*100 if r['Clicks']>0 else 0,axis=1)

lp25=ad25.groupby('Final URL').agg({'Cost':'sum','Conversions':'sum','Clicks':'sum'}).sort_values('Cost',ascending=False)
lp25=lp25[lp25['Cost']>0]
lp25['CPA']=lp25.apply(lambda r:r['Cost']/r['Conversions'] if r['Conversions']>0 else 999,axis=1)
lp25['CVR']=lp25.apply(lambda r:r['Conversions']/r['Clicks']*100 if r['Clicks']>0 else 0,axis=1)

# Ad strength
str26=ad26[ad26['Ad strength'].notna()&(ad26['Ad strength']!=' --')]['Ad strength'].value_counts()
str25=ad25[ad25['Ad strength'].notna()&(ad25['Ad strength']!=' --')]['Ad strength'].value_counts()

print(f"2025: L{S25:,.0f} | {V25:.0f} cv | L{CPA25:.0f} CPA | NB L{NBCPA25:.0f}")
print(f"2026: L{S26:,.0f} | {V26:.0f} cv | L{CPA26:.0f} CPA | NB L{NBCPA26:.0f}")
print(f"QS weighted avg: 2025={avg_qs_w25:.1f}, 2026={avg_qs_w:.1f}")
print("Data loaded.")

# ══════════════════════════════════════════════════════════════════
# CHARTS
# ══════════════════════════════════════════════════════════════════
print("Generating charts...")

# Chart 1: QS distribution 2026
fig,ax=plt.subplots(figsize=(7,3.5))
qs_dist=scored26.groupby('QS')['Cost'].sum()
qs_range=range(1,11)
vals=[qs_dist.get(q,0) for q in qs_range]
colors=[CR if q<=3 else CY if q<=5 else CG for q in qs_range]
bars=ax.bar(qs_range,vals,color=colors,width=0.7)
ax.set_xlabel('Quality Score',fontsize=10); ax.set_ylabel('Spend (\u00a3)',fontsize=10)
ax.set_xticks(range(1,11))
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f'\u00a3{x:,.0f}'))
for bar,val in zip(bars,vals):
    if val>0: ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+200,f'\u00a3{val:,.0f}',ha='center',fontsize=8,color=CN)
ax.set_title('2026 Keyword Spend by Quality Score',fontsize=11,color=CN,fontweight='bold',loc='left')
plt.tight_layout()
plt.savefig(f"{CHARTS}/qs_dist_2026.png",dpi=200,bbox_inches='tight'); plt.close()

# Chart 2: Ad strength distribution 2026
fig,ax=plt.subplots(figsize=(5,3.5))
strength_order=['Excellent','Good','Average','Poor','Pending']
str_vals=[str26.get(s,0) for s in strength_order]
str_colors=[CG,CB,CY,CR,'#CCCCCC']
bars=ax.barh(range(len(strength_order)),str_vals,color=str_colors,height=0.5)
ax.set_yticks(range(len(strength_order))); ax.set_yticklabels(strength_order,fontsize=10)
ax.invert_yaxis()
for bar,val in zip(bars,str_vals): ax.text(bar.get_width()+1,bar.get_y()+bar.get_height()/2,str(val),va='center',fontsize=10,color=CN,fontweight='bold')
ax.set_xlabel('Number of Ads',fontsize=10)
ax.set_title('2026 Ad Strength',fontsize=11,color=CN,fontweight='bold',loc='left')
plt.tight_layout()
plt.savefig(f"{CHARTS}/ad_strength_2026.png",dpi=200,bbox_inches='tight'); plt.close()

# Chart 3: Landing page CPA YoY comparison
# Find common pages between years (top pages by combined spend)
all_urls=set(lp25[lp25['Cost']>500].index) | set(lp26[lp26['Cost']>500].index)
lp_combined=[]
for u in all_urls:
    short=str(u).replace('https://dentalbydesign.co.uk','').replace('/promo/google/','/')[:30]
    if short=='' or short=='/': short='(homepage)'
    cpa_25=lp25.loc[u,'CPA'] if u in lp25.index and lp25.loc[u,'Conversions']>0 else 0
    cpa_26=lp26.loc[u,'CPA'] if u in lp26.index and lp26.loc[u,'Conversions']>0 else 0
    if cpa_25>0 or cpa_26>0:
        lp_combined.append((short,cpa_25,cpa_26))
lp_combined.sort(key=lambda x:-max(x[1],x[2]))
lp_combined=lp_combined[:7]

fig,ax=plt.subplots(figsize=(10,4.5))
labels_lp=[x[0] for x in lp_combined]
v25_lp=[x[1] for x in lp_combined]
v26_lp=[x[2] for x in lp_combined]
y=np.arange(len(labels_lp)); w=0.35
b1=ax.barh(y-w/2,v25_lp,w,label='2025',color=CB,alpha=0.8)
b2=ax.barh(y+w/2,v26_lp,w,label='2026',color=CR,alpha=0.8)
ax.set_yticks(y); ax.set_yticklabels(labels_lp,fontsize=9)
ax.invert_yaxis(); ax.legend(fontsize=9)
for bar,val in zip(b1,v25_lp):
    if val>0: ax.text(bar.get_width()+1,bar.get_y()+bar.get_height()/2,f'\u00a3{val:.0f}',va='center',fontsize=9,color=CN)
for bar,val in zip(b2,v26_lp):
    if val>0: ax.text(bar.get_width()+1,bar.get_y()+bar.get_height()/2,f'\u00a3{val:.0f}',va='center',fontsize=9,color=CR)
ax.set_xlabel('CPA (\u00a3)',fontsize=10)
ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f'\u00a3{x:.0f}'))
ax.set_title('Landing Page CPA: 2025 vs 2026',fontsize=11,color=CN,fontweight='bold',loc='left')
plt.tight_layout()
plt.savefig(f"{CHARTS}/lp_cpa_yoy.png",dpi=200,bbox_inches='tight'); plt.close()

# Chart 4: Conversion tracking breakdown (pie)
fig,ax=plt.subplots(figsize=(5,4))
labels=['Real Leads\n(6,182)','Funnel\n(815)','Soft Signals\n(445)','Ghost/KMG\n(1,511)']
sizes=[6182,815,445,1511]
colors_pie=[CG,CB,CY,CR]
explode=(0,0,0.05,0.1)
ax.pie(sizes,explode=explode,labels=labels,colors=colors_pie,autopct='%1.0f%%',startangle=90,textprops={'fontsize':10})
ax.set_title('Reported Conversions: 8,953',fontsize=11,color=CN,fontweight='bold')
plt.tight_layout()
plt.savefig(f"{CHARTS}/conv_breakdown.png",dpi=200,bbox_inches='tight'); plt.close()

# Chart 5: Campaign structure (horizontal bars showing active vs paused spend)
fig,ax=plt.subplots(figsize=(8,3))
camp_labels=['PMax','Search Intent','Brand','Demand Gen']
camp_spend_26=[93550,74024,4133,0]
camp_spend_25=[67586,14977,7663,3065]
x=np.arange(len(camp_labels)); w=0.35
b1=ax.barh(x-w/2,camp_spend_25,w,label='2025',color=CB,alpha=0.8)
b2=ax.barh(x+w/2,camp_spend_26,w,label='2026 (Q1)',color=CR,alpha=0.8)
ax.set_yticks(x); ax.set_yticklabels(camp_labels,fontsize=10)
ax.invert_yaxis(); ax.legend(fontsize=9)
ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f'\u00a3{x:,.0f}'))
for bar,val in zip(b1,camp_spend_25):
    if val>0: ax.text(bar.get_width()+500,bar.get_y()+bar.get_height()/2,f'\u00a3{val:,.0f}',va='center',fontsize=9,color=CN)
for bar,val in zip(b2,camp_spend_26):
    if val>0: ax.text(bar.get_width()+500,bar.get_y()+bar.get_height()/2,f'\u00a3{val:,.0f}',va='center',fontsize=9,color=CR)
ax.set_title('Active Campaign Spend (2025 vs 2026)',fontsize=11,color=CN,fontweight='bold',loc='left')
plt.tight_layout()
plt.savefig(f"{CHARTS}/camp_spend_yoy.png",dpi=200,bbox_inches='tight'); plt.close()

print("Charts done.")

# ══════════════════════════════════════════════════════════════════
# PPTX HELPERS (same as Report 1)
# ══════════════════════════════════════════════════════════════════
prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)

def rect(sl,l,t,w,h,fc,lc=None,lw=None):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h);s.fill.solid();s.fill.fore_color.rgb=fc
    if lc:s.line.color.rgb=lc;s.line.width=lw or Pt(1)
    else:s.line.fill.background()
    return s

def txt(sl,l,t,w,h,c,sz=11,co=BLACK,b=False,a=PP_ALIGN.LEFT):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    r=tf.paragraphs[0].add_run();r.text=c;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
    tf.paragraphs[0].alignment=a;return bx

def multirun(sl,l,t,w,h,runs,a=PP_ALIGN.LEFT):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=a
    for tx,sz,co,b in runs:
        r=p.add_run();r.text=tx;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
    return bx

def multipara(sl,l,t,w,h,pdata):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    for i,(tx,sz,co,b,sp) in enumerate(pdata):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run();r.text=tx;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
        if sp:p.space_after=Pt(sp)
    return bx

def bullets(sl,l,t,w,h,items,sz=11,co=BLACK):
    return multipara(sl,l,t,w,h,[(i,sz,co,False,5) for i in items])

def top_bar(sl,h=Inches(0.07)):
    for i,c in enumerate([BLUE,RED,YELLOW,GREEN]):rect(sl,Inches(i*3.333),Inches(0),Inches(3.333),h,c)
def bot_bar(sl,y=Inches(6.92)):
    for p,c in zip([Inches(0.60),Inches(3.63),Inches(6.67),Inches(9.70)],[BLUE,RED,YELLOW,GREEN]):rect(sl,p,y,Inches(3.03),Inches(0.03),c)
def bot_bar_title(sl):
    for i,c in enumerate([BLUE,RED,YELLOW,GREEN]):rect(sl,Inches(i*3.333),Inches(7.0),Inches(3.333),Inches(0.04),c)
def sidebar(sl): rect(sl,Inches(0),Inches(0.07),Inches(0.12),Inches(7.43),BLUE)
def logo(sl,l=Inches(0.60),t=Inches(0.50),s=Inches(0.65)):
    try:sl.shapes.add_picture(LOGO,l,t,s,s)
    except:pass
def footer(sl,n):
    bot_bar(sl)
    try:sl.shapes.add_picture(LOGO,Inches(0.60),Inches(7.0),Inches(0.22),Inches(0.22))
    except:pass
    txt(sl,Inches(0.90),Inches(7.0),Inches(6),Inches(0.25),"Christopher Hoole  |  christopherhoole.com  |  Confidential",11,NAVY,True)
    txt(sl,Inches(12.23),Inches(7.0),Inches(0.50),Inches(0.25),str(n),11,NAVY,a=PP_ALIGN.RIGHT)
def stitle(sl,t,sub=None):
    txt(sl,Inches(0.60),Inches(0.30),Inches(7),Inches(0.50),t,28,NAVY,True)
    if sub:txt(sl,Inches(0.60),Inches(0.85),Inches(9),Inches(0.30),sub,11,RED)
def badge(sl,t):
    rect(sl,Inches(9.13),Inches(0.30),Inches(3.60),Inches(0.45),WHITE,BLUE,Pt(1))
    txt(sl,Inches(9.23),Inches(0.32),Inches(3.40),Inches(0.40),t,11,BLUE,True,PP_ALIGN.CENTER)
def mcard(sl,l,t,v,lab,ac,w=Inches(2.80),h=Inches(1.05)):
    rect(sl,l,t,w,h,WHITE,BORDER,Pt(0.75));rect(sl,l,t,Inches(0.06),h,ac)
    txt(sl,l+Inches(0.20),t+Inches(0.10),w-Inches(0.30),Inches(0.50),v,22,ac,True)
    txt(sl,l+Inches(0.20),t+Inches(0.60),w-Inches(0.30),Inches(0.35),lab,11,BLACK)
def minicard(sl,l,t,v,lab,ac):
    rect(sl,l,t,Inches(1.83),Inches(1.15),WHITE);rect(sl,l,t,Inches(0.06),Inches(1.15),ac)
    txt(sl,l+Inches(0.15),t+Inches(0.12),Inches(1.58),Inches(0.45),v,22,ac,True)
    txt(sl,l+Inches(0.15),t+Inches(0.60),Inches(1.58),Inches(0.30),lab,11,BLACK)
def ibox(sl,l,t,w,h,ti,body,ac,bg=None):
    rect(sl,l,t,w,h,bg or GREY_BG);rect(sl,l,t,Inches(0.06),h,ac)
    txt(sl,l+Inches(0.25),t+Inches(0.10),w-Inches(0.40),Inches(0.30),ti,14,NAVY,True)
    txt(sl,l+Inches(0.25),t+Inches(0.40),w-Inches(0.40),h-Inches(0.50),body,11,BLACK)
def rbox(sl,l,t,w,h,c):
    rect(sl,l,t,w,h,BLUE_TINT,BLUE,Pt(1))
    txt(sl,l+Inches(0.25),t+Inches(0.05),w-Inches(0.50),h-Inches(0.10),c,11,NAVY,True)
def add_table(sl,l,t,w,h,hdrs,rows,cw=None):
    ts=sl.shapes.add_table(len(rows)+1,len(hdrs),l,t,w,h);tb=ts.table
    if cw:
        for i,c in enumerate(cw):tb.columns[i].width=Inches(c)
    for ci,hd in enumerate(hdrs):
        c=tb.cell(0,ci);c.text=hd;c.fill.solid();c.fill.fore_color.rgb=NAVY
        for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.font.color.rgb=WHITE;p.font.bold=True;p.font.name=FONT
    for ri,rd in enumerate(rows):
        for ci,v in enumerate(rd):
            c=tb.cell(ri+1,ci);c.text=str(v);c.fill.solid();c.fill.fore_color.rgb=GREY_BG if ri%2==0 else WHITE
            for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.font.color.rgb=BLACK;p.font.name=FONT
    return ts
def img(sl,path,l,t,w,h):
    try:sl.shapes.add_picture(path,l,t,w,h)
    except Exception as e:print(f"  Warning: {e}")

# ══════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════
print("Building slides...")

# ── S1: TITLE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=GREY_BG
top_bar(sl);sidebar(sl);logo(sl)
txt(sl,Inches(0.60),Inches(1.40),Inches(5.74),Inches(1.80),"Account Structure\n& Issues",44,NAVY,True)
rect(sl,Inches(0.60),Inches(3.30),Inches(2.50),Inches(0.05),BLUE)
txt(sl,Inches(0.60),Inches(3.55),Inches(5.50),Inches(0.50),"Dental by Design  |  Prodent Group",22,BLUE)
multipara(sl,Inches(0.60),Inches(4.30),Inches(5.50),Inches(1.00),[
    ("Christopher Hoole",11,BLACK,True,2),("Google Ads Specialist  |  April 2026",11,BLACK,False,2),
    ("christopherhoole.com",11,BLUE,False,0)])
rect(sl,Inches(0.60),Inches(5.50),Inches(4.50),Inches(0.50),WHITE,BLUE,Pt(1))
txt(sl,Inches(0.70),Inches(5.52),Inches(4.30),Inches(0.45),"Data: 2025 (full year) + 2026 (Q1, Jan\u2013Apr)",11,BLUE,True)

# Hero — lead with the structural problem
rect(sl,Inches(6.80),Inches(0.50),Inches(5.90),Inches(2.50),WHITE)
rect(sl,Inches(6.80),Inches(0.50),Inches(0.08),Inches(2.50),RED)
txt(sl,Inches(7.10),Inches(0.60),Inches(5.40),Inches(0.35),"Non-Brand CPA Trend",14,NAVY)
txt(sl,Inches(7.10),Inches(1.00),Inches(5.40),Inches(0.80),f"\u00a3{NBCPA25:.0f} \u2192 \u00a3{NBCPA26:.0f}",44,RED,True)
txt(sl,Inches(7.10),Inches(1.80),Inches(5.40),Inches(0.50),f"+{(NBCPA26-NBCPA25)/NBCPA25*100:.0f}% increase year-on-year while spend nearly doubled",11,BLACK)

minicard(sl,Inches(6.80),Inches(3.20),"1,511","Ghost Conversions (KMG)",RED)
minicard(sl,Inches(8.83),Inches(3.20),"47","Total Campaigns",YELLOW)
minicard(sl,Inches(10.87),Inches(3.20),f"{avg_qs_w:.1f}","Avg QS (Weighted)",BLUE)
bot_bar_title(sl)


# ── S2: EXECUTIVE SUMMARY ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06));stitle(sl,"Executive Summary");badge(sl,"2025 + 2026")

mcard(sl,Inches(0.60),Inches(1.20),f"\u00a3{NBCPA25:.0f}\u2192\u00a3{NBCPA26:.0f}","Non-Brand CPA (+{:.0f}%)".format((NBCPA26-NBCPA25)/NBCPA25*100),RED)
mcard(sl,Inches(3.70),Inches(1.20),"1,511","Ghost Conversions (KMG Legacy)",YELLOW)
mcard(sl,Inches(6.80),Inches(1.20),"90%","Ads Rated Poor or Average",RED)
mcard(sl,Inches(9.90),Inches(1.20),f"{avg_qs_w:.1f}/10","Weighted Avg Quality Score",BLUE)

txt(sl,Inches(0.60),Inches(2.60),Inches(6),Inches(0.35),"Key Issues Identified",14,NAVY,True)
bullets(sl,Inches(0.70),Inches(3.00),Inches(11.50),Inches(2.60),[
    "1,511 ghost conversions from legacy KMG landing pages are inflating reported performance. Giulio confirmed these should be ignored \u2014 but they're still set as Primary conversions feeding the bidding algorithm.",
    f"Non-brand CPA has risen {(NBCPA26-NBCPA25)/NBCPA25*100:.0f}% year-on-year (\u00a3{NBCPA25:.0f} \u2192 \u00a3{NBCPA26:.0f}) while monthly spend nearly doubled from ~\u00a322K to ~\u00a358K.",
    "Dengro Offline Leads (4,663) may include non-ad-sourced leads, further inflating Google Ads attribution. Pending verification.",
    f"90% of ads are rated Poor or Average strength. The main Search campaign ({c26[c26['Campaign']=='Dental Implants Intent']['Cost'].sum():,.0f} spend) has Poor-rated ads.",
    f"Ad relevance is Below Average on 71% of scored keywords. Expected CTR is Below Average on 63%. These are the primary Quality Score drivers.",
])

rbox(sl,Inches(0.60),Inches(5.85),Inches(12.13),Inches(0.70),
    "The core issue: the account is scaling spend aggressively but the foundations (conversion tracking, ad quality, keyword Quality Scores) haven't kept pace. Higher spend on weak foundations = rising CPA.")
footer(sl,2)


# ── S3: CONVERSION TRACKING — THE CRITICAL FINDING ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Conversion Tracking Audit","1,511 ghost conversions + 445 soft signals are inflating reported performance")
badge(sl,"Critical Issue")

img(sl,f"{CHARTS}/conv_breakdown.png",Inches(0.60),Inches(1.30),Inches(4.50),Inches(3.50))

# Classification table
txt(sl,Inches(5.40),Inches(1.20),Inches(7.30),Inches(0.30),"Conversion Classification",14,NAVY,True)
add_table(sl,Inches(5.40),Inches(1.55),Inches(7.30),Inches(2.80),
    ["Category","Conversions","% of Total","Status"],
    [["Real Leads (forms, calls, quiz)","6,182","69%","Legitimate"],
     ["Funnel (Dengro bookings + purchases)","815","9%","Legitimate (downstream)"],
     ["Soft Signals (clicks, not actions)","445","5%","Should be Secondary"],
     ["Ghost (KMG legacy \u2014 confirmed ignore)","1,511","17%","Should be removed"]],
    cw=[2.8,1.0,0.8,1.8])

ibox(sl,Inches(5.40),Inches(4.55),Inches(7.30),Inches(1.80),
    "Impact on CPA",
    f"Reported CPA: \u00a3{441000/8953:.0f} (using all 8,953 conversions)\n"
    f"Real leads CPA: \u00a3{441000/6182:.0f} (excluding ghost + soft signals)\n"
    f"Leads + funnel CPA: \u00a3{441000/(6182+815):.0f} (including bookings)\n\n"
    "The bidding algorithm is optimising toward 8,953 events \u2014 including 1,511 that don't exist and 445 that aren't real actions. "
    "This is why CPA appears lower than it actually is.",
    RED,RED_TINT)

ibox(sl,Inches(0.60),Inches(5.00),Inches(4.50),Inches(1.30),
    "Giulio Confirmed",
    "KMG OLD LP actions: \"data that have been preserved but target old lp\" \u2014 confirmed ghost data.\n"
    "Dengro Offline Lead: Lead from CRM.\n"
    "Dengro Offline Booking/Purchase: Booked/completed patient.\n"
    "Every booking has an arbitrary \u00a3300 value.",
    GREEN,GREEN_TINT)
footer(sl,3)


# ── S4: DENGRO & ROAS — THE VALUE PUZZLE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"The ROAS Puzzle","Target ROAS bidding is based on arbitrary \u00a3300 booking values \u2014 not real revenue")
badge(sl,"Bid Strategy")

# Current setup
txt(sl,Inches(0.60),Inches(1.20),Inches(5.80),Inches(0.30),"Current Bid Strategy Setup",14,NAVY,True)
add_table(sl,Inches(0.60),Inches(1.55),Inches(5.80),Inches(1.80),
    ["Campaign","Bid Strategy","Target ROAS","2026 CPA"],
    [["Dental Implants \u2014 PMax","Max Conv Value (tROAS)","30%",f"\u00a3{c26[c26['Campaign']=='Dental Implants - Pmax']['Cost'].sum()/max(c26[c26['Campaign']=='Dental Implants - Pmax']['Conversions'].sum(),0.1):.0f}"],
     ["Dental Implants Intent","Max Conv Value (tROAS)","60%",f"\u00a3{c26[c26['Campaign']=='Dental Implants Intent']['Cost'].sum()/max(c26[c26['Campaign']=='Dental Implants Intent']['Conversions'].sum(),0.1):.0f}"],
     ["Brand","Max Conv Value (tROAS)","10%",f"\u00a3{c26[c26['Campaign']=='Brand']['Cost'].sum()/max(c26[c26['Campaign']=='Brand']['Conversions'].sum(),0.1):.0f}"]],
    cw=[2.2,1.5,0.8,0.8])

ibox(sl,Inches(0.60),Inches(3.60),Inches(5.80),Inches(2.70),
    "Why This Is a Problem",
    "The algorithm is told to maximise conversion VALUE, not conversion VOLUME. "
    "With \u00a3300 arbitrary booking values, a 30% ROAS target means the algorithm wants \u00a3300 in value for every \u00a31,000 spent \u2014 "
    "i.e. one booking per \u00a31,000 spend.\n\n"
    "But only 621 of 8,953 reported conversions are actual bookings. "
    "The algorithm is chasing a small pool of high-value events while spending broadly. "
    "This explains why PMax CPA rose from \u00a326 to \u00a352 \u2014 it's optimising for bookings, not leads, "
    "but the booking signal is sparse and unreliable with an arbitrary value.",
    RED,RED_TINT)

ibox(sl,Inches(6.80),Inches(1.20),Inches(5.80),Inches(2.10),
    "The Dengro Pipeline",
    "Dengro Offline Lead: 4,663 (all CRM leads \u2014 may include non-ad sources)\n"
    "Dengro Offline Booking: 621 (\u00a3300 arbitrary value each = \u00a370,910)\n"
    "Dengro Offline Purchase: 194 (completed treatment = \u00a373,106)\n\n"
    "Lead \u2192 Booking conversion rate: 13.3% (621/4,663)\n"
    "Booking \u2192 Purchase conversion rate: 31.2% (194/621)",
    BLUE,BLUE_TINT)

ibox(sl,Inches(6.80),Inches(3.60),Inches(5.80),Inches(2.70),
    "Recommendation",
    "Option A: Switch to Maximise Conversions with Target CPA. "
    "Optimise for LEAD VOLUME at a controllable cost per lead. "
    "This gives direct control over acquisition cost.\n\n"
    "Option B: Keep Target ROAS but replace the \u00a3300 arbitrary value with actual treatment revenue from Dengro. "
    "If a booking for a single implant is worth \u00a31,695 and a Vivo Bridge is worth \u00a315,990, "
    "the algorithm can intelligently bid more for high-value treatments.\n\n"
    "Either option is better than the current setup.",
    GREEN,GREEN_TINT)
footer(sl,4)


# ── S5: ACCOUNT STRUCTURE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Account Structure","47 campaigns, 4 active \u2014 a legacy of multiple agencies and strategy changes")
badge(sl,"Structure")

img(sl,f"{CHARTS}/camp_spend_yoy.png",Inches(0.60),Inches(1.30),Inches(7.50),Inches(2.80))

# Active campaigns table
txt(sl,Inches(8.40),Inches(1.20),Inches(4.30),Inches(0.30),"Currently Active (4 Campaigns)",14,NAVY,True)
add_table(sl,Inches(8.40),Inches(1.55),Inches(4.30),Inches(2.50),
    ["Campaign","2025 CPA","2026 CPA","Change"],
    [["PMax",f"\u00a3{c25[c25['Campaign']=='Dental Implants - Pmax']['Cost'].sum()/max(c25[c25['Campaign']=='Dental Implants - Pmax']['Conversions'].sum(),0.1):.0f}",f"\u00a3{c26[c26['Campaign']=='Dental Implants - Pmax']['Cost'].sum()/max(c26[c26['Campaign']=='Dental Implants - Pmax']['Conversions'].sum(),0.1):.0f}","+100%"],
     ["Search Intent",f"\u00a3{c25[c25['Campaign']=='Dental Implants Intent']['Cost'].sum()/max(c25[c25['Campaign']=='Dental Implants Intent']['Conversions'].sum(),0.1):.0f}",f"\u00a3{c26[c26['Campaign']=='Dental Implants Intent']['Cost'].sum()/max(c26[c26['Campaign']=='Dental Implants Intent']['Conversions'].sum(),0.1):.0f}","+54%"],
     ["Brand","\u00a317","\u00a318","Stable"],
     ["Demand Gen","\u00a3180","\u2014","Not serving"]],
    cw=[1.2,0.8,0.8,0.7])

# The mess below
txt(sl,Inches(0.60),Inches(4.30),Inches(12),Inches(0.30),"The 43 Paused Campaigns",14,NAVY,True)

ibox(sl,Inches(0.60),Inches(4.65),Inches(3.80),Inches(1.70),
    "Legacy Agencies",
    "KMG (Kau Media Group): 6 campaigns, L87K spend\n"
    "Multiple naming conventions and strategies\n"
    "KMG conversion actions still firing as Primary",
    RED,RED_TINT)

ibox(sl,Inches(4.70),Inches(4.65),Inches(3.80),Inches(1.70),
    "Historical Experiments",
    "Device-split campaigns ('Mobile Only', 'Desktop')\n"
    "Geo-split campaigns ('London', 'Geo')\n"
    "Service campaigns (Invisalign, Veneers, Full Arch, Teeth In A Day)\n"
    "Dating back to 2022",
    YELLOW)

ibox(sl,Inches(8.80),Inches(4.65),Inches(3.90),Inches(1.70),
    "Recommendation",
    "Remove all campaigns paused 6+ months with no active learnings. "
    "This doesn't affect performance but makes the account manageable. "
    "Clean structure is essential for multi-clinic expansion.",
    GREEN,GREEN_TINT)
footer(sl,5)


# ── S6: AD GROUP STRUCTURE (2026) ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Ad Group Structure (2026)","Search Intent campaign: 8 ad groups with wide CPA variation")
badge(sl,"2025 vs 2026")

# 2026 ad groups
ag26=ct(ld(D26,"Ad group report 2026.csv"),'Ad group')
ag26=ag26[~ag26['Ad group status'].astype(str).str.startswith('Total')]
ag26=ag26[ag26['Cost']>0].sort_values('Cost',ascending=False)

add_table(sl,Inches(0.60),Inches(1.30),Inches(12.10),Inches(3.80),
    ["Campaign","Ad Group","Spend","Conv","CPA","CVR"],
    [[r['Campaign'][:25],r['Ad group'][:28],f"\u00a3{r['Cost']:,.0f}",f"{r['Conversions']:.0f}",
      f"\u00a3{r['Cost']/r['Conversions']:.0f}" if r['Conversions']>0 else "N/A",
      f"{r['Conversions']/r['Clicks']*100:.1f}%" if r['Clicks']>0 else "N/A"]
     for _,r in ag26.iterrows()],
    cw=[2.2,2.5,1.2,0.8,0.8,0.8])

ibox(sl,Inches(0.60),Inches(5.30),Inches(5.80),Inches(1.20),
    "Key Finding",
    f"The 'Dental Implants' ad group (\u00a326.6K) runs at \u00a3111 CPA \u2014 the worst in the account. "
    f"Meanwhile 'Financing' (\u00a32.9K) runs at \u00a334 CPA with 24.7% CVR. "
    "Budget is flowing to the least efficient ad groups.",
    RED,RED_TINT)

ibox(sl,Inches(6.80),Inches(5.30),Inches(5.90),Inches(1.20),
    "Recommendation",
    "The 'Financing' and 'Implants Cost' ad groups convert better because they match cost-conscious search intent with the right landing page. "
    "Restructure ad groups around patient intent, not just keyword categories.",
    GREEN,GREEN_TINT)
footer(sl,6)


# ── S7: QUALITY SCORES (both years) ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Quality Score Analysis",f"Weighted average QS: {avg_qs_w25:.1f} (2025) \u2192 {avg_qs_w:.1f} (2026) \u2014 declining")
badge(sl,"2025 vs 2026")

img(sl,f"{CHARTS}/qs_dist_2026.png",Inches(0.60),Inches(1.30),Inches(6.50),Inches(3.00))

# QS components — BOTH YEARS
txt(sl,Inches(7.40),Inches(1.20),Inches(5.30),Inches(0.30),"QS Components (2025 vs 2026)",14,NAVY,True)
add_table(sl,Inches(7.40),Inches(1.55),Inches(5.30),Inches(2.80),
    ["Component","2025 Below","2025 Above","2026 Below","2026 Above"],
    [["Ad Relevance","57 (74%)","12 (16%)","58 (71%)","14 (17%)"],
     ["Expected CTR","51 (66%)","7 (9%)","52 (63%)","8 (10%)"],
     ["Landing Page","3 (4%)","55 (71%)","3 (4%)","58 (71%)"]],
    cw=[1.3,0.9,0.9,0.9,0.9])

ibox(sl,Inches(7.40),Inches(4.55),Inches(5.30),Inches(1.00),
    "Consistent Pattern Both Years",
    "Landing pages: strong (71% Above Average both years). "
    "Ad relevance: weak (74% Below 2025, 71% Below 2026). "
    "The ads don't match what people search for.",
    RED,RED_TINT)

# Low QS keywords with spend (2026)
low_qs=scored26[scored26['QS']<=3].sort_values('Cost',ascending=False).head(5)
txt(sl,Inches(0.60),Inches(4.50),Inches(6),Inches(0.30),"Worst QS Keywords (QS 1-3) \u2014 2026 Spend",12,NAVY,True)
add_table(sl,Inches(0.60),Inches(4.80),Inches(6.50),Inches(1.70),
    ["Keyword","QS","Spend","CPA","Ad Rel."],
    [[r['Keyword'][:35],f"{int(r['QS'])}",f"\u00a3{r['Cost']:,.0f}",
      f"\u00a3{r['Cost']/r['Conversions']:.0f}" if r['Conversions']>0 else "N/A",
      str(r['Ad relevance'])[:15]]
     for _,r in low_qs.iterrows()],
    cw=[2.5,0.4,0.9,0.7,1.5])

# Note about 2025 QS
txt(sl,Inches(0.60),Inches(6.55),Inches(12.10),Inches(0.25),
    "Note: In 2025, most low-QS keywords (1-5) had zero spend because the Search Intent campaign was only active late in the year. The QS issues have emerged as this campaign scaled in 2026.",
    10,GREY_MID)
footer(sl,7)


# ── S8: AD COPY AUDIT ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Ad Copy Audit","90% of ads rated Poor or Average \u2014 the primary driver of low Quality Scores")
badge(sl,"Current State")

img(sl,f"{CHARTS}/ad_strength_2026.png",Inches(0.60),Inches(1.30),Inches(4.50),Inches(3.00))

# Ad strength table — current state (not YoY since it's the same ads)
txt(sl,Inches(5.40),Inches(1.20),Inches(7.30),Inches(0.30),"Current Ad Strength (All Ads in Account)",14,NAVY,True)
add_table(sl,Inches(5.40),Inches(1.55),Inches(5.50),Inches(2.00),
    ["Strength","Count","% of Total","Impact"],
    [["Excellent","4","2%","Keep \u2014 best performers"],
     ["Good","16","8%","Keep \u2014 performing well"],
     ["Average","105","50%","Rewrite \u2014 opportunity to improve"],
     ["Poor","68","32%","Rewrite urgently \u2014 dragging QS down"],
     ["Pending","19","9%","New ads \u2014 monitor"]],
    cw=[1.0,0.7,0.7,2.5])

# Performance by ad strength
ibox(sl,Inches(5.40),Inches(3.80),Inches(7.30),Inches(1.00),
    "Critical Finding",
    "The main Search Intent campaign (\u00a315K spend in 2025, \u00a374K in 2026) has Poor-rated ads. "
    "This directly causes Below Average ad relevance on 71% of keywords. "
    "Poor ads = low QS = higher CPCs = higher CPA. This has been the case across both years.",
    RED,RED_TINT)

# The spend impact
ibox(sl,Inches(0.60),Inches(5.00),Inches(5.50),Inches(1.50),
    "Spend on Poor/Average Ads",
    f"2025: \u00a3{S25:,.0f} total spend \u2014 same ads running throughout\n"
    f"2026: \u00a3{S26:,.0f} total spend \u2014 same ads, but at nearly double the budget\n\n"
    "The ads haven't been rewritten as spend has scaled. "
    "The same Poor-rated ads that ran at \u00a315K/month are now running at \u00a358K/month.",
    YELLOW)

rbox(sl,Inches(6.40),Inches(5.00),Inches(6.30),Inches(0.60),
    "Recommendation: Rewrite all Poor and Average ads. Match headlines to search intent. Test 2-3 RSA variants per ad group.")

txt(sl,Inches(6.40),Inches(5.70),Inches(6.30),Inches(0.30),
    "This is the #1 lever for reducing CPA. Better ads = higher QS = lower CPCs = more patients for the same budget.",
    11,NAVY,True)
footer(sl,8)


# ── S9: LANDING PAGE PERFORMANCE (YoY chart) ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Landing Page Performance","CPA varies 7x+ across landing pages \u2014 and has worsened year-on-year on the top pages")
badge(sl,"2025 vs 2026")

img(sl,f"{CHARTS}/lp_cpa_yoy.png",Inches(0.60),Inches(1.30),Inches(8.00),Inches(3.80))

ibox(sl,Inches(8.90),Inches(1.30),Inches(3.80),Inches(1.60),
    "Best Performer",
    "/dental-implants-offer:\n2025: \u00a367 CPA, 6.8% CVR\n2026: \u00a334 CPA, 24.7% CVR\n\n"
    "Improved as Financing ad group sends cost-conscious traffic here.",
    GREEN,GREEN_TINT)

ibox(sl,Inches(8.90),Inches(3.10),Inches(3.80),Inches(2.00),
    "Worst Performer",
    "/teeth-implants (highest spend):\n2025: \u00a382 CPA, 9.3% CVR\n2026: \u00a3111 CPA, 8.5% CVR\n\n"
    "CPA risen 35% YoY despite receiving \u00a331K in 2026 \u2014 the most of any page.",
    RED,RED_TINT)

ibox(sl,Inches(0.60),Inches(5.30),Inches(5.80),Inches(1.20),
    "The Good News",
    "QS data confirms landing pages score well (71% Above Average both years). "
    "The pages aren't the problem \u2014 it's which traffic reaches which page.",
    BLUE,BLUE_TINT)

rbox(sl,Inches(6.80),Inches(5.30),Inches(5.90),Inches(0.60),
    "Recommendation: Route traffic to the best-converting pages. /dental-implants-offer at \u00a334 CPA should receive more traffic than /teeth-implants at \u00a3111.")
footer(sl,9)


# ── S10: CHANGE HISTORY ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Change History","2,543 changes \u2014 89% by one person, 24 by auto-apply")
badge(sl,"Account Management")

# User breakdown
txt(sl,Inches(0.60),Inches(1.20),Inches(5.80),Inches(0.30),"Who Manages the Account",14,NAVY,True)
add_table(sl,Inches(0.60),Inches(1.55),Inches(5.80),Inches(2.80),
    ["User","Changes","Role"],
    [["giulio.zanchetta@majestico.it","2,264 (89%)","In-house (Giulio)"],
     ["taranjeet@kaumediagroup.com","166 (7%)","Previous agency (KMG)"],
     ["nav@kaumediagroup.com","40","Previous agency (KMG)"],
     ["Recommendations Auto-Apply","24","Google system"],
     ["junad@kaumediagroup.com","22","Previous agency (KMG)"],
     ["sarahm@kaumediagroup.com","14","Previous agency (KMG)"],
     ["Google Ads System","4","Google system"]],
    cw=[2.5,1.0,1.5])

# Change categories
txt(sl,Inches(6.80),Inches(1.20),Inches(5.90),Inches(0.30),"What's Being Changed",14,NAVY,True)
add_table(sl,Inches(6.80),Inches(1.55),Inches(5.90),Inches(2.80),
    ["Category","Count","Assessment"],
    [["Keywords added/paused","1,315","Highest activity \u2014 frequent keyword management"],
     ["Campaign changes","715","Many campaigns created and paused over time"],
     ["Targeting","420","Location, device, audience changes"],
     ["Ad changes","418","Ad creation and modification"],
     ["Bid strategy","223","Multiple strategy changes over time"],
     ["Budget changes","198","Frequent budget adjustments"],
     ["Assets (PMax)","197","PMax asset group management"],
     ["Negative keywords","182","Negative keyword additions"]],
    cw=[2.0,0.8,2.5])

ibox(sl,Inches(0.60),Inches(4.60),Inches(5.80),Inches(1.80),
    "Key Observation",
    "Giulio is doing the work of a full-time PPC specialist alongside his other responsibilities. "
    "2,264 changes shows significant effort and engagement with the account. "
    "The issue isn't lack of activity \u2014 it's that the foundations (conversion tracking, ad quality, "
    "negative keywords) need specialist attention that goes beyond day-to-day management.",
    BLUE,BLUE_TINT)

ibox(sl,Inches(6.80),Inches(4.60),Inches(5.90),Inches(1.80),
    "Auto-Apply: 24 Changes",
    "Google auto-removed keywords including [low cost dental implants], [cheapest dental implants], "
    "[all on 6 dental implants], [implant teeth cost], and [replace teeth]. "
    "Several of these are high-intent keywords that should be in the account. "
    "Auto-apply is removing keywords without human oversight \u2014 recommend disabling.",
    RED,RED_TINT)
footer(sl,10)


# ── S11: ISSUES SUMMARY (3 priority cards) ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Issues Summary \u2014 Prioritised")
badge(sl,"12 Issues Identified")

# CRITICAL card
rect(sl,Inches(0.60),Inches(1.20),Inches(3.80),Inches(5.20),WHITE,BORDER,Pt(0.75))
rect(sl,Inches(0.60),Inches(1.20),Inches(3.80),Pt(5),RED)
rect(sl,Inches(0.60),Inches(1.20),Inches(0.06),Inches(5.20),RED)
txt(sl,Inches(0.80),Inches(1.35),Inches(3.40),Inches(0.35),"Critical (Fix Immediately)",14,RED,True)
multipara(sl,Inches(0.80),Inches(1.80),Inches(3.40),Inches(4.40),[
    ("1. Ghost conversions (KMG)",11,NAVY,True,3),
    ("1,511 phantom conversions misleading the bidding algorithm. Confirmed by Giulio as legacy data.",10,BLACK,False,10),
    ("2. ROAS on arbitrary \u00a3300 values",11,NAVY,True,3),
    ("Algorithm chasing sparse booking signals instead of optimising for lead volume.",10,BLACK,False,10),
    ("3. 90% of ads Poor/Average",11,NAVY,True,3),
    ("Directly causing Below Average ad relevance on 71% of keywords. Same ads, never rewritten.",10,BLACK,False,10),
    ("4. Dengro Lead attribution",11,NAVY,True,3),
    ("May include non-ad leads, inflating Google Ads credit. Pending verification.",10,BLACK,False,0),
])

# IMPORTANT card
rect(sl,Inches(4.70),Inches(1.20),Inches(3.80),Inches(5.20),WHITE,BORDER,Pt(0.75))
rect(sl,Inches(4.70),Inches(1.20),Inches(3.80),Pt(5),YELLOW)
rect(sl,Inches(4.70),Inches(1.20),Inches(0.06),Inches(5.20),YELLOW)
txt(sl,Inches(4.90),Inches(1.35),Inches(3.40),Inches(0.35),"Important (Fix in 30 Days)",14,YELLOW,True)
multipara(sl,Inches(4.90),Inches(1.80),Inches(3.40),Inches(4.40),[
    ("5. No device bid adjustments",11,NAVY,True,3),
    ("Desktop \u00a367\u2192\u00a392 vs mobile \u00a342\u2192\u00a356. Gap widening.",10,BLACK,False,10),
    ("6. No ad schedule adjustments",11,NAVY,True,3),
    ("7am: \u00a356\u2192\u00a374 CPA (worst both years). No scheduling exists.",10,BLACK,False,10),
    ("7. No day-of-week adjustments",11,NAVY,True,3),
    ("Sunday: \u00a345\u2192\u00a368 CPA, 4.1% CVR (lowest). Emerging problem.",10,BLACK,False,10),
    ("8. 47 campaigns (43 paused)",11,NAVY,True,3),
    ("Structural confusion from legacy agencies. Remove all paused 6+ months.",10,BLACK,False,10),
    ("9. Auto-apply enabled",11,NAVY,True,3),
    ("24 keyword removals without human review. Disable immediately.",10,BLACK,False,0),
])

# MONITOR card
rect(sl,Inches(8.80),Inches(1.20),Inches(3.90),Inches(5.20),WHITE,BORDER,Pt(0.75))
rect(sl,Inches(8.80),Inches(1.20),Inches(3.90),Pt(5),BLUE)
rect(sl,Inches(8.80),Inches(1.20),Inches(0.06),Inches(5.20),BLUE)
txt(sl,Inches(9.00),Inches(1.35),Inches(3.50),Inches(0.35),"Monitor (Ongoing)",14,BLUE,True)
multipara(sl,Inches(9.00),Inches(1.80),Inches(3.50),Inches(4.40),[
    ("10. 784 broad match negatives",11,NAVY,True,3),
    ("Account-level lists may block legitimate traffic. Full audit needed. Replace with phrase/exact.",10,BLACK,False,10),
    ("11. Landing page CPA variance",11,NAVY,True,3),
    ("7x+ range (\u00a318\u2013\u00a3142 in 2026). Optimise traffic routing to best-converting pages.",10,BLACK,False,10),
    ("12. Quality Score declining",11,NAVY,True,3),
    ("Weighted avg: 6.9\u21926.4. Ad relevance is the driver. Ad rewrites will address this.",10,BLACK,False,0),
])

rbox(sl,Inches(0.60),Inches(6.15),Inches(12.10),Inches(0.40),
    "The 4 Critical issues affect every conversion and every pound spent. Fixing them is the foundation for everything else.")
footer(sl,11)


# ── S12: NEXT STEPS ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=GREY_BG
top_bar(sl);sidebar(sl);logo(sl)
txt(sl,Inches(0.60),Inches(1.40),Inches(5.50),Inches(0.80),"What's Next",44,NAVY,True)
rect(sl,Inches(0.60),Inches(2.20),Inches(2.50),Inches(0.05),BLUE)
txt(sl,Inches(0.60),Inches(2.40),Inches(5.50),Inches(0.40),"These issues become the foundation for Report 3",14,NAVY)

wins=[
    ("1","Remove KMG ghost conversion actions","Stop 1,511 phantom conversions from misleading the algorithm",RED),
    ("2","Move soft signals to Secondary","Online Booking Click, Phone Click, WhatsApp Click \u2014 not real leads",RED),
    ("3","Verify Dengro Offline Lead attribution","Confirm whether non-ad leads are being attributed to Google Ads",YELLOW),
    ("4","Rewrite all Poor/Average ads","The #1 lever for improving QS and reducing CPA",BLUE),
    ("5","Evaluate bid strategy: tCPA vs tROAS","Replace arbitrary L300 values with either real revenue or target CPA",BLUE),
    ("6","Clean up 43 paused campaigns","Remove structural clutter from legacy agencies",GREEN),
]

for i,(num,action,detail,color) in enumerate(wins):
    y=Inches(2.85+i*0.60)
    rect(sl,Inches(0.60),y,Inches(5.50),Inches(0.50),WHITE)
    rect(sl,Inches(0.60),y,Inches(0.06),Inches(0.50),color)
    multirun(sl,Inches(0.85),y+Inches(0.08),Inches(5.00),Inches(0.38),[
        (f"{num}.  ",11,color,True),(action,11,NAVY,True)])

rect(sl,Inches(6.80),Inches(0.50),Inches(5.90),Inches(2.50),WHITE)
rect(sl,Inches(6.80),Inches(0.50),Inches(0.08),Inches(2.50),GREEN)
txt(sl,Inches(7.10),Inches(0.60),Inches(5.40),Inches(0.35),"Coming Next",14,NAVY)
txt(sl,Inches(7.10),Inches(1.00),Inches(5.40),Inches(0.80),"Report 3:\nRestructure Plan",40,GREEN,True)
txt(sl,Inches(7.10),Inches(1.90),Inches(5.40),Inches(0.60),
    "Proposed campaign structure, bid strategy, conversion tracking fixes, ad copy plan, and projected CPA targets.",11,BLACK)

minicard(sl,Inches(6.80),Inches(3.20),"4","Critical Issues",RED)
minicard(sl,Inches(8.83),Inches(3.20),"5","Important Issues",YELLOW)
minicard(sl,Inches(10.87),Inches(3.20),"Report 3","Coming Next",GREEN)
bot_bar_title(sl)


# ══════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════
output=os.path.join(REPORTS,'02_account_structure_report_v3.pptx')
prs.save(output)
print(f"\nSaved to: {output}")
print(f"Slides: {len(prs.slides)}")
print("Done!")
