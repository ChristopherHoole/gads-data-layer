"""
Dental by Design - Report 3: Restructure & Growth Plan
Design system: exact match to Objection Experts winning reports.
Both 2025 and 2026 data. Aspirational tone — this is the "hire me" report.
"""
import os, sys, warnings
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

warnings.filterwarnings('ignore')

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
D26 = os.path.join(SCRIPT_DIR, 'data', '2026')
CHARTS = os.path.join(SCRIPT_DIR, 'charts')
REPORTS = os.path.join(SCRIPT_DIR, 'reports')
LOGO = os.path.join(SCRIPT_DIR, 'act_logo_official.png')
os.makedirs(CHARTS, exist_ok=True)

# ── Design tokens ──
BLUE=RGBColor(0x42,0x85,0xF4); RED=RGBColor(0xEA,0x43,0x35)
YELLOW=RGBColor(0xFB,0xBC,0x05); GREEN=RGBColor(0x34,0xA8,0x53)
NAVY=RGBColor(0x1A,0x23,0x7E); BLACK=RGBColor(0x1A,0x1A,0x1A)
GREY_BG=RGBColor(0xF5,0xF6,0xFA); WHITE=RGBColor(0xFF,0xFF,0xFF)
BORDER=RGBColor(0xE2,0xE8,0xF0)
BLUE_TINT=RGBColor(0xE8,0xF0,0xFE); RED_TINT=RGBColor(0xFC,0xE8,0xE6)
GREEN_TINT=RGBColor(0xE6,0xF4,0xEA); GREY_MID=RGBColor(0xCB,0xD5,0xE1)
FONT='Calibri'
CB='#4285F4'; CR='#EA4335'; CY='#FBBC05'; CG='#34A853'; CN='#1A237E'

plt.rcParams.update({'font.family':'Calibri','font.size':11,
    'axes.spines.top':False,'axes.spines.right':False,
    'axes.edgecolor':'#E2E8F0','axes.linewidth':0.8,
    'figure.facecolor':'white','axes.facecolor':'white',
    'grid.color':'#F0F0F0','grid.linewidth':0.5})

# ── Verified headline numbers (hardcoded from analysis) ──
S25=265869; V25=6070; CPA25=44; NBCPA25=46
S26=175467; V26=2889; CPA26=61; NBCPA26=64
NB26_SPEND=171334; NB26_CONV=2660
MONTHLY_RUN=53990
TARGET_CPA=50

print("Building Report 3...")

# ══════════════════════════════════════════════════════════════════
# CHARTS
# ══════════════════════════════════════════════════════════════════
print("Generating charts...")

# Chart 1: Projected results — before vs after
fig,ax=plt.subplots(figsize=(8,4))
metrics=['Non-Brand CPA','Monthly Leads\n(non-brand)','Monthly Waste']
current=[64, int(NB26_CONV/3.25), int(59349/3.25)]
projected=[50, int(NB26_SPEND/TARGET_CPA/3.25), int(59349*0.4/3.25)]
x=np.arange(len(metrics)); w=0.35
b1=ax.bar(x-w/2,current,w,label='Current (2026)',color=CR,alpha=0.8)
b2=ax.bar(x+w/2,projected,w,label='Projected (90 days)',color=CG,alpha=0.8)
for bar,val in zip(b1,current):
    lbl = f'\u00a3{val}' if val < 200 else f'{val:,}'
    ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+5,lbl,ha='center',fontsize=11,color=CN,fontweight='bold')
for bar,val in zip(b2,projected):
    lbl = f'\u00a3{val}' if val < 200 else f'{val:,}'
    ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+5,lbl,ha='center',fontsize=11,color=CG,fontweight='bold')
ax.set_xticks(x); ax.set_xticklabels(metrics,fontsize=10)
ax.legend(fontsize=10); ax.set_ylim(0,max(current)*1.25)
plt.tight_layout()
plt.savefig(f"{CHARTS}/r3_projections.png",dpi=200,bbox_inches='tight'); plt.close()

# Chart 2: Implementation timeline
fig,ax=plt.subplots(figsize=(10,3))
phases=[
    ('Week 1','Conversion tracking\nfix + neg keywords',CR),
    ('Week 2-4','Ad rewrites +\nbid adjustments',CY),
    ('Month 2','Geo restructure +\nkeyword optimization',CB),
    ('Month 3','Scale + multi-clinic\nrollout planning',CG),
]
for i,(label,desc,color) in enumerate(phases):
    ax.barh(0,1,left=i,color=color,height=0.6,edgecolor='white',linewidth=2)
    ax.text(i+0.5,0.05,label,ha='center',va='center',fontsize=11,color='white',fontweight='bold')
    ax.text(i+0.5,-0.55,desc,ha='center',va='center',fontsize=9,color=CN)
ax.set_xlim(-0.1,4.1); ax.set_ylim(-1,0.5)
ax.axis('off')
ax.set_title('90-Day Implementation Roadmap',fontsize=12,color=CN,fontweight='bold',loc='left')
plt.tight_layout()
plt.savefig(f"{CHARTS}/r3_timeline.png",dpi=200,bbox_inches='tight'); plt.close()

# Chart 3: Device bid adjustment visual
fig,ax=plt.subplots(figsize=(6,3))
devices=['Mobile','Desktop','Tablet']
current_cpa=[56,92,73]
adj=['No change','-25% to -35%','-20% to -30%']
colors=[CG,CR,CY]
bars=ax.barh(range(3),current_cpa,color=colors,height=0.5)
ax.set_yticks(range(3)); ax.set_yticklabels(devices,fontsize=10)
ax.invert_yaxis()
for bar,cpa,a in zip(bars,current_cpa,adj):
    ax.text(bar.get_width()+1,bar.get_y()+bar.get_height()/2,f'\u00a3{cpa} CPA \u2192 {a}',va='center',fontsize=10,color=CN)
ax.set_xlabel('Current CPA (\u00a3)',fontsize=10)
ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f'\u00a3{x:.0f}'))
ax.set_xlim(0,140)
ax.set_title('Proposed Device Bid Adjustments',fontsize=11,color=CN,fontweight='bold',loc='left')
plt.tight_layout()
plt.savefig(f"{CHARTS}/r3_device_adj.png",dpi=200,bbox_inches='tight'); plt.close()

print("Charts done.")

# ══════════════════════════════════════════════════════════════════
# PPTX HELPERS (same as Reports 1+2)
# ══════════════════════════════════════════════════════════════════
prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)

def rect(sl,l,t,w,h,fc,lc=None,lw=None):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h);s.fill.solid();s.fill.fore_color.rgb=fc
    if lc:s.line.color.rgb=lc;s.line.width=lw or Pt(1)
    else:s.line.fill.background()
def txt(sl,l,t,w,h,c,sz=11,co=BLACK,b=False,a=PP_ALIGN.LEFT):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    r=tf.paragraphs[0].add_run();r.text=c;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
    tf.paragraphs[0].alignment=a;return bx
def multirun(sl,l,t,w,h,runs,a=PP_ALIGN.LEFT):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=a
    for tx,sz,co,b in runs:
        r=p.add_run();r.text=tx;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
def multipara(sl,l,t,w,h,pdata):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    for i,(tx,sz,co,b,sp) in enumerate(pdata):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run();r.text=tx;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
        if sp:p.space_after=Pt(sp)
def bullets(sl,l,t,w,h,items,sz=11,co=BLACK):
    multipara(sl,l,t,w,h,[(i,sz,co,False,5) for i in items])
def top_bar(sl,h=Inches(0.07)):
    for i,c in enumerate([BLUE,RED,YELLOW,GREEN]):rect(sl,Inches(i*3.333),Inches(0),Inches(3.333),h,c)
def bot_bar(sl,y=Inches(6.92)):
    for p,c in zip([Inches(0.60),Inches(3.63),Inches(6.67),Inches(9.70)],[BLUE,RED,YELLOW,GREEN]):rect(sl,p,y,Inches(3.03),Inches(0.03),c)
def bot_bar_title(sl):
    for i,c in enumerate([BLUE,RED,YELLOW,GREEN]):rect(sl,Inches(i*3.333),Inches(7.0),Inches(3.333),Inches(0.04),c)
def sidebar(sl): rect(sl,Inches(0),Inches(0.07),Inches(0.12),Inches(7.43),BLUE)
def logo(sl,l=Inches(0.60),t=Inches(0.50),s=Inches(0.65)):
    try:sl.shapes.add_picture(LOGO,l,t,s,s)
    except:pass
def footer(sl,n):
    bot_bar(sl)
    try:sl.shapes.add_picture(LOGO,Inches(0.60),Inches(7.0),Inches(0.22),Inches(0.22))
    except:pass
    txt(sl,Inches(0.90),Inches(7.0),Inches(6),Inches(0.25),"Christopher Hoole  |  christopherhoole.com  |  Confidential",11,NAVY,True)
    txt(sl,Inches(12.23),Inches(7.0),Inches(0.50),Inches(0.25),str(n),11,NAVY,a=PP_ALIGN.RIGHT)
def stitle(sl,t,sub=None):
    txt(sl,Inches(0.60),Inches(0.30),Inches(7),Inches(0.50),t,28,NAVY,True)
    if sub:txt(sl,Inches(0.60),Inches(0.85),Inches(9),Inches(0.30),sub,11,RED)
def badge(sl,t):
    rect(sl,Inches(9.13),Inches(0.30),Inches(3.60),Inches(0.45),WHITE,BLUE,Pt(1))
    txt(sl,Inches(9.23),Inches(0.32),Inches(3.40),Inches(0.40),t,11,BLUE,True,PP_ALIGN.CENTER)
def mcard(sl,l,t,v,lab,ac,w=Inches(2.80),h=Inches(1.05)):
    rect(sl,l,t,w,h,WHITE,BORDER,Pt(0.75));rect(sl,l,t,Inches(0.06),h,ac)
    txt(sl,l+Inches(0.20),t+Inches(0.10),w-Inches(0.30),Inches(0.50),v,22,ac,True)
    txt(sl,l+Inches(0.20),t+Inches(0.60),w-Inches(0.30),Inches(0.35),lab,11,BLACK)
def minicard(sl,l,t,v,lab,ac):
    rect(sl,l,t,Inches(1.83),Inches(1.15),WHITE);rect(sl,l,t,Inches(0.06),Inches(1.15),ac)
    txt(sl,l+Inches(0.15),t+Inches(0.12),Inches(1.58),Inches(0.45),v,22,ac,True)
    txt(sl,l+Inches(0.15),t+Inches(0.60),Inches(1.58),Inches(0.30),lab,11,BLACK)
def ibox(sl,l,t,w,h,ti,body,ac,bg=None):
    rect(sl,l,t,w,h,bg or GREY_BG);rect(sl,l,t,Inches(0.06),h,ac)
    txt(sl,l+Inches(0.25),t+Inches(0.10),w-Inches(0.40),Inches(0.30),ti,14,NAVY,True)
    txt(sl,l+Inches(0.25),t+Inches(0.40),w-Inches(0.40),h-Inches(0.50),body,11,BLACK)
def rbox(sl,l,t,w,h,c):
    rect(sl,l,t,w,h,BLUE_TINT,BLUE,Pt(1))
    txt(sl,l+Inches(0.25),t+Inches(0.05),w-Inches(0.50),h-Inches(0.10),c,11,NAVY,True)
def add_table(sl,l,t,w,h,hdrs,rows,cw=None):
    ts=sl.shapes.add_table(len(rows)+1,len(hdrs),l,t,w,h);tb=ts.table
    if cw:
        for i,c in enumerate(cw):tb.columns[i].width=Inches(c)
    for ci,hd in enumerate(hdrs):
        c=tb.cell(0,ci);c.text=hd;c.fill.solid();c.fill.fore_color.rgb=NAVY
        for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.font.color.rgb=WHITE;p.font.bold=True;p.font.name=FONT
    for ri,rd in enumerate(rows):
        for ci,v in enumerate(rd):
            c=tb.cell(ri+1,ci);c.text=str(v);c.fill.solid();c.fill.fore_color.rgb=GREY_BG if ri%2==0 else WHITE
            for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.font.color.rgb=BLACK;p.font.name=FONT
    return ts
def img(sl,path,l,t,w,h):
    try:sl.shapes.add_picture(path,l,t,w,h)
    except Exception as e:print(f"  Warning: {e}")
def ncircle(sl,l,t,num,col):
    ns=sl.shapes.add_shape(MSO_SHAPE.OVAL,l,t,Inches(0.45),Inches(0.45))
    ns.fill.solid();ns.fill.fore_color.rgb=col;ns.line.fill.background()
    tf=ns.text_frame;tf.paragraphs[0].alignment=PP_ALIGN.CENTER;tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    r=tf.paragraphs[0].add_run();r.text=str(num);r.font.size=Pt(16);r.font.color.rgb=WHITE;r.font.bold=True;r.font.name=FONT


# ══════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════
print("Building slides...")

# ── S1: TITLE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=GREY_BG
top_bar(sl);sidebar(sl);logo(sl)
txt(sl,Inches(0.60),Inches(1.40),Inches(5.74),Inches(1.80),"Restructure &\nGrowth Plan",44,NAVY,True)
rect(sl,Inches(0.60),Inches(3.30),Inches(2.50),Inches(0.05),BLUE)
txt(sl,Inches(0.60),Inches(3.55),Inches(5.50),Inches(0.50),"Dental by Design  |  Prodent Group",22,BLUE)
multipara(sl,Inches(0.60),Inches(4.30),Inches(5.50),Inches(1.00),[
    ("Christopher Hoole",11,BLACK,True,2),("Google Ads Specialist  |  April 2026",11,BLACK,False,2),
    ("christopherhoole.com",11,BLUE,False,0)])
rect(sl,Inches(0.60),Inches(5.50),Inches(4.50),Inches(0.50),WHITE,BLUE,Pt(1))
txt(sl,Inches(0.70),Inches(5.52),Inches(4.30),Inches(0.45),"Based on data: 2025 (full year) + 2026 (Q1)",11,BLUE,True)

# Hero
rect(sl,Inches(6.80),Inches(0.50),Inches(5.90),Inches(2.50),WHITE)
rect(sl,Inches(6.80),Inches(0.50),Inches(0.08),Inches(2.50),GREEN)
txt(sl,Inches(7.10),Inches(0.60),Inches(5.40),Inches(0.35),"Target Non-Brand CPA",14,NAVY)
txt(sl,Inches(7.10),Inches(1.00),Inches(5.40),Inches(0.80),f"\u00a3{TARGET_CPA}",44,GREEN,True)
txt(sl,Inches(7.10),Inches(1.80),Inches(5.40),Inches(0.50),f"Down from \u00a3{NBCPA26} \u2014 a {(NBCPA26-TARGET_CPA)/NBCPA26*100:.0f}% reduction within 90 days",11,BLACK)

minicard(sl,Inches(6.80),Inches(3.20),"+29%","More Leads (Same Budget)",GREEN)
minicard(sl,Inches(8.83),Inches(3.20),"90 Days","To Full Impact",BLUE)
minicard(sl,Inches(10.87),Inches(3.20),f"\u00a3{12000:,}/mo","Potential Monthly Saving",YELLOW)
bot_bar_title(sl)


# ── S2: EXECUTIVE SUMMARY ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06));stitle(sl,"Executive Summary");badge(sl,"Restructure Plan")

mcard(sl,Inches(0.60),Inches(1.20),f"\u00a3{TARGET_CPA}","Target Non-Brand CPA",GREEN)
mcard(sl,Inches(3.70),Inches(1.20),"+236/mo","Additional Leads (Same Budget)",BLUE)
mcard(sl,Inches(6.80),Inches(1.20),"90 Days","Time to Full Impact",YELLOW)
mcard(sl,Inches(9.90),Inches(1.20),f"\u00a3{MONTHLY_RUN:,.0f}","Current Monthly Spend",NAVY)

txt(sl,Inches(0.60),Inches(2.60),Inches(6),Inches(0.35),"Restructure Plan Summary",14,NAVY,True)
bullets(sl,Inches(0.70),Inches(3.00),Inches(11.50),Inches(2.60),[
    "Week 1: Fix conversion tracking (remove 1,511 KMG ghosts, move soft signals to Secondary), implement negative keyword overhaul (9-list system), disable auto-apply.",
    "Weeks 2-4: Rewrite all Poor/Average ads (the #1 CPA lever), implement device + schedule + day-of-week bid adjustments, evaluate bid strategy (tCPA vs tROAS with real values).",
    "Month 2: Restructure geo targeting from radius to postcodes, keyword optimization based on first month of clean data, landing page traffic routing improvements.",
    "Month 3: Scale winning campaigns, establish per-clinic benchmarks, begin multi-clinic rollout planning. Full reporting framework.",
])

rbox(sl,Inches(0.60),Inches(5.85),Inches(12.13),Inches(0.70),
    f"Conservative projection: Reducing non-brand CPA from \u00a3{NBCPA26} to \u00a3{TARGET_CPA} within 90 days. "
    f"At the current \u00a3{MONTHLY_RUN:,.0f}/month spend, this means 236 additional patient leads per month or \u00a312,000/month in savings.")
footer(sl,2)


# ── S3: CONVERSION TRACKING FIX ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Step 1: Fix Conversion Tracking","The foundation \u2014 everything else depends on clean data")
badge(sl,"Week 1")

add_table(sl,Inches(0.60),Inches(1.30),Inches(7.50),Inches(3.30),
    ["Conversion Action","Current","Proposed","Why"],
    [["KMG OLD LP | Book A Free Consult","Primary","Remove","Ghost data \u2014 confirmed by Giulio. 1,250 phantom conversions."],
     ["KMG OLD LP | Website Phone Call","Primary","Remove","Ghost data \u2014 legacy landing page. 57 phantom conversions."],
     ["KMG | Call Extension","Primary","Remove","Legacy agency call extension. 204 phantom conversions."],
     ["Online Booking Click","Primary","Secondary","Click, not a booking. 264 soft signals misleading algorithm."],
     ["Phone Click","Primary","Secondary","Click, not a call. 123 soft signals."],
     ["WhatsApp Click","Primary","Secondary","Click, not a conversation. 48 soft signals."],
     ["Conversation Started","Primary","Secondary","Chat opened, not a lead. 10 soft signals."],
     ["Dengro Offline Lead","Primary","Keep (verify)","Real CRM lead \u2014 but verify it only counts ad-sourced leads."],
     ["Dengro Offline Booking","Primary","Keep","Real booking. \u00a3300 arbitrary value \u2014 replace with actual treatment value."],
     ["Calls from the ads","Primary","Keep","Real phone calls. 965 legitimate conversions."],
    ],
    cw=[2.0,0.7,0.8,3.5])

ibox(sl,Inches(8.40),Inches(1.30),Inches(4.30),Inches(1.50),
    "Impact",
    "Removing 1,511 ghost + 445 soft signal conversions will cause reported CPA to rise from \u00a349 to ~\u00a371. "
    "This isn't a deterioration \u2014 it's seeing the real numbers for the first time.",
    RED,RED_TINT)

ibox(sl,Inches(8.40),Inches(3.00),Inches(4.30),Inches(1.60),
    "ROAS Values",
    "Replace the arbitrary \u00a3300 booking value with actual treatment revenue:\n"
    "\u2022 Single implant: \u00a31,695\n"
    "\u2022 All-on-4: \u00a39,995\n"
    "\u2022 Vivo Bridge: \u00a315,990\n\n"
    "This lets the algorithm bid more for high-value treatments.",
    GREEN,GREEN_TINT)

rbox(sl,Inches(0.60),Inches(4.85),Inches(12.10),Inches(0.50),
    "This is a 30-minute fix with outsized impact. Clean tracking = better algorithm decisions = lower CPA.")
footer(sl,3)


# ── S4: NEGATIVE KEYWORD OVERHAUL ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Step 2: Negative Keyword Overhaul","Replace 784 broad match negatives with a structured 9-list system")
badge(sl,"Week 1")

# Current vs proposed
ibox(sl,Inches(0.60),Inches(1.30),Inches(5.80),Inches(2.50),
    "Current Setup",
    "2,192 total negative keywords across the account:\n"
    "\u2022 784 broad match at account level (dangerous \u2014 may block converting traffic)\n"
    "\u2022 770 phrase match\n"
    "\u2022 425 exact match\n"
    "\u2022 Scattered across 35 campaigns with no structure\n\n"
    "Auto-apply has removed 24 keywords including high-intent terms like "
    "[low cost dental implants] and [cheapest dental implants].",
    RED,RED_TINT)

ibox(sl,Inches(6.80),Inches(1.30),Inches(5.80),Inches(2.50),
    "Proposed: 9-List System",
    "9 shared account-level lists organised by word count + match type:\n\n"
    "\u2022 1 word \u2014 phrase + exact\n"
    "\u2022 2 words \u2014 phrase + exact\n"
    "\u2022 3 words \u2014 phrase + exact\n"
    "\u2022 4 words \u2014 exact only\n"
    "\u2022 5+ words \u2014 exact only\n\n"
    "All new negatives added as phrase or exact match only. No broad match. "
    "Weekly search term reviews to keep lists current.",
    GREEN,GREEN_TINT)

rbox(sl,Inches(0.60),Inches(4.05),Inches(12.10),Inches(0.50),
    f"The 9-list system blocked \u00a3{59349:,} of search term waste in Q1 2026. A structured negative keyword framework is the most cost-effective optimization available.")

# Key categories to block
txt(sl,Inches(0.60),Inches(4.80),Inches(12),Inches(0.30),"Priority Negative Keyword Categories",14,NAVY,True)
add_table(sl,Inches(0.60),Inches(5.10),Inches(12.10),Inches(1.40),
    ["Category","Examples","2025 Waste","2026 Waste","Match Type"],
    [["Price Research (no location)","implant cost uk, how much is dental implants","\u00a320,540","\u00a314,211","Phrase"],
     ["Wrong Treatment","implant dentures, dental bridge, dental crown","\u00a33,112","\u00a33,661","Phrase"],
     ["Competitor / Tourism","brighton implant clinic, turkey teeth","\u00a31,724","\u00a31,390","Exact"],
     ["NHS / Free","nhs dental implants, free dental implants","\u00a31,378","\u00a3801","Phrase"]],
    cw=[2.0,3.0,1.0,1.0,1.0])
footer(sl,4)


# ── S5: AD COPY REWRITE PLAN ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Step 3: Ad Copy Rewrite","90% of ads rated Poor or Average \u2014 the #1 lever for reducing CPA")
badge(sl,"Weeks 2-4")

mcard(sl,Inches(0.60),Inches(1.20),"90%","Ads Rated Poor/Average",RED,Inches(3.00))
mcard(sl,Inches(3.90),Inches(1.20),"71%","Keywords with Below Avg Ad Relevance",YELLOW,Inches(3.00))
mcard(sl,Inches(7.20),Inches(1.20),"6.4","Weighted Avg QS (down from 6.9)",BLUE,Inches(3.00))
mcard(sl,Inches(10.50),Inches(1.20),"#1","CPA Reduction Lever",GREEN,Inches(2.20))

txt(sl,Inches(0.60),Inches(2.60),Inches(12),Inches(0.30),"Proposed Ad Testing Framework",14,NAVY,True)
add_table(sl,Inches(0.60),Inches(2.95),Inches(12.10),Inches(2.00),
    ["Ad Group","Current Strength","Angle 1 (Price)","Angle 2 (Speed)","Angle 3 (Quality)"],
    [["Dental Implants","Poor","From \u00a31,695. 0% Finance.","Permanent Teeth in 48 Hours.","4.8\u2605 Rating. 399+ Reviews."],
     ["Implants Cost","Poor","Up to 60% Less Than Other Clinics.","Same-Day Consultation Available.","In-House Lab. Premium Materials."],
     ["All on 4","Poor","Full Arch From \u00a39,995.","New Teeth in 48 Hours. Walk Out Smiling.","CQC Registered. Expert Surgeons."],
     ["Financing","Average","0% Finance. From \u00a347/month.","Book Free Consultation Today.","Trusted by 399+ Patients."],
     ["Invisalign","Average","Invisalign From \u00a32,500.","Free Scan. Results in 6 Months.","Diamond Apex Provider."]],
    cw=[1.5,1.0,2.5,2.5,2.5])

ibox(sl,Inches(0.60),Inches(5.20),Inches(5.80),Inches(1.30),
    "Testing Plan",
    "2-3 RSA variants per ad group, each testing a different messaging angle. "
    "Run for 4-6 weeks, then keep winner and test new challenger. "
    "Every headline must match the keyword intent in that ad group.",
    BLUE,BLUE_TINT)

ibox(sl,Inches(6.80),Inches(5.20),Inches(5.80),Inches(1.30),
    "Expected Impact",
    "Improving ad relevance from Below Average to Average on the top 20 keywords "
    "would lift QS by 2-3 points, reducing CPC by an estimated 20-30%. "
    "At \u00a3{:,}/month spend, that's \u00a310-16K/month in CPC savings.".format(MONTHLY_RUN),
    GREEN,GREEN_TINT)
footer(sl,5)


# ── S6: BID STRATEGY RECOMMENDATION ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Step 4: Bid Strategy","Switch from arbitrary ROAS to Target CPA \u2014 or use real treatment values")
badge(sl,"Weeks 2-4")

# Current vs options
add_table(sl,Inches(0.60),Inches(1.30),Inches(12.10),Inches(1.80),
    ["Campaign","Current Strategy","Current Target","Proposed","Rationale"],
    [["PMax","Max Conv Value (tROAS)","30% ROAS (\u00a3300 arb.)","tCPA \u00a345 OR tROAS with real values","Arbitrary values mislead the algorithm"],
     ["Search Intent","Max Conv Value (tROAS)","60% ROAS (\u00a3300 arb.)","tCPA \u00a365 (reduce to \u00a350 by month 3)","Start conservative, tighten gradually"],
     ["Brand","Max Conv Value (tROAS)","10% ROAS","Keep or switch to Max Clicks","Brand always converts well \u2014 focus on IS"]],
    cw=[1.5,1.8,2.0,2.5,3.0])

ibox(sl,Inches(0.60),Inches(3.40),Inches(5.80),Inches(3.00),
    "Option A: Target CPA (Recommended)",
    "Switch all campaigns to Maximise Conversions with Target CPA.\n\n"
    "\u2022 PMax: Start at \u00a345 (current \u00a352), tighten to \u00a340 as QS improves\n"
    "\u2022 Search: Start at \u00a365 (current \u00a394), tighten to \u00a350 over 8 weeks\n"
    "\u2022 Never change target by more than 15-20% at a time\n\n"
    "Advantage: Simple, direct control over cost per lead. "
    "The algorithm optimises for lead volume at your price.",
    GREEN,GREEN_TINT)

ibox(sl,Inches(6.80),Inches(3.40),Inches(5.80),Inches(3.00),
    "Option B: Target ROAS with Real Values",
    "Keep Target ROAS but replace \u00a3300 arbitrary values with actual treatment revenue "
    "from Dengro.\n\n"
    "\u2022 Single implant booking: \u00a31,695\n"
    "\u2022 All-on-4 booking: \u00a39,995\n"
    "\u2022 Vivo Bridge booking: \u00a315,990\n\n"
    "Advantage: Algorithm bids more aggressively for high-value treatments. "
    "Risk: Requires reliable Dengro-to-Google Ads attribution.",
    BLUE,BLUE_TINT)
footer(sl,6)


# ── S7: DEVICE + SCHEDULE + DAY ADJUSTMENTS ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Step 5: Bid Adjustments","Device, schedule, and day-of-week \u2014 no adjustments currently exist anywhere")
badge(sl,"Weeks 2-4")

img(sl,f"{CHARTS}/r3_device_adj.png",Inches(0.60),Inches(1.30),Inches(5.50),Inches(2.50))

# Day of week adjustments
txt(sl,Inches(6.50),Inches(1.20),Inches(6.20),Inches(0.30),"Day of Week (2026 Data)",14,NAVY,True)
add_table(sl,Inches(6.50),Inches(1.55),Inches(6.20),Inches(2.20),
    ["Day","2025 CPA","2026 CPA","Proposed Adj."],
    [["Monday","\u00a341","\u00a364","-5%"],
     ["Tuesday","\u00a344","\u00a362","No change"],
     ["Wednesday","\u00a344","\u00a361","No change"],
     ["Thursday","\u00a344","\u00a358","+5%"],
     ["Friday","\u00a346","\u00a356","+10%"],
     ["Saturday","\u00a343","\u00a356","+10%"],
     ["Sunday","\u00a345","\u00a368","-15%"]],
    cw=[1.2,0.8,0.8,1.0])

# Schedule adjustments
txt(sl,Inches(0.60),Inches(4.10),Inches(12),Inches(0.30),"Hour of Day Adjustments",14,NAVY,True)
ibox(sl,Inches(0.60),Inches(4.45),Inches(5.80),Inches(2.00),
    "Increase Bids (+10 to +15%)",
    "Hours 16-17 (4-5pm): Best CPA at \u00a346-48 both years. "
    "High-volume, high-efficiency slot. Late afternoon is when patients are ready to book.",
    GREEN,GREEN_TINT)

ibox(sl,Inches(6.80),Inches(4.45),Inches(5.80),Inches(2.00),
    "Reduce Bids (-15 to -20%)",
    "Hours 7-8 (7-8am): Worst CPA at \u00a356-74 across both years. "
    "Early commute clicks are expensive and convert poorly. "
    "Reduce bids during this window.",
    RED,RED_TINT)
footer(sl,7)


# ── S8: GEO RESTRUCTURE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Step 6: Geographic Restructure","Replace radius targeting with postcodes to unlock bid control per area")
badge(sl,"Month 2")

# Current vs proposed
ibox(sl,Inches(0.60),Inches(1.30),Inches(5.80),Inches(2.50),
    "Current: Radius Targeting",
    "Multiple overlapping radiuses around the clinic:\n"
    "\u2022 30km ring: \u00a359 CPA (2026)\n"
    "\u2022 35km ring: \u00a393 CPA (2026)\n\n"
    "Problems:\n"
    "\u2022 No bid adjustments possible on radius targets\n"
    "\u2022 Wider rings reach patients less likely to travel\n"
    "\u2022 Different campaigns target different radiuses \u2014 fragmented data",
    RED,RED_TINT)

ibox(sl,Inches(6.80),Inches(1.30),Inches(5.80),Inches(2.50),
    "Proposed: Postcode Targeting",
    "Individual London postcodes with bid adjustments per area:\n\n"
    "\u2022 Bid +20% for W6, W12, W3 (near clinic, low CPA)\n"
    "\u2022 No adjustment for mid-range postcodes\n"
    "\u2022 Bid -20% for outer postcodes (N13 \u00a3169, TW2 \u00a3137)\n\n"
    "All campaigns share the same postcode targets \u2014 one unified geographic view. "
    "Data from the matched locations report shows CPA ranges from \u00a314 to \u00a3169 across postcodes.",
    GREEN,GREEN_TINT)

# Best/worst postcodes reminder
txt(sl,Inches(0.60),Inches(4.10),Inches(6),Inches(0.30),"Best Postcodes (2026, 3+ conv)",12,GREEN,True)
add_table(sl,Inches(0.60),Inches(4.40),Inches(5.80),Inches(1.40),
    ["Postcode","CPA","Conv"],
    [["SM4","\u00a314","5"],["EN4","\u00a314","3"],["CM18","\u00a315","3"],
     ["WD7","\u00a318","4"],["SW13","\u00a319","3"]],
    cw=[2.0,1.0,1.0])

txt(sl,Inches(6.80),Inches(4.10),Inches(6),Inches(0.30),"Worst Postcodes (2026, 3+ conv)",12,RED,True)
add_table(sl,Inches(6.80),Inches(4.40),Inches(5.80),Inches(1.40),
    ["Postcode","CPA","Conv"],
    [["N13","\u00a3169","3"],["TW2","\u00a3137","5"],["WD23","\u00a3123","5"],
     ["SM3","\u00a3123","4"],["N17","\u00a3120","3"]],
    cw=[2.0,1.0,1.0])
footer(sl,8)


# ── S9: MULTI-CLINIC SCALE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Multi-Clinic Growth","How the optimised Dental by Design account becomes the blueprint for the partner network")
badge(sl,"Month 3+")

# Flagship → Partner flow
rect(sl,Inches(0.60),Inches(1.30),Inches(5.00),Inches(3.00),WHITE,BORDER,Pt(0.75))
rect(sl,Inches(0.60),Inches(1.30),Inches(0.08),Inches(3.00),GREEN)
txt(sl,Inches(0.90),Inches(1.45),Inches(4.50),Inches(0.40),"Dental by Design",22,GREEN,True)
txt(sl,Inches(0.90),Inches(1.90),Inches(4.50),Inches(0.25),"Flagship Clinic \u2014 Testing Ground",14,NAVY,True)
bullets(sl,Inches(0.90),Inches(2.25),Inches(4.50),Inches(1.80),[
    "Highest budget = fastest learning",
    "Test ad copy angles, keywords, bid strategies here",
    "Establish CPA benchmarks per treatment type",
    "Proven winning strategies become the playbook",
],12)

# Arrow
txt(sl,Inches(5.80),Inches(2.50),Inches(0.80),Inches(0.50),"\u25B6",28,BLUE,a=PP_ALIGN.CENTER)

# Partner clinics
rect(sl,Inches(6.80),Inches(1.30),Inches(5.80),Inches(3.00),WHITE,BORDER,Pt(0.75))
rect(sl,Inches(6.80),Inches(1.30),Inches(0.08),Inches(3.00),BLUE)
txt(sl,Inches(7.10),Inches(1.45),Inches(5.20),Inches(0.40),"Partner Clinics",22,BLUE,True)
txt(sl,Inches(7.10),Inches(1.90),Inches(5.20),Inches(0.25),"UK Network \u2014 Scaled Learnings",14,NAVY,True)
bullets(sl,Inches(7.10),Inches(2.25),Inches(5.20),Inches(1.80),[
    "Separate campaigns per clinic \u2014 tailored messaging + local targeting",
    "Shared negative keyword lists \u2014 one block benefits all clinics",
    "Performance-based budgets \u2014 best clinics earn more spend",
    "Postcode targeting per clinic location \u2014 unified geographic view",
],12)

# What scales
txt(sl,Inches(0.60),Inches(4.60),Inches(12),Inches(0.30),"What Scales from Flagship to Network",14,NAVY,True)
mcard(sl,Inches(0.60),Inches(5.00),"9-List","Neg Keyword Framework",GREEN,Inches(2.40))
mcard(sl,Inches(3.20),Inches(5.00),"Ad Copy","Winning Angles",BLUE,Inches(2.40))
mcard(sl,Inches(5.80),Inches(5.00),"Bid Strategy","tCPA Benchmarks",YELLOW,Inches(2.40))
mcard(sl,Inches(8.40),Inches(5.00),"Reporting","Monthly Framework",GREEN,Inches(2.20))
mcard(sl,Inches(10.80),Inches(5.00),"Postcodes","Geo Structure",BLUE,Inches(1.90))
footer(sl,9)


# ── S10: PROJECTED RESULTS ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Projected Results","Conservative projections based on the data in Reports 1 and 2")
badge(sl,"90-Day Target")

img(sl,f"{CHARTS}/r3_projections.png",Inches(0.60),Inches(1.30),Inches(7.50),Inches(3.50))

# Projection cards
ibox(sl,Inches(8.40),Inches(1.30),Inches(4.30),Inches(1.50),
    "Scenario 1: More Leads",
    f"Keep spend at \u00a3{MONTHLY_RUN:,.0f}/month.\n"
    f"Reduce non-brand CPA from \u00a3{NBCPA26} to \u00a3{TARGET_CPA}.\n"
    f"Result: 236 additional leads per month (+29%).",
    GREEN,GREEN_TINT)

ibox(sl,Inches(8.40),Inches(3.00),Inches(4.30),Inches(1.50),
    "Scenario 2: Same Leads, Less Spend",
    f"Maintain current lead volume.\n"
    f"Reduce CPA from \u00a3{NBCPA26} to \u00a3{TARGET_CPA}.\n"
    f"Result: \u00a312,000/month saving (\u00a3144K/year).",
    BLUE,BLUE_TINT)

rbox(sl,Inches(0.60),Inches(5.10),Inches(12.10),Inches(0.60),
    "These projections are conservative. They're based on the waste identified in Report 1 and the structural issues in Report 2. "
    "The actual improvement depends on how aggressively the changes are implemented.")

txt(sl,Inches(0.60),Inches(5.90),Inches(12.10),Inches(0.50),
    f"Basis: 2025 non-brand CPA was \u00a3{NBCPA25}. The account achieved this before PMax scaled and Search Intent expanded. "
    f"A target of \u00a3{TARGET_CPA} is between the 2025 level and the current \u00a3{NBCPA26} \u2014 achievable with the optimisations proposed.",
    10,GREY_MID)
footer(sl,10)


# ── S11: IMPLEMENTATION TIMELINE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Implementation Timeline")
badge(sl,"90-Day Roadmap")

img(sl,f"{CHARTS}/r3_timeline.png",Inches(0.60),Inches(1.10),Inches(12.10),Inches(2.20))

# 4 phase cards
phases=[
    ("Week 1","Fix Foundations",RED,[
        "Remove KMG ghost conversions",
        "Move soft signals to Secondary",
        "Negative keyword 9-list overhaul",
        "Disable auto-apply (except ad rotation)",
    ]),
    ("Weeks 2-4","Optimise",YELLOW,[
        "Rewrite all Poor/Average ads",
        "Implement device bid adjustments",
        "Implement schedule + day-of-week adjustments",
        "Evaluate tCPA vs tROAS bid strategy",
    ]),
    ("Month 2","Restructure",BLUE,[
        "Geo: radius \u2192 postcode targeting",
        "Keyword optimization (first month clean data)",
        "Landing page traffic routing",
        "Remove 43 paused campaigns",
    ]),
    ("Month 3","Scale",GREEN,[
        "Scale winning campaigns",
        "Establish per-treatment benchmarks",
        "Multi-clinic rollout planning",
        "Full reporting framework",
    ]),
]

for i,(label,title,color,items) in enumerate(phases):
    cx=Inches(0.60+i*3.10)
    rect(sl,cx,Inches(3.50),Inches(2.90),Inches(3.00),WHITE,BORDER,Pt(0.75))
    rect(sl,cx,Inches(3.50),Inches(2.90),Pt(4),color)
    rect(sl,cx,Inches(3.50),Inches(0.06),Inches(3.00),color)
    txt(sl,cx+Inches(0.15),Inches(3.60),Inches(2.60),Inches(0.25),label,11,color,True)
    txt(sl,cx+Inches(0.15),Inches(3.85),Inches(2.60),Inches(0.30),title,14,NAVY,True)
    bullets(sl,cx+Inches(0.15),Inches(4.25),Inches(2.60),Inches(2.10),items,10)
footer(sl,11)


# ── S12: NEXT STEPS (closing slide) ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=GREY_BG
top_bar(sl);sidebar(sl);logo(sl)
txt(sl,Inches(0.60),Inches(1.40),Inches(5.50),Inches(0.80),"Next Steps",44,NAVY,True)
rect(sl,Inches(0.60),Inches(2.20),Inches(2.50),Inches(0.05),BLUE)

steps=[
    ("1","Grant me editor access to the account","I can implement the Week 1 fixes in a single session",GREEN),
    ("2","Conversion tracking fix + negative keyword overhaul in Week 1","The two highest-impact changes \u2014 immediate CPA improvement",BLUE),
    ("3","Ad rewrites + bid adjustments in Weeks 2-4","The #1 CPA reduction lever, backed by data from all 3 reports",YELLOW),
    ("4","Monthly reporting from Day 1","Full transparency on what changed, why, and the measured impact",RED),
]

for i,(num,action,detail,color) in enumerate(steps):
    y=Inches(2.55+i*0.85)
    rect(sl,Inches(0.60),y,Inches(5.50),Inches(0.70),WHITE)
    rect(sl,Inches(0.60),y,Inches(0.06),Inches(0.70),color)
    ncircle(sl,Inches(0.80),y+Inches(0.10),num,color)
    txt(sl,Inches(1.40),y+Inches(0.05),Inches(4.50),Inches(0.30),action,12,NAVY,True)
    txt(sl,Inches(1.40),y+Inches(0.35),Inches(4.50),Inches(0.30),detail,10,BLACK)

# Right: contact card
rect(sl,Inches(6.80),Inches(0.50),Inches(5.90),Inches(2.50),WHITE)
rect(sl,Inches(6.80),Inches(0.50),Inches(0.08),Inches(2.50),BLUE)
txt(sl,Inches(7.10),Inches(0.60),Inches(5.40),Inches(0.35),"Ready to Start Immediately",14,NAVY)
txt(sl,Inches(7.10),Inches(1.00),Inches(5.40),Inches(0.80),"Christopher Hoole",40,BLUE,True)
multipara(sl,Inches(7.10),Inches(1.80),Inches(5.40),Inches(1.00),[
    ("chris@christopherhoole.com",14,BLUE,False,4),
    ("christopherhoole.com",14,BLUE,False,4),
    ("+44 7999 500 184",14,BLACK,False,0)])

minicard(sl,Inches(6.80),Inches(3.20),f"\u00a3{TARGET_CPA}","Target CPA",GREEN)
minicard(sl,Inches(8.83),Inches(3.20),"90 Days","To Full Impact",BLUE)
minicard(sl,Inches(10.87),Inches(3.20),"+29%","More Leads",RED)

rect(sl,Inches(6.80),Inches(4.60),Inches(5.90),Inches(1.20),BLUE_TINT,BLUE,Pt(1))
txt(sl,Inches(7.05),Inches(4.65),Inches(5.40),Inches(1.10),
    "\u201CThe data is clear. The opportunities are identified. "
    "The plan is specific. I'm ready to implement Week 1 the day I start.\u201D\n\n"
    "\u2014 Christopher Hoole",11,NAVY,True)
bot_bar_title(sl)


# ══════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════
output=os.path.join(REPORTS,'03_restructure_report.pptx')
prs.save(output)
print(f"\nSaved to: {output}")
print(f"Slides: {len(prs.slides)}")
print("Done!")
