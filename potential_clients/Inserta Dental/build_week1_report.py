"""Build Week 1 weekly report for Dental by Design.

Follows DBD_REPORT_TEMPLATE_SPEC.md exactly. 11 slides.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from lxml import etree

# === Palette ===
NAVY = RGBColor(0x1A, 0x23, 0x7E)
GBLUE = RGBColor(0x42, 0x85, 0xF4)
GRED = RGBColor(0xEA, 0x43, 0x35)
GGREEN = RGBColor(0x34, 0xA8, 0x53)
GYELLOW = RGBColor(0xFB, 0xBC, 0x05)
BODY = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x66, 0x66, 0x66)
LIGHTGREY = RGBColor(0xE0, 0xE0, 0xE0)
PALEGREY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
LOGO = r"potential_clients/Inserta Dental/act_logo_official.png"

# === Presentation 16:9 ===
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW = prs.slide_width
SH = prs.slide_height


def add_text(slide, left, top, width, height, text, *, size=14, bold=False,
             color=BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
             font=FONT):
    """Add a textbox and return the paragraph for further tweaking."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb, p, run


def add_multirun(slide, left, top, width, height, runs, *, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP):
    """runs = list of dicts: text, size, bold, color, italic (optional)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, r in enumerate(runs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = r["text"]
        run.font.name = FONT
        run.font.size = Pt(r.get("size", 14))
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
        run.font.color.rgb = r.get("color", BODY)
        if "space_after" in r:
            p.space_after = Pt(r["space_after"])
    return tb


def add_rect(slide, left, top, width, height, *, fill=None, line=None,
             line_width=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_width)
    shp.text_frame.margin_left = Emu(0)
    shp.text_frame.margin_right = Emu(0)
    shp.text_frame.margin_top = Emu(0)
    shp.text_frame.margin_bottom = Emu(0)
    return shp


def add_line(slide, x1, y1, x2, y2, color=LIGHTGREY, width=1.0, dash=False):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if dash:
        # MSO_LINE_DASH_STYLE.DASH == 7
        lnL = conn.line._get_or_add_lnL() if hasattr(conn.line, "_get_or_add_lnL") else None
        # Use oxml to set dash
        sppr = conn.line._get_or_add_ln()
        prstDash = etree.SubElement(sppr, qn("a:prstDash"))
        prstDash.set("val", "dash")
    return conn


def add_footer(slide, page_num):
    """Standard footer: [logo] Christopher Hoole | christopherhoole.com | Confidential   [N]"""
    # Logo bottom-left
    slide.shapes.add_picture(LOGO, Inches(0.3), Inches(7.1), height=Inches(0.3))
    # Footer text
    add_text(slide, Inches(0.75), Inches(7.1), Inches(10), Inches(0.3),
             "Christopher Hoole  |  christopherhoole.com  |  Confidential",
             size=10, color=MUTED)
    # Page number
    add_text(slide, Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.3),
             str(page_num), size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_title(slide, title, context=None, supporting=None):
    """Main navy title at 44pt + optional context (GBlue italic) + supporting line."""
    add_text(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
             title, size=32, bold=True, color=NAVY)
    y = Inches(1.05)
    if context:
        add_text(slide, Inches(0.5), y, Inches(12.3), Inches(0.35),
                 context, size=14, italic=True, color=GBLUE)
        y = Inches(1.35)
    if supporting:
        add_text(slide, Inches(0.5), y, Inches(12.3), Inches(0.35),
                 supporting, size=13, color=BODY)


def stat_callout(slide, left, top, width, value, label, *, color=NAVY,
                 sub=None, value_size=40, label_size=13):
    """Vertical stack: big number, label, optional sublabel. Returns shape."""
    # Value
    add_text(slide, left, top, width, Inches(0.8),
             value, size=value_size, bold=True, color=color, align=PP_ALIGN.LEFT)
    # Label
    add_text(slide, left, top + Inches(0.75), width, Inches(0.3),
             label, size=label_size, color=BODY, align=PP_ALIGN.LEFT)
    if sub:
        add_text(slide, left, top + Inches(1.05), width, Inches(0.3),
                 sub, size=10, color=MUTED, align=PP_ALIGN.LEFT)


def four_stats(slide, top, stats, *, left=Inches(0.5), right=Inches(12.83)):
    """stats = list of 4 dicts: value, label, color (default NAVY), sub (optional)."""
    total = right - left
    col_w = total // 4
    for i, s in enumerate(stats):
        stat_callout(slide, left + col_w * i, top, col_w - Inches(0.2),
                     s["value"], s["label"],
                     color=s.get("color", NAVY),
                     sub=s.get("sub"),
                     value_size=s.get("value_size", 40))


def key_finding(slide, left, top, width, height, label, content,
                *, label_color=GBLUE):
    """Key Finding / Recommendation block."""
    add_text(slide, left, top, width, Inches(0.3),
             label, size=13, bold=True, color=label_color)
    add_text(slide, left, top + Inches(0.3), width, height - Inches(0.3),
             content, size=12, color=BODY)


# =========================================================================
# SLIDE 1 — TITLE
# =========================================================================
s = prs.slides.add_slide(BLANK)

# Top-left stat 1
stat_callout(s, Inches(0.7), Inches(0.6), Inches(5.5),
             "25", "Ad groups (was 4)", color=NAVY, value_size=72, label_size=16)
# Top-right stat 2
stat_callout(s, Inches(7.1), Inches(0.6), Inches(5.5),
             "~675", "New keywords added", color=NAVY, value_size=72, label_size=16)

# Centre: main title + subtitle
add_text(s, Inches(0.5), Inches(2.8), Inches(12.3), Inches(0.8),
         "Week 1 Report", size=48, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
         "13–17 April 2026", size=28, color=GBLUE, align=PP_ALIGN.CENTER)

# Bottom stats
stat_callout(s, Inches(0.7), Inches(4.6), Inches(5.5),
             "75", "New RSAs written", color=NAVY, value_size=72, label_size=16)
stat_callout(s, Inches(7.1), Inches(4.6), Inches(5.5),
             "~5,000", "Structured negatives", color=NAVY, value_size=72, label_size=16)

# Bottom footer area: client name + CH details + logo
add_text(s, Inches(0.5), Inches(6.55), Inches(8), Inches(0.35),
         "Dental by Design  |  Prodent Group", size=14, bold=True, color=BODY)
add_text(s, Inches(0.5), Inches(6.85), Inches(8), Inches(0.3),
         "Christopher Hoole  ·  Google Ads Specialist  ·  christopherhoole.com",
         size=11, color=MUTED)
add_text(s, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
         "Data: 13–17 April 2026", size=10, color=MUTED)
# Logo bottom right
s.shapes.add_picture(LOGO, Inches(12.3), Inches(6.85), height=Inches(0.5))


# =========================================================================
# SLIDE 2 — EXECUTIVE SUMMARY
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Executive Summary",
          context="Week 1: 13–17 April 2026",
          supporting="Structural reset across three live campaigns; early signal positive on Brand and PMax.")

# 4 big stats — all account-level
four_stats(s, Inches(1.8), [
    {"value": "£8,462", "label": "Week 1 spend", "color": NAVY},
    {"value": "87", "label": "Dengro Offline Leads", "color": GGREEN},
    {"value": "101", "label": "Primary conversions", "color": GGREEN},
    {"value": "£97", "label": "Cost per lead", "color": NAVY},
])

# Status pills row
pill_y = Inches(3.5)
pill_defs = [
    ("Brand: CPC at 8-week low", GGREEN),
    ("PMax: CPA improving", GGREEN),
    ("DII: Learning phase (rebuild)", GYELLOW),
]
pill_x = Inches(0.5)
for label, col in pill_defs:
    w = Inches(3.95)
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pill_x, pill_y, w, Inches(0.5))
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = col
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = FONT; r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = WHITE
    pill_x += w + Inches(0.1)

# Core Insight
add_text(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.35),
         "Core Insight", size=14, bold=True, color=GBLUE)
core = ("Week 1 was structural. Three live campaigns moved off Max Conversion Value bidding onto controllable targets. "
        "Around 5,000 structured negatives replaced 2,192 unstructured broad-match entries, and the Dental Implants Intent campaign "
        "was rebuilt from 4 legacy ad groups into 25 intent-grouped ad groups. Early signal is positive — Brand CPC is at its "
        "8-week low and PMax CPA improved week-on-week — but DII needs a further week of stabilisation and Monday's "
        "end-to-end tracking audit is the gate to any CPA-based decisions from here.")
add_text(s, Inches(0.5), Inches(4.65), Inches(12.3), Inches(2.0),
         core, size=13, color=BODY)

add_footer(s, 2)


# =========================================================================
# SLIDE 3 — WORK DELIVERED
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Work Delivered",
          context="Tue 14 → Fri 17 April (4 working days post-signing)")

# Horizontal timeline
days = [
    ("Tue 14 Apr", "Negatives Overhaul",
     "2,192 unstructured negatives → ~4,000+ structured across 11 shared lists (grew to ~5,000 by Fri). ~£63k of 60-day waste blocked at auction source."),
    ("Wed 15 Apr", "Account-Wide Reset",
     "Brand rebuilt (tIS 90% / £1.50 CPC cap). Auto-apply disabled. PMax → Target CPA £60. DII → Target CPA £75. 270-postcode geo + 42-cell schedule rolled to all 3 campaigns. 10 new RSAs."),
    ("Thu 16 Apr", "DII Rebuild (Part 1)",
     "20 new ad groups built on Dental Implants Intent. 537 keywords, 60 RSAs. 4 legacy ad groups paused. Website USP audit. 18 questions sent to Giulio + Tommaso."),
    ("Fri 17 Apr", "Audit + Completion",
     "173 DBD-blocking negatives removed (zirconia, same-day teeth, Hammersmith, crowns). AG21-25 built — DII now 25/25 complete. Ad Strength audit across 63 RSAs. Client replies received."),
]

# Timeline
tl_y = Inches(1.95)
tl_left = Inches(0.6)
tl_right = Inches(12.73)
tl_w = tl_right - tl_left
col_w = tl_w // 4

# Horizontal baseline line
add_line(s, tl_left + col_w // 2, tl_y, tl_right - col_w // 2 + Emu(1),
         tl_y, color=GBLUE, width=2.5)

for i, (day, title, body) in enumerate(days):
    cx = tl_left + col_w * i + col_w // 2
    # Circle marker
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              cx - Inches(0.18), tl_y - Inches(0.18),
                              Inches(0.36), Inches(0.36))
    circ.shadow.inherit = False
    circ.fill.solid()
    circ.fill.fore_color.rgb = GBLUE
    circ.line.color.rgb = WHITE
    circ.line.width = Pt(2)
    tf = circ.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i + 1)
    r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = WHITE

    # Day label above
    add_text(s, cx - col_w // 2 + Inches(0.2), tl_y - Inches(0.8),
             col_w - Inches(0.4), Inches(0.4),
             day, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Title below marker
    add_text(s, cx - col_w // 2 + Inches(0.2), tl_y + Inches(0.4),
             col_w - Inches(0.4), Inches(0.4),
             title, size=14, bold=True, color=GBLUE, align=PP_ALIGN.CENTER)

    # Body
    add_text(s, cx - col_w // 2 + Inches(0.2), tl_y + Inches(0.85),
             col_w - Inches(0.4), Inches(4),
             body, size=11, color=BODY, align=PP_ALIGN.LEFT)

# Key Finding at bottom
key_finding(s, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.6),
            "Key Finding",
            "Three campaigns moved from Max Conversion Value bidding (optimising toward the arbitrary £300 booking value) onto Target CPA / Target Impression Share — putting real cost control in place for the first time.")

add_footer(s, 3)


# =========================================================================
# SLIDE 4 — BRAND CAMPAIGN
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Brand Campaign",
          context="Week 1: 13–17 April 2026",
          supporting="Bid strategy rebuilt Wednesday 15 April — Target Impression Share 90% / Max CPC £1.50.")

four_stats(s, Inches(1.75), [
    {"value": "£280", "label": "Week spend", "color": NAVY},
    {"value": "12.0", "label": "Primary conv.", "color": GGREEN},
    {"value": "£23.39", "label": "Cost / conv.", "color": GGREEN},
    {"value": "£2.16", "label": "Avg CPC", "color": GGREEN, "sub": "Lowest of 8 weeks"},
])

# Chart — 8-week Brand CPC line
chart_data = CategoryChartData()
chart_data.categories = ["23 Feb", "2 Mar", "9 Mar", "16 Mar",
                         "23 Mar", "30 Mar", "6 Apr", "13 Apr"]
chart_data.add_series("Brand Avg CPC (£)",
                      [3.09, 2.41, 3.49, 4.00, 3.24, 2.31, 2.90, 2.16])
chart_x = Inches(0.5); chart_y = Inches(3.65)
chart_w = Inches(7.8); chart_h = Inches(3.3)
gframe = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, chart_x, chart_y,
                            chart_w, chart_h, chart_data)
chart = gframe.chart
chart.has_title = False
chart.has_legend = False
# Style line
plot = chart.plots[0]
series = plot.series[0]
line = series.format.line
line.color.rgb = GBLUE
line.width = Pt(2.5)
# Marker color
from pptx.enum.chart import XL_MARKER_STYLE
series.marker.style = XL_MARKER_STYLE.CIRCLE
series.marker.size = 7
series.marker.format.fill.solid()
series.marker.format.fill.fore_color.rgb = GBLUE
series.marker.format.line.color.rgb = GBLUE
# Data labels
plot.has_data_labels = True
dl = plot.data_labels
dl.font.size = Pt(9)
dl.font.name = FONT
dl.font.color.rgb = BODY
dl.position = XL_LABEL_POSITION.ABOVE
dl.number_format = '"£"0.00'

# Key Finding + Recommendation (right column)
key_finding(s, Inches(8.6), Inches(3.65), Inches(4.25), Inches(1.5),
            "Key Finding",
            "Brand CPC is at its lowest of 8 tracked weeks. The £1.50 Max CPC cap is biting without suppressing volume — impression share held at 92%+ and conversion volume tracked with spend.")
key_finding(s, Inches(8.6), Inches(5.3), Inches(4.25), Inches(1.5),
            "Recommendation",
            "Hold the £1.50 cap through Week 2. Evaluate lowering toward £1.20 once the 14-day bid-strategy learning phase completes.")

add_footer(s, 4)


# =========================================================================
# SLIDE 5 — PERFORMANCE MAX
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Performance Max",
          context="Week 1: 13–17 April 2026",
          supporting="Bid strategy switched Wed 15 April — Max Conversion Value (tROAS) → Target CPA £60.")

four_stats(s, Inches(1.75), [
    {"value": "£5,960", "label": "Week spend", "color": NAVY},
    {"value": "81.5", "label": "Primary conv.", "color": NAVY},
    {"value": "£73.13", "label": "Cost / conv.", "color": GGREEN, "sub": "Down from £85.86"},
    {"value": "48,349", "label": "Impressions", "color": NAVY},
])

# Daily CPA bar chart Mon-Fri
chart_data = CategoryChartData()
chart_data.categories = ["Mon 13", "Tue 14", "Wed 15", "Thu 16", "Fri 17"]
chart_data.add_series("Daily PMax CPA (£)",
                      [118.68, 55.85, 25.18, 76.29, 79.22])
chart_x = Inches(0.5); chart_y = Inches(3.65)
chart_w = Inches(7.8); chart_h = Inches(3.3)
gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_x, chart_y,
                            chart_w, chart_h, chart_data)
chart = gframe.chart
chart.has_title = False
chart.has_legend = False
plot = chart.plots[0]
plot.gap_width = 80
# Colour bars: red for pre-switch, green for post-switch
series = plot.series[0]
fills = [GRED, GYELLOW, GGREEN, GGREEN, GGREEN]
for i, pt in enumerate(series.points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = fills[i]
    pt.format.line.fill.background()
plot.has_data_labels = True
dl = plot.data_labels
dl.font.size = Pt(10)
dl.font.name = FONT
dl.font.color.rgb = BODY
dl.font.bold = True
dl.position = XL_LABEL_POSITION.OUTSIDE_END
dl.number_format = '"£"0'

# Annotation: Target CPA £60 applied Wed afternoon
# Vertical marker between Tue and Wed (between categories 2 and 3)
# Chart occupies Inches(0.5) to Inches(8.3), plot area approx 0.9 to 8.1
# 5 bars; dividing line between bar 2 and bar 3 at approx 1/2.5 of chart
ann_x = Inches(0.5 + 7.8 * (2.5 / 5))  # midpoint between bar 2 and 3
add_line(s, ann_x, Inches(3.85), ann_x, Inches(6.8),
         color=NAVY, width=1.5, dash=True)
add_text(s, ann_x - Inches(1.7), Inches(6.75), Inches(3.4), Inches(0.25),
         "Target CPA £60 applied Wed 15 Apr",
         size=9, bold=True, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

# Key Finding + Flag
key_finding(s, Inches(8.6), Inches(3.65), Inches(4.25), Inches(1.6),
            "Key Finding",
            "CPA improved week-on-week despite the mid-week bid-strategy switch and a persistent \"asset groups limited by policy\" flag restricting reach. Wednesday's £25 CPA shows the ceiling potential once policy constraints are cleared.")
key_finding(s, Inches(8.6), Inches(5.4), Inches(4.25), Inches(1.5),
            "Flagged for Week 2",
            "Rebuild PMax asset groups with policy-safe copy aligned to Giulio's verified claims (no \"99.12%\", no \"Save 60%\", no \"Top 10 Europe\"). Scheduled Wed 22 April.",
            label_color=GYELLOW)

add_footer(s, 5)


# =========================================================================
# SLIDE 6 — DENTAL IMPLANTS INTENT
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Dental Implants Intent",
          context="Week 1: 13–17 April 2026",
          supporting="Campaign fully rebuilt Thu–Fri. Entered learning phase Wed PM; meaningful signal needs 5–7 days.")

four_stats(s, Inches(1.75), [
    {"value": "£2,222", "label": "Week spend", "color": NAVY},
    {"value": "25", "label": "Ad groups (was 4)", "color": NAVY},
    {"value": "~675", "label": "Keywords (was 27)", "color": NAVY},
    {"value": "75", "label": "RSAs (was 6)", "color": NAVY},
])

# Before / After visual: 4 horizontal bars comparing before/after
visual_x = Inches(0.5)
visual_y = Inches(3.55)
visual_w = Inches(7.8)

add_text(s, visual_x, visual_y, visual_w, Inches(0.3),
         "Structural rebuild — before / after", size=13, bold=True, color=NAVY)

rows = [
    ("Ad groups",  4, 25),
    ("Keywords", 27, 675),
    ("RSAs",      6, 75),
    ("Headlines written", 0, 1050),
]
row_y = visual_y + Inches(0.45)
max_after = max(r[2] for r in rows)
bar_area_x = visual_x + Inches(1.9)
bar_area_w = Inches(4.4)
row_h = Inches(0.5)
row_gap = Inches(0.08)

for label, before, after in rows:
    # Label
    add_text(s, visual_x, row_y + Inches(0.08), Inches(1.85), Inches(0.35),
             label, size=11, bold=True, color=BODY)
    # Before bar (grey, smaller)
    if before > 0:
        bw = int(bar_area_w * (before / max_after))
        add_rect(s, bar_area_x, row_y, bw, Inches(0.2), fill=LIGHTGREY)
    add_text(s, bar_area_x + Inches(0.05), row_y - Inches(0.01),
             Inches(1.0), Inches(0.22),
             f"{before}", size=9, color=MUTED)
    # After bar (gblue)
    aw = int(bar_area_w * (after / max_after))
    add_rect(s, bar_area_x, row_y + Inches(0.22), aw, Inches(0.25), fill=GBLUE)
    add_text(s, bar_area_x + aw + Inches(0.05), row_y + Inches(0.22),
             Inches(1.6), Inches(0.25),
             f"{after:,}", size=10, bold=True, color=GBLUE)
    row_y += row_h + row_gap

# Legend
leg_y = row_y + Inches(0.05)
add_rect(s, bar_area_x, leg_y + Inches(0.05), Inches(0.18), Inches(0.12), fill=LIGHTGREY)
add_text(s, bar_area_x + Inches(0.25), leg_y, Inches(0.8), Inches(0.25),
         "Before", size=10, color=MUTED)
add_rect(s, bar_area_x + Inches(1.0), leg_y + Inches(0.05), Inches(0.18), Inches(0.12), fill=GBLUE)
add_text(s, bar_area_x + Inches(1.25), leg_y, Inches(0.8), Inches(0.25),
         "After", size=10, color=GBLUE, bold=True)

# Key finding + rec
key_finding(s, Inches(8.6), Inches(3.65), Inches(4.25), Inches(1.7),
            "Key Finding",
            "The campaign entered learning Wed PM after a full structural rebuild. The five days of data (2.5 on the new structure) show a high nominal CPA — expected noise, not signal. All 27 legacy keywords are covered by the new 25-ad-group structure.")
key_finding(s, Inches(8.6), Inches(5.45), Inches(4.25), Inches(1.5),
            "Recommendation",
            "Hold the new structure untouched through Wed 22 April. Ad Strength audit on AG22-25 resumes Tuesday once Google clears the new ads from Pending review.")

add_footer(s, 6)


# =========================================================================
# SLIDE 7 — NUMBERS THAT MATTER
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "The Numbers That Matter",
          context="Business-level outcomes — Week 1",
          supporting="Full end-to-end tracking audit scheduled Mon 20 April — the gate to all CPA-based decisions.")

# Funnel horizontal — 4 stages
stages = [
    ("£8,462", "Spend", NAVY),
    ("2,339", "Clicks", GBLUE),
    ("87", "Dengro Leads", GGREEN),
    ("[audit]", "Bookings", MUTED),
    ("7", "Purchases", GGREEN),
]

fn_y = Inches(1.9)
fn_x = Inches(0.5)
fn_w = Inches(12.3)
fn_h = Inches(1.5)
stage_w = fn_w // len(stages)
gap = Inches(0.12)
box_w = stage_w - gap

for i, (val, label, col) in enumerate(stages):
    x = fn_x + stage_w * i + gap // 2
    box = s.shapes.add_shape(MSO_SHAPE.CHEVRON, x, fn_y, box_w, fn_h)
    box.shadow.inherit = False
    box.fill.solid()
    box.fill.fore_color.rgb = col
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.35)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = val
    r.font.name = FONT; r.font.size = Pt(24); r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.name = FONT; r2.font.size = Pt(12); r2.font.color.rgb = WHITE

# Note below funnel
add_text(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.55),
         "Dengro Offline Booking fired 0 times in this week's conversion export. Monday's tracking audit will confirm whether this is a reporting gap or a genuine zero-booking week. No CPA-based decisions will be drawn before the audit closes.",
         size=11, italic=True, color=MUTED)

# Table — Leads by campaign
tbl_y = Inches(4.4)
add_text(s, Inches(0.5), tbl_y, Inches(12.3), Inches(0.3),
         "Dengro Offline Leads — by campaign", size=13, bold=True, color=NAVY)

rows = 4; cols = 4
tbl_top = tbl_y + Inches(0.4)
tbl = s.shapes.add_table(rows, cols, Inches(0.5), tbl_top,
                         Inches(12.3), Inches(1.7)).table
headers = ["Campaign", "Spend", "Dengro Leads", "Spend per Lead"]
data = [
    ["Performance Max", "£5,960", "78.5", "£75.92"],
    ["Dental Implants Intent", "£2,222", "6.5", "£341.83"],
    ["Brand", "£280", "2.0", "£140.24"],
]
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    tf = cell.text_frame; tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.paragraphs[0].text = ""
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
    r = tf.paragraphs[0].add_run(); r.text = h
    r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = WHITE

for i, row in enumerate(data):
    for j, val in enumerate(row):
        cell = tbl.cell(i + 1, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else PALEGREY
        tf = cell.text_frame
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.paragraphs[0].text = ""
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
        rr = tf.paragraphs[0].add_run(); rr.text = val
        rr.font.name = FONT; rr.font.size = Pt(12); rr.font.color.rgb = BODY

add_footer(s, 7)


# =========================================================================
# SLIDE 8 — 8-WEEK CONTEXT
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "8-Week Context",
          context="23 February – 17 April 2026",
          supporting="Account trajectory into Week 1. The engagement began 14 April (Week 8 on these charts).")

# Two line charts side by side
cats_8w = ["23 Feb", "2 Mar", "9 Mar", "16 Mar", "23 Mar", "30 Mar", "6 Apr", "13 Apr"]
spend = [17568, 14856, 8393, 13318, 11558, 9997, 10291, 8462]
allconv = [519, 630, 420, 605, 441, 555, 425, 231]

# Left chart — Weekly spend
d1 = CategoryChartData()
d1.categories = cats_8w
d1.add_series("Weekly spend (£)", spend)
gframe = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.5), Inches(2.1),
                            Inches(6.15), Inches(4.5), d1)
c = gframe.chart; c.has_title = True
c.chart_title.text_frame.text = "Weekly spend (£)"
for p in c.chart_title.text_frame.paragraphs:
    for r in p.runs:
        r.font.name = FONT; r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = NAVY
c.has_legend = False
plot = c.plots[0]
ser = plot.series[0]
ser.format.line.color.rgb = GBLUE; ser.format.line.width = Pt(2.5)
ser.marker.style = XL_MARKER_STYLE.CIRCLE; ser.marker.size = 6
ser.marker.format.fill.solid(); ser.marker.format.fill.fore_color.rgb = GBLUE
ser.marker.format.line.color.rgb = GBLUE
plot.has_data_labels = True
dl = plot.data_labels
dl.font.size = Pt(9); dl.font.name = FONT; dl.font.color.rgb = BODY
dl.position = XL_LABEL_POSITION.ABOVE
dl.number_format = '"£"#,##0'

# Right chart — All conversions
d2 = CategoryChartData()
d2.categories = cats_8w
d2.add_series("Weekly all conversions", allconv)
gframe = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(6.85), Inches(2.1),
                            Inches(6.15), Inches(4.5), d2)
c = gframe.chart; c.has_title = True
c.chart_title.text_frame.text = "Weekly all conversions"
for p in c.chart_title.text_frame.paragraphs:
    for r in p.runs:
        r.font.name = FONT; r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = NAVY
c.has_legend = False
plot = c.plots[0]
ser = plot.series[0]
ser.format.line.color.rgb = GGREEN; ser.format.line.width = Pt(2.5)
ser.marker.style = XL_MARKER_STYLE.CIRCLE; ser.marker.size = 6
ser.marker.format.fill.solid(); ser.marker.format.fill.fore_color.rgb = GGREEN
ser.marker.format.line.color.rgb = GGREEN
plot.has_data_labels = True
dl = plot.data_labels
dl.font.size = Pt(9); dl.font.name = FONT; dl.font.color.rgb = BODY
dl.position = XL_LABEL_POSITION.ABOVE
dl.number_format = '0'

# Key finding
key_finding(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.55),
            "Key Finding",
            "Spend is down 52% vs 8-week peak as wasted impressions are removed at auction source. Conversion volume fell proportionally during structural disruption — efficiency (cost per lead) is the metric to watch from Week 3 onward, once the rebuild stabilises.")

add_footer(s, 8)


# =========================================================================
# SLIDE 9 — CLIENT CLARIFICATIONS & ACTIONS
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Client Clarifications & Actions",
          context="Giulio's replies — Fri 17 April 16:47")

# Two-column table
rows_data = [
    ("Consultation: Free + £25 booking deposit",
     "Update ad copy to \"Free Consultation (£25 booking deposit)\" — W2 Mon"),
    ("\"25+ years expertise\" = team-wide across dental field",
     "Keep as-is"),
    ("\"99.12% success rate\" — estimate, not tracked",
     "Remove from all ads — W2 Tue"),
    ("\"Save 60%\" — estimate, no source",
     "Remove — W2 Tue"),
    ("\"Top 10 Implant Provider in Europe\" — self-claim, not an award",
     "Reword as a factual descriptor"),
    ("\"Best Implant Clinic 2024\" — clinic was a finalist",
     "Reword to \"2024 Finalist\""),
    ("Lead → Booking: 13%  ·  Booking → Treatment: 35%",
     "Use in CPA / cost-per-booking modelling"),
    ("~10,000 patients  ·  20,000+ implants placed",
     "Replace outdated \"2,500+ patients\" claim with verified figures"),
    ("Advertising focus: dental implants only (single + double arches)",
     "Scope confirmed — cosmetic / general deferred"),
]

tbl_top = Inches(1.6)
n = len(rows_data) + 1
tbl = s.shapes.add_table(n, 2, Inches(0.5), tbl_top,
                         Inches(12.3), Inches(5.4)).table
tbl.columns[0].width = Inches(6.5)
tbl.columns[1].width = Inches(5.8)
# Header
for j, h in enumerate(["Giulio confirmed", "Action taken / scheduled"]):
    cell = tbl.cell(0, j)
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    tf = cell.text_frame
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    tf.paragraphs[0].text = ""
    r = tf.paragraphs[0].add_run(); r.text = h
    r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = WHITE

for i, (a, b) in enumerate(rows_data):
    for j, val in enumerate([a, b]):
        cell = tbl.cell(i + 1, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else PALEGREY
        tf = cell.text_frame
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
        tf.word_wrap = True
        tf.paragraphs[0].text = ""
        rr = tf.paragraphs[0].add_run(); rr.text = val
        rr.font.name = FONT; rr.font.size = Pt(11); rr.font.color.rgb = BODY

add_footer(s, 9)


# =========================================================================
# SLIDE 10 — WEEK 2 PLAN
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_title(s, "Week 2 Plan",
          context="20 – 24 April 2026")

w2 = [
    ("Mon 20 Apr", "Conversion Tracking Audit",
     "End-to-end verification on all Search campaign landing pages. Confirmation that Lead / Booking / Purchase fire reliably — the gate to all CPA analysis going forward."),
    ("Tue 21 Apr", "Claim Corrections + Ad Strength",
     "Remove \"99.12%\" and \"Save 60%\" from all ads. Reword Top 10 / Best Clinic claims per Giulio's replies. AG22-25 Ad Strength audit once Google clears Pending. Second pass on AG1-21."),
    ("Wed 22 Apr", "PMax Rebuild",
     "New asset groups with policy-safe copy aligned to Giulio's verified claims. Target: clear the \"asset groups limited by policy\" flag that has restricted reach since engagement started."),
    ("Thu 23 Apr", "Ad Extensions",
     "Sitelinks, callouts, and structured snippets deployed across all 3 live campaigns. First time the account has had a complete extensions suite."),
    ("Fri 24 Apr", "Week 2 Performance Review + Report",
     "7-day comparison vs Week 1. Signal emerging on Brand (mature) and PMax (settling). DII still in learning but trajectory visible. Week 2 report delivered EoD."),
]

tbl_top = Inches(1.75)
n = len(w2) + 1
tbl = s.shapes.add_table(n, 3, Inches(0.5), tbl_top,
                         Inches(12.3), Inches(5.15)).table
tbl.columns[0].width = Inches(1.6)
tbl.columns[1].width = Inches(3.0)
tbl.columns[2].width = Inches(7.7)
# Header
for j, h in enumerate(["Day", "Focus", "Expected output"]):
    cell = tbl.cell(0, j)
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    tf = cell.text_frame
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    tf.paragraphs[0].text = ""
    r = tf.paragraphs[0].add_run(); r.text = h
    r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = WHITE

for i, (day, focus, out) in enumerate(w2):
    for j, val in enumerate([day, focus, out]):
        cell = tbl.cell(i + 1, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else PALEGREY
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
        tf.paragraphs[0].text = ""
        rr = tf.paragraphs[0].add_run(); rr.text = val
        rr.font.name = FONT
        if j == 0:
            rr.font.size = Pt(12); rr.font.bold = True; rr.font.color.rgb = NAVY
        elif j == 1:
            rr.font.size = Pt(12); rr.font.bold = True; rr.font.color.rgb = GBLUE
        else:
            rr.font.size = Pt(11); rr.font.color.rgb = BODY

add_footer(s, 10)


# =========================================================================
# SLIDE 11 — CLOSER
# =========================================================================
s = prs.slides.add_slide(BLANK)

add_text(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(0.9),
         "Questions?", size=54, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.5),
         "Happy to walk through any section of this report or align on Week 2 priorities over a call.",
         size=16, color=GBLUE, align=PP_ALIGN.CENTER, italic=True)

# Contact details
add_text(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.4),
         "Christopher Hoole  ·  Google Ads Specialist",
         size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.4),
         "christopherhoole.com   ·   chrishoole101@gmail.com   ·   +44 7451 252857",
         size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Logo bottom right
s.shapes.add_picture(LOGO, Inches(12.3), Inches(6.7), height=Inches(0.55))


# =========================================================================
# SAVE
# =========================================================================
OUT = r"potential_clients/Inserta Dental/End-of-week reports/DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v1.pptx"
import os
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
