"""
Dental by Design - Report 1: Waste Spend Analysis (v4)
v4: Fix slide 11 with matched locations postcode data, fix slide 13 for both years,
update slide 14 quick win #1 to negative keyword overhaul.
"""
import os, sys, warnings, json
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

warnings.filterwarnings('ignore')

# ── Paths ──
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
D25 = os.path.join(SCRIPT_DIR, 'data', '2025')
D26 = os.path.join(SCRIPT_DIR, 'data', '2026')
DOLD = os.path.join(SCRIPT_DIR, 'data', 'old')
CHARTS = os.path.join(SCRIPT_DIR, 'charts')
REPORTS = os.path.join(SCRIPT_DIR, 'reports')
LOGO = os.path.join(SCRIPT_DIR, 'act_logo_official.png')
os.makedirs(CHARTS, exist_ok=True)
os.makedirs(REPORTS, exist_ok=True)

# ── Design tokens ──
BLUE=RGBColor(0x42,0x85,0xF4); RED=RGBColor(0xEA,0x43,0x35)
YELLOW=RGBColor(0xFB,0xBC,0x05); GREEN=RGBColor(0x34,0xA8,0x53)
NAVY=RGBColor(0x1A,0x23,0x7E); BLACK=RGBColor(0x1A,0x1A,0x1A)
GREY_BG=RGBColor(0xF5,0xF6,0xFA); WHITE=RGBColor(0xFF,0xFF,0xFF)
BORDER=RGBColor(0xE2,0xE8,0xF0)
BLUE_TINT=RGBColor(0xE8,0xF0,0xFE); RED_TINT=RGBColor(0xFC,0xE8,0xE6)
GREEN_TINT=RGBColor(0xE6,0xF4,0xEA); GREY_MID=RGBColor(0xCB,0xD5,0xE1)
FONT='Calibri'
CB='#4285F4'; CR='#EA4335'; CY='#FBBC05'; CG='#34A853'; CN='#1A237E'

plt.rcParams.update({'font.family':'Calibri','font.size':11,
    'axes.spines.top':False,'axes.spines.right':False,
    'axes.edgecolor':'#E2E8F0','axes.linewidth':0.8,
    'figure.facecolor':'white','axes.facecolor':'white',
    'grid.color':'#F0F0F0','grid.linewidth':0.5})

def ld(folder,f): return pd.read_csv(f"{folder}/{f}",skiprows=2,thousands=',')
def ct(df,col): return df[df[col].notna()&~df[col].astype(str).str.startswith('Total')&(df[col]!=' --')]

# ══════════════════════════════════════════════════════════════════
# DATA LOADING
# ══════════════════════════════════════════════════════════════════
print("Loading data...")

# 2025
c25=ct(ld(D25,"Campaign report 2025.csv"),'Campaign')
st25=ct(ld(D25,"Search terms report 2025.csv"),'Search term')
dv25=ct(ld(D25,"Device report 2025.csv"),'Device')
dow25=ct(ld(D25,"Campaign report by day of week 2025.csv"),'Day of the week')
hod25=ct(ld(D25,"Campaign report by hour of day 2025.csv"),'Hour of the day')
kw25=ct(ld(D25,"Search keyword report 2025.csv"),'Keyword')
kw25=kw25[~kw25['Match type'].astype(str).str.startswith('Total')]
geo25=ct(ld(D25,"Location report 2025.csv"),'Location')
ml25=ct(ld(D25,"Matched locations report 2025.csv"),'Matched location')

# 2026
c26=ct(ld(D26,"Campaign report 2026.csv"),'Campaign')
st26=ct(ld(D26,"Search terms report 2026.csv"),'Search term')
dv26=ct(ld(D26,"Device report 2026.csv"),'Device')
dow26=ct(ld(D26,"Campaign report by day of week 2026.csv"),'Day of the week')
hod26=ct(ld(D26,"Campaign report by hour of day 2026.csv"),'Hour of the day')
kw26=ct(ld(D26,"Search keyword report 2026.csv"),'Keyword')
kw26=kw26[~kw26['Match type'].astype(str).str.startswith('Total')]
geo26=ct(ld(D26,"Location report 2026.csv"),'Location')
ml26=ct(ld(D26,"Matched locations report 2026.csv"),'Matched location')

# Verified headline numbers
S25=c25['Cost'].sum(); V25=c25['Conversions'].sum(); CPA25=S25/V25
S26=c26['Cost'].sum(); V26=c26['Conversions'].sum(); CPA26=S26/V26
NB25=c25[c25['Campaign']!='Brand']['Cost'].sum()
NV25=c25[c25['Campaign']!='Brand']['Conversions'].sum()
NBCPA25=NB25/NV25
NB26=c26[c26['Campaign']!='Brand']['Cost'].sum()
NV26=c26[c26['Campaign']!='Brand']['Conversions'].sum()
NBCPA26=NB26/NV26

# Search term waste
ZS25=st25[st25['Conversions']==0]['Cost'].sum()
ZP25=ZS25/st25['Cost'].sum()*100
ZS26=st26[st26['Conversions']==0]['Cost'].sum()
ZP26=ZS26/st26['Cost'].sum()*100

# Combined
TOTAL_S=S25+S26; TOTAL_V=V25+V26; TOTAL_CPA=TOTAL_S/TOTAL_V
TOTAL_ZS=ZS25+ZS26

print(f"2025: L{S25:,.0f} spend, {V25:,.0f} cv, L{CPA25:.0f} CPA, NB CPA L{NBCPA25:.0f}")
print(f"2026: L{S26:,.0f} spend, {V26:,.0f} cv, L{CPA26:.0f} CPA, NB CPA L{NBCPA26:.0f}")
print(f"ST waste 2025: L{ZS25:,.0f} ({ZP25:.1f}%), 2026: L{ZS26:,.0f} ({ZP26:.1f}%)")

# Categorize 2026 search term waste
big26=st26[(st26['Conversions']==0)&(st26['Cost']>=1)].copy()
def categorize(t):
    t=str(t).lower()
    if any(c in t for c in ['brighton implant','damira','bupa','mydentist','dental prime','turkey','abroad','dental tourism','smile direct']): return 'Competitor / Dental Tourism'
    if any(w in t for w in ['nhs','free dental','free implant','free consultation']): return 'NHS / Free Treatment'
    if any(w in t for w in ['how to','what is','can i','should i','is it','why do','how long','reviews','blog','reddit','forum','pros and cons']): return 'Informational / Research'
    if any(w in t for w in ['denture','crown','filling','extraction','root canal','braces','orthodont','retainer','wisdom','whitening','bonding']) and 'implant' not in t: return 'Wrong Treatment Type'
    if any(w in t for w in ['manchester','birmingham','leeds','liverpool','glasgow','edinburgh','cardiff','belfast','nottingham','sheffield','newcastle','brighton','bristol','kent','essex','surrey']) and 'london' not in t: return 'Wrong Location'
    if any(w in t for w in ['cost','price','cheap','cheapest','affordable','low cost','how much']) and not any(w in t for w in ['london','near me','clinic','dentist','hammersmith','book','consult']): return 'Price Research (No Location)'
    return 'Other Non-Converting'

big26['Category']=big26['Search term'].apply(categorize)
cat26=big26.groupby('Category').agg({'Cost':'sum','Search term':'count','Clicks':'sum'}).sort_values('Cost',ascending=False)

# 2025 waste categories
big25=st25[(st25['Conversions']==0)&(st25['Cost']>=1)].copy()
big25['Category']=big25['Search term'].apply(categorize)
cat25=big25.groupby('Category').agg({'Cost':'sum','Search term':'count','Clicks':'sum'}).sort_values('Cost',ascending=False)

# Device aggregation
def agg_device(dv):
    g=dv.groupby('Device').agg({'Cost':'sum','Clicks':'sum','Conversions':'sum'}).sort_values('Cost',ascending=False)
    r={}
    for d,row in g.iterrows():
        if row['Conversions']>0:
            r[d]={'spend':row['Cost'],'conv':row['Conversions'],'cpa':row['Cost']/row['Conversions'],'cvr':row['Conversions']/row['Clicks']*100 if row['Clicks']>0 else 0}
    return r

dev25=agg_device(dv25); dev26=agg_device(dv26)

# Day of week
def agg_dow(dow):
    g=dow.groupby('Day of the week').agg({'Cost':'sum','Clicks':'sum','Conversions':'sum'}).reset_index()
    r={}
    for _,row in g.iterrows():
        if row['Conversions']>0:
            r[row['Day of the week']]={'spend':row['Cost'],'conv':row['Conversions'],'cpa':row['Cost']/row['Conversions'],'cvr':row['Conversions']/row['Clicks']*100 if row['Clicks']>0 else 0}
    return r

days25=agg_dow(dow25); days26=agg_dow(dow26)

# Hour of day
def agg_hod(hod):
    hod=hod.copy()
    hod['Hour of the day']=pd.to_numeric(hod['Hour of the day'],errors='coerce')
    hod=hod.dropna(subset=['Hour of the day'])
    g=hod.groupby('Hour of the day').agg({'Cost':'sum','Clicks':'sum','Conversions':'sum'}).reset_index().sort_values('Hour of the day')
    g['CPA']=g.apply(lambda r:r['Cost']/r['Conversions'] if r['Conversions']>0 else 0,axis=1)
    return g

hours25=agg_hod(hod25); hours26=agg_hod(hod26)

# Keywords
def agg_kw(kw):
    g=kw.groupby('Keyword').agg({'Cost':'sum','Conversions':'sum','Clicks':'sum'}).reset_index()
    z=g[(g['Conversions']==0)&(g['Cost']>0)]
    h=g[g['Conversions']>0].copy()
    h['CPA']=h['Cost']/h['Conversions']
    bad=h[h['CPA']>150]
    return {'zero_count':len(z),'zero_spend':z['Cost'].sum(),'high_count':len(bad),'high_spend':bad['Cost'].sum(),'high_df':bad}

kwd25=agg_kw(kw25); kwd26=agg_kw(kw26)

print("Data loaded and verified.")

# ══════════════════════════════════════════════════════════════════
# CHART GENERATION
# ══════════════════════════════════════════════════════════════════
print("Generating charts...")

# CHART 1: YoY CPA comparison (the killer chart)
fig,ax=plt.subplots(figsize=(8,4))
metrics=['Account CPA','Non-Brand CPA','PMax CPA','Search CPA']
v25_vals=[CPA25,NBCPA25]
v26_vals=[CPA26,NBCPA26]
# Get PMax and Search CPAs
for c in ['Dental Implants - Pmax','Dental Implants Intent']:
    for camp_df,yr in [(c25,'2025'),(c26,'2026')]:
        sub=camp_df[camp_df['Campaign']==c] if c in camp_df['Campaign'].values else pd.DataFrame()
        if yr=='2025':
            v25_vals.append(sub['Cost'].sum()/max(sub['Conversions'].sum(),0.1) if len(sub)>0 and sub['Cost'].sum()>0 else 0)
        else:
            v26_vals.append(sub['Cost'].sum()/max(sub['Conversions'].sum(),0.1) if len(sub)>0 and sub['Cost'].sum()>0 else 0)

# Filter out zero entries
valid=[(m,a,b) for m,a,b in zip(metrics,v25_vals,v26_vals) if a>0 and b>0]
ms=[v[0] for v in valid]; a25=[v[1] for v in valid]; a26=[v[2] for v in valid]
x=np.arange(len(ms)); w=0.35
b1=ax.bar(x-w/2,a25,w,label='2025',color=CB,alpha=0.8)
b2=ax.bar(x+w/2,a26,w,label='2026 (Q1)',color=CR,alpha=0.8)
for bar,val in zip(b1,a25): ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+1,f'\u00a3{val:.0f}',ha='center',fontsize=10,color=CN,fontweight='bold')
for bar,val in zip(b2,a26): ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+1,f'\u00a3{val:.0f}',ha='center',fontsize=10,color=CR,fontweight='bold')
ax.set_xticks(x); ax.set_xticklabels(ms,fontsize=10)
ax.set_ylabel('Cost Per Conversion (\u00a3)',fontsize=10)
ax.legend(fontsize=10); ax.set_ylim(0,max(a26)*1.2)
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f'\u00a3{x:.0f}'))
plt.tight_layout()
plt.savefig(f"{CHARTS}/yoy_cpa.png",dpi=200,bbox_inches='tight'); plt.close()

# CHART 2: Search term waste YoY
fig,ax=plt.subplots(figsize=(6,4))
cats=['2025','2026 (Q1)']
waste_vals=[ZS25,ZS26]
waste_pcts=[ZP25,ZP26]
bars=ax.bar(cats,waste_vals,color=[CB,CR],width=0.5)
for bar,val,pct in zip(bars,waste_vals,waste_pcts):
    ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+1000,f'\u00a3{val:,.0f}\n({pct:.1f}%)',ha='center',fontsize=11,color=CN,fontweight='bold')
ax.set_ylabel('Zero-Conversion Search Term Spend (\u00a3)',fontsize=10)
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f'\u00a3{x:,.0f}'))
ax.set_ylim(0,max(waste_vals)*1.25)
plt.tight_layout()
plt.savefig(f"{CHARTS}/st_waste_yoy.png",dpi=200,bbox_inches='tight'); plt.close()

# CHART 3: Waste categories 2026
fig,ax=plt.subplots(figsize=(7,4))
cats_list=cat26.index.tolist(); vals=cat26['Cost'].tolist()
colors=[CR if v>10000 else CY if v>3000 else CB for v in vals]
bars=ax.barh(range(len(cats_list)),vals,color=colors,height=0.6)
ax.set_yticks(range(len(cats_list))); ax.set_yticklabels(cats_list,fontsize=10)
ax.invert_yaxis()
for bar,val in zip(bars,vals): ax.text(bar.get_width()+300,bar.get_y()+bar.get_height()/2,f'\u00a3{val:,.0f}',va='center',fontsize=9,color=CN)
ax.set_xlim(0,max(vals)*1.3)
ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f'\u00a3{x:,.0f}'))
ax.set_xlabel('Wasted Spend (\u00a3)',fontsize=10)
plt.tight_layout()
plt.savefig(f"{CHARTS}/waste_cats_2026.png",dpi=200,bbox_inches='tight'); plt.close()

# CHART 4: Device CPA YoY
fig,ax=plt.subplots(figsize=(7,4))
devs=['Mobile Phones','Computers','Tablets']
d25v=[dev25.get(d,{}).get('cpa',0) for d in devs]
d26v=[dev26.get(d,{}).get('cpa',0) for d in devs]
x=np.arange(len(devs)); w=0.35
b1=ax.bar(x-w/2,d25v,w,label='2025',color=CB,alpha=0.8)
b2=ax.bar(x+w/2,d26v,w,label='2026 (Q1)',color=CR,alpha=0.8)
for bar,val in zip(b1,d25v): ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+1,f'\u00a3{val:.0f}',ha='center',fontsize=10,color=CN,fontweight='bold')
for bar,val in zip(b2,d26v): ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+1,f'\u00a3{val:.0f}',ha='center',fontsize=10,color=CR,fontweight='bold')
ax.set_xticks(x); ax.set_xticklabels(devs,fontsize=10)
ax.set_ylabel('CPA (\u00a3)',fontsize=10); ax.legend(fontsize=10)
ax.set_ylim(0,max(d26v)*1.2)
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f'\u00a3{x:.0f}'))
plt.tight_layout()
plt.savefig(f"{CHARTS}/device_yoy.png",dpi=200,bbox_inches='tight'); plt.close()

# CHART 5: Day of week CPA 2026
fig,ax=plt.subplots(figsize=(8,4))
day_order=['Monday','Tuesday','Wednesday','Thursday','Friday','Saturday','Sunday']
cpas26=[days26.get(d,{}).get('cpa',0) for d in day_order]
avg26=sum(cpas26)/len(cpas26)
colors_d=[CR if c>avg26*1.05 else CG if c<avg26*0.95 else CB for c in cpas26]
bars=ax.bar(range(7),cpas26,color=colors_d,width=0.6)
ax.set_xticks(range(7)); ax.set_xticklabels([d[:3] for d in day_order],fontsize=10)
ax.axhline(y=avg26,color=CN,linestyle='--',linewidth=1,alpha=0.7)
ax.text(6.5,avg26+0.5,f'Avg: \u00a3{avg26:.0f}',fontsize=9,color=CN,ha='right')
for bar,cpa in zip(bars,cpas26): ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+0.3,f'\u00a3{cpa:.0f}',ha='center',fontsize=10,color=CN,fontweight='bold')
ax.set_ylabel('CPA (\u00a3)',fontsize=10); ax.set_ylim(0,max(cpas26)*1.15)
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f'\u00a3{x:.0f}'))
ax.set_title('2026 (Q1)',fontsize=12,color=CN,fontweight='bold',loc='left')
plt.tight_layout()
plt.savefig(f"{CHARTS}/dow_2026.png",dpi=200,bbox_inches='tight'); plt.close()

# CHART 6: Hour of day CPA 2026
fig,ax=plt.subplots(figsize=(10,4))
h26=hours26; avg_h=h26['Cost'].sum()/h26['Conversions'].sum()
colors_h=[CR if c>avg_h*1.1 else CG if c<avg_h*0.9 else CB for c in h26['CPA'].values]
ax.bar(h26['Hour of the day'].astype(int),h26['CPA'],color=colors_h,width=0.7)
ax.axhline(y=avg_h,color=CN,linestyle='--',linewidth=1,alpha=0.7)
ax.text(23,avg_h+1,f'Avg: \u00a3{avg_h:.0f}',fontsize=9,color=CN,ha='right')
ax.set_xlabel('Hour of Day',fontsize=10); ax.set_ylabel('CPA (\u00a3)',fontsize=10)
ax.set_xticks(range(24)); ax.set_xticklabels([f'{h:02d}' for h in range(24)],fontsize=8)
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f'\u00a3{x:.0f}'))
ax.set_title('2026 (Q1)',fontsize=12,color=CN,fontweight='bold',loc='left')
plt.tight_layout()
plt.savefig(f"{CHARTS}/hod_2026.png",dpi=200,bbox_inches='tight'); plt.close()

# CHART 7: PMax waste rate YoY
fig,ax=plt.subplots(figsize=(5,3.5))
pm25=st25[st25['Campaign type']=='Performance Max']
pm26=st26[st26['Campaign type']=='Performance Max']
pm25_pct=pm25[pm25['Conversions']==0]['Cost'].sum()/pm25['Cost'].sum()*100
pm26_pct=pm26[pm26['Conversions']==0]['Cost'].sum()/pm26['Cost'].sum()*100
bars=ax.bar(['2025','2026 (Q1)'],[pm25_pct,pm26_pct],color=[CY,CR],width=0.4)
for bar,val in zip(bars,[pm25_pct,pm26_pct]): ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+0.5,f'{val:.1f}%',ha='center',fontsize=12,color=CN,fontweight='bold')
ax.set_ylabel('% of PMax ST Spend with 0 Conversions',fontsize=9)
ax.set_ylim(0,75)
ax.set_title('PMax Search Term Waste Rate',fontsize=11,color=CN,fontweight='bold',loc='left')
plt.tight_layout()
plt.savefig(f"{CHARTS}/pmax_waste_yoy.png",dpi=200,bbox_inches='tight'); plt.close()

# CHART 8: Day of week YoY comparison
fig,ax=plt.subplots(figsize=(9,4))
day_order=['Monday','Tuesday','Wednesday','Thursday','Friday','Saturday','Sunday']
cpas25_d=[days25.get(d,{}).get('cpa',0) for d in day_order]
cpas26_d=[days26.get(d,{}).get('cpa',0) for d in day_order]
x=np.arange(7); w=0.35
b1=ax.bar(x-w/2,cpas25_d,w,label='2025',color=CB,alpha=0.8)
b2=ax.bar(x+w/2,cpas26_d,w,label='2026 (Q1)',color=CR,alpha=0.8)
for bar,val in zip(b1,cpas25_d): ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+0.3,f'\u00a3{val:.0f}',ha='center',fontsize=9,color=CN,fontweight='bold')
for bar,val in zip(b2,cpas26_d): ax.text(bar.get_x()+bar.get_width()/2,bar.get_height()+0.3,f'\u00a3{val:.0f}',ha='center',fontsize=9,color=CR,fontweight='bold')
ax.set_xticks(x); ax.set_xticklabels([d[:3] for d in day_order],fontsize=10)
ax.set_ylabel('CPA (\u00a3)',fontsize=10); ax.legend(fontsize=10)
ax.set_ylim(0,max(cpas26_d)*1.2)
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f'\u00a3{x:.0f}'))
plt.tight_layout()
plt.savefig(f"{CHARTS}/dow_yoy.png",dpi=200,bbox_inches='tight'); plt.close()

# CHART 9: Hour of day YoY (overlay lines)
fig,ax=plt.subplots(figsize=(10,4))
h25=hours25; h26=hours26
ax.plot(h25['Hour of the day'].astype(int),h25['CPA'],color=CB,linewidth=2,marker='o',markersize=4,label='2025',alpha=0.8)
ax.plot(h26['Hour of the day'].astype(int),h26['CPA'],color=CR,linewidth=2,marker='s',markersize=4,label='2026 (Q1)',alpha=0.8)
ax.set_xlabel('Hour of Day',fontsize=10); ax.set_ylabel('CPA (\u00a3)',fontsize=10)
ax.set_xticks(range(24)); ax.set_xticklabels([f'{h:02d}' for h in range(24)],fontsize=8)
ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f'\u00a3{x:.0f}'))
ax.legend(fontsize=10); ax.grid(axis='y',alpha=0.3)
plt.tight_layout()
plt.savefig(f"{CHARTS}/hod_yoy.png",dpi=200,bbox_inches='tight'); plt.close()

# CHART 10: Waste categories YoY comparison
all_cats=sorted(set(cat25.index.tolist()+cat26.index.tolist()))
fig,ax=plt.subplots(figsize=(8,4.5))
y=np.arange(len(all_cats)); w=0.35
v25c=[cat25.loc[c,'Cost'] if c in cat25.index else 0 for c in all_cats]
v26c=[cat26.loc[c,'Cost'] if c in cat26.index else 0 for c in all_cats]
b1=ax.barh(y-w/2,v25c,w,label='2025',color=CB,alpha=0.8)
b2=ax.barh(y+w/2,v26c,w,label='2026 (Q1)',color=CR,alpha=0.8)
ax.set_yticks(y); ax.set_yticklabels(all_cats,fontsize=9)
ax.invert_yaxis(); ax.legend(fontsize=10)
ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_:f'\u00a3{x:,.0f}'))
ax.set_xlabel('Wasted Spend (\u00a3)',fontsize=10)
plt.tight_layout()
plt.savefig(f"{CHARTS}/waste_cats_yoy.png",dpi=200,bbox_inches='tight'); plt.close()

print("Charts generated.")

# ══════════════════════════════════════════════════════════════════
# PPTX HELPERS
# ══════════════════════════════════════════════════════════════════
prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)

def rect(sl,l,t,w,h,fc,lc=None,lw=None):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    s.fill.solid();s.fill.fore_color.rgb=fc
    if lc:s.line.color.rgb=lc;s.line.width=lw or Pt(1)
    else:s.line.fill.background()
    return s

def txt(sl,l,t,w,h,c,sz=11,co=BLACK,b=False,a=PP_ALIGN.LEFT):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    r=tf.paragraphs[0].add_run();r.text=c;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
    tf.paragraphs[0].alignment=a;return bx

def multirun(sl,l,t,w,h,runs,a=PP_ALIGN.LEFT):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.alignment=a
    for tx,sz,co,b in runs:
        r=p.add_run();r.text=tx;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
    return bx

def multipara(sl,l,t,w,h,pdata):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    for i,(tx,sz,co,b,sp) in enumerate(pdata):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run();r.text=tx;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
        if sp:p.space_after=Pt(sp)
    return bx

def bullets(sl,l,t,w,h,items,sz=11,co=BLACK):
    return multipara(sl,l,t,w,h,[(i,sz,co,False,5) for i in items])

def top_bar(sl,h=Inches(0.07)):
    for i,c in enumerate([BLUE,RED,YELLOW,GREEN]): rect(sl,Inches(i*3.333),Inches(0),Inches(3.333),h,c)

def bot_bar(sl,y=Inches(6.92)):
    for p,c in zip([Inches(0.60),Inches(3.63),Inches(6.67),Inches(9.70)],[BLUE,RED,YELLOW,GREEN]):rect(sl,p,y,Inches(3.03),Inches(0.03),c)

def bot_bar_title(sl):
    for i,c in enumerate([BLUE,RED,YELLOW,GREEN]):rect(sl,Inches(i*3.333),Inches(7.0),Inches(3.333),Inches(0.04),c)

def sidebar(sl): rect(sl,Inches(0),Inches(0.07),Inches(0.12),Inches(7.43),BLUE)

def logo(sl,l=Inches(0.60),t=Inches(0.50),s=Inches(0.65)):
    try:sl.shapes.add_picture(LOGO,l,t,s,s)
    except:pass

def footer(sl,n):
    bot_bar(sl)
    try:sl.shapes.add_picture(LOGO,Inches(0.60),Inches(7.0),Inches(0.22),Inches(0.22))
    except:pass
    txt(sl,Inches(0.90),Inches(7.0),Inches(6),Inches(0.25),"Christopher Hoole  |  christopherhoole.com  |  Confidential",11,NAVY,True)
    txt(sl,Inches(12.23),Inches(7.0),Inches(0.50),Inches(0.25),str(n),11,NAVY,a=PP_ALIGN.RIGHT)

def stitle(sl,t,sub=None):
    txt(sl,Inches(0.60),Inches(0.30),Inches(7),Inches(0.50),t,28,NAVY,True)
    if sub:txt(sl,Inches(0.60),Inches(0.85),Inches(9),Inches(0.30),sub,11,RED)

def badge(sl,t):
    rect(sl,Inches(9.13),Inches(0.30),Inches(3.60),Inches(0.45),WHITE,BLUE,Pt(1))
    txt(sl,Inches(9.23),Inches(0.32),Inches(3.40),Inches(0.40),t,11,BLUE,True,PP_ALIGN.CENTER)

def mcard(sl,l,t,v,lab,ac,w=Inches(2.80),h=Inches(1.05)):
    rect(sl,l,t,w,h,WHITE,BORDER,Pt(0.75));rect(sl,l,t,Inches(0.06),h,ac)
    txt(sl,l+Inches(0.20),t+Inches(0.10),w-Inches(0.30),Inches(0.50),v,22,ac,True)
    txt(sl,l+Inches(0.20),t+Inches(0.60),w-Inches(0.30),Inches(0.35),lab,11,BLACK)

def minicard(sl,l,t,v,lab,ac):
    rect(sl,l,t,Inches(1.83),Inches(1.15),WHITE);rect(sl,l,t,Inches(0.06),Inches(1.15),ac)
    txt(sl,l+Inches(0.15),t+Inches(0.12),Inches(1.58),Inches(0.45),v,22,ac,True)
    txt(sl,l+Inches(0.15),t+Inches(0.60),Inches(1.58),Inches(0.30),lab,11,BLACK)

def ibox(sl,l,t,w,h,ti,body,ac,bg=None):
    rect(sl,l,t,w,h,bg or GREY_BG);rect(sl,l,t,Inches(0.06),h,ac)
    txt(sl,l+Inches(0.25),t+Inches(0.10),w-Inches(0.40),Inches(0.30),ti,14,NAVY,True)
    txt(sl,l+Inches(0.25),t+Inches(0.40),w-Inches(0.40),h-Inches(0.50),body,11,BLACK)

def rbox(sl,l,t,w,h,c):
    rect(sl,l,t,w,h,BLUE_TINT,BLUE,Pt(1))
    txt(sl,l+Inches(0.25),t+Inches(0.05),w-Inches(0.50),h-Inches(0.10),c,11,NAVY,True)

def add_table(sl,l,t,w,h,hdrs,rows,cw=None):
    ts=sl.shapes.add_table(len(rows)+1,len(hdrs),l,t,w,h);tb=ts.table
    if cw:
        for i,c in enumerate(cw):tb.columns[i].width=Inches(c)
    for ci,hd in enumerate(hdrs):
        c=tb.cell(0,ci);c.text=hd;c.fill.solid();c.fill.fore_color.rgb=NAVY
        for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.font.color.rgb=WHITE;p.font.bold=True;p.font.name=FONT
    for ri,rd in enumerate(rows):
        for ci,v in enumerate(rd):
            c=tb.cell(ri+1,ci);c.text=str(v);c.fill.solid();c.fill.fore_color.rgb=GREY_BG if ri%2==0 else WHITE
            for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.font.color.rgb=BLACK;p.font.name=FONT
    return ts

def img(sl,path,l,t,w,h):
    try:sl.shapes.add_picture(path,l,t,w,h)
    except Exception as e:print(f"  Warning: {e}")

def ncircle(sl,l,t,num,col):
    ns=sl.shapes.add_shape(MSO_SHAPE.OVAL,l,t,Inches(0.45),Inches(0.45))
    ns.fill.solid();ns.fill.fore_color.rgb=col;ns.line.fill.background()
    tf=ns.text_frame;tf.paragraphs[0].alignment=PP_ALIGN.CENTER;tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    r=tf.paragraphs[0].add_run();r.text=str(num);r.font.size=Pt(16);r.font.color.rgb=WHITE;r.font.bold=True;r.font.name=FONT


# ══════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════
print("Building slides...")

# ── S1: TITLE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=GREY_BG
top_bar(sl);sidebar(sl);logo(sl)
txt(sl,Inches(0.60),Inches(1.40),Inches(5.74),Inches(1.80),"Waste Spend\nAnalysis",44,NAVY,True)
rect(sl,Inches(0.60),Inches(3.30),Inches(2.50),Inches(0.05),BLUE)
txt(sl,Inches(0.60),Inches(3.55),Inches(5.50),Inches(0.50),"Dental by Design  |  Prodent Group",22,BLUE)
multipara(sl,Inches(0.60),Inches(4.30),Inches(5.50),Inches(1.00),[
    ("Christopher Hoole",11,BLACK,True,2),("Google Ads Specialist  |  April 2026",11,BLACK,False,2),
    ("christopherhoole.com",11,BLUE,False,0)])
rect(sl,Inches(0.60),Inches(5.50),Inches(4.50),Inches(0.50),WHITE,BLUE,Pt(1))
txt(sl,Inches(0.70),Inches(5.52),Inches(4.30),Inches(0.45),"Data: 2025 (full year) + 2026 (Q1, Jan\u2013Apr)",11,BLUE,True)

# Hero — lead with waste spend
rect(sl,Inches(6.80),Inches(0.50),Inches(5.90),Inches(2.50),WHITE)
rect(sl,Inches(6.80),Inches(0.50),Inches(0.08),Inches(2.50),RED)
txt(sl,Inches(7.10),Inches(0.60),Inches(5.40),Inches(0.35),"Identified Waste (Search Terms)",14,NAVY)
txt(sl,Inches(7.10),Inches(1.00),Inches(5.40),Inches(0.80),f"\u00a3{TOTAL_ZS:,.0f}",44,RED,True)
txt(sl,Inches(7.10),Inches(1.80),Inches(5.40),Inches(0.50),f"2025: \u00a3{ZS25:,.0f} ({ZP25:.0f}%)  |  2026 Q1: \u00a3{ZS26:,.0f} ({ZP26:.0f}%)  |  Waste rate increasing",11,BLACK)

minicard(sl,Inches(6.80),Inches(3.20),f"\u00a3{NBCPA25:.0f}\u2192\u00a3{NBCPA26:.0f}","Non-Brand CPA (YoY)",RED)
minicard(sl,Inches(8.83),Inches(3.20),f"\u00a3{TOTAL_S:,.0f}","Total Spend Analysed",BLUE)
minicard(sl,Inches(10.87),Inches(3.20),f"{TOTAL_V:,.0f}","Total Conversions",GREEN)
bot_bar_title(sl)


# ── S2: EXECUTIVE SUMMARY ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06));stitle(sl,"Executive Summary");badge(sl,"2025 + 2026 (Q1)")

mcard(sl,Inches(0.60),Inches(1.20),f"\u00a3{TOTAL_S:,.0f}","Total Spend Analysed",NAVY)
mcard(sl,Inches(3.70),Inches(1.20),f"\u00a3{TOTAL_ZS:,.0f}","Search Term Waste (0 Conv)",RED)
mcard(sl,Inches(6.80),Inches(1.20),f"\u00a3{NBCPA25:.0f} \u2192 \u00a3{NBCPA26:.0f}","Non-Brand CPA (YoY)",YELLOW)
mcard(sl,Inches(9.90),Inches(1.20),f"{TOTAL_V:,.0f}","Total Conversions",GREEN)

txt(sl,Inches(0.60),Inches(2.60),Inches(6),Inches(0.35),"Key Findings",14,NAVY,True)
bullets(sl,Inches(0.70),Inches(3.00),Inches(11.50),Inches(2.60),[
    f"Non-brand CPA has increased {(NBCPA26-NBCPA25)/NBCPA25*100:.0f}% year-on-year (\u00a3{NBCPA25:.0f} \u2192 \u00a3{NBCPA26:.0f}) while monthly spend has nearly doubled.",
    f"\u00a3{ZS26:,.0f} spent on zero-conversion search terms in Q1 2026 alone \u2014 {ZP26:.0f}% of all tracked search term spend generated no patient enquiries.",
    f"PMax search term waste has risen from {pm25[pm25['Conversions']==0]['Cost'].sum()/pm25['Cost'].sum()*100:.0f}% to {pm26[pm26['Conversions']==0]['Cost'].sum()/pm26['Cost'].sum()*100:.0f}% year-on-year \u2014 the campaign is scaling but efficiency is declining.",
    f"Desktop CPA is \u00a3{dev26['Computers']['cpa']:.0f} vs mobile \u00a3{dev26['Mobile Phones']['cpa']:.0f} ({(dev26['Computers']['cpa']-dev26['Mobile Phones']['cpa'])/dev26['Mobile Phones']['cpa']*100:.0f}% higher) with no device bid adjustments in place.",
    f"17 keywords are running at CPA above \u00a3150 in 2026, spending \u00a3{kwd26['high_spend']:,.0f} at 2-5x the account average.",
])

rbox(sl,Inches(0.60),Inches(5.85),Inches(12.13),Inches(0.70),
    "This report identifies two types of waste: money that generated nothing (zero-conversion spend) and money that could work harder (high-CPA inefficiency). Both represent growth opportunity.")
footer(sl,2)


# ── S3: YoY CPA TREND ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06));stitle(sl,"Year-on-Year CPA Comparison","Non-brand CPA has increased across every campaign type and device");badge(sl,"2025 vs 2026")

img(sl,f"{CHARTS}/yoy_cpa.png",Inches(0.60),Inches(1.30),Inches(7.50),Inches(4.00))

# Before/after cards
rect(sl,Inches(8.40),Inches(1.30),Inches(4.30),Inches(1.80),WHITE,BORDER,Pt(0.75))
rect(sl,Inches(8.40),Inches(1.30),Inches(0.06),Inches(1.80),GREEN)
txt(sl,Inches(8.65),Inches(1.40),Inches(3.90),Inches(0.25),"2025 (Full Year)",11,GREEN,True)
txt(sl,Inches(8.65),Inches(1.70),Inches(3.90),Inches(0.30),f"\u00a3{NBCPA25:.0f} Non-Brand CPA",18,GREEN,True)
txt(sl,Inches(8.65),Inches(2.05),Inches(3.90),Inches(0.80),f"\u00a3{S25:,.0f} total spend\n{V25:,.0f} conversions\n15 active campaigns",11,BLACK)

rect(sl,Inches(8.40),Inches(3.30),Inches(4.30),Inches(1.80),WHITE,BORDER,Pt(0.75))
rect(sl,Inches(8.40),Inches(3.30),Inches(0.06),Inches(1.80),RED)
txt(sl,Inches(8.65),Inches(3.40),Inches(3.90),Inches(0.25),"2026 (Q1, Jan\u2013Apr)",11,RED,True)
txt(sl,Inches(8.65),Inches(3.70),Inches(3.90),Inches(0.30),f"\u00a3{NBCPA26:.0f} Non-Brand CPA",18,RED,True)
txt(sl,Inches(8.65),Inches(4.05),Inches(3.90),Inches(0.80),f"\u00a3{S26:,.0f} total spend (annualised \u00a3{S26*3:,.0f})\n{V26:,.0f} conversions\n4 active campaigns",11,BLACK)

ibox(sl,Inches(0.60),Inches(5.55),Inches(12.10),Inches(1.00),
    "The Pattern",
    f"Spend has nearly doubled but CPA has risen {(NBCPA26-NBCPA25)/NBCPA25*100:.0f}%. "
    "The account is scaling aggressively but efficiency is declining \u2014 meaning more money is being spent to acquire each patient. "
    "This is the core issue the optimisations in this report address.",
    BLUE,BLUE_TINT)
footer(sl,3)


# ── S4: SEARCH TERM WASTE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Search Term Waste",f"\u00a3{ZS26:,.0f} spent on zero-conversion search terms in Q1 2026 ({ZP26:.0f}% of search term spend)")
badge(sl,"2026 Focus")

# Top 15 wasted terms 2026
top_waste_26=st26[st26['Conversions']==0].nlargest(15,'Cost')
add_table(sl,Inches(0.60),Inches(1.30),Inches(7.80),Inches(5.00),
    ["Search Term","Spend","Clicks","Campaign"],
    [[r['Search term'][:40],f"\u00a3{r['Cost']:,.0f}",f"{r['Clicks']:.0f}",r['Campaign'][:25]] for _,r in top_waste_26.iterrows()],
    cw=[3.2,0.8,0.6,3.2])

img(sl,f"{CHARTS}/st_waste_yoy.png",Inches(8.70),Inches(1.30),Inches(4.00),Inches(2.80))

ibox(sl,Inches(8.70),Inches(4.30),Inches(4.00),Inches(2.00),
    "Year-on-Year Trend",
    f"2025: \u00a3{ZS25:,.0f} wasted ({ZP25:.0f}% of ST spend)\n"
    f"2026: \u00a3{ZS26:,.0f} wasted ({ZP26:.0f}% of ST spend)\n\n"
    "The waste rate is increasing \u2014 nearly half of all search term spend now generates zero patient enquiries.",
    RED,RED_TINT)
footer(sl,4)


# ── S5: WASTE CATEGORIES ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Waste by Category","Zero-conversion search terms grouped by intent type \u2014 2025 vs 2026")
badge(sl,"2025 vs 2026")

img(sl,f"{CHARTS}/waste_cats_yoy.png",Inches(0.60),Inches(1.30),Inches(7.50),Inches(4.00))

txt(sl,Inches(8.30),Inches(1.30),Inches(4.40),Inches(0.30),"2025 vs 2026 Totals",14,NAVY,True)
cat_items=[]
for cat in ['Other Non-Converting','Price Research (No Location)','Wrong Treatment Type','Competitor / Dental Tourism','NHS / Free Treatment','Informational / Research','Wrong Location']:
    v25=cat25.loc[cat,'Cost'] if cat in cat25.index else 0
    v26=cat26.loc[cat,'Cost'] if cat in cat26.index else 0
    cat_items.append(f"{cat}:  \u00a3{v25:,.0f} \u2192 \u00a3{v26:,.0f}")
bullets(sl,Inches(8.30),Inches(1.70),Inches(4.40),Inches(3.60),cat_items,10)

oc25=cat25.loc['Other Non-Converting','Cost'] if 'Other Non-Converting' in cat25.index else 0
oc26=cat26.loc['Other Non-Converting','Cost'] if 'Other Non-Converting' in cat26.index else 0
pr25=cat25.loc['Price Research (No Location)','Cost'] if 'Price Research (No Location)' in cat25.index else 0
pr26=cat26.loc['Price Research (No Location)','Cost'] if 'Price Research (No Location)' in cat26.index else 0
rbox(sl,Inches(0.60),Inches(5.55),Inches(12.13),Inches(0.80),
    f"\"Other Non-Converting\" (\u00a3{oc25:,.0f} in 2025, \u00a3{oc26:,.0f} in 2026) includes relevant dental terms that didn't convert \u2014 "
    "these need landing page and ad copy work, not just negative keywords. "
    f"\"Price Research\" (\u00a3{pr25:,.0f} \u2192 \u00a3{pr26:,.0f}) are blockable with negative keywords.")
footer(sl,5)


# ── S6: PMAX WASTE DEEP-DIVE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
pm26_spend=c26[c26['Campaign']=='Dental Implants - Pmax']['Cost'].sum()
pm26_conv=c26[c26['Campaign']=='Dental Implants - Pmax']['Conversions'].sum()
pm26_cpa=pm26_spend/pm26_conv
pm26_zero=st26[(st26['Campaign type']=='Performance Max')&(st26['Conversions']==0)]['Cost'].sum()
pm26_zero_pct=pm26_zero/st26[st26['Campaign type']=='Performance Max']['Cost'].sum()*100
stitle(sl,"PMax Campaign \u2014 Rising Waste",f"60% of PMax search term spend in 2026 generated zero conversions, up from 53% in 2025")
badge(sl,"Performance Max")

mcard(sl,Inches(0.60),Inches(1.20),f"\u00a3{pm26_spend:,.0f}","2026 PMax Spend",BLUE,Inches(3.00))
mcard(sl,Inches(3.90),Inches(1.20),f"\u00a3{pm26_cpa:.0f}","2026 PMax CPA",YELLOW,Inches(3.00))
mcard(sl,Inches(7.20),Inches(1.20),f"\u00a3{pm26_zero:,.0f}","2026 Zero-Conv ST Spend",RED,Inches(3.00))
mcard(sl,Inches(10.50),Inches(1.20),f"{pm26_zero_pct:.0f}%","Waste Rate",RED,Inches(2.20))

img(sl,f"{CHARTS}/pmax_waste_yoy.png",Inches(0.60),Inches(2.60),Inches(4.50),Inches(3.00))

ibox(sl,Inches(5.40),Inches(2.60),Inches(7.30),Inches(1.40),
    "Why PMax Waste Is High",
    "Performance Max campaigns show ads across Search, Display, YouTube, Gmail, and Maps. "
    "Many search terms triggered by PMax are broad and irrelevant \u2014 but unlike Search campaigns, "
    "you have limited control over which terms PMax targets. Negative keywords can only be applied at account level for PMax.",
    BLUE,BLUE_TINT)

ibox(sl,Inches(5.40),Inches(4.20),Inches(7.30),Inches(1.40),
    "The Concern",
    f"PMax CPA has risen from \u00a326 in 2025 to \u00a3{pm26_cpa:.0f} in 2026 \u2014 a {(pm26_cpa-26)/26*100:.0f}% increase. "
    "As the campaign scales, it's reaching into less relevant search territory. "
    "The account-level negative keyword lists (currently 784 broad match) need a full audit to protect PMax from irrelevant traffic.",
    RED,RED_TINT)
footer(sl,6)


# ── S7: KEYWORD INEFFICIENCY ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Keyword Inefficiency",f"Zero-conv waste: \u00a3{kwd25['zero_spend']:,.0f} (2025) + \u00a3{kwd26['zero_spend']:,.0f} (2026)  |  High-CPA keywords: \u00a3{kwd25['high_spend']:,.0f} (2025) + \u00a3{kwd26['high_spend']:,.0f} (2026)")
badge(sl,"High-CPA Keywords")

# Two sections: zero-conv + high-CPA
txt(sl,Inches(0.60),Inches(1.20),Inches(5.80),Inches(0.30),"Angle 1: Zero Conversions",14,NAVY,True)
zero_kw_26=kw26.groupby('Keyword').agg({'Cost':'sum','Conversions':'sum','Clicks':'sum'}).reset_index()
zero_kw_26=zero_kw_26[(zero_kw_26['Conversions']==0)&(zero_kw_26['Cost']>0)].nlargest(8,'Cost')
add_table(sl,Inches(0.60),Inches(1.55),Inches(5.80),Inches(2.50),
    ["Keyword","Spend","Clicks"],
    [[r['Keyword'][:45],f"\u00a3{r['Cost']:,.0f}",f"{r['Clicks']:.0f}"] for _,r in zero_kw_26.iterrows()],
    cw=[3.8,1.0,1.0])

txt(sl,Inches(6.80),Inches(1.20),Inches(5.80),Inches(0.30),"Angle 2: CPA Above \u00a3150",14,NAVY,True)
high_kw=kw26.groupby('Keyword').agg({'Cost':'sum','Conversions':'sum'}).reset_index()
high_kw=high_kw[high_kw['Conversions']>0].copy()
high_kw['CPA']=high_kw['Cost']/high_kw['Conversions']
high_kw=high_kw[high_kw['CPA']>150].nlargest(8,'Cost')
add_table(sl,Inches(6.80),Inches(1.55),Inches(5.80),Inches(2.50),
    ["Keyword","Spend","Conv","CPA"],
    [[r['Keyword'][:40],f"\u00a3{r['Cost']:,.0f}",f"{r['Conversions']:.0f}",f"\u00a3{r['CPA']:.0f}"] for _,r in high_kw.iterrows()],
    cw=[3.0,1.0,0.7,1.1])

ibox(sl,Inches(0.60),Inches(4.30),Inches(5.80),Inches(2.00),
    "Zero-Conv Keywords: Both Years",
    f"2025: {kwd25['zero_count']} keywords, \u00a3{kwd25['zero_spend']:,.0f} wasted\n"
    f"2026: {kwd26['zero_count']} keywords, \u00a3{kwd26['zero_spend']:,.0f} wasted\n\n"
    "These should be paused and budget redirected to converting keywords.",
    RED,RED_TINT)

ibox(sl,Inches(6.80),Inches(4.30),Inches(5.80),Inches(2.00),
    "High-CPA Keywords: Both Years",
    f"2025: {kwd25['high_count']} keywords at \u00a3150+ CPA, \u00a3{kwd25['high_spend']:,.0f} spent\n"
    f"2026: {kwd26['high_count']} keywords at \u00a3150+ CPA, \u00a3{kwd26['high_spend']:,.0f} spent\n\n"
    "The problem is growing \u2014 3.6x more spend on high-CPA keywords in 2026.",
    YELLOW)
footer(sl,7)


# ── S8: DEVICE PERFORMANCE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Device Performance","No device bid adjustments despite significant and widening CPA gaps")
badge(sl,"2025 vs 2026")

img(sl,f"{CHARTS}/device_yoy.png",Inches(0.60),Inches(1.30),Inches(7.00),Inches(3.50))

add_table(sl,Inches(7.90),Inches(1.30),Inches(4.80),Inches(1.80),
    ["Device","2025 CPA","2026 CPA","Change","Bid Adj."],
    [["Mobile",f"\u00a3{dev25['Mobile Phones']['cpa']:.0f}",f"\u00a3{dev26['Mobile Phones']['cpa']:.0f}",f"+{(dev26['Mobile Phones']['cpa']-dev25['Mobile Phones']['cpa'])/dev25['Mobile Phones']['cpa']*100:.0f}%","None"],
     ["Desktop",f"\u00a3{dev25['Computers']['cpa']:.0f}",f"\u00a3{dev26['Computers']['cpa']:.0f}",f"+{(dev26['Computers']['cpa']-dev25['Computers']['cpa'])/dev25['Computers']['cpa']*100:.0f}%","None"],
     ["Tablet",f"\u00a3{dev25['Tablets']['cpa']:.0f}",f"\u00a3{dev26['Tablets']['cpa']:.0f}",f"+{(dev26['Tablets']['cpa']-dev25['Tablets']['cpa'])/dev25['Tablets']['cpa']*100:.0f}%","None"]],
    cw=[0.9,0.8,0.8,0.8,0.7])

ibox(sl,Inches(7.90),Inches(3.30),Inches(4.80),Inches(1.20),
    "Desktop Premium Widening",
    f"Desktop CPA is {(dev26['Computers']['cpa']-dev26['Mobile Phones']['cpa'])/dev26['Mobile Phones']['cpa']*100:.0f}% higher than mobile in 2026 (was {(dev25['Computers']['cpa']-dev25['Mobile Phones']['cpa'])/dev25['Mobile Phones']['cpa']*100:.0f}% in 2025). The gap is growing.",
    RED,RED_TINT)

rbox(sl,Inches(0.60),Inches(5.20),Inches(12.10),Inches(0.60),
    "Recommendation: Implement device bid adjustments. Reduce desktop bids by 20-25% and tablet bids by 30-40% to align spend with conversion efficiency. No adjustments exist anywhere in the account.")
footer(sl,8)


# ── S9: DAY OF WEEK ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
sun25=days25.get('Sunday',{}).get('cpa',0)
sun26=days26.get('Sunday',{}).get('cpa',0)
fri26=days26.get('Friday',{}).get('cpa',0)
stitle(sl,"Day of Week Performance",f"Sunday CPA: \u00a3{sun25:.0f} (2025) \u2192 \u00a3{sun26:.0f} (2026) \u2014 emerging as the clear waste day")
badge(sl,"2025 vs 2026")

img(sl,f"{CHARTS}/dow_yoy.png",Inches(0.60),Inches(1.30),Inches(7.50),Inches(3.50))

ibox(sl,Inches(8.40),Inches(1.30),Inches(4.30),Inches(1.30),
    "2025: Tight Spread",
    f"Range: \u00a3{min(d['cpa'] for d in days25.values()):.0f}\u2013\u00a3{max(d['cpa'] for d in days25.values()):.0f}\n"
    "Day-of-week barely mattered. \u00a35 spread across all days.",
    GREEN,GREEN_TINT)

ibox(sl,Inches(8.40),Inches(2.80),Inches(4.30),Inches(1.30),
    "2026: Sunday Emerging",
    f"Range: \u00a3{min(d['cpa'] for d in days26.values()):.0f}\u2013\u00a3{max(d['cpa'] for d in days26.values()):.0f}\n"
    f"Sunday at \u00a3{sun26:.0f} is now {(sun26-fri26)/fri26*100:.0f}% above Friday. CVR {days26.get('Sunday',{}).get('cvr',0):.1f}% (lowest).",
    RED,RED_TINT)

rbox(sl,Inches(0.60),Inches(5.10),Inches(12.10),Inches(0.60),
    f"Recommendation: Reduce Sunday bids by 15-20%. Sunday spent \u00a3{days26.get('Sunday',{}).get('spend',0):,.0f} in Q1 2026 at the worst CPA and CVR. No day-of-week bid adjustments exist.")
footer(sl,9)


# ── S10: HOUR OF DAY ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Hour of Day Performance","7am worst in both years: \u00a356 (2025) \u2192 \u00a374 (2026) \u2014 +32% increase")
badge(sl,"2025 vs 2026")

img(sl,f"{CHARTS}/hod_yoy.png",Inches(0.60),Inches(1.30),Inches(12.10),Inches(3.50))

ibox(sl,Inches(0.60),Inches(5.00),Inches(3.80),Inches(1.40),
    "2025 Range",
    "Best: \u00a339 CPA (5am, 4pm, 1am)\nWorst: \u00a356 CPA (7am)\nSpread: \u00a317",
    GREEN,GREEN_TINT)

ibox(sl,Inches(4.70),Inches(5.00),Inches(3.80),Inches(1.40),
    "2026 Range",
    "Best: \u00a342 CPA (3am)\nWorst: \u00a374 CPA (7am)\nSpread: \u00a332 \u2014 nearly doubled",
    RED,RED_TINT)

ibox(sl,Inches(8.80),Inches(5.00),Inches(3.90),Inches(1.40),
    "Recommendation",
    "Implement ad schedule bid adjustments. Reduce 7-8am by 15-20%, boost 4-5pm by 10%. No scheduling exists in the account.",
    BLUE,BLUE_TINT)
footer(sl,10)


# ── S11: GEO PERFORMANCE (using matched locations postcode data) ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))

# Aggregate matched locations by postcode
mlg25=ml25.groupby('Matched location').agg({'Cost':'sum','Clicks':'sum','Conversions':'sum'}).sort_values('Cost',ascending=False)
mlg26=ml26.groupby('Matched location').agg({'Cost':'sum','Clicks':'sum','Conversions':'sum'}).sort_values('Cost',ascending=False)
ml_zero25=mlg25[(mlg25['Conversions']==0)&(mlg25['Cost']>0)]
ml_zero26=mlg26[(mlg26['Conversions']==0)&(mlg26['Cost']>0)]

stitle(sl,"Geographic Performance by Postcode",
    f"Zero-conv postcodes: \u00a3{ml_zero25['Cost'].sum():,.0f} waste (2025) + \u00a3{ml_zero26['Cost'].sum():,.0f} waste (2026)")
badge(sl,"2025 vs 2026")

# Top spending postcodes 2025 — left table
pc25_top=mlg25[mlg25['Conversions']>=3].copy()
pc25_top['CPA']=pc25_top['Cost']/pc25_top['Conversions']
worst25=pc25_top.nlargest(6,'CPA')
best25=pc25_top.nsmallest(6,'CPA')

txt(sl,Inches(0.60),Inches(1.20),Inches(6),Inches(0.30),"2025 \u2014 Worst CPA Postcodes (3+ conv)",12,NAVY,True)
add_table(sl,Inches(0.60),Inches(1.50),Inches(5.80),Inches(2.10),
    ["Postcode","Spend","Conv","CPA"],
    [[loc[:25],f"\u00a3{r['Cost']:,.0f}",f"{r['Conversions']:.0f}",f"\u00a3{r['CPA']:.0f}"] for loc,r in worst25.iterrows()],
    cw=[2.2,1.2,0.7,0.8])

# Top spending postcodes 2026 — right table
pc26_top=mlg26[mlg26['Conversions']>=3].copy()
pc26_top['CPA']=pc26_top['Cost']/pc26_top['Conversions']
worst26=pc26_top.nlargest(6,'CPA')

txt(sl,Inches(6.80),Inches(1.20),Inches(6),Inches(0.30),"2026 \u2014 Worst CPA Postcodes (3+ conv)",12,NAVY,True)
add_table(sl,Inches(6.80),Inches(1.50),Inches(5.80),Inches(2.10),
    ["Postcode","Spend","Conv","CPA"],
    [[loc[:25],f"\u00a3{r['Cost']:,.0f}",f"{r['Conversions']:.0f}",f"\u00a3{r['CPA']:.0f}"] for loc,r in worst26.iterrows()],
    cw=[2.2,1.2,0.7,0.8])

# Best CPAs 2026
best26=pc26_top.nsmallest(6,'CPA')
txt(sl,Inches(0.60),Inches(3.80),Inches(6),Inches(0.30),"2026 \u2014 Best CPA Postcodes (3+ conv)",12,GREEN,True)
add_table(sl,Inches(0.60),Inches(4.10),Inches(5.80),Inches(2.10),
    ["Postcode","Spend","Conv","CPA"],
    [[loc[:25],f"\u00a3{r['Cost']:,.0f}",f"{r['Conversions']:.0f}",f"\u00a3{r['CPA']:.0f}"] for loc,r in best26.iterrows()],
    cw=[2.2,1.2,0.7,0.8])

ibox(sl,Inches(6.80),Inches(3.80),Inches(5.80),Inches(2.50),
    "Recommendation: Postcode Targeting",
    f"The data shows CPA ranging from \u00a3{best26['CPA'].min():.0f} to \u00a3{worst26['CPA'].max():.0f} across London postcodes. "
    "Replace radius targeting with individual postcodes to unlock bid adjustments per area. "
    "Bid higher for postcodes near the clinic (W6, W12, W3) and lower for outer areas with high CPAs. "
    "All campaigns should share the same postcode targets \u2014 one unified geographic view across the account.",
    GREEN,GREEN_TINT)
footer(sl,11)


# ── S12: NEGATIVE KEYWORDS & AUTO-APPLY ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Negative Keyword Gaps","784 account-level negatives use broad match \u2014 a risk to legitimate converting traffic")
badge(sl,"Protection Gaps")

mcard(sl,Inches(0.60),Inches(1.20),"2,192","Total Negative Keywords",BLUE,Inches(3.00))
mcard(sl,Inches(3.90),Inches(1.20),"784","Broad Match (Risky)",RED,Inches(3.00))
mcard(sl,Inches(7.20),Inches(1.20),"770","Phrase Match",GREEN,Inches(3.00))
mcard(sl,Inches(10.50),Inches(1.20),"425","Exact Match",YELLOW,Inches(2.20))

ibox(sl,Inches(0.60),Inches(2.60),Inches(5.80),Inches(2.00),
    "The Broad Match Problem",
    "784 negative keywords at account level use broad match. "
    "Broad match negatives can block legitimate searches. "
    "For example, a broad negative for 'bridge' would block 'dental bridge implants' \u2014 "
    "a high-intent search. These need a full audit to ensure they're not blocking revenue.",
    RED,RED_TINT)

ibox(sl,Inches(6.80),Inches(2.60),Inches(5.80),Inches(2.00),
    "Recommended Framework",
    "Replace broad match negatives with phrase or exact match. "
    "Implement a structured 9-list system organised by word count and match type. "
    "Weekly search term reviews to keep the lists current. "
    "This gives precise control without accidentally blocking converting traffic.",
    GREEN,GREEN_TINT)

ibox(sl,Inches(0.60),Inches(4.85),Inches(12.10),Inches(1.50),
    "Auto-Apply: 5 Recommendations Enabled",
    "Google is automatically applying 5 recommendation types, including 'Remove conflicting negative keywords' (~20 times historically). "
    "This means Google has been removing negative keywords it considers conflicting \u2014 potentially re-opening waste that was deliberately blocked. "
    "Last auto-applied: 19 Feb 2026. Additionally, 'Remove redundant keywords' has been applied ~170 times \u2014 pausing keywords without human review. "
    "Recommendation: Disable all auto-apply except 'Use optimised ad rotation' (which is standard best practice).",
    YELLOW)
footer(sl,12)


# ── S13: TOTAL WASTE SUMMARY ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Total Waste Summary")
badge(sl,"Two Angles")

# Angle 1
rect(sl,Inches(0.60),Inches(1.10),Inches(5.80),Inches(0.80),WHITE)
rect(sl,Inches(0.60),Inches(1.10),Inches(0.08),Inches(0.80),RED)
txt(sl,Inches(0.90),Inches(1.15),Inches(5.30),Inches(0.35),"Angle 1: Money That Generated Nothing",14,NAVY,True)
txt(sl,Inches(0.90),Inches(1.50),Inches(5.30),Inches(0.30),f"\u00a3{ZS26:,.0f} in Q1 2026 (search terms with zero conversions)",11,RED,True)

waste_a1=[
    ("Zero-conv search terms (2026)",f"\u00a3{ZS26:,.0f}",f"{ZP26:.0f}% of ST spend",RED),
    ("Zero-conv keywords (2026)",f"\u00a3{kwd26['zero_spend']:,.0f}",f"{kwd26['zero_count']} keywords",YELLOW),
    ("Zero-conv search terms (2025)",f"\u00a3{ZS25:,.0f}",f"{ZP25:.0f}% of ST spend",RED),
    ("Zero-conv keywords (2025)",f"\u00a3{kwd25['zero_spend']:,.0f}",f"{kwd25['zero_count']} keywords",YELLOW),
]
for i,(label,value,detail,color) in enumerate(waste_a1):
    y=Inches(2.05+i*0.60)
    rect(sl,Inches(0.60),y,Inches(5.80),Inches(0.50),WHITE,BORDER,Pt(0.5))
    rect(sl,Inches(0.60),y,Inches(0.06),Inches(0.50),color)
    multirun(sl,Inches(0.85),y+Inches(0.08),Inches(5.30),Inches(0.35),[
        (label,10,BLACK,False),("  ",10,BLACK,False),(value,11,color,True),("  ",10,BLACK,False),(detail,10,GREY_MID,False)])

# Angle 2
rect(sl,Inches(6.80),Inches(1.10),Inches(5.80),Inches(0.80),WHITE)
rect(sl,Inches(6.80),Inches(1.10),Inches(0.08),Inches(0.80),YELLOW)
txt(sl,Inches(7.10),Inches(1.15),Inches(5.30),Inches(0.35),"Angle 2: Money That Could Work Harder",14,NAVY,True)
txt(sl,Inches(7.10),Inches(1.50),Inches(5.30),Inches(0.30),f"High-CPA keywords: \u00a3{kwd25['high_spend']:,.0f} (2025) + \u00a3{kwd26['high_spend']:,.0f} (2026) at 2-5x account average",11,YELLOW,True)

waste_a2=[
    (f"Desktop: \u00a3{dev25['Computers']['cpa']:.0f}\u2192\u00a3{dev26['Computers']['cpa']:.0f} vs mobile \u00a3{dev25['Mobile Phones']['cpa']:.0f}\u2192\u00a3{dev26['Mobile Phones']['cpa']:.0f}","Gap widening","No bid adjustments",RED),
    (f"Sunday: \u00a3{days25.get('Sunday',{}).get('cpa',0):.0f}\u2192\u00a3{sun26:.0f} CPA (worst day both years)","Emerging problem","No schedule adjustments",YELLOW),
    (f"High-CPA keywords: {kwd25['high_count']}\u2192{kwd26['high_count']} at \u00a3150+",f"\u00a3{kwd25['high_spend']+kwd26['high_spend']:,.0f} total","Bid reduction needed",RED),
    (f"Geo: worst postcodes \u00a3120-169 CPA vs best \u00a314-29","5-10x range","Switch to postcodes",YELLOW),
]
for i,(label,value,detail,color) in enumerate(waste_a2):
    y=Inches(2.05+i*0.60)
    rect(sl,Inches(6.80),y,Inches(5.80),Inches(0.50),WHITE,BORDER,Pt(0.5))
    rect(sl,Inches(6.80),y,Inches(0.06),Inches(0.50),color)
    multirun(sl,Inches(7.05),y+Inches(0.08),Inches(5.30),Inches(0.35),[
        (label,10,BLACK,False),("  ",10,BLACK,False),(value,10,color,True),("  ",10,BLACK,False),(detail,10,GREY_MID,False)])

rbox(sl,Inches(0.60),Inches(4.60),Inches(12.10),Inches(0.50),
    f"Combined: \u00a3{ZS25+ZS26+kwd25['zero_spend']+kwd26['zero_spend']:,.0f} in zero-conversion waste across the data period, plus significant efficiency gaps in device, schedule, geo, and keyword bidding.")

# Conversion tracking note
ibox(sl,Inches(0.60),Inches(5.30),Inches(12.10),Inches(1.20),
    "Important Note: Conversion Tracking",
    "These numbers use the conversions reported in Google Ads, which include Dengro offline leads (4,663), form submissions, phone calls, quiz completions, and several soft signals (booking clicks, phone clicks, WhatsApp clicks). "
    "If soft signals are excluded, the true CPA is higher and the efficiency gaps are wider. "
    "A conversion tracking audit is recommended as part of Report 2: Account Structure & Issues.",
    BLUE,BLUE_TINT)
footer(sl,13)


# ── S14: QUICK WINS ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=GREY_BG
top_bar(sl);sidebar(sl);logo(sl)
txt(sl,Inches(0.60),Inches(1.40),Inches(5.50),Inches(0.80),"Quick Wins",44,NAVY,True)
rect(sl,Inches(0.60),Inches(2.20),Inches(2.50),Inches(0.05),BLUE)
txt(sl,Inches(0.60),Inches(2.40),Inches(5.50),Inches(0.40),"Immediate actions to reduce waste and improve efficiency",14,NAVY)

wins=[
    ("1","Total negative keyword overhaul","Replace broad match with structured 9-list system (phrase + exact by word count)",RED),
    ("2","Add negative keywords for top waste categories","Block price research, competitor, wrong treatment, and wrong location searches",RED),
    ("3","Implement device bid adjustments","Reduce desktop bids 20-25%, tablet 30-40%",BLUE),
    ("4","Implement day-of-week bid adjustments","Reduce Sunday bids 15-20%",BLUE),
    ("5","Implement hour-of-day bid adjustments","Reduce 7-8am bids 15-20%, boost 4-5pm",BLUE),
    ("6","Review 17 high-CPA keywords (\u00a3150+)","Reduce bids or improve landing pages for these keywords",YELLOW),
    ("7","Disable auto-apply for negative keyword removal","Stop Google removing your negative keyword protection",GREEN),
    ("8","Tighten geo targeting: 35km radius \u2192 postcodes","Unlock bid control per postcode area",GREEN),
]

for i,(num,action,detail,color) in enumerate(wins):
    y=Inches(2.85+i*0.52)
    rect(sl,Inches(0.60),y,Inches(5.50),Inches(0.45),WHITE)
    rect(sl,Inches(0.60),y,Inches(0.06),Inches(0.45),color)
    multirun(sl,Inches(0.85),y+Inches(0.06),Inches(5.00),Inches(0.35),[
        (f"{num}.  ",11,color,True),(action,11,NAVY,True)])

# Right: what's next
rect(sl,Inches(6.80),Inches(0.50),Inches(5.90),Inches(2.50),WHITE)
rect(sl,Inches(6.80),Inches(0.50),Inches(0.08),Inches(2.50),GREEN)
txt(sl,Inches(7.10),Inches(0.60),Inches(5.40),Inches(0.35),"What's Next",14,NAVY)
txt(sl,Inches(7.10),Inches(1.00),Inches(5.40),Inches(0.80),"Report 2:\nAccount Structure",40,GREEN,True)
txt(sl,Inches(7.10),Inches(1.90),Inches(5.40),Inches(0.60),
    "Conversion tracking audit, bid strategy analysis, Quality Scores, ad copy, and the structural issues driving CPA higher.",11,BLACK)

minicard(sl,Inches(6.80),Inches(3.20),f"\u00a3{NBCPA25:.0f}\u2192\u00a3{NBCPA26:.0f}","CPA Trend (YoY)",RED)
minicard(sl,Inches(8.83),Inches(3.20),"8","Quick Win Actions",BLUE)
minicard(sl,Inches(10.87),Inches(3.20),"Report 2","Coming Next",GREEN)
bot_bar_title(sl)


# ══════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════
output=os.path.join(REPORTS,'01_waste_spend_report_v4.pptx')
prs.save(output)
print(f"\nSaved to: {output}")
print(f"Slides: {len(prs.slides)}")
print("Done!")
