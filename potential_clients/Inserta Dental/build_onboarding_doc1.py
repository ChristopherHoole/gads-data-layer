"""
Dental by Design — Document 1: Business Discovery Questionnaire
Same design system as the audit reports.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
REPORTS = os.path.join(SCRIPT_DIR, 'reports')
LOGO = os.path.join(SCRIPT_DIR, 'act_logo_official.png')
os.makedirs(REPORTS, exist_ok=True)

BLUE=RGBColor(0x42,0x85,0xF4); RED=RGBColor(0xEA,0x43,0x35)
YELLOW=RGBColor(0xFB,0xBC,0x05); GREEN=RGBColor(0x34,0xA8,0x53)
NAVY=RGBColor(0x1A,0x23,0x7E); BLACK=RGBColor(0x1A,0x1A,0x1A)
GREY_BG=RGBColor(0xF5,0xF6,0xFA); WHITE=RGBColor(0xFF,0xFF,0xFF)
BORDER=RGBColor(0xE2,0xE8,0xF0)
BLUE_TINT=RGBColor(0xE8,0xF0,0xFE); RED_TINT=RGBColor(0xFC,0xE8,0xE6)
GREEN_TINT=RGBColor(0xE6,0xF4,0xEA); GREY_MID=RGBColor(0xCB,0xD5,0xE1)
FONT='Calibri'

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)

def rect(sl,l,t,w,h,fc,lc=None,lw=None):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h);s.fill.solid();s.fill.fore_color.rgb=fc
    if lc:s.line.color.rgb=lc;s.line.width=lw or Pt(1)
    else:s.line.fill.background()
def txt(sl,l,t,w,h,c,sz=11,co=BLACK,b=False,a=PP_ALIGN.LEFT):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    r=tf.paragraphs[0].add_run();r.text=c;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
    tf.paragraphs[0].alignment=a;return bx
def multipara(sl,l,t,w,h,pdata):
    bx=sl.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    for i,(tx,sz,co,b,sp) in enumerate(pdata):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run();r.text=tx;r.font.size=Pt(sz);r.font.color.rgb=co;r.font.bold=b;r.font.name=FONT
        if sp:p.space_after=Pt(sp)
def bullets(sl,l,t,w,h,items,sz=11,co=BLACK):
    multipara(sl,l,t,w,h,[(i,sz,co,False,5) for i in items])
def top_bar(sl,h=Inches(0.07)):
    for i,c in enumerate([BLUE,RED,YELLOW,GREEN]):rect(sl,Inches(i*3.333),Inches(0),Inches(3.333),h,c)
def bot_bar(sl,y=Inches(6.92)):
    for p,c in zip([Inches(0.60),Inches(3.63),Inches(6.67),Inches(9.70)],[BLUE,RED,YELLOW,GREEN]):rect(sl,p,y,Inches(3.03),Inches(0.03),c)
def bot_bar_title(sl):
    for i,c in enumerate([BLUE,RED,YELLOW,GREEN]):rect(sl,Inches(i*3.333),Inches(7.0),Inches(3.333),Inches(0.04),c)
def sidebar(sl): rect(sl,Inches(0),Inches(0.07),Inches(0.12),Inches(7.43),BLUE)
def logo(sl,l=Inches(0.60),t=Inches(0.50),s=Inches(0.65)):
    try:sl.shapes.add_picture(LOGO,l,t,s,s)
    except:pass
def footer(sl,n):
    bot_bar(sl)
    try:sl.shapes.add_picture(LOGO,Inches(0.60),Inches(7.0),Inches(0.22),Inches(0.22))
    except:pass
    txt(sl,Inches(0.90),Inches(7.0),Inches(6),Inches(0.25),"Christopher Hoole  |  christopherhoole.com  |  Confidential",11,NAVY,True)
    txt(sl,Inches(12.23),Inches(7.0),Inches(0.50),Inches(0.25),str(n),11,NAVY,a=PP_ALIGN.RIGHT)
def stitle(sl,t,sub=None):
    txt(sl,Inches(0.60),Inches(0.30),Inches(7),Inches(0.50),t,28,NAVY,True)
    if sub:txt(sl,Inches(0.60),Inches(0.85),Inches(9),Inches(0.30),sub,11,GREY_MID)
def badge(sl,t):
    rect(sl,Inches(9.13),Inches(0.30),Inches(3.60),Inches(0.45),WHITE,BLUE,Pt(1))
    txt(sl,Inches(9.23),Inches(0.32),Inches(3.40),Inches(0.40),t,11,BLUE,True,PP_ALIGN.CENTER)
def ibox(sl,l,t,w,h,ti,body,ac,bg=None):
    rect(sl,l,t,w,h,bg or GREY_BG);rect(sl,l,t,Inches(0.06),h,ac)
    txt(sl,l+Inches(0.25),t+Inches(0.10),w-Inches(0.40),Inches(0.30),ti,14,NAVY,True)
    txt(sl,l+Inches(0.25),t+Inches(0.40),w-Inches(0.40),h-Inches(0.50),body,11,BLACK)

def question_section(sl, left, top, width, title, questions, color):
    """Build a question section with title and numbered questions."""
    rect(sl,left,top,width,Inches(0.35),color)
    txt(sl,left+Inches(0.15),top+Inches(0.02),width-Inches(0.30),Inches(0.35),title,13,WHITE,True)
    y = top + Inches(0.45)
    for i, q in enumerate(questions):
        q_height = Inches(0.35) if len(q) < 80 else Inches(0.50)
        txt(sl,left+Inches(0.10),y,width-Inches(0.20),q_height,f"{i+1}.  {q}",11,BLACK)
        y += q_height
    return y


# ══════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════
print("Building Document 1...")

# ── S1: TITLE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=GREY_BG
top_bar(sl);sidebar(sl);logo(sl)
txt(sl,Inches(0.60),Inches(1.40),Inches(5.74),Inches(1.80),"Business Discovery\nQuestionnaire",44,NAVY,True)
rect(sl,Inches(0.60),Inches(3.30),Inches(2.50),Inches(0.05),BLUE)
txt(sl,Inches(0.60),Inches(3.55),Inches(5.50),Inches(0.50),"Dental by Design  |  Prodent Group",22,BLUE)
multipara(sl,Inches(0.60),Inches(4.30),Inches(5.50),Inches(1.00),[
    ("Christopher Hoole",11,BLACK,True,2),("Google Ads Specialist  |  April 2026",11,BLACK,False,2),
    ("christopherhoole.com",11,BLUE,False,0)])
rect(sl,Inches(0.60),Inches(5.50),Inches(4.50),Inches(0.50),WHITE,BLUE,Pt(1))
txt(sl,Inches(0.70),Inches(5.52),Inches(4.30),Inches(0.45),"Onboarding  |  Document 1 of 2",11,BLUE,True)

rect(sl,Inches(6.80),Inches(0.50),Inches(5.90),Inches(2.50),WHITE)
rect(sl,Inches(6.80),Inches(0.50),Inches(0.08),Inches(2.50),GREEN)
txt(sl,Inches(7.10),Inches(0.60),Inches(5.40),Inches(0.35),"Purpose",14,NAVY)
txt(sl,Inches(7.10),Inches(1.00),Inches(5.40),Inches(1.80),
    "This questionnaire helps me understand your business, patients, goals, and competitive landscape in detail.\n\n"
    "I've pre-filled answers where I already have data from the audit reports. "
    "On the onboarding call we'll validate these and fill in the gaps.\n\n"
    "The answers will directly inform campaign structure, ad copy, targeting, and bidding strategy.",
    12,BLACK)

from pptx.util import Inches as I
minicard_data = [("10","Sections to Cover",BLUE),("3","Patient Personas",GREEN),("Pre-filled","From Audit Data",YELLOW)]
for i,(v,l,c) in enumerate(minicard_data):
    left=Inches(6.80+i*2.03)
    rect(sl,left,Inches(3.20),Inches(1.83),Inches(1.15),WHITE)
    rect(sl,left,Inches(3.20),Inches(0.06),Inches(1.15),c)
    txt(sl,left+Inches(0.15),Inches(3.32),Inches(1.58),Inches(0.45),v,22,c,True)
    txt(sl,left+Inches(0.15),Inches(3.80),Inches(1.58),Inches(0.30),l,11,BLACK)

bot_bar_title(sl)


# ── S2: BUSINESS DEEP-DIVE ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Business Deep-Dive","Understanding how your business works, your revenue model, and what makes you different")
badge(sl,"Section 1 of 10")

ibox(sl,Inches(0.60),Inches(1.20),Inches(5.80),Inches(2.80),
    "What We Already Know (From Audit)",
    "\u2022 Integrated model: Clinic (DBD) + Lab (Vivo) + Manufacturer (Inserta)\n"
    "\u2022 Competitive advantage: \u201Cup to 60% less\u201D due to in-house production\n"
    "\u2022 Signature product: Vivo Bridge (permanent teeth in 48 hours)\n"
    "\u2022 4.8\u2605 rating, 399+ Google reviews\n"
    "\u2022 Prices: Single implant from \u00a31,695, Full Arch from \u00a39,995, Vivo Bridge \u00a315,990+\n"
    "\u2022 0% finance available\n"
    "\u2022 CQC registered, Hammersmith location (W6)",
    GREEN,GREEN_TINT)

ibox(sl,Inches(6.80),Inches(1.20),Inches(5.80),Inches(2.80),
    "Questions for the Call",
    "1. What is the average patient lifetime value by treatment type?\n"
    "2. What percentage of consultations convert to booked treatment?\n"
    "3. Average time from first enquiry to treatment completion?\n"
    "4. How does the in-house lab/manufacturing keep costs down day-to-day?\n"
    "5. How many treatment chairs / surgeons / capacity per day?\n"
    "6. What's the current monthly revenue attributed to Google Ads patients?\n"
    "7. Who handles inbound enquiries? How quickly are leads followed up?",
    BLUE,BLUE_TINT)

ibox(sl,Inches(0.60),Inches(4.20),Inches(12.10),Inches(2.20),
    "Answers (To Complete on Call)",
    "\n\n\n\n\n",
    GREY_MID,WHITE)
footer(sl,2)


# ── S3: CLIENT PERSONAS ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Client Personas","Proposed 3 personas based on account data \u2014 validate and refine on the call")
badge(sl,"Section 2 of 10")

# Persona 1
ibox(sl,Inches(0.60),Inches(1.20),Inches(3.80),Inches(3.80),
    "Persona 1: Price-Conscious Researcher",
    "Age: 45-65. Lost teeth or facing extraction.\n\n"
    "Searches: \u201Cdental implants cost\u201D, \u201Ccheapest implants\u201D, \u201Cpayment plans\u201D\n\n"
    "Pain point: Shocked by implant costs. Comparing prices.\n\n"
    "Decision trigger: Quality at a price they can justify + finance.\n\n"
    "Why DBD wins: \u201CUp to 60% less\u201D, 0% finance, in-house lab.\n\n"
    "Data: Financing ad group = \u00a334 CPA, 24.7% CVR (proven).",
    GREEN,GREEN_TINT)

# Persona 2
ibox(sl,Inches(4.70),Inches(1.20),Inches(3.80),Inches(3.80),
    "Persona 2: Quality-Driven Urgent Case",
    "Age: 50-70. Failed dental work or deteriorating health.\n\n"
    "Searches: \u201Call on 4\u201D, \u201Cfull mouth restoration\u201D, \u201Cteeth in a day\u201D\n\n"
    "Pain point: Wants it done properly this time. Willing to pay.\n\n"
    "Decision trigger: Trust (reviews, CQC), speed (48hrs).\n\n"
    "Why DBD wins: 4.8\u2605, Vivo Bridge 48hrs, expert surgeons.\n\n"
    "Data: All-on-4 ad group = \u00a3106-128 CPA (high value).",
    BLUE,BLUE_TINT)

# Persona 3
ibox(sl,Inches(8.80),Inches(1.20),Inches(3.90),Inches(3.80),
    "Persona 3: Cosmetic Improver",
    "Age: 30-50. Functional teeth but wants a better smile.\n\n"
    "Searches: \u201CInvisalign London\u201D, \u201Cveneers near me\u201D, \u201Csmile makeover\u201D\n\n"
    "Pain point: Self-conscious. Wants transformation.\n\n"
    "Decision trigger: Before/after photos, ease of treatment.\n\n"
    "Why DBD wins: Diamond Apex Invisalign, premium location.\n\n"
    "Data: Invisalign = \u00a345 CPA, 10.1% CVR.",
    YELLOW)

ibox(sl,Inches(0.60),Inches(5.20),Inches(12.10),Inches(1.20),
    "Questions for the Call",
    "1. Do these 3 personas match who actually walks through your door?  2. Which persona has the highest lifetime value?  "
    "3. Which has the highest consultation-to-treatment rate?  4. Are there patient types you actively DON'T want?  "
    "5. What's the typical age range of paying implant patients?  6. Do family members often research on behalf of the patient?",
    BLUE,BLUE_TINT)
footer(sl,3)


# ── S4: KEY SERVICES ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Key Services to Promote","Which treatments should we prioritise in Google Ads?")
badge(sl,"Section 3 of 10")

from pptx.util import Inches
# Pre-filled table
ts=sl.shapes.add_table(8,5,Inches(0.60),Inches(1.20),Inches(7.50),Inches(3.20))
tb=ts.table
hdrs=["Service","Current Ad Spend","Current CPA","Priority?","Notes"]
cws=[1.8,1.2,0.9,0.8,2.8]
for i,cw in enumerate(cws): tb.columns[i].width=Inches(cw)
for ci,hd in enumerate(hdrs):
    c=tb.cell(0,ci);c.text=hd;c.fill.solid();c.fill.fore_color.rgb=NAVY
    for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.font.color.rgb=WHITE;p.font.bold=True;p.font.name=FONT
rows=[
    ["Dental Implants (single)","\u00a374K (2026 Q1)","\u00a394 (Search)","Primary?","Highest volume. Needs restructure."],
    ["Full Arch / All-on-4","\u00a316.5K (2026 Q1)","\u00a3106-128","Primary?","Highest patient value. Worth the CPA?"],
    ["Financing / Payment Plans","\u00a32.9K (2026 Q1)","\u00a334","Scale up?","Best CPA. Proven demand."],
    ["Vivo Bridge (48hrs)","\u2014 (bundled)","\u2014","Discuss","Separate campaign or USP within implants?"],
    ["Invisalign","\u00a33.8K (2026 Q1)","\u00a345","Secondary?","Good CPA. Currently paused."],
    ["General Dentistry","\u00a30 (paused)","\u00a3118 (when active)","Discuss","Gateway to implant patients?"],
    ["Emergency Dental","\u00a30","\u2014","Discuss","Worth advertising?"],
]
for ri,rd in enumerate(rows):
    for ci,v in enumerate(rd):
        c=tb.cell(ri+1,ci);c.text=v;c.fill.solid()
        c.fill.fore_color.rgb=GREY_BG if ri%2==0 else WHITE
        for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.font.color.rgb=BLACK;p.font.name=FONT

ibox(sl,Inches(8.40),Inches(1.20),Inches(4.30),Inches(3.20),
    "Questions for the Call",
    "1. Rank these services by profitability \u2014 which puts the most revenue per chair-hour?\n\n"
    "2. Should Vivo Bridge be a separate campaign or a USP within implants?\n\n"
    "3. Appetite to restart Invisalign?\n\n"
    "4. General dentistry / emergency \u2014 worth advertising?\n\n"
    "5. Any NEW services launching?\n\n"
    "6. Which services have the longest waiting list?",
    BLUE,BLUE_TINT)

ibox(sl,Inches(0.60),Inches(4.60),Inches(12.10),Inches(1.80),
    "Answers (To Complete on Call)",
    "\n\n\n\n",
    GREY_MID,WHITE)
footer(sl,4)


# ── S5: BUSINESS GOALS ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Business Goals","What does success look like?")
badge(sl,"Section 4 of 10")

ibox(sl,Inches(0.60),Inches(1.20),Inches(5.80),Inches(3.00),
    "Questions for the Call",
    "1. Primary goal right now \u2014 more leads, better quality leads, or lower cost per lead?\n\n"
    "2. Is there a monthly patient acquisition target? (e.g. \u201C50 new implant patients/month\u201D)\n\n"
    "3. What does \u201Csuccess\u201D look like in 90 days? (Specific number.)\n\n"
    "4. Are there revenue targets tied to marketing?\n\n"
    "5. Is \u00a354K/month the budget ceiling, or would you scale up if CPA hits target?\n\n"
    "6. What's the growth plan for the next 12 months?",
    BLUE,BLUE_TINT)

ibox(sl,Inches(6.80),Inches(1.20),Inches(5.80),Inches(3.00),
    "What We Already Know",
    "Current state (2026 Q1):\n"
    "\u2022 Monthly spend: ~\u00a354,000\n"
    "\u2022 Non-brand CPA: \u00a364\n"
    "\u2022 ~820 non-brand leads/month\n"
    "\u2022 Real booking CPA: \u00a3418 (PMax)\n"
    "\u2022 Lead-to-booking rate: 13.3%\n\n"
    "Our proposed targets (from Report 3):\n"
    "\u2022 Non-brand CPA target: \u00a350\n"
    "\u2022 +236 additional leads/month OR \u00a312K/month saving",
    GREEN,GREEN_TINT)

ibox(sl,Inches(0.60),Inches(4.40),Inches(12.10),Inches(2.00),
    "Answers (To Complete on Call)",
    "\n\n\n\n",
    GREY_MID,WHITE)
footer(sl,5)


# ── S6: SEASONALITY ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Seasonality & Busy/Quiet Periods","Understanding demand patterns to optimise budget pacing")
badge(sl,"Section 5 of 10")

ibox(sl,Inches(0.60),Inches(1.20),Inches(5.80),Inches(2.50),
    "Questions for the Call",
    "1. Which months are busiest for implant enquiries? (Jan/New Year resolutions? Post-summer?)\n\n"
    "2. Is there a seasonal pattern to consultations vs treatments? (e.g. enquire in Jan, book in Mar)\n\n"
    "3. Any quiet periods where we should reduce spend? (August? Christmas?)\n\n"
    "4. Do you run seasonal promotions? (January offers, Black Friday?)\n\n"
    "5. How far in advance do patients typically book?",
    BLUE,BLUE_TINT)

ibox(sl,Inches(6.80),Inches(1.20),Inches(5.80),Inches(2.50),
    "What the Data Shows",
    "From the account data (2025-2026):\n\n"
    "\u2022 Spend ramped significantly from Sep 2025 (\u00a320K/month) to Jan 2026 (\u00a355K/month)\n"
    "\u2022 CPA was lowest in Sep-Oct 2025 (when PMax launched)\n"
    "\u2022 CPA has been rising since Jan 2026 as spend scaled\n"
    "\u2022 No clear seasonal pattern visible yet in the ad data \u2014 but the business will know their demand cycles better than the data shows",
    GREEN,GREEN_TINT)

ibox(sl,Inches(0.60),Inches(3.90),Inches(12.10),Inches(2.50),
    "Answers (To Complete on Call)",
    "\n\n\n\n\n\n",
    GREY_MID,WHITE)
footer(sl,6)


# ── S7: TARGET KPIs ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Target KPIs","Pre-filled from audit reports \u2014 validate and agree on the call")
badge(sl,"Section 6 of 10")

ts=sl.shapes.add_table(7,4,Inches(0.60),Inches(1.20),Inches(7.50),Inches(2.80))
tb=ts.table
hdrs=["KPI","Current (2026 Q1)","Proposed Target","Basis"]
cws=[2.0,1.5,1.5,2.5]
for i,cw in enumerate(cws): tb.columns[i].width=Inches(cw)
for ci,hd in enumerate(hdrs):
    c=tb.cell(0,ci);c.text=hd;c.fill.solid();c.fill.fore_color.rgb=NAVY
    for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.font.color.rgb=WHITE;p.font.bold=True;p.font.name=FONT
kpi_rows=[
    ["Non-brand CPA","\u00a364","\u00a350","Report 3 projection"],
    ["Monthly spend","~\u00a354,000","Maintain","\u2014"],
    ["Monthly non-brand leads","~820","~1,050 (+29%)","Same spend, lower CPA"],
    ["Cost per booked patient","\u00a3418 (PMax)","\u00a3200 target","Search migration"],
    ["Lead-to-booking rate","13.3%","20%+ target","Better lead quality"],
    ["Quality Score (weighted)","6.4","7.5+","Ad rewrites"],
]
for ri,rd in enumerate(kpi_rows):
    for ci,v in enumerate(rd):
        c=tb.cell(ri+1,ci);c.text=v;c.fill.solid()
        c.fill.fore_color.rgb=GREY_BG if ri%2==0 else WHITE
        for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.font.color.rgb=BLACK;p.font.name=FONT

ibox(sl,Inches(8.40),Inches(1.20),Inches(4.30),Inches(2.80),
    "Questions for the Call",
    "1. Do these targets feel right? Too aggressive or too conservative?\n\n"
    "2. Is there a hard CPA ceiling where a lead stops being profitable?\n\n"
    "3. What's an acceptable cost per BOOKED patient?\n\n"
    "4. How do you currently measure Google Ads success? (Dengro? Manual? Gut feel?)",
    BLUE,BLUE_TINT)

ibox(sl,Inches(0.60),Inches(4.20),Inches(12.10),Inches(2.20),
    "Answers (To Complete on Call)",
    "\n\n\n\n\n",
    GREY_MID,WHITE)
footer(sl,7)


# ── S8: COMPETITORS ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Competitors","Who are you competing against for patients?")
badge(sl,"Section 7 of 10")

ibox(sl,Inches(0.60),Inches(1.20),Inches(5.80),Inches(2.50),
    "Competitors Seen in Search Terms",
    "These competitor names appeared in the search term data:\n\n"
    "\u2022 Brighton Implant Clinic (\u00a3139 spent on this term, 0 conversions)\n"
    "\u2022 Harley Street Dental Studio (\u00a347 spent)\n"
    "\u2022 Dental Prime Stratford (\u00a3123 spent)\n"
    "\u2022 Marylebone Implant Centre (\u00a340 spent)\n"
    "\u2022 Evodental (\u00a339 spent)\n\n"
    "All zero conversions \u2014 people searching for competitors and clicking your ads.",
    RED,RED_TINT)

ibox(sl,Inches(6.80),Inches(1.20),Inches(5.80),Inches(2.50),
    "Questions for the Call",
    "1. Who are your top 3 competitors for dental implants in London?\n\n"
    "2. Are there competitors specifically in the Hammersmith area?\n\n"
    "3. Do any competitors offer a \u201Cteeth in a day\u201D type USP? Is Vivo Bridge unique?\n\n"
    "4. Are you losing patients to dental tourism (Turkey etc.)?\n\n"
    "5. Any competitors you'd want to conquest-target (bid on their brand terms)?",
    BLUE,BLUE_TINT)

ibox(sl,Inches(0.60),Inches(3.90),Inches(12.10),Inches(2.50),
    "Answers (To Complete on Call)",
    "\n\n\n\n\n\n",
    GREY_MID,WHITE)
footer(sl,8)


# ── S9: GEOGRAPHIC TARGETING ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Geographic Targeting","Understanding where your patients come from")
badge(sl,"Section 8 of 10")

ibox(sl,Inches(0.60),Inches(1.20),Inches(5.80),Inches(2.50),
    "What We Already Know",
    "Current setup: Radius targeting (30km, 35km around Hammersmith)\n\n"
    "CPA by radius (2026):\n"
    "\u2022 30km ring: \u00a359 CPA\n"
    "\u2022 35km ring: \u00a393 CPA\n\n"
    "Best postcodes: SM4 \u00a314, EN4 \u00a314, WD7 \u00a318\n"
    "Worst postcodes: N13 \u00a3169, TW2 \u00a3137, WD23 \u00a3123\n\n"
    "Recommendation: Switch from radius to postcode targeting for bid control.",
    GREEN,GREEN_TINT)

ibox(sl,Inches(6.80),Inches(1.20),Inches(5.80),Inches(2.50),
    "Questions for the Call",
    "1. How far do patients actually travel to the clinic? Is 35km realistic?\n\n"
    "2. Are there specific areas you know your patients come from?\n\n"
    "3. Can you provide a full download of patient postcodes from Dengro? (We'll create a heatmap.)\n\n"
    "4. Are there areas you explicitly DON'T want to target?\n\n"
    "5. Are there areas with high no-show rates?",
    BLUE,BLUE_TINT)

ibox(sl,Inches(0.60),Inches(3.90),Inches(12.10),Inches(2.50),
    "Answers (To Complete on Call)",
    "\n\n\n\n\n\n",
    GREY_MID,WHITE)
footer(sl,9)


# ── S10: DECISION MAKERS & COMMUNICATION ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Decision Makers & Communication","Who approves what, and how do we stay in touch?")
badge(sl,"Section 9 of 10")

ibox(sl,Inches(0.60),Inches(1.20),Inches(5.80),Inches(2.50),
    "Questions for the Call",
    "1. Who approves ad copy changes? (Can I change ads without sign-off?)\n\n"
    "2. Who approves budget changes? (If I want to shift budget between campaigns.)\n\n"
    "3. Who's my day-to-day contact? (Giulio for technical, Tommaso for strategic?)\n\n"
    "4. Preferred communication channel? (Email? WhatsApp? Slack?)\n\n"
    "5. How quickly do you expect responses to questions?",
    BLUE,BLUE_TINT)

ibox(sl,Inches(6.80),Inches(1.20),Inches(5.80),Inches(2.50),
    "What We Know So Far",
    "\u2022 Tommaso Grassi: Commercial Director \u2014 strategic decisions\n"
    "\u2022 Giulio Zanchetta: Marketing / Technical \u2014 manages the Google Ads account (89% of change history)\n"
    "\u2022 Giulio's email: giulio.zanchetta@majestico.it \u2014 is Majestico a separate entity or part of Prodent Group?\n"
    "\u2022 Ellie Carvell: Recruitment Specialist (Inserta Dental)\n\n"
    "Questions to clarify:\n"
    "\u2022 What is Giulio's role exactly? Full-time DBD or shared with Majestico?\n"
    "\u2022 Who else needs to be kept in the loop on reporting?",
    GREEN,GREEN_TINT)

ibox(sl,Inches(0.60),Inches(3.90),Inches(12.10),Inches(2.50),
    "Answers (To Complete on Call)",
    "\n\n\n\n\n\n",
    GREY_MID,WHITE)
footer(sl,10)


# ── S11: REPORTING ──
sl=prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid();sl.background.fill.fore_color.rgb=WHITE
top_bar(sl,Inches(0.06))
stitle(sl,"Reporting Expectations","Proposed reporting cadence \u2014 validate on the call")
badge(sl,"Section 10 of 10")

ibox(sl,Inches(0.60),Inches(1.20),Inches(5.80),Inches(3.00),
    "Proposed: Weekly Report (Every Friday)",
    "\u2022 What I worked on each morning and afternoon throughout the week\n\n"
    "\u2022 Key metrics: spend, leads, CPA, conversion rate\n\n"
    "\u2022 Last 8 weeks week-on-week comparison table\n\n"
    "\u2022 Changes made and rationale\n\n"
    "\u2022 Issues flagged and actions planned for next week",
    GREEN,GREEN_TINT)

ibox(sl,Inches(6.80),Inches(1.20),Inches(5.80),Inches(3.00),
    "Proposed: Monthly Report (First Week of Month)",
    "\u2022 Full month performance vs KPI targets\n\n"
    "\u2022 Month-on-month and year-on-year comparisons\n\n"
    "\u2022 Campaign-level and treatment-level breakdown\n\n"
    "\u2022 Search term audit findings\n\n"
    "\u2022 Recommendations and strategic plan for next month\n\n"
    "\u2022 Budget utilisation and pacing review",
    GREEN,GREEN_TINT)

ibox(sl,Inches(0.60),Inches(4.40),Inches(12.10),Inches(2.00),
    "Questions for the Call",
    "1. Does this weekly + monthly cadence work?  "
    "2. Who should receive the reports? (Tommaso, Giulio, anyone else?)  "
    "3. Do you want a call to walk through the monthly report?  "
    "4. Is there a specific dashboard or format you currently use? (Looker Studio? Dengro?)  "
    "5. Are there specific metrics Tommaso watches most closely?",
    BLUE,BLUE_TINT)
footer(sl,11)


# ══════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════
output=os.path.join(REPORTS,'Onboarding Doc 1 - Business Discovery Questionnaire.pptx')
prs.save(output)
print(f"Saved to: {output}")
print(f"Slides: {len(prs.slides)}")
print("Done!")
