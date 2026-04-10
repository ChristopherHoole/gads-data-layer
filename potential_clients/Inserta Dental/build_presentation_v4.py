"""
Christopher Hoole — Dental by Design Interview Presentation v4
Restructured narrative: lead with THEM, not me.
New slide 2 "Their Opportunity" — show homework before credentials.
13 min presenting, 17 min conversation.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Design system (unchanged) ──
BLUE    = RGBColor(0x42, 0x85, 0xF4)
RED     = RGBColor(0xEA, 0x43, 0x35)
YELLOW  = RGBColor(0xFB, 0xBC, 0x05)
GREEN   = RGBColor(0x34, 0xA8, 0x53)
NAVY    = RGBColor(0x1A, 0x23, 0x7E)
BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
GREY_BG = RGBColor(0xF5, 0xF6, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BORDER  = RGBColor(0xE2, 0xE8, 0xF0)
BLUE_TINT = RGBColor(0xE8, 0xF0, 0xFE)
RED_TINT  = RGBColor(0xFC, 0xE8, 0xE6)
GREEN_TINT = RGBColor(0xE6, 0xF4, 0xEA)
GREY_MID  = RGBColor(0xCB, 0xD5, 0xE1)

FONT = 'Calibri'
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.normpath(os.path.join(SCRIPT_DIR, '..', '..'))
ACT_LOGO = os.path.join(PROJECT_ROOT, 'potential_clients', 'objection_experts', 'charts', 'act_logo.png')
ACT_SCREENSHOT = os.path.join(PROJECT_ROOT, 'act_dashboard', 'prototypes', 'screenshots', 'v6-light-full.png')

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers (compact) ──

def rect(sl, l, t, w, h, fc, lc=None, lw=None):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if lc: s.line.color.rgb = lc; s.line.width = lw or Pt(1)
    else: s.line.fill.background()
    return s

def txt(sl, l, t, w, h, c, sz=11, co=BLACK, b=False, a=PP_ALIGN.LEFT):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = c; r.font.size = Pt(sz)
    r.font.color.rgb = co; r.font.bold = b; r.font.name = FONT
    tf.paragraphs[0].alignment = a
    return bx

def multirun(sl, l, t, w, h, runs, a=PP_ALIGN.LEFT):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = a
    for tx, sz, co, b in runs:
        r = p.add_run(); r.text = tx; r.font.size = Pt(sz)
        r.font.color.rgb = co; r.font.bold = b; r.font.name = FONT
    return bx

def multipara(sl, l, t, w, h, pdata):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, (tx, sz, co, b, sp) in enumerate(pdata):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = tx; r.font.size = Pt(sz)
        r.font.color.rgb = co; r.font.bold = b; r.font.name = FONT
        if sp: p.space_after = Pt(sp)
    return bx

def bullets(sl, l, t, w, h, items, sz=11, co=BLACK):
    return multipara(sl, l, t, w, h, [(i, sz, co, False, 4) for i in items])

def top_bar(sl, h=Inches(0.07)):
    for i, c in enumerate([BLUE, RED, YELLOW, GREEN]):
        rect(sl, Inches(i*3.333), Inches(0), Inches(3.333), h, c)

def bot_bar(sl, y=Inches(6.92)):
    for p, c in zip([Inches(0.60),Inches(3.63),Inches(6.67),Inches(9.70)], [BLUE,RED,YELLOW,GREEN]):
        rect(sl, p, y, Inches(3.03), Inches(0.03), c)

def bot_bar_title(sl):
    for i, c in enumerate([BLUE, RED, YELLOW, GREEN]):
        rect(sl, Inches(i*3.333), Inches(7.0), Inches(3.333), Inches(0.04), c)

def sidebar(sl):
    rect(sl, Inches(0), Inches(0.07), Inches(0.12), Inches(7.43), BLUE)

def logo(sl, l=Inches(0.60), t=Inches(0.50), s=Inches(0.65)):
    try: sl.shapes.add_picture(ACT_LOGO, l, t, s, s)
    except: pass

def footer(sl, n):
    bot_bar(sl)
    try: sl.shapes.add_picture(ACT_LOGO, Inches(0.60), Inches(7.0), Inches(0.22), Inches(0.22))
    except: pass
    txt(sl, Inches(0.90), Inches(7.0), Inches(6), Inches(0.25),
        "Christopher Hoole  |  christopherhoole.com  |  Confidential", 11, NAVY, True)
    txt(sl, Inches(12.23), Inches(7.0), Inches(0.50), Inches(0.25), str(n), 11, NAVY, a=PP_ALIGN.RIGHT)

def stitle(sl, t, sub=None):
    txt(sl, Inches(0.60), Inches(0.30), Inches(7), Inches(0.50), t, 28, NAVY, True)
    if sub: txt(sl, Inches(0.60), Inches(0.85), Inches(9), Inches(0.30), sub, 11, RED)

def badge(sl, t):
    rect(sl, Inches(9.13), Inches(0.30), Inches(3.60), Inches(0.45), WHITE, BLUE, Pt(1))
    txt(sl, Inches(9.23), Inches(0.32), Inches(3.40), Inches(0.40), t, 11, BLUE, True, PP_ALIGN.CENTER)

def mcard(sl, l, t, v, lab, ac, w=Inches(2.80), h=Inches(1.05)):
    rect(sl, l, t, w, h, WHITE, BORDER, Pt(0.75))
    rect(sl, l, t, Inches(0.06), h, ac)
    txt(sl, l+Inches(0.20), t+Inches(0.10), w-Inches(0.30), Inches(0.50), v, 22, ac, True)
    txt(sl, l+Inches(0.20), t+Inches(0.60), w-Inches(0.30), Inches(0.35), lab, 11, BLACK)

def minicard(sl, l, t, v, lab, ac):
    rect(sl, l, t, Inches(1.83), Inches(1.15), WHITE)
    rect(sl, l, t, Inches(0.06), Inches(1.15), ac)
    txt(sl, l+Inches(0.15), t+Inches(0.12), Inches(1.58), Inches(0.45), v, 22, ac, True)
    txt(sl, l+Inches(0.15), t+Inches(0.60), Inches(1.58), Inches(0.30), lab, 11, BLACK)

def ibox(sl, l, t, w, h, ti, body, ac, bg=None):
    rect(sl, l, t, w, h, bg or GREY_BG)
    rect(sl, l, t, Inches(0.06), h, ac)
    txt(sl, l+Inches(0.25), t+Inches(0.10), w-Inches(0.40), Inches(0.30), ti, 14, NAVY, True)
    txt(sl, l+Inches(0.25), t+Inches(0.40), w-Inches(0.40), h-Inches(0.50), body, 11, BLACK)

def rbox(sl, l, t, w, h, c):
    rect(sl, l, t, w, h, BLUE_TINT, BLUE, Pt(1))
    txt(sl, l+Inches(0.25), t+Inches(0.05), w-Inches(0.50), h-Inches(0.10), c, 11, NAVY, True)

def add_table(sl, l, t, w, h, hdrs, rows):
    ts = sl.shapes.add_table(len(rows)+1, len(hdrs), l, t, w, h)
    tb = ts.table; cw = int(w/len(hdrs))
    for i in range(len(hdrs)): tb.columns[i].width = cw
    for ci, hd in enumerate(hdrs):
        c = tb.cell(0, ci); c.text = hd; c.fill.solid(); c.fill.fore_color.rgb = NAVY
        for p in c.text_frame.paragraphs: p.font.size=Pt(11); p.font.color.rgb=WHITE; p.font.bold=True; p.font.name=FONT
    for ri, rd in enumerate(rows):
        for ci, v in enumerate(rd):
            c = tb.cell(ri+1, ci); c.text = str(v); c.fill.solid()
            c.fill.fore_color.rgb = GREY_BG if ri%2==0 else WHITE
            for p in c.text_frame.paragraphs: p.font.size=Pt(11); p.font.color.rgb=BLACK; p.font.name=FONT
    return ts

def num_circle(sl, l, t, num, col):
    ns = sl.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Inches(0.45), Inches(0.45))
    ns.fill.solid(); ns.fill.fore_color.rgb = col; ns.line.fill.background()
    tf = ns.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run(); r.text = str(num); r.font.size = Pt(16)
    r.font.color.rgb = WHITE; r.font.bold = True; r.font.name = FONT


# ══════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE — set the tone, about THEM
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = GREY_BG

top_bar(sl)
sidebar(sl)
logo(sl)

txt(sl, Inches(0.60), Inches(1.40), Inches(5.74), Inches(1.80),
    "DRIVING PATIENT\nACQUISITION", 44, NAVY, True)
rect(sl, Inches(0.60), Inches(3.30), Inches(2.50), Inches(0.05), BLUE)
txt(sl, Inches(0.60), Inches(3.55), Inches(5.50), Inches(0.50),
    "Dental by Design  |  Prodent Group", 22, BLUE)

multipara(sl, Inches(0.60), Inches(4.30), Inches(5.50), Inches(1.00), [
    ("Christopher Hoole", 11, BLACK, True, 2),
    ("Google Ads Specialist  |  April 2026", 11, BLACK, False, 2),
    ("christopherhoole.com", 11, BLUE, False, 0),
])

rect(sl, Inches(0.60), Inches(5.50), Inches(4.50), Inches(0.50), WHITE, BLUE, Pt(1))
txt(sl, Inches(0.70), Inches(5.52), Inches(4.30), Inches(0.45),
    "Performance Marketing Specialist  |  Interview Presentation", 11, BLUE, True)

# Right hero
rect(sl, Inches(6.80), Inches(0.50), Inches(5.90), Inches(2.50), WHITE)
rect(sl, Inches(6.80), Inches(0.50), Inches(0.08), Inches(2.50), GREEN)
txt(sl, Inches(7.10), Inches(0.70), Inches(5.40), Inches(0.35), "Specialist Focus", 14, NAVY)
txt(sl, Inches(7.10), Inches(1.10), Inches(5.40), Inches(0.80),
    "Dental & Aesthetics PPC", 40, GREEN, True)
txt(sl, Inches(7.10), Inches(1.90), Inches(5.40), Inches(0.30),
    "Direct experience managing campaigns for dental clinics and cosmetic practices", 11, BLACK)

minicard(sl, Inches(6.80), Inches(3.20), "16 Years", "Hands-On Google Ads", BLUE)
minicard(sl, Inches(8.83), Inches(3.20), "100+", "Accounts Managed", YELLOW)
minicard(sl, Inches(10.87), Inches(3.20), "\u00a350M+", "Total Spend", GREEN)

bot_bar_title(sl)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2: THEIR OPPORTUNITY — NEW. Lead with them, not you.
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE

top_bar(sl, Inches(0.06))
stitle(sl, "Your Competitive Advantage",
       "You have assets most competitors don't \u2014 let me turn them into patients.")
badge(sl, "About Dental by Design")

# 3 big advantage cards — full width, generous spacing
advantages = [
    ("Integrated Model", "Clinic + Lab + Manufacturer",
     "You control the entire chain \u2014 quality AND pricing. No other London dental clinic can say this. "
     "In PPC, this means ad copy that competitors simply cannot match.",
     GREEN),
    ("\u201CUp to 60% Less\u201D", "Transparent Pricing Advantage",
     "This is the most powerful ad copy hook in your market. Price-sensitive patients searching for implants "
     "will click \u2014 and your quality keeps them. This messaging should be in every headline.",
     BLUE),
    ("Vivo Bridge \u2014 48 Hours", "Permanent Teeth, Fast",
     "\u201CPermanent teeth in 48 hours\u201D is an incredible USP. Speed + quality + price = unbeatable value proposition. "
     "This is your premium product and it deserves its own campaign.",
     RED),
]

for i, (title_t, subtitle_t, body_t, color) in enumerate(advantages):
    y = Inches(1.30 + i * 1.65)
    rect(sl, Inches(0.60), y, Inches(12.10), Inches(1.45), WHITE, BORDER, Pt(0.75))
    rect(sl, Inches(0.60), y, Inches(0.08), Inches(1.45), color)

    txt(sl, Inches(0.90), y + Inches(0.12), Inches(4.00), Inches(0.40), title_t, 22, color, True)
    txt(sl, Inches(0.90), y + Inches(0.55), Inches(4.00), Inches(0.25), subtitle_t, 11, GREY_MID, True)
    txt(sl, Inches(5.20), y + Inches(0.15), Inches(7.20), Inches(1.15), body_t, 12, BLACK)

# Supporting stats
mcard(sl, Inches(0.60), Inches(6.00), "4.8\u2605", "399+ Google Reviews", YELLOW, Inches(3.80))
mcard(sl, Inches(4.70), Inches(6.00), "\u00a31,695+", "Single Implant (from)", GREEN, Inches(3.80))
mcard(sl, Inches(8.80), Inches(6.00), "UK Network", "Flagship + Partner Clinics", BLUE, Inches(3.80))

footer(sl, 2)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3: THE PROBLEM I SOLVE
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE

top_bar(sl, Inches(0.06))
stitle(sl, "The Problem I Solve",
       "Every dental clinic I've worked with has the same three problems.")
badge(sl, "Why You're Hiring")

problems = [
    ("Wasted Spend on Irrelevant Clicks",
     "Without proper negative keyword management, dental implant campaigns attract searches for "
     "DIY guides, dental jobs, dental tourism, and NHS alternatives. "
     "In my last audit, 42.7% of budget was going to non-converting traffic.",
     RED, "42.7%", "Average Waste I Find"),
    ("No Target CPA Guardrails",
     "When bid strategies run without a target CPA, the algorithm has permission to spend anything per lead. "
     "CPAs spike from \u00a330 one week to \u00a365 the next. "
     "Setting the right guardrails gives you predictable, controllable costs.",
     YELLOW, "\u00a335\u2192\u00a365", "Typical CPA Volatility"),
    ("Lead Quality Issues",
     "Not all form fills are real patients. When conversion tracking counts email clicks or irrelevant actions as leads, "
     "the algorithm optimises for the wrong thing. "
     "Clean tracking = better algorithm decisions = more real patient enquiries.",
     BLUE, "3x", "Quality Improvement"),
]

for i, (title_t, body_t, color, stat, stat_label) in enumerate(problems):
    y = Inches(1.30 + i * 1.70)

    # Main content card
    rect(sl, Inches(0.60), y, Inches(9.30), Inches(1.50), WHITE, BORDER, Pt(0.75))
    rect(sl, Inches(0.60), y, Inches(0.08), Inches(1.50), color)
    txt(sl, Inches(0.90), y + Inches(0.12), Inches(8.80), Inches(0.35), title_t, 16, NAVY, True)
    txt(sl, Inches(0.90), y + Inches(0.50), Inches(8.80), Inches(0.90), body_t, 11, BLACK)

    # Stat card to the right
    mcard(sl, Inches(10.20), y + Inches(0.15), stat, stat_label, color, Inches(2.50), Inches(1.05))

rbox(sl, Inches(0.60), Inches(6.25), Inches(12.13), Inches(0.40),
     "These are the three things I fix first. Every time. The methodology is proven.")

footer(sl, 3)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4: CASE STUDY — THE PROOF (3 hero stats only)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE

top_bar(sl, Inches(0.06))
stitle(sl, "Case Study \u2014 The Proof",
       "A UK professional services business \u2014 the exact methodology I would apply to your account")
badge(sl, "Anonymised Client")

# 3 hero stat cards — big and prominent
stats = [
    ("42.7%", "Of Budget Wasted", "Identified in first audit \u2014 \u00a37,836 of \u00a318,338 spent on non-converting traffic", RED),
    ("35-48%", "CPA Reduction", "Target CPA brought down from \u00a338 to \u00a320-25 through structured optimization", GREEN),
    ("Week 1", "Delivery", "Full audit report with findings, quick wins, and restructure plan \u2014 delivered in 5 days", BLUE),
]

for i, (val, label, desc, color) in enumerate(stats):
    col = Inches(0.60 + i * 4.10)
    w = Inches(3.80)

    rect(sl, col, Inches(1.30), w, Inches(3.20), WHITE, BORDER, Pt(0.75))
    rect(sl, col, Inches(1.30), Inches(0.08), Inches(3.20), color)

    txt(sl, col + Inches(0.30), Inches(1.50), w - Inches(0.50), Inches(1.00), val, 44, color, True)
    txt(sl, col + Inches(0.30), Inches(2.50), w - Inches(0.50), Inches(0.40), label, 18, NAVY, True)
    txt(sl, col + Inches(0.30), Inches(3.00), w - Inches(0.50), Inches(1.20), desc, 11, BLACK)

# The one-liner
rect(sl, Inches(0.60), Inches(4.80), Inches(12.13), Inches(1.00), NAVY)
txt(sl, Inches(1.00), Inches(4.90), Inches(11.33), Inches(0.80),
    "This is exactly the audit I would run on your account in week 1.",
    22, WHITE, True, PP_ALIGN.CENTER)

footer(sl, 4)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5: DENTAL & AESTHETICS EXPERIENCE
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE

top_bar(sl, Inches(0.06))
stitle(sl, "Dental & Aesthetics Experience",
       "Direct experience in the same market dynamics as Dental by Design")
badge(sl, "Sector Experience")

# Left: Cosmetic Dental
rect(sl, Inches(0.60), Inches(1.30), Inches(5.80), Inches(3.00), GREY_BG)
rect(sl, Inches(0.60), Inches(1.30), Inches(0.06), Inches(3.00), BLUE)
txt(sl, Inches(0.85), Inches(1.40), Inches(5.20), Inches(0.35),
    "Cosmetic Dental Clinic", 18, NAVY, True)
txt(sl, Inches(0.85), Inches(1.78), Inches(5.20), Inches(0.25),
    "via PodDigital Agency", 11, GREY_MID)
bullets(sl, Inches(0.85), Inches(2.10), Inches(5.20), Inches(2.00), [
    "Google Ads for implants, veneers, and cosmetic dental treatments",
    "Competitive local market with CPCs of \u00a35-15+",
    "Confidential \u2014 active agency client (can discuss approach in conversation)",
], 12)

# Right: Medizen
rect(sl, Inches(6.80), Inches(1.30), Inches(5.80), Inches(3.00), GREY_BG)
rect(sl, Inches(6.80), Inches(1.30), Inches(0.06), Inches(3.00), GREEN)
txt(sl, Inches(7.10), Inches(1.40), Inches(5.20), Inches(0.35),
    "Medizen Aesthetics", 18, NAVY, True)
txt(sl, Inches(7.10), Inches(1.78), Inches(5.20), Inches(0.25),
    "Birmingham \u2014 medically-led cosmetic clinic", 11, GREY_MID)
bullets(sl, Inches(7.10), Inches(2.10), Inches(5.20), Inches(2.00), [
    "Injectable treatments, laser procedures, aesthetic services",
    "High-value patient leads in a competitive B2C market",
    "CPA reduction through audience and landing page optimization",
], 12)

# Bottom: the connector statement
ibox(sl, Inches(0.60), Inches(4.60), Inches(3.80), Inches(1.50),
     "Same Patient Journey",
     "Search intent \u2192 ad click \u2192 landing page \u2192 enquiry \u2192 booking.",
     BLUE, BLUE_TINT)
ibox(sl, Inches(4.70), Inches(4.60), Inches(3.80), Inches(1.50),
     "Same Cost Dynamics",
     "High CPCs, long decision cycle, price-sensitive patients.",
     RED, RED_TINT)
ibox(sl, Inches(8.80), Inches(4.60), Inches(3.80), Inches(1.50),
     "Same Success Metrics",
     "Patient bookings, not clicks. Quality over volume.",
     GREEN, GREEN_TINT)

footer(sl, 5)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6: WEEK 1 PLAN — visual timeline
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE

top_bar(sl, Inches(0.06))
stitle(sl, "Week 1 Plan")
badge(sl, "First 5 Days")

days = [
    ("Days 1-2", "Full Account Audit", BLUE, [
        "Google Ads + Meta account structure deep-dive",
        "Bid strategies, conversion tracking, Quality Scores",
        "Benchmark CPL, CPA, conversion rates per clinic",
    ]),
    ("Days 3-4", "Competitive Analysis", RED, [
        "Competitor bidding \u2014 keywords, ad copy, landing pages",
        "Gaps in the dental implant keyword space",
        "Your USP positioning vs the competition",
    ]),
    ("Day 5", "Quick Wins Live", GREEN, [
        "Fix conversion tracking issues",
        "Add negative keywords to block waste",
        "Set target CPA / ROAS guardrails",
    ]),
]

# Connecting arrow line
rect(sl, Inches(1.80), Inches(1.70), Inches(9.80), Inches(0.04), BORDER)

for i, (day_l, day_t, color, items) in enumerate(days):
    cx = Inches(0.60 + i * 4.10)
    cw = Inches(3.80)

    # Circle on timeline
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(1.60), Inches(1.55), Inches(0.35), Inches(0.35))
    circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()

    # Day label
    rect(sl, cx, Inches(2.10), Inches(1.40), Inches(0.35), color)
    txt(sl, cx, Inches(2.12), Inches(1.40), Inches(0.35), day_l, 12, WHITE, True, PP_ALIGN.CENTER)

    # Title
    txt(sl, cx, Inches(2.55), cw, Inches(0.35), day_t, 16, NAVY, True)

    # Content card
    rect(sl, cx, Inches(2.95), cw, Inches(2.10), GREY_BG)
    rect(sl, cx, Inches(2.95), Inches(0.06), Inches(2.10), color)
    bullets(sl, cx + Inches(0.20), Inches(3.10), cw - Inches(0.35), Inches(1.80), items, 11)

# Deliverable — the strongest line
rect(sl, Inches(0.60), Inches(5.40), Inches(12.13), Inches(1.10), BLUE_TINT, BLUE, Pt(1))
txt(sl, Inches(0.85), Inches(5.45), Inches(11.63), Inches(0.50),
    "Full audit report with findings and 90-day roadmap by Friday.", 22, NAVY, True)
txt(sl, Inches(0.85), Inches(5.95), Inches(11.63), Inches(0.40),
    "Delivered before you've finished onboarding with any other candidate.", 12, BLACK)

footer(sl, 6)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7: 90-DAY ROADMAP
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE

top_bar(sl, Inches(0.06))
stitle(sl, "90-Day Roadmap")
badge(sl, "Months 1-3")

# Hero stat
rect(sl, Inches(0.60), Inches(1.00), Inches(12.13), Inches(1.00), WHITE)
rect(sl, Inches(0.60), Inches(1.00), Inches(0.08), Inches(1.00), GREEN)
txt(sl, Inches(0.90), Inches(1.02), Inches(5.00), Inches(0.70),
    "20-30% CPL Reduction", 40, GREEN, True)
txt(sl, Inches(6.50), Inches(1.15), Inches(6.00), Inches(0.50),
    "Within 90 days, based on experience with similar dental and aesthetics accounts", 14, BLACK)

months = [
    ("Month 1", "Restructure", BLUE, [
        "Restructure campaigns by treatment type",
        "Launch 9-list negative keyword framework",
        "Fix conversion tracking \u2014 real leads only",
    ]),
    ("Month 2", "Optimise", YELLOW, [
        "Landing page + conversion rate optimization",
        "Expand keywords from search term mining",
        "Begin multi-clinic rollout",
    ]),
    ("Month 3", "Scale", GREEN, [
        "Scale winners, cut underperformers",
        "Establish per-clinic benchmarks",
        "Full reporting + months 4-6 roadmap",
    ]),
]

for i, (lab, ti, co, its) in enumerate(months):
    cx = Inches(0.60 + i * 4.10)
    cw = Inches(3.80)

    rect(sl, cx, Inches(2.30), cw, Inches(4.10), GREY_BG)
    rect(sl, cx, Inches(2.30), Inches(0.06), Inches(4.10), co)
    txt(sl, cx+Inches(0.20), Inches(2.45), cw-Inches(0.30), Inches(0.25), lab, 11, co, True)
    txt(sl, cx+Inches(0.20), Inches(2.70), cw-Inches(0.30), Inches(0.35), ti, 16, NAVY, True)
    bullets(sl, cx+Inches(0.20), Inches(3.15), cw-Inches(0.35), Inches(2.80), its, 11)

footer(sl, 7)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8: MULTI-CLINIC SCALE
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE

top_bar(sl, Inches(0.06))
stitle(sl, "Multi-Clinic Scale",
       "How insights from Dental by Design flow across your partner network")
badge(sl, "Network Strategy")

# Flagship card (left, prominent)
rect(sl, Inches(0.60), Inches(1.30), Inches(5.00), Inches(3.80), WHITE, BORDER, Pt(0.75))
rect(sl, Inches(0.60), Inches(1.30), Inches(0.08), Inches(3.80), GREEN)
txt(sl, Inches(0.90), Inches(1.45), Inches(4.50), Inches(0.40),
    "Dental by Design", 22, GREEN, True)
txt(sl, Inches(0.90), Inches(1.90), Inches(4.50), Inches(0.30),
    "Flagship Clinic \u2014 Testing Ground", 14, NAVY, True)
bullets(sl, Inches(0.90), Inches(2.30), Inches(4.50), Inches(2.50), [
    "Highest budget = most data = fastest learning",
    "Test ad copy angles, keywords, bid strategies here first",
    "Establish benchmarks for CPL, conversion rate, Quality Score",
    "Winning strategies proven with real patient enquiry data",
], 12)

# Arrow (visual connector)
rect(sl, Inches(5.80), Inches(2.80), Inches(1.00), Inches(0.06), BLUE)
# Arrow head (small triangle approximation)
txt(sl, Inches(6.50), Inches(2.60), Inches(0.50), Inches(0.50), "\u25B6", 22, BLUE)

# Partner clinics card (right)
rect(sl, Inches(7.10), Inches(1.30), Inches(5.50), Inches(3.80), WHITE, BORDER, Pt(0.75))
rect(sl, Inches(7.10), Inches(1.30), Inches(0.08), Inches(3.80), BLUE)
txt(sl, Inches(7.40), Inches(1.45), Inches(5.00), Inches(0.40),
    "Partner Clinics", 22, BLUE, True)
txt(sl, Inches(7.40), Inches(1.90), Inches(5.00), Inches(0.30),
    "UK Network \u2014 Scaled Learnings", 14, NAVY, True)
bullets(sl, Inches(7.40), Inches(2.30), Inches(5.00), Inches(2.50), [
    "Separate campaigns per clinic \u2014 tailored messaging + targeting",
    "Shared negative keyword lists \u2014 one block benefits all clinics",
    "Performance-based budgets \u2014 best clinics earn more spend",
    "Localised ad copy with proven angles from flagship testing",
], 12)

# 3 key points at bottom
mcard(sl, Inches(0.60), Inches(5.40), "Separate", "Campaigns Per Clinic", GREEN, Inches(3.80))
mcard(sl, Inches(4.70), Inches(5.40), "Shared", "Negative Keyword Lists", BLUE, Inches(3.80))
mcard(sl, Inches(8.80), Inches(5.40), "Performance", "Based Budgets", RED, Inches(3.80))

footer(sl, 8)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9: ACT — MY EDGE (screenshot + 3 bullets + quote)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE

top_bar(sl, Inches(0.06))
stitle(sl, "ACT \u2014 My Competitive Edge",
       "A proprietary Google Ads optimization engine I built from the ground up")
badge(sl, "Proprietary Platform")

# Left: 3 bullets only + quote
txt(sl, Inches(0.60), Inches(1.30), Inches(5.80), Inches(0.35),
    "What ACT Does", 14, NAVY, True)

act_items = [
    ("30+ automated checks", "running overnight across Account, Campaign, Ad Group, Keyword, Ad, and Shopping levels", BLUE),
    ("7 bid strategy rule sets", "covering Manual CPC, tCPA, tROAS, Max Conversions, Max Clicks, PMax, and Shopping", GREEN),
    ("9 negative keyword lists", "auto-sorted by word count and match type \u2014 systematic protection against wasted spend", RED),
]

for i, (ti, bd, co) in enumerate(act_items):
    y = Inches(1.80 + i * 0.85)
    rect(sl, Inches(0.60), y, Inches(5.80), Inches(0.75), WHITE, BORDER, Pt(0.5))
    rect(sl, Inches(0.60), y, Inches(0.06), Inches(0.75), co)
    multirun(sl, Inches(0.85), y + Inches(0.10), Inches(5.30), Inches(0.55), [
        (ti + "  \u2014  ", 13, NAVY, True),
        (bd, 11, BLACK, False),
    ])

# Quote
rect(sl, Inches(0.60), Inches(4.60), Inches(5.80), Inches(1.80), NAVY)
txt(sl, Inches(0.85), Inches(4.75), Inches(5.30), Inches(1.50),
    "\u201CYour campaigns are being optimised 24/7, not just when I log in. "
    "Every action is logged, explainable, and reversible. "
    "No other candidate has built their own optimization platform.\u201D",
    14, WHITE)

# Right: Screenshot
txt(sl, Inches(6.80), Inches(1.30), Inches(5.80), Inches(0.35),
    "ACT Morning Review Dashboard", 14, NAVY, True)
rect(sl, Inches(6.80), Inches(1.75), Inches(5.80), Inches(4.65), GREY_BG, BORDER, Pt(1))
try:
    sl.shapes.add_picture(ACT_SCREENSHOT, Inches(6.90), Inches(1.85), Inches(5.60), Inches(4.45))
except:
    txt(sl, Inches(7.50), Inches(3.80), Inches(4.50), Inches(1),
        "[ACT Dashboard Screenshot]", 18, GREY_MID, a=PP_ALIGN.CENTER)

footer(sl, 9)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10: WHY ME — 5 one-liners, no descriptions
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE

top_bar(sl, Inches(0.06))
stitle(sl, "Why Me")
badge(sl, "Summary")

reasons = [
    ("Hands-on specialist \u2014 I work inside the platforms daily, not a strategist who delegates.", BLUE),
    ("Direct dental and aesthetics PPC experience \u2014 same market, same dynamics.", GREEN),
    ("Built my own Google Ads optimization platform (ACT) \u2014 no other candidate has this.", RED),
    ("Proven audit methodology \u2014 42.7% waste found, CPA reduced 35-48%.", YELLOW),
    ("I measure success by patient bookings, not clicks.", BLUE),
]

for i, (reason, color) in enumerate(reasons):
    y = Inches(1.30 + i * 1.05)
    rect(sl, Inches(0.60), y, Inches(12.10), Inches(0.85), WHITE, BORDER, Pt(0.75))
    rect(sl, Inches(0.60), y, Inches(0.06), Inches(0.85), color)
    num_circle(sl, Inches(0.85), y + Inches(0.18), i + 1, color)
    txt(sl, Inches(1.50), y + Inches(0.20), Inches(10.90), Inches(0.50), reason, 16, NAVY)

footer(sl, 10)


# ══════════════════════════════════════════════════════════════════
# SLIDE 11: NEXT STEPS — close with confidence
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
sl.background.fill.solid(); sl.background.fill.fore_color.rgb = GREY_BG

top_bar(sl)
sidebar(sl)
logo(sl)

txt(sl, Inches(0.60), Inches(1.40), Inches(5.50), Inches(0.80),
    "NEXT STEPS", 44, NAVY, True)
rect(sl, Inches(0.60), Inches(2.20), Inches(2.50), Inches(0.05), BLUE)

steps = [
    ("1", "Grant me access to your Google Ads and Meta accounts", "Read-only access is all I need to run the full audit", GREEN),
    ("2", "I'll deliver a full audit report within the first week", "Waste analysis, structure review, and restructure plan", BLUE),
    ("3", "You'll see measurable improvements within 30 days", "Quick wins in week 1, structured optimization months 1-3", RED),
]

for i, (num, main, sub, color) in enumerate(steps):
    y = Inches(2.60 + i * 1.15)
    rect(sl, Inches(0.60), y, Inches(5.50), Inches(1.00), WHITE)
    rect(sl, Inches(0.60), y, Inches(0.06), Inches(1.00), color)
    num_circle(sl, Inches(0.85), y + Inches(0.22), num, color)
    txt(sl, Inches(1.50), y + Inches(0.12), Inches(4.40), Inches(0.35), main, 14, NAVY, True)
    txt(sl, Inches(1.50), y + Inches(0.50), Inches(4.40), Inches(0.40), sub, 11, BLACK)

# Right: contact
rect(sl, Inches(6.80), Inches(0.50), Inches(5.90), Inches(2.50), WHITE)
rect(sl, Inches(6.80), Inches(0.50), Inches(0.08), Inches(2.50), BLUE)
txt(sl, Inches(7.10), Inches(0.60), Inches(5.40), Inches(0.35), "Ready to Start Immediately", 14, NAVY)
txt(sl, Inches(7.10), Inches(1.00), Inches(5.40), Inches(0.80), "Christopher Hoole", 40, BLUE, True)
multipara(sl, Inches(7.10), Inches(1.80), Inches(5.40), Inches(1.00), [
    ("chris@christopherhoole.com", 14, BLUE, False, 4),
    ("christopherhoole.com", 14, BLUE, False, 4),
    ("+44 7999 500 184", 14, BLACK, False, 0),
])

minicard(sl, Inches(6.80), Inches(3.20), "16 Years", "Hands-On Experience", BLUE)
minicard(sl, Inches(8.83), Inches(3.20), "Dental PPC", "Direct Experience", GREEN)
minicard(sl, Inches(10.87), Inches(3.20), "ACT Engine", "Proprietary Platform", RED)

rect(sl, Inches(6.80), Inches(4.60), Inches(5.90), Inches(1.20), BLUE_TINT, BLUE, Pt(1))
txt(sl, Inches(7.05), Inches(4.65), Inches(5.40), Inches(1.10),
    "\u201CI don't just manage campaigns. I've built a system to continuously optimise them. "
    "I have direct dental PPC experience. And I'm ready to start immediately.\u201D\n\n"
    "\u2014 Christopher Hoole", 11, NAVY, True)

bot_bar_title(sl)


# ── Save ──
out = os.path.join(SCRIPT_DIR, 'Christopher_Hoole_Dental_by_Design_Presentation_v4.pptx')
prs.save(out)
print(f"Saved to: {out}")
print(f"Slides: {len(prs.slides)}")
