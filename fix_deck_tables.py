# -*- coding: utf-8 -*-
"""Remove leftover tables from slides 2 and 3 of call deck."""
import sys, io, shutil
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

SRC = r'C:\Users\User\Desktop\gads-data-layer\potential_clients\Inserta Dental\End-of-week reports\DentalByDesign.co.uk - Call Deck 21 April 2026.pptx'
TMP = SRC.replace('.pptx', '_TMP.pptx')

p = Presentation(SRC)

removed = 0
for idx in [1, 2]:  # slides 2, 3
    s = p.slides[idx]
    # Find and remove table shapes
    shapes_to_remove = []
    for sh in s.shapes:
        if hasattr(sh, 'has_table') and sh.has_table:
            shapes_to_remove.append(sh)
    for sh in shapes_to_remove:
        sh._element.getparent().remove(sh._element)
        removed += 1
        print(f'Removed table on slide {idx+1}: {sh.name}')

print(f'Total removed: {removed}')

p.save(TMP)
shutil.move(TMP, SRC)
print(f'Saved: {SRC}')
