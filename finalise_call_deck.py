# -*- coding: utf-8 -*-
"""Final pass: add question content to slide 3, save clean copy."""
import sys, io, shutil
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = r'C:\Users\User\Desktop\gads-data-layer\potential_clients\Inserta Dental\End-of-week reports\DentalByDesign.co.uk - Call Deck 21 April 2026.pptx'
TMP = r'C:\Users\User\Desktop\gads-data-layer\potential_clients\Inserta Dental\End-of-week reports\DentalByDesign.co.uk - Call Deck 21 April 2026 FINAL.pptx'

p = Presentation(SRC)
s3 = p.slides[2]

# Add question content text box over empty container at (63, 148) 1164x283
tb = s3.shapes.add_textbox(Emu(63*9525), Emu(150*9525), Emu(1164*9525), Emu(283*9525))
tf = tb.text_frame
tf.word_wrap = True
tf.margin_left = Emu(120000)
tf.margin_top = Emu(80000)
tf.margin_right = Emu(120000)

NAVY = RGBColor(0x1A, 0x3E, 0x72)
RED = RGBColor(0xEA, 0x43, 0x35)

questions = [
    ('1.  Bridges — standard dental bridge vs. Vivo bridge?',
     "We're currently blocking all \"bridge\" search terms as out-of-scope. If Vivo bridges are a service you want to advertise, we split them out from generic bridge negatives and open targeted traffic."),
    ('2.  Lead quality since negatives rolled out — what are you seeing?',
     'Volume will be lower by design (we\'ve blocked freebie + NHS seekers). Any early signal from your team on consultation show-rate, treatment value, or close rate?'),
    ('3.  Search campaign priorities — where do we focus the rebuild?',
     'Search CPL is still high (£58-£200 recent days) — tracking audit + bespoke LPs queued for Week 2. Specific sub-segment (location, service, device) you want prioritised first, or let the data decide?'),
]

first = True
for q, detail in questions:
    # Question headline
    if first:
        para = tf.paragraphs[0]
        first = False
    else:
        para = tf.add_paragraph()
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = q
    run.font.name = 'Calibri'
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RED
    para.space_after = Pt(2)

    # Detail
    para2 = tf.add_paragraph()
    para2.alignment = PP_ALIGN.LEFT
    run2 = para2.add_run()
    run2.text = '     ' + detail
    run2.font.name = 'Calibri'
    run2.font.size = Pt(11)
    run2.font.color.rgb = NAVY
    para2.space_after = Pt(10)

# Save clean file (new zip, no dup entries)
p.save(TMP)
# Replace original with clean version
shutil.move(TMP, SRC)
print(f'Saved final clean deck: {SRC}')

# Verify
p2 = Presentation(SRC)
print(f'Slides: {len(p2.slides)}')
for i, s in enumerate(p2.slides):
    print(f'--- Slide {i+1}: {sum(1 for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip())} text shapes')
