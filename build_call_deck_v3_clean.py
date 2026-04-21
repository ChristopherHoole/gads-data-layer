# -*- coding: utf-8 -*-
"""Clean rebuild of call deck from v8 — removes tables first, then overlays.
Eliminates zip duplicate entries by building in one pass."""
import sys, io, shutil, copy, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

SRC = r'C:\Users\User\Desktop\gads-data-layer\potential_clients\Inserta Dental\End-of-week reports\DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v8.pptx'
DST = r'C:\Users\User\Desktop\gads-data-layer\potential_clients\Inserta Dental\End-of-week reports\DentalByDesign.co.uk - Call Deck 21 April 2026 v3.pptx'

# Clean slate — remove any existing broken file
if os.path.exists(DST):
    try:
        os.remove(DST)
    except PermissionError:
        print(f'WARNING: {DST} is open. Saving to _NEW suffix.')
        DST = DST.replace('.pptx', '_NEW.pptx')
shutil.copy(SRC, DST)

p = Presentation(DST)

# Keep slides 2, 3, 10, 11 (0-idx: 1, 2, 9, 10). Delete others.
def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])

keep = {1, 2, 9, 10}
to_delete = [i for i in range(len(p.slides)) if i not in keep]
for i in sorted(to_delete, reverse=True):
    delete_slide(p, i)

assert len(p.slides) == 4

# ---------- REMOVE TABLES from kept slides (before any editing) ----------
for idx in range(len(p.slides)):
    s = p.slides[idx]
    to_remove = [sh for sh in s.shapes if hasattr(sh, 'has_table') and sh.has_table]
    for sh in to_remove:
        sh._element.getparent().remove(sh._element)
        print(f'Removed table from slide {idx+1}: {sh.name}')

# ---------- Helpers ----------
def find_shape(slide, text_needle):
    for sh in slide.shapes:
        if sh.has_text_frame and text_needle in sh.text_frame.text:
            return sh
    return None

def set_single(shape, new_text):
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = new_text
        for r in list(p0.runs[1:]):
            r._r.getparent().remove(r._r)
    else:
        tf.text = new_text
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)

def set_multi(shape, lines):
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    ref_run = tf.paragraphs[0].runs[0] if (tf.paragraphs and tf.paragraphs[0].runs) else None
    for para in list(tf.paragraphs):
        para._p.getparent().remove(para._p)
    for line in lines:
        p_el = etree.SubElement(tf._txBody, qn('a:p'))
        r_el = etree.SubElement(p_el, qn('a:r'))
        if ref_run is not None:
            rpr = ref_run._r.find(qn('a:rPr'))
            if rpr is not None:
                r_el.append(copy.deepcopy(rpr))
        t_el = etree.SubElement(r_el, qn('a:t'))
        t_el.text = line

NAVY = RGBColor(0x1A, 0x3E, 0x72)
RED = RGBColor(0xEA, 0x43, 0x35)

def add_textbox(slide, left_px, top_px, width_px, height_px, lines, size=12, color=NAVY, bold=False, space_after=4):
    tb = slide.shapes.add_textbox(
        Emu(left_px * 9525), Emu(top_px * 9525),
        Emu(width_px * 9525), Emu(height_px * 9525))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(120000)
    tf.margin_top = Emu(80000)
    tf.margin_right = Emu(120000)
    first = True
    for line in lines:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = line
        run.font.name = 'Calibri'
        run.font.size = Pt(size)
        run.font.color.rgb = color
        if bold:
            run.font.bold = True
        para.space_after = Pt(space_after)
    return tb

# =============================================================
# SLIDE 1 — Week 1 Highlights
# =============================================================
s1 = p.slides[0]
set_single(find_shape(s1, 'Executive Summary'), 'Week 1 Highlights')
set_single(find_shape(s1, 'Week 1 · 13-17 Apr'), 'Since 14 April · 8 days')
set_single(find_shape(s1, '£8,462'), '£12.3k')
set_single(find_shape(s1, 'Week 1 Spend'), 'Spend (14-21 Apr)')
set_single(find_shape(s1, '87'), '5,000+')
set_single(find_shape(s1, 'Dengro Offline Leads'), 'Structured negatives')
set_single(find_shape(s1, '101'), '−62%')
set_single(find_shape(s1, 'Primary Conversions'), 'PMax CPA day-1 drop')
set_single(find_shape(s1, '£97'), '−49%')
set_single(find_shape(s1, 'Cost per Lead'), 'Brand CPC drop')
set_single(find_shape(s1, 'Key Wins This Week'), 'Two Big Headlines')

wins = find_shape(s1, 'Brand CPC at 8-week low')
set_multi(wins, [
    '1.  CPA is trending in the right direction — PMax dropped from £67 → £25 within 24h of the bid strategy flip (Max Conv Value → Target CPA). Brand CPA is now £11.29 (was £60+ pre-engagement).',
    '2.  Freebie & NHS-seeker traffic has been cut off — 5,000+ structured negatives live, 95 more pushed today. Volume of converts is lower by design, lead quality is rising as junk is blocked at auction source.',
    'Reporting caveat: primary-conversion CPA numbers are not fully trusted until the tracking audit completes this week. All-conversion CPL (£29.50 blended) includes verified phone + offline leads.',
])

summary = find_shape(s1, 'Week 1 was structural')
set_multi(summary, [
    'Week 1 was foundation. Week 2 is execution.',
    'Tracking audit, search-term completion, PMax rebuild and landing pages all queued this week. Bespoke LP designs will be with you by Friday for sign-off.',
])

for sh in s1.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip() == '2' and (sh.width or 0) < 500000:
        set_single(sh, '1')
        break

# =============================================================
# SLIDE 2 — What We Delivered (table removed; overlay achievements)
# =============================================================
s2 = p.slides[1]
set_single(find_shape(s2, 'Work Delivered'), 'What We Delivered — Week 1')
set_single(find_shape(s2, 'Tue 14 → Fri 17 April · 4 working days post-signing · 28 hours'),
           '14 – 21 April · 8 calendar days · ~40 hours')
set_single(find_shape(s2, 'Foundation Week'), 'Foundation Week')
set_single(find_shape(s2, 'Day-by-Day Activity'), 'Major Achievements')

# Overlay achievements in the now-empty space where table was
add_textbox(s2, 58, 155, 1158, 270, [
    '•  PMax bid strategy flipped (Max Conv Value → Target CPA £60) — CPA −62% day-1',
    '•  Brand rebuilt end-to-end: Target IS 90% + £1.50 max CPC cap — CPC −49%, CPA £11.29',
    '•  Dental Implants Intent rebuilt: 25 new ad groups, 537 keywords, 75 RSAs (1,125 headlines)',
    '•  Negative keyword structure: 2,192 loose → 5,000+ organised across 11 shared lists',
    '•  Deep audit unblocked 173 self-blocked terms (zirconia £9,990–15,990 premium, "same day teeth" USP, Hammersmith catchment)',
    '•  266-postcode Greater London geo master list replaces loose 25-mile radius (was stretching to Reading/Romford)',
    '•  Google auto-apply recommendations disabled account-wide — no more silent strategy reversals',
    '•  Website USP audit delivered; 2 live policy/accuracy risks flagged (free consult messaging, 25+ yrs claim)',
    '•  95 further wasted queries pushed as negatives today — invisalign, veneers, whitening, bonding, bridges, crowns',
], size=11, color=NAVY, space_after=5)

rt = find_shape(s2, 'Positive keywords added')
set_multi(rt, [
    'Structured negatives live: 5,000+   ·   RSAs written: 75 (1,125 headlines, 300 descriptions)',
    'New ad groups built: 25   ·   Keywords added: 537   ·   Geo targets live: 810',
    '42-cell ad schedule deployed across 3 campaigns   ·   Negatives pushed today: 95',
])

set_single(find_shape(s2, 'Key Findings'), "What's Working")
kf = find_shape(s2, 'All 3 live campaigns moved from Max Conversion Value')
set_multi(kf, [
    'PMax CPA is trending in the right direction — down from £67 to £25 within 24h of bid change.',
    'Lead quality is up now that freebie & NHS seekers are blocked at auction source.',
    'Brand is printing — 100% impression share, £1.61 CPC, £11.29 CPA.',
    'Search still needs work — tracking audit + bespoke LPs are this week\'s focus.',
])

for sh in s2.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip() == '3' and (sh.width or 0) < 500000:
        set_single(sh, '2')
        break

# =============================================================
# SLIDE 3 — Questions for You (table removed; overlay 2 questions)
# =============================================================
s3 = p.slides[2]
set_single(find_shape(s3, 'Questions for Client'), 'Questions for You')
set_single(find_shape(s3, 'To refine strategy'),
           'Two clarifications to sharpen strategy')
set_single(find_shape(s3, 'Your Input'), 'Your Input')
set_single(find_shape(s3, 'Budget Questions (Urgent)'), 'Questions')

# Question 1 header + body
add_textbox(s3, 65, 155, 1160, 30, [
    '1.  Bridges — what\'s the difference between a standard dental bridge and a Vivo bridge?',
], size=14, color=RED, bold=True)
add_textbox(s3, 85, 195, 1140, 60, [
    'We\'re currently blocking all "bridge" search terms as out-of-scope. If Vivo bridges are a service you want to advertise, we split them out from generic bridge negatives and open targeted traffic.',
], size=12, color=NAVY)

# Question 2 header + body
add_textbox(s3, 65, 290, 1160, 30, [
    '2.  Lead quality since negatives rolled out — what are you seeing from your side?',
], size=14, color=RED, bold=True)
add_textbox(s3, 85, 330, 1140, 60, [
    'Volume will be lower by design (we\'ve blocked freebie + NHS seekers). Early signal on consultation show-rate, treatment value, or close rate from your team?',
], size=12, color=NAVY)

# Footer cards
set_single(find_shape(s3, 'How to Reply'), 'Next Steps')
hr = find_shape(s3, 'Reply to this email with answers')
set_multi(hr, [
    'Today\'s call covers these 2 questions + any of yours.',
    'Week 2 plan locks in after this call — see next slide for priorities.',
])
set_single(find_shape(s3, 'Why These Matter'), 'Why These Matter')
wm = find_shape(s3, 'Every piece of clarification')
set_multi(wm, [
    'The bridge question unlocks potential new service traffic we\'re currently blocking.',
    'Lead quality feedback tells us if negatives are working as intended — or over-blocking.',
])

for sh in s3.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip() == '10' and (sh.width or 0) < 500000:
        set_single(sh, '3')
        break

# =============================================================
# SLIDE 4 — This Week's Priorities
# =============================================================
s4 = p.slides[3]
set_single(find_shape(s4, 'Week 2 Plan'), "This Week's Priorities")
set_single(find_shape(s4, '5 Working Days'), 'Week 2 · Execution')

set_single(find_shape(s4, 'Monday 20 Apr — Tracking Audit'), '1.  Finish Search Term Project')
set_single(find_shape(s4, 'Verify form_submit fires on all landing pages'),
           'Push remaining automated negatives (Leak-exact 71, Leak-phrase 55, Outside area 5, plus manual review of ~166 ambiguous). Complete Pass 3 phrase-suggestion rollout for ongoing protection.')

set_single(find_shape(s4, 'Tue 21 – Thu 23 Apr — Execution'), '2.  Rebuild & Improve PMax')
set_single(find_shape(s4, "Apply Giulio's claim corrections"),
           'Review asset group performance, tune audience signals, audit shopping feed integration, optimise creative mix. Target: sustain the CPA drop and drive volume.')

set_single(find_shape(s4, 'Friday 24 Apr — Review + Report'), '3.  Full Tracking & LP Work')
set_single(find_shape(s4, '7-day performance comparison vs Week 1'),
           'End-to-end tracking audit (form submits, phone, offline leads). Then bespoke LP designs for all Implant services — with you by Friday for sign-off. LP build follows tracking.')

bottom = find_shape(s4, 'Week 2 is about execution: tracking audit unblocks')
set_multi(bottom, [
    'Sequence matters: search-term cleanup → PMax optimisation → tracking audit → LP designs for sign-off → LP build.',
    'LPs come AFTER the structural work. You\'ll have LP designs for sign-off by Friday.',
])

for sh in s4.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip() == '11' and (sh.width or 0) < 500000:
        set_single(sh, '4')
        break

# Save (single pass, no moves)
p.save(DST)
print(f'Saved clean: {DST}')
print(f'Slides: {len(Presentation(DST).slides)}')
