# -*- coding: utf-8 -*-
"""Build 3-slide DBD call deck for 21 April 2026 call.
Clones v8 template, keeps slides 2/3/10, edits text, renumbers."""
import sys, io, shutil, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

SRC = r'C:\Users\User\Desktop\gads-data-layer\potential_clients\Inserta Dental\End-of-week reports\DentalByDesign.co.uk - Week 1 Report 13-17 April 2026-v8.pptx'
DST = r'C:\Users\User\Desktop\gads-data-layer\potential_clients\Inserta Dental\End-of-week reports\DentalByDesign.co.uk - Call Deck 21 April 2026.pptx'

shutil.copy(SRC, DST)

p = Presentation(DST)

# Keep slides 2, 3, 10 (0-idx: 1, 2, 9). Delete rest in reverse order.
# python-pptx delete hack: remove from XML + rels
def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])

keep = {1, 2, 9}  # 0-indexed
to_delete = [i for i in range(len(p.slides)) if i not in keep]
# Reverse so indexes stay valid
for i in sorted(to_delete, reverse=True):
    delete_slide(p, i)

print(f'Slides remaining: {len(p.slides)}')

# ---------- Helper: replace text in a shape preserving formatting ----------
def set_text(shape, new_text):
    """Replace text keeping first run's formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        tf.text = new_text
        return
    p0 = tf.paragraphs[0]
    # Keep first run's props
    if p0.runs:
        r0 = p0.runs[0]
        # clear other runs in para
        for r in list(p0.runs[1:]):
            r._r.getparent().remove(r._r)
        r0.text = new_text
    else:
        p0.text = new_text
    # Remove other paragraphs
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)

def set_multiline(shape, lines):
    """Replace text frame content with multiple lines, preserving first-run format."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    # Copy props from first run if exists
    ref_run = p0.runs[0] if p0.runs else None
    # Nuke all paragraphs
    for para in list(tf.paragraphs):
        para._p.getparent().remove(para._p)
    # Add new paragraphs
    from pptx.oxml.ns import qn
    from lxml import etree
    for i, line in enumerate(lines):
        p_el = etree.SubElement(tf._txBody, qn('a:p'))
        r_el = etree.SubElement(p_el, qn('a:r'))
        if ref_run is not None:
            rpr_src = ref_run._r.find(qn('a:rPr'))
            if rpr_src is not None:
                r_el.append(copy.deepcopy(rpr_src))
        t_el = etree.SubElement(r_el, qn('a:t'))
        t_el.text = line

# Find shapes by their current text content
def find_by_text(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            return sh
    return None

# ---------- SLIDE 1 (was v8 slide 2: Executive Summary) ----------
s1 = p.slides[0]

# Title + subtitle pill
set_text(find_by_text(s1, 'Executive Summary'), 'Week 1 Highlights')
set_text(find_by_text(s1, 'Week 1 · 13-17 Apr'), 'Since 14 April · 7 days')

# 4 stat cards — values + labels
# Card 1: £8,462 / Week 1 Spend
set_text(find_by_text(s1, '£8,462'), '£24.9k')
set_text(find_by_text(s1, 'Week 1 Spend'), 'Total spend (6-21 Apr)')

# Card 2: 87 / Dengro Offline Leads
set_text(find_by_text(s1, '87'), '95')
set_text(find_by_text(s1, 'Dengro Offline Leads'), 'Negatives pushed today')

# Card 3: 101 / Primary Conversions
set_text(find_by_text(s1, '101'), '847')
set_text(find_by_text(s1, 'Primary Conversions'), 'All-conversion total')

# Card 4: £97 / Cost per Lead
set_text(find_by_text(s1, '£97'), '£29.50')
set_text(find_by_text(s1, 'Cost per Lead'), 'Blended CPL (all conv.)')

# Key Wins This Week section
set_text(find_by_text(s1, 'Key Wins This Week'), 'Key Wins — First 7 Days')
wins_shape = find_by_text(s1, 'Brand CPC at 8-week low')
if wins_shape:
    set_multiline(wins_shape, [
        'Freebie & NHS seekers eliminated — lead quality rising as junk traffic is blocked at source.',
        'PMax stabilising at £31 CPL vs £45+ pre-engagement — Max Conv Value → Target CPA fix paying off.',
        '95 high-waste search queries blocked today alone (invisalign, veneers, whitening, composite bonding, etc. — services not advertised).',
        'Services scope clarified: implants-only ad spend, everything else negatived.',
        'Full tracking audit + bespoke landing pages queued for Week 2 — Search CPL will follow PMax once LP work lands.',
    ])

# Summary paragraph at bottom
summary_shape = find_by_text(s1, 'Week 1 was structural')
if summary_shape:
    set_multiline(summary_shape, [
        'Week 1 was foundation. Week 2 is optimisation.',
        'The engine is now clean enough to start real work. Focus moves to tracking accuracy, landing pages, and Search campaign rebuild.',
    ])

# Page number
set_text(find_by_text(s1, '2') if find_by_text(s1, '\n2') is None else None, '1')
# Actually find page num by size (47x23)
for sh in s1.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip() == '2' and (sh.width or 0) < 500000:
        set_text(sh, '1')
        break

# ---------- SLIDE 2 (was v8 slide 3: Work Delivered) ----------
s2 = p.slides[1]

set_text(find_by_text(s2, 'Work Delivered'), 'What We Delivered — Week 1')
set_text(find_by_text(s2, 'Tue 14 → Fri 17 April · 4 working days post-signing · 28 hours'),
         '14 – 21 April · 7 calendar days · ~40+ hours')
set_text(find_by_text(s2, 'Foundation Week'), 'Highlights')

# Day-by-Day Activity — replace with achievement bullets
set_text(find_by_text(s2, 'Day-by-Day Activity'), 'Major Achievements')

# Find the big day-by-day content shape (the one with lots of text after "Day-by-Day")
# It's the shape at (56, 148) width 1162 — large content block
day_shape = None
for sh in s2.shapes:
    if sh.has_text_frame and (sh.width or 0) > 10000000 and (sh.top or 0) > 1300000 and (sh.top or 0) < 2000000:
        day_shape = sh
        break

if day_shape:
    set_multiline(day_shape, [
        '•  3 live campaigns moved from Max Conv Value → Target CPA / Target Impression Share bidding',
        '•  Brand CPC cut from £3.29 → £1.61 (−49%) via tCPA + £1.50 max CPC cap',
        '•  2,192 loose negatives restructured into 5,000+ structured negatives across 11 shared lists',
        '•  Services & ad scope locked: implants-only (single + double arches), everything else blocked',
        '•  95 wasted queries pushed as negatives TODAY (veneers, invisalign, whitening, bonding, bridges, crowns, cosmetic dentistry)',
        '•  Account reporting infrastructure built — daily visibility into waste, leaks, and lead quality',
    ])

# Running Totals card
rt_shape = find_by_text(s2, 'Positive keywords added')
if rt_shape:
    set_multiline(rt_shape, [
        'Structured negatives live: 5,000+  ·  New RSAs written: 75 (1,050 headlines, 300 descriptions)',
        'Ad groups with positive keywords added: 25  ·  Total negatives pushed this session: 95',
        'Campaigns stabilised: 3  ·  Shared negative lists restructured: 11',
    ])

# Key Findings card
set_text(find_by_text(s2, 'Key Findings'), 'What\'s Working')
kf_shape = find_by_text(s2, 'All 3 live campaigns moved from Max Conversion Value')
if kf_shape:
    set_multiline(kf_shape, [
        'PMax CPL trending in the right direction — £31 blended, down from £45+ under old bid strategy.',
        'Lead quality is up now that freebie/NHS seekers are blocked. Volume down, but each lead is real.',
        'Brand campaign is printing — 100% IS, 0.22% optimisation headroom, £1.61 CPC.',
        'Search still struggling — tracking audit + bespoke LPs needed before we can optimise properly.',
    ])

# Page number
for sh in s2.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip() == '3' and (sh.width or 0) < 500000:
        set_text(sh, '2')
        break

# ---------- SLIDE 3 (was v8 slide 10: Questions for Client) ----------
s3 = p.slides[2]

set_text(find_by_text(s3, 'Questions for Client'), 'Questions for You')
set_text(find_by_text(s3, 'To refine strategy'),
         'Three clarifications I need to sharpen strategy')
set_text(find_by_text(s3, 'Your Input'), 'Your Input')
set_text(find_by_text(s3, 'Budget Questions (Urgent)'), 'Questions')

# The big questions content block — let me find it and replace
q_shape = None
for sh in s3.shapes:
    if sh.has_text_frame and (sh.top or 0) > 1200000 and (sh.top or 0) < 1700000 and (sh.width or 0) > 9000000:
        q_shape = sh
        break
# If not found by geometry, try finding by top position
if not q_shape:
    for sh in s3.shapes:
        if sh.has_text_frame and sh.name.endswith('389;p10'):
            q_shape = sh
            break

# Fallback — find shape at (63, 148) width 1000+
if not q_shape:
    for sh in s3.shapes:
        l = sh.left or 0
        t = sh.top or 0
        w = sh.width or 0
        if 500000 < l < 700000 and 1300000 < t < 1500000 and w > 9000000:
            q_shape = sh
            break

if q_shape:
    set_multiline(q_shape, [
        '1.  Bridges — what\'s the difference between a standard dental bridge and a Vivo bridge?',
        '    Right now we\'re blocking all "bridge" search terms as out-of-scope. If Vivo bridges are',
        '    a service you want to advertise, we need to split them out from generic bridge negatives',
        '    and open up targeted traffic.',
        '',
        '2.  Lead quality since negatives rolled out — on your side, are the leads coming through',
        '    feeling more qualified? Volume will be lower by design (we\'ve blocked freebie + NHS seekers).',
        '    Any early signal from your team on consultation show-rate, treatment value, or close rate?',
        '',
        '3.  Search campaign priorities — Search CPL is still high (£58-£200 on recent days) because',
        '    tracking fires before real intent and LPs aren\'t optimised. Tracking audit + bespoke LPs',
        '    are queued for Week 2. Is there a specific Search sub-segment (location, service, device)',
        '    you want us to prioritise rebuilding first, or should we let the data decide?',
    ])

# How to Reply footer card
set_text(find_by_text(s3, 'How to Reply'), 'Next Steps')
hr_shape = find_by_text(s3, 'Reply to this email with answers')
if hr_shape:
    set_multiline(hr_shape, [
        'Today\'s call covers these 3 questions + any of yours.',
        'Week 2 plan locks in after this call — tracking audit, LP work, next negative sweep.',
    ])

# Why These Matter
set_text(find_by_text(s3, 'Why These Matter'), 'Why These Matter')
wm_shape = find_by_text(s3, 'Every piece of clarification')
if wm_shape:
    set_multiline(wm_shape, [
        'The bridge question unlocks potential new service traffic we\'re currently blocking.',
        'Lead quality feedback tells us if negatives are working as intended — or if we\'re over-blocking.',
        'Search direction decides where the Week 2 rebuild effort goes first.',
    ])

# Page number
for sh in s3.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip() == '10' and (sh.width or 0) < 500000:
        set_text(sh, '3')
        break

p.save(DST)
print(f'✓ Saved: {DST}')
print(f'  Slides: {len(p.slides)}')
