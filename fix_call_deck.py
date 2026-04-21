# -*- coding: utf-8 -*-
"""Fix call deck:
1. Slide 1: first card value "1" → "£24.9k"
2. Slide 2: add achievement bullets text box over empty container
"""
import sys, io, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

DST = r'C:\Users\User\Desktop\gads-data-layer\potential_clients\Inserta Dental\End-of-week reports\DentalByDesign.co.uk - Call Deck 21 April 2026.pptx'

p = Presentation(DST)

# ---------- Fix slide 1 card 1 ----------
s1 = p.slides[0]
for sh in s1.shapes:
    if not sh.has_text_frame:
        continue
    l = sh.left or 0
    t = sh.top or 0
    w = sh.width or 0
    # Card 1 value: (76, 124) width 239 — in EMU ~723900, 1181100, 2276475
    if abs(l - 76*9525) < 50000 and abs(t - 124*9525) < 50000 and sh.text_frame.text.strip() == '1':
        # Replace text while keeping run formatting
        tf = sh.text_frame
        p0 = tf.paragraphs[0]
        if p0.runs:
            p0.runs[0].text = '£24.9k'
            for r in list(p0.runs[1:]):
                r._r.getparent().remove(r._r)
        else:
            tf.text = '£24.9k'
        for para in list(tf.paragraphs[1:]):
            para._p.getparent().remove(para._p)
        print('Fixed slide 1 card 1: 1 → £24.9k')
        break

# ---------- Add achievements text to slide 2 ----------
s2 = p.slides[1]

# Find the "Running Totals" text shape for formatting reference
ref_shape = None
for sh in s2.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip().startswith('Structured negatives'):
        ref_shape = sh
        break

# Add a new text box at (56, 148) with width 1162 height 249 (EMU conversion)
# Original empty container sits there — our text will overlay it
from pptx.util import Emu
left = Emu(56 * 9525)
top = Emu(150 * 9525)
width = Emu(1162 * 9525)
height = Emu(249 * 9525)

tb = s2.shapes.add_textbox(left, top, width, height)
tf = tb.text_frame
tf.word_wrap = True
tf.margin_left = Emu(120000)
tf.margin_top = Emu(80000)
tf.margin_right = Emu(120000)

bullets = [
    '3 live campaigns moved from Max Conv Value → Target CPA / Target Impression Share bidding',
    'Brand CPC cut from £3.29 → £1.61 (−49%) via tCPA + £1.50 max CPC cap',
    '2,192 loose negatives restructured into 5,000+ structured across 11 shared lists',
    'Services & ad scope locked: implants-only (single + double arches); all other services blocked',
    '95 wasted queries pushed as negatives today (veneers, invisalign, whitening, bonding, bridges, crowns, cosmetic dentistry)',
    'Account reporting infrastructure built — daily visibility into waste, leaks, lead quality',
]

# Format: navy bullets, Calibri 12pt
NAVY = RGBColor(0x1A, 0x3E, 0x72)
for i, line in enumerate(bullets):
    if i == 0:
        para = tf.paragraphs[0]
    else:
        para = tf.add_paragraph()
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = '•  ' + line
    run.font.name = 'Calibri'
    run.font.size = Pt(13)
    run.font.color.rgb = NAVY
    para.space_after = Pt(6)

print('Added 6 achievement bullets to slide 2')

p.save(DST)
print('Saved.')
